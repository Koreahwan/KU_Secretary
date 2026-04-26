from __future__ import annotations

from collections import Counter
from copy import deepcopy
from dataclasses import asdict, dataclass, field
from datetime import datetime, timedelta, timezone
from hashlib import sha1, sha256
import html
import json
import logging
import mimetypes
import os
from pathlib import Path
import re
import struct
import threading
import time
from typing import Any, Callable
from urllib.parse import parse_qsl, unquote, urlencode, urlparse, urlunparse
import zlib
from zoneinfo import ZoneInfo

from dateutil import parser as dt_parser
from dateutil.rrule import rrulestr
from dotenv import dotenv_values
import requests

from ku_secretary.assistant import (
    execute_assistant_plan,
    get_capability,
    plan_assistant_request,
)
from ku_secretary.browser_session import sanitize_browser_session_result
from ku_secretary.briefing_relay import build_signed_briefing_delivery_request
from ku_secretary.config import Settings
from ku_secretary.connectors.llm import LLMClient, LLMConfig
from ku_secretary.connectors.portal import (
    PortalNoticeFetchError,
    parse_csv_file,
    parse_ics_file,
    parse_ics_url,
)
from ku_secretary.connectors.ku_notices import (
    KuNoticeFetchError,
    KuNoticeFetchResult,
    fetch_ku_notice_feed,
)
from ku_secretary.connectors.ku_openapi import (
    KU_OPENAPI_BUILDING_SOURCE,
    KU_OPENAPI_OFFICIAL_BUILDING_URL,
    KU_OPENAPI_OFFICIAL_TIMETABLE_URL,
    KU_OPENAPI_TIMETABLE_SOURCE,
    KUOpenAPIBuildingCatalogError,
    KUOpenAPIBuildingCatalogMalformedPayload,
    KUOpenAPITimetableMalformedPayload,
    KUOpenAPITimetableUnsupported,
    fetch_ku_openapi_building_catalog,
    fetch_ku_openapi_timetable,
    ku_openapi_timetable_configured,
    ku_openapi_uses_official_catalog_mode,
)
from ku_secretary.connectors.ku_library import (
    get_library_seats,
    list_known_libraries,
)
from ku_secretary.connectors.ku_kupid_timetable import (
    KUPID_SSO_TIMETABLE_SOURCE,
    fetch_kupid_sso_timetable,
)
from ku_secretary.connectors.ku_portal import (
    KU_PORTAL_LOGIN_URL,
    KU_PORTAL_PROVIDER,
    KU_PORTAL_SCHOOL_SLUG,
    UOS_TIMETABLE_TITLE,
    KU_WISE_INDEX_URL,
    fetch_ku_portal_timetable,
)
from ku_secretary.connectors.seoul_air import SeoulAirQualityClient
from ku_secretary.connectors.telegram import TelegramBotClient, normalize_updates
from ku_secretary.connectors.uclass import (
    MoodleWSClient,
    extract_course_index,
    extract_material_candidates,
    extract_material_candidates_from_course_contents,
    normalize_action_events,
    normalize_assignments,
    normalize_forum_notifications,
    normalize_notifications,
    request_moodle_mobile_launch_token,
    request_moodle_ws_token,
)
from ku_secretary.connectors.weather_kma import (
    KMAWeatherClient,
    resolve_weather_location_query,
)
from ku_secretary.day_agenda_state import (
    DayAgendaState,
    build_day_agenda_state,
)
from ku_secretary.db import (
    Database,
    attach_provenance,
    normalize_datetime,
    normalize_provenance,
    normalize_course_alias,
    now_utc_iso,
)
from ku_secretary.models import Task
from ku_secretary.onboarding import (
    MOODLE_ONBOARDING_SESSION_KIND,
    build_public_moodle_connect_url,
    normalize_public_moodle_connect_base_url,
    onboarding_allowed_school_slugs,
    school_entry_allowed_for_onboarding,
    visible_onboarding_school_entries,
)
from ku_secretary.ops_health_state import (
    _build_notice_feed_health as _build_notice_feed_health_impl,
    _build_telegram_listener_health as _build_telegram_listener_health_impl,
    _build_telegram_send_health as _build_telegram_send_health_impl,
    _build_uclass_sync_health as _build_uclass_sync_health_impl,
    _build_ku_official_api_health as _build_ku_official_api_health_impl,
    _build_weather_sync_health as _build_weather_sync_health_impl,
    build_beta_ops_health_report as _build_beta_ops_health_report_state,
    ops_surface_state as _ops_surface_state_impl,
)
from ku_secretary.publish.dashboard import render_dashboard_snapshot
from ku_secretary.school_support import school_support_summary
from ku_secretary.secret_store import SecretStoreError, StoredSecretRef, default_secret_store
from ku_secretary.storage import (
    backups_dir as storage_backups_dir,
    dashboard_dir as storage_dashboard_dir,
    materials_dir as storage_materials_dir,
    resolve_storage_root,
)
from ku_secretary.telegram_setup_state import (
    TelegramSetupState,
    build_telegram_setup_state,
    chat_lms_connection_snapshot as _chat_lms_connection_snapshot,
    looks_like_auth_or_session_issue as _looks_like_auth_or_session_issue,
    looks_like_secure_storage_missing as _looks_like_secure_storage_missing,
    sync_dashboard_source_card as _sync_dashboard_source_card,
)


logger = logging.getLogger(__name__)
TELEGRAM_MENU_RETRY_COOLDOWN_SECONDS = 300
UOS_ONLINE_CLASS_SCHOOL_SLUG = "ku_online_class"
TELEGRAM_UCLASS_NOTICE_STALE_HOURS = 48
TELEGRAM_UOS_NOTICE_LIMIT = 10
TELEGRAM_ASSISTANT_CHAT_ACTION = "typing"
TELEGRAM_ASSISTANT_CHAT_ACTION_INTERVAL_SEC = 4.0
TELEGRAM_LMS_COURSE_SCAN_LIMIT = 20
TELEGRAM_LMS_BOARD_SCAN_LIMIT_PER_COURSE = 8
TELEGRAM_LMS_BOARD_POST_LIMIT_PER_BOARD = 5
TELEGRAM_LMS_ASSIGNMENT_DISPLAY_LIMIT = 16
TELEGRAM_LMS_ASSIGNMENT_HINT_DISPLAY_LIMIT = 12
TELEGRAM_LMS_ASSIGNMENT_BOARD_DETAIL_LIMIT_PER_BOARD = 3
TELEGRAM_LMS_SUBMITTED_DISPLAY_LIMIT = 20
TELEGRAM_LMS_MATERIAL_DISPLAY_LIMIT_PER_COURSE = 5
TELEGRAM_LMS_MESSAGE_SOFT_LIMIT = 3600
TELEGRAM_ASSIGNMENTS_CACHE_TTL_SECONDS = 180
TELEGRAM_ASSIGNMENTS_CACHE_JOB_PREFIX = "telegram_assignments_cache"
PORTAL_SECURE_STORAGE_MISSING_REASON = "KU portal session missing from secure storage; reconnect required"
UCLASS_SECURE_STORAGE_MISSING_REASON = "UClass token missing from secure storage; reconnect required"
UCLASS_RECONNECT_REQUIRED_REASON = "UClass token expired or unavailable; reconnect required"
KU_PORTAL_BROWSER_TIMETABLE_SOURCE = "ku_portal_browser"
_DAY_BRIEF_CACHE_LOCK = threading.Lock()
_DAY_BRIEF_CACHE: dict[str, dict[str, Any]] = {}
WORD_RE = re.compile(r"[A-Za-z][A-Za-z0-9_-]{2,}")
STOPWORDS = {
    "the",
    "and",
    "for",
    "with",
    "that",
    "this",
    "from",
    "your",
    "are",
    "was",
    "were",
    "have",
    "has",
    "had",
    "will",
    "about",
    "into",
    "after",
    "before",
    "during",
    "when",
    "where",
    "what",
    "which",
    "they",
    "them",
    "their",
    "you",
    "his",
    "her",
    "she",
    "him",
    "its",
    "all",
    "can",
    "could",
    "should",
    "would",
    "class",
    "lecture",
    "slide",
    "slides",
}
LOCATION_BUILDING_RE = re.compile(
    r"^\s*(?P<building_no>\d+)\s*-\s*(?P<room>[A-Za-z]?\d+)\s*(?:,\s*(?P<extra>.+))?$"
)
LOW_SIGNAL_BRIEF_PATTERNS = [
    re.compile(r"^token$", re.IGNORECASE),
    re.compile(r"^<answer>$", re.IGNORECASE),
    re.compile(r"^no additional(?: material details detected| update)?\.?$", re.IGNORECASE),
    re.compile(r"^review updates and schedule next step\.?$", re.IGNORECASE),
    re.compile(r"^what is the core argument in .+\?$", re.IGNORECASE),
]
MATERIAL_WEEK_MARKER_PATTERNS = (
    re.compile(r"(?<!\d)(\d{1,2})\s*주차"),
    re.compile(r"제\s*(\d{1,2})\s*주"),
    re.compile(r"\bweek\s*(\d{1,2})\b", re.IGNORECASE),
    re.compile(r"\bweek[-_]?(\d{1,2})\b", re.IGNORECASE),
    re.compile(r"\bw(\d{1,2})\b", re.IGNORECASE),
)
MATERIAL_BRIEF_INVALID_ITEM_PATTERNS = [
    re.compile(r"로그인 페이지", re.IGNORECASE),
    re.compile(r"실제 (?:수업|강의) (?:내용|주제).*(?:포함되어 있지 않|없)", re.IGNORECASE),
    re.compile(r"파일(?:이)? 첨부되지 않았", re.IGNORECASE),
    re.compile(r"구체적인 .*(?:생성|확인)할 수 없", re.IGNORECASE),
    re.compile(r"실제 (?:수업|강의) 자료.*제공", re.IGNORECASE),
    re.compile(r"해당 url 로 접속", re.IGNORECASE),
]
MATERIAL_GENERIC_HTML_FILENAME_RE = re.compile(
    r"^(?:index|view|resource|file)(?:_\d+)?\.(?:php|html?|aspx?)$",
    re.IGNORECASE,
)
MATERIAL_HTML_LOGIN_MARKERS = (
    re.compile(r'name=["\']logintoken["\']', re.IGNORECASE),
    re.compile(r'name=["\']username["\']', re.IGNORECASE),
    re.compile(r'type=["\']password["\']', re.IGNORECASE),
    re.compile(r'action=["\'][^"\']*/login/index\.php', re.IGNORECASE),
)
MATERIAL_HTML_LOGIN_TEXT_HINTS = (
    "login",
    "로그인",
    "username",
    "password",
    "아이디",
    "비밀번호",
    "remember username",
    "온라인강의실",
)
BRIEF_PROVENANCE_SOURCE_LABELS = {
    "llm_inferred": "AI",
    "portal_csv": "포털",
    "portal_ics": "포털",
    "portal_ics_url": "포털",
    "portal_uos_timetable": "포털",
    "telegram_draft": "Telegram",
    "uclass_html": "UClass",
    "uclass_ws": "UClass",
}
UOS_NOTICE_FEEDS = {
    "general": {
        "label": "학교 일반공지",
        "board_id": "566",
        # legacy fields retained for downstream metadata payloads that key on
        # `list_id` / `menuid`. We mirror the KU board id so existing readers
        # still get a stable identifier without referring to the old uos.ac.kr
        # endpoints.
        "list_id": "566",
        "menuid": "566",
    },
    "academic": {
        "label": "학교 학사공지",
        "board_id": "567",
        "list_id": "567",
        "menuid": "567",
    },
}
MATERIAL_BRIEF_SYSTEM_PROMPT = (
    "You are an academic assistant for a Korean university student. "
    "Write in concise Korean. Summarize the class material into exactly 3 bullet points and 1 action item. "
    "Focus on lecture topics, exam-important concepts, and what the student should prepare next. "
    "If a file is attached, use it as the primary source and use the text excerpt only as fallback context. "
    "Format:\n"
    "- ...\n"
    "- ...\n"
    "- ...\n"
    "Action: ..."
)
MATERIAL_BRIEF_VERSION = 2
COURSE_DAY_SUMMARY_VERSION = 1
COURSE_DAY_SUMMARY_SYSTEM_PROMPT = (
    "You are an academic assistant for a Korean university student. "
    "You will receive already-summarized class materials for one course on one day. "
    "Synthesize them into one combined course summary in concise Korean. "
    "Return JSON only with this schema: "
    '{"short_summary":"...","long_bullets":["...","...","..."],"review":"..."}. '
    "Rules: `short_summary` must be 1-2 short sentences for a Telegram preview. "
    "`long_bullets` must contain 2-3 distinct bullets for /todaysummary. "
    "`review` must be one short action sentence for what to review next. "
    "Do not mention filenames unless needed to disambiguate topics. "
    "Merge overlapping material points instead of repeating them."
)
MATERIAL_BRIEF_ATTACHMENT_SUFFIXES = {
    ".csv",
    ".doc",
    ".docx",
    ".md",
    ".pdf",
    ".ppt",
    ".pptx",
    ".rtf",
    ".txt",
}
MATERIAL_DEADLINE_SYSTEM_PROMPT = (
    "You extract assignment deadlines from Korean university class materials. "
    "Return JSON only in this schema: "
    '{"tasks":[{"title":"...","due_at":"2026-03-10T23:59:00+09:00","evidence":"..."}]}. '
    "Include only actual student deliverables with a submission deadline. "
    "Ignore lecture meeting times, class dates, and vague schedules. "
    "If no deadline exists, return {\"tasks\":[]}."
)
TASK_MERGE_CACHE_VERSION = 1
TASK_MERGE_CACHE_JOB_NAME = "task_merge_cache"
TASK_MERGE_LLM_MAX_COMPONENTS = 12
TASK_MERGE_LLM_MAX_COMPONENT_SIZE = 6
TASK_MERGE_CANDIDATE_DUE_WINDOW_HOURS = 36
TASK_MERGE_SYSTEM_PROMPT = (
    "You identify duplicate university assignments for a student reminder app. "
    "Return JSON only in this schema: "
    '{"groups":[{"ids":["...","..."],"merged_title":"...","confidence":"high|medium|low","reason":"..."}]}. '
    "Merge only when items are clearly the same deliverable for the same course. "
    "Do not merge different weeks, different parts, different files, or merely related topics. "
    "If nothing should merge, return {\"groups\":[]}."
)
TASK_MERGE_TOKEN_RE = re.compile(r"[0-9A-Za-z\u3131-\u318E\uAC00-\uD7A3]{2,}")
TASK_MERGE_STOPWORDS = {
    "assignment",
    "class",
    "course",
    "task",
    "과제",
    "과목",
    "강의",
    "수업",
    "제출",
    "안내",
    "자료",
    "파일",
}
MATERIAL_DEADLINE_SCAN_VERSION = 2
DAY_BRIEF_ARTIFACT_LOOKBACK_DAYS = 120
DAY_BRIEF_NOTIFICATION_LOOKBACK_DAYS = 45
TELEGRAM_DAY_ARTIFACT_LIMIT = 360
TELEGRAM_DAY_NOTIFICATION_LIMIT = 120
TELEGRAM_DAY_OPEN_TASK_LIMIT = 180
TELEGRAM_DAY_SUMMARY_ARTIFACT_LIMIT = 480
TELEGRAM_DAY_SUMMARY_NOTIFICATION_LIMIT = 120
TELEGRAM_DAY_SUMMARY_OPEN_TASK_LIMIT = 180
SCHEDULED_BRIEFING_ARTIFACT_LIMIT = 480
SCHEDULED_BRIEFING_NOTIFICATION_LIMIT = 150
SCHEDULED_BRIEFING_OPEN_TASK_LIMIT = 240
DAY_BRIEF_CACHE_MAX_ENTRIES = 128
MATERIAL_TASK_HINT_KEYWORDS = (
    "과제",
    "제출",
    "마감",
    "기한",
    "assignment",
    "deadline",
    "due",
    "homework",
    "hw",
    "project",
    "report",
    "quiz",
    "레포트",
    "보고서",
)
WEATHER_SNAPSHOT_MAX_AGE_MINUTES = 90
WEATHER_REGION_RESET_TOKENS = {
    "default",
    "reset",
    "clear",
    "기본",
    "기본값",
    "초기화",
    "삭제",
    "해제",
}
MATERIAL_DATE_PATTERNS = (
    re.compile(
        r"(?P<year>20\d{2})\s*[./-]\s*(?P<month>\d{1,2})\s*[./-]\s*(?P<day>\d{1,2})"
    ),
    re.compile(
        r"(?:(?P<year>20\d{2})\s*년\s*)?(?P<month>\d{1,2})\s*월\s*(?P<day>\d{1,2})\s*일"
    ),
    re.compile(r"(?<!\d)(?P<month>\d{1,2})\s*/\s*(?P<day>\d{1,2})(?!\d)"),
)
MATERIAL_TIME_PATTERNS = (
    re.compile(
        r"(?P<ampm>오전|오후|am|pm|AM|PM)?\s*(?P<hour>\d{1,2})\s*:\s*(?P<minute>\d{1,2})"
    ),
    re.compile(
        r"(?P<ampm>오전|오후|am|pm|AM|PM)?\s*(?P<hour>\d{1,2})\s*시\s*(?P<minute>\d{1,2})?\s*분?"
    ),
    re.compile(r"(?P<ampm>오전|오후|am|pm|AM|PM)\s*(?P<hour>\d{1,2})(?:\s*시)?"),
)
MATERIAL_HTML_SECTION_PATTERNS = (
    re.compile(
        r'<div[^>]*class="[^"]*article-content[^"]*"[^>]*>([\s\S]*?)<div[^>]*class="[^"]*article-buttons[^"]*"',
        re.I,
    ),
    re.compile(
        r'<div[^>]*class="[^"]*text_to_html[^"]*"[^>]*>([\s\S]*?)</div>',
        re.I,
    ),
    re.compile(r"<main[^>]*>([\s\S]*?)</main>", re.I),
    re.compile(r"<body[^>]*>([\s\S]*?)</body>", re.I),
)
MATERIAL_HTML_TITLE_RE = re.compile(r"<title[^>]*>([\s\S]*?)</title>", re.I)
MATERIAL_HTML_TAG_RE = re.compile(r"<[^>]+>")
MATERIAL_HTML_DROP_BLOCK_RE = re.compile(
    r"<(script|style|svg|noscript)[^>]*>[\s\S]*?</\1>",
    re.I,
)
MATERIAL_BLOCK_TAG_RE = re.compile(
    r"</?(?:p|div|li|ul|ol|table|tr|td|th|section|article|main|header|footer|aside|h[1-6])[^>]*>",
    re.I,
)
MATERIAL_BREAK_TAG_RE = re.compile(r"<br\s*/?>", re.I)
MATERIAL_CONTROL_CHAR_RE = re.compile(r"[\x00-\x08\x0b-\x1f\x7f]")
HWP_PARA_TEXT_TAG = 67


@dataclass
class PipelineSummary:
    ok: bool
    stats: dict[str, Any]
    errors: list[str]


@dataclass(frozen=True)
class PortalTimetableAdapter:
    provider: str
    school_slug: str
    job_name: str
    no_target_reason: str
    resolve_targets: Callable[[Settings, Database], list[dict[str, Any]]]
    sync_target: Callable[..., dict[str, Any]]


def _append_token_to_url(url: str, token: str) -> str:
    parsed = urlparse(url)
    query_items = parse_qsl(parsed.query, keep_blank_values=True)
    has_token = any(key == "token" for key, _ in query_items)
    if not has_token:
        query_items.append(("token", token))
    return urlunparse(
        (
            parsed.scheme,
            parsed.netloc,
            parsed.path,
            parsed.params,
            urlencode(query_items, doseq=True),
            parsed.fragment,
        )
    )


def _safe_filename(name: str) -> str:
    base = str(name or "").replace("\\", "/").split("/")[-1]
    keep = []
    for ch in base:
        if ch.isalnum() or ch in {"-", "_", ".", " "}:
            keep.append(ch)
        else:
            keep.append("_")
    cleaned = re.sub(r"_+", "_", "".join(keep)).strip().replace(" ", "_")
    return cleaned or "material.bin"


def _safe_path_component(name: str, fallback: str) -> str:
    sanitized = _safe_filename(name).strip("._")
    if not sanitized:
        return fallback
    return sanitized


def _sha256_file(path: Path) -> str:
    digest = sha256()
    with path.open("rb") as fp:
        while True:
            chunk = fp.read(1024 * 1024)
            if not chunk:
                break
            digest.update(chunk)
    return digest.hexdigest()


def _filename_from_content_disposition(value: str | None) -> str | None:
    if not value:
        return None

    def _repair_text(raw_text: str) -> str:
        candidate = raw_text.strip()
        if not candidate:
            return candidate
        try:
            repaired = candidate.encode("latin-1").decode("utf-8")
        except Exception:
            return candidate
        return repaired.strip() or candidate

    parts = [part.strip() for part in value.split(";") if part.strip()]
    for part in parts:
        lowered = part.lower()
        if lowered.startswith("filename*="):
            raw = part.split("=", 1)[1].strip().strip('"').strip("'")
            if "''" in raw:
                raw = raw.split("''", 1)[1]
            decoded = unquote(raw).strip()
            return _safe_filename(decoded) if decoded else None
    for part in parts:
        lowered = part.lower()
        if lowered.startswith("filename="):
            raw = part.split("=", 1)[1].strip().strip('"').strip("'")
            decoded = _repair_text(raw)
            return _safe_filename(decoded) if decoded else None
    return None


def _mime_extension(content_type: str | None) -> str | None:
    if not content_type:
        return None
    mime = str(content_type).split(";", 1)[0].strip().lower()
    if not mime:
        return None
    if mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return ".pptx"
    if mime == "application/vnd.ms-powerpoint":
        return ".ppt"
    guessed = mimetypes.guess_extension(mime)
    if guessed in {".jpe"}:
        return ".jpg"
    return guessed


def _filename_from_url_path(url: str, default: str = "material") -> str:
    parsed = urlparse(url)
    tail = Path(parsed.path).name
    return tail or default


def _normalize_content_type(content_type: str | None) -> str:
    return str(content_type or "").split(";", 1)[0].strip().lower()


def _parse_material_error_payload(
    payload: bytes,
    *,
    content_type: str | None,
) -> dict[str, Any] | None:
    normalized_content_type = _normalize_content_type(content_type)
    stripped = payload.lstrip()
    looks_json = "json" in normalized_content_type or stripped.startswith(b"{")
    if not looks_json:
        return None

    text: str | None = None
    for encoding in ("utf-8-sig", "utf-8", "cp949", "latin-1"):
        try:
            text = payload.decode(encoding)
            break
        except Exception:
            continue
    if text is None:
        return None
    try:
        parsed = json.loads(text)
    except Exception:
        return None
    if not isinstance(parsed, dict):
        return None
    if not any(
        key in parsed for key in ("error", "errorcode", "message", "exception", "debuginfo")
    ):
        return None
    return parsed


def _material_error_message_from_payload(
    payload: bytes,
    *,
    content_type: str | None,
    resolved_url: str | None = None,
) -> str | None:
    if _looks_like_uclass_login_html_payload(
        payload,
        content_type=content_type,
        resolved_url=resolved_url,
    ):
        return "uclass returned login page html instead of material content"
    parsed = _parse_material_error_payload(payload, content_type=content_type)
    if not isinstance(parsed, dict):
        return None

    details: list[str] = []
    for key in ("errorcode", "error", "message", "exception"):
        value = str(parsed.get(key) or "").strip()
        if value and value not in details:
            details.append(value)
    detail_text = " / ".join(details) if details else "unknown json error payload"
    return f"uclass returned error payload: {detail_text}"


def _material_error_message_from_file(path: Path, *, content_type: str | None) -> str | None:
    normalized_content_type = _normalize_content_type(content_type)
    try:
        if "json" not in normalized_content_type and path.stat().st_size > 16384:
            return None
        payload = path.read_bytes()
    except Exception:
        return None
    return _material_error_message_from_payload(payload, content_type=content_type)


def _decode_material_payload_text(payload: bytes, *, max_chars: int = 16000) -> str:
    if not payload:
        return ""
    try:
        return payload[:max_chars].decode("utf-8", errors="ignore")
    except Exception:
        return ""


def _looks_like_uclass_login_html_payload(
    payload: bytes,
    *,
    content_type: str | None,
    resolved_url: str | None = None,
) -> bool:
    normalized_content_type = _normalize_content_type(content_type)
    text = _decode_material_payload_text(payload)
    lowered = text.lower()
    if (
        "html" not in normalized_content_type
        and "<html" not in lowered
        and "<!doctype html" not in lowered
        and "<form" not in lowered
    ):
        return False
    resolved_path = urlparse(str(resolved_url or "")).path.lower()
    if "/login/index.php" in resolved_path:
        return True
    structure_hits = sum(1 for pattern in MATERIAL_HTML_LOGIN_MARKERS if pattern.search(text))
    text_hits = sum(1 for hint in MATERIAL_HTML_LOGIN_TEXT_HINTS if hint in lowered)
    return structure_hits >= 2 or (structure_hits >= 1 and text_hits >= 3)


def _current_uclass_download_token(settings: Settings) -> str:
    resolved = str(getattr(settings, "_uclass_resolved_token", "") or "").strip()
    if resolved:
        return resolved
    return str(getattr(settings, "uclass_wstoken", "") or "").strip()


def _rewrite_uclass_session_download_url(url: str) -> str:
    parsed = urlparse(url)
    marker = "/webservice/pluginfile.php/"
    if marker not in parsed.path:
        return url
    return urlunparse(
        (
            parsed.scheme,
            parsed.netloc,
            parsed.path.replace(marker, "/pluginfile.php/", 1),
            parsed.params,
            parsed.query,
            parsed.fragment,
        )
    )


def _resolve_download_filename(
    *,
    url: str,
    content_disposition: str | None,
    content_type: str | None,
    fallback_name: str,
) -> str:
    candidate = (
        _filename_from_content_disposition(content_disposition)
        or _safe_filename(_filename_from_url_path(url, default=fallback_name))
    )
    if Path(candidate).suffix:
        return _safe_filename(candidate)
    ext = _mime_extension(content_type) or ""
    return _safe_filename(f"{candidate}{ext}") if ext else _safe_filename(candidate)


def _resolve_collision_target(target: Path, content_hash: str) -> Path:
    if not target.exists():
        return target
    try:
        if _sha256_file(target) == content_hash:
            return target
    except Exception:
        pass
    stem = target.stem
    suffix = target.suffix
    for idx in range(1, 2000):
        candidate = target.with_name(f"{stem}_{idx}{suffix}")
        if not candidate.exists():
            return candidate
        try:
            if _sha256_file(candidate) == content_hash:
                return candidate
        except Exception:
            continue
    return target.with_name(f"{stem}_{sha1(content_hash.encode('utf-8')).hexdigest()[:8]}{suffix}")


def _download_material_response(
    url: str,
    token: str | None,
    retries: int,
    backoff_sec: float,
) -> tuple[bytes, dict[str, str], str]:
    request_url = _append_token_to_url(url, token) if token else url
    attempt = 0
    while True:
        attempt += 1
        try:
            response = requests.get(request_url, timeout=60)
            status = response.status_code
            if status in {408, 429} or status >= 500:
                raise requests.HTTPError(
                    f"transient http status: {status}",
                    response=response,
                )
            response.raise_for_status()
            headers = {str(k).lower(): str(v) for k, v in dict(response.headers).items()}
            return response.content, headers, str(response.url or request_url)
        except Exception as exc:
            transient = isinstance(
                exc,
                (
                    requests.Timeout,
                    requests.ConnectionError,
                    requests.HTTPError,
                ),
            )
            if isinstance(exc, requests.HTTPError):
                response = exc.response
                status = response.status_code if response is not None else 0
                transient = status in {408, 425, 429} or status >= 500
            if not transient or attempt >= max(retries, 1):
                raise
            sleep_for = backoff_sec * (2 ** (attempt - 1))
            logger.warning(
                "retrying material download",
                extra={"url": url, "attempt": attempt, "sleep_sec": sleep_for},
            )
            time.sleep(sleep_for)


def _download_material_response_via_session(
    *,
    url: str,
    session: requests.Session,
    retries: int,
    backoff_sec: float,
) -> tuple[bytes, dict[str, str], str]:
    attempt = 0
    while True:
        attempt += 1
        try:
            response = session.get(url, timeout=60, allow_redirects=True)
            status = response.status_code
            if status in {408, 429} or status >= 500:
                raise requests.HTTPError(
                    f"transient http status: {status}",
                    response=response,
                )
            response.raise_for_status()
            headers = {str(k).lower(): str(v) for k, v in dict(response.headers).items()}
            return response.content, headers, str(response.url or url)
        except Exception as exc:
            transient = isinstance(
                exc,
                (
                    requests.Timeout,
                    requests.ConnectionError,
                    requests.HTTPError,
                ),
            )
            if isinstance(exc, requests.HTTPError):
                response = exc.response
                status = response.status_code if response is not None else 0
                transient = status in {408, 425, 429} or status >= 500
            if not transient or attempt >= max(retries, 1):
                raise
            sleep_for = backoff_sec * (2 ** (attempt - 1))
            logger.warning(
                "retrying session material download",
                extra={"url": url, "attempt": attempt, "sleep_sec": sleep_for},
            )
            time.sleep(sleep_for)


def _should_use_uclass_session_download(settings: Settings, url: str) -> bool:
    session = getattr(settings, "_uclass_html_session", None)
    if session is None:
        return False
    ws_base = str(getattr(settings, "uclass_ws_base", "") or "").strip()
    if not ws_base:
        return False
    try:
        ws_netloc = urlparse(ws_base).netloc
        url_netloc = urlparse(url).netloc
    except Exception:
        return False
    return bool(ws_netloc and ws_netloc == url_netloc)


def _download_material(
    db: Database,
    settings: Settings,
    external_id: str,
    url: str,
    target: Path,
    *,
    owner_id: int,
) -> tuple[str, str, bool, dict[str, Any]]:
    existing = db.get_artifact(
        external_id=external_id,
        source="uclass",
        user_id=owner_id,
    )
    invalid_existing_path: Path | None = None
    if existing and existing.content_hash and existing.icloud_path:
        existing_path = Path(existing.icloud_path)
        if existing_path.exists():
            existing_meta = _json_load(existing.metadata_json)
            existing_content_type = existing_meta.get("content_type")
            existing_error = _material_error_message_from_file(
                existing_path,
                content_type=str(existing_content_type or ""),
            )
            if existing_error:
                invalid_existing_path = existing_path
                logger.warning(
                    "discarding invalid existing material artifact",
                    extra={
                        "external_id": external_id,
                        "path": str(existing_path),
                        "error": existing_error,
                    },
                )
            else:
                return (
                    str(existing_path),
                    existing.content_hash,
                    False,
                    {
                        "original_url": str(existing_meta.get("original_url") or url),
                        "resolved_filename": str(
                            existing_meta.get("resolved_filename") or existing_path.name
                        ),
                        "content_type": existing_content_type,
                    },
                )

    if _should_use_uclass_session_download(settings, url):
        payload, headers, final_url = _download_material_response_via_session(
            url=_rewrite_uclass_session_download_url(url),
            session=getattr(settings, "_uclass_html_session"),
            retries=settings.uclass_download_retries,
            backoff_sec=settings.uclass_download_backoff_sec,
        )
    else:
        token = _current_uclass_download_token(settings)
        payload, headers, final_url = _download_material_response(
            url=url,
            token=token,
            retries=settings.uclass_download_retries,
            backoff_sec=settings.uclass_download_backoff_sec,
        )
        content_type = headers.get("content-type")
        payload_error = _material_error_message_from_payload(
            payload,
            content_type=content_type,
            resolved_url=final_url,
        )
        if payload_error:
            try:
                fresh_token = _resolve_uclass_token(settings, prefer_static=False)
            except Exception:
                fresh_token = ""
            if fresh_token and fresh_token != token:
                setattr(settings, "_uclass_resolved_token", fresh_token)
                payload, headers, final_url = _download_material_response(
                    url=url,
                    token=fresh_token,
                    retries=settings.uclass_download_retries,
                    backoff_sec=settings.uclass_download_backoff_sec,
                )
                content_type = headers.get("content-type")
                payload_error = _material_error_message_from_payload(
                    payload,
                    content_type=content_type,
                    resolved_url=final_url,
                )
            if payload_error:
                raise RuntimeError(payload_error)
    content_type = headers.get("content-type")
    payload_error = _material_error_message_from_payload(
        payload,
        content_type=content_type,
        resolved_url=final_url,
    )
    if payload_error:
        raise RuntimeError(payload_error)
    resolved_filename = _resolve_download_filename(
        url=final_url or url,
        content_disposition=headers.get("content-disposition"),
        content_type=content_type,
        fallback_name=target.name,
    )
    target = target.with_name(_safe_filename(resolved_filename))
    digest = sha256(payload).hexdigest()
    target.parent.mkdir(parents=True, exist_ok=True)
    if invalid_existing_path and invalid_existing_path.exists():
        try:
            invalid_existing_path.unlink()
        except Exception:
            logger.warning(
                "failed to remove invalid existing material artifact",
                extra={
                    "external_id": external_id,
                    "path": str(invalid_existing_path),
                },
            )
    target = _resolve_collision_target(target, digest)
    if target.exists() and _sha256_file(target) == digest:
        return (
            str(target),
            digest,
            False,
            {
                "original_url": url,
                "resolved_filename": target.name,
                "content_type": content_type,
            },
        )
    target.write_bytes(payload)
    return (
        str(target),
        digest,
        True,
        {
            "original_url": url,
            "resolved_filename": target.name,
            "content_type": content_type,
        },
    )


def _json_load(value: str | None) -> dict[str, Any]:
    if not value:
        return {}
    try:
        parsed = json.loads(value)
    except Exception:
        return {}
    return parsed if isinstance(parsed, dict) else {}


def _sync_dashboard_cursor(
    db: Database,
    job_name: str,
    *,
    status: str,
    new_items: int = 0,
    action_required: int = 0,
    last_error: str | None = None,
    cursor_payload: dict[str, Any] | None = None,
    user_id: int | None = None,
) -> dict[str, Any]:
    payload = dict(cursor_payload or {})
    previous = db.get_sync_state(job_name, user_id=user_id)
    previous_cursor = _json_load(previous.last_cursor_json)
    previous_meta = (
        previous_cursor.get("_sync_dashboard")
        if isinstance(previous_cursor.get("_sync_dashboard"), dict)
        else {}
    )
    last_success_at = str(previous_meta.get("last_success_at") or "").strip() or None
    if status == "success":
        last_success_at = now_utc_iso()
    payload["_sync_dashboard"] = {
        "status": status,
        "new_items": max(int(new_items), 0),
        "action_required": max(int(action_required), 0),
        "last_error": str(last_error or "").strip() or None,
        "last_success_at": last_success_at,
    }
    return payload


def _record_sync_dashboard_state(
    db: Database,
    job_name: str,
    *,
    status: str,
    new_items: int = 0,
    action_required: int = 0,
    last_error: str | None = None,
    cursor_payload: dict[str, Any] | None = None,
    last_run_at: str | None = None,
    user_id: int | None = None,
) -> None:
    db.update_sync_state(
        job_name,
        last_run_at=last_run_at or now_utc_iso(),
        last_cursor_json=_sync_dashboard_cursor(
            db=db,
            job_name=job_name,
            status=status,
            new_items=new_items,
            action_required=action_required,
            last_error=last_error,
            cursor_payload=cursor_payload,
            user_id=user_id,
        ),
        user_id=user_id,
    )


def _provenance_brief_tag(metadata_json: str | dict[str, Any] | None, *, fallback_source: str) -> str:
    provenance = normalize_provenance(metadata_json, fallback_source=fallback_source)
    source = str(provenance.get("source") or "").strip().lower()
    label = BRIEF_PROVENANCE_SOURCE_LABELS.get(source) or str(
        provenance.get("source_label") or ""
    ).strip()
    if not label or label.lower() == "unknown":
        return ""
    return f"[{label}]"


def _ws_required(required_ws: set[str], wsfunction: str) -> bool:
    return wsfunction in required_ws


def _call_optional_ws(
    ws_status: dict[str, dict[str, Any]],
    wsfunction: str,
    enabled: bool,
    required_ws: set[str],
    loader: Callable[[], Any],
    log_skip: bool = True,
) -> Any:
    status = ws_status.setdefault(
        wsfunction, {"ok": 0, "failed": 0, "skipped": 0, "last_error": None}
    )
    if not enabled:
        status["skipped"] += 1
        if log_skip:
            logger.info("uclass ws skipped", extra={"wsfunction": wsfunction})
        return None
    try:
        payload = loader()
        status["ok"] += 1
        logger.info("uclass ws success", extra={"wsfunction": wsfunction})
        return payload
    except Exception as exc:
        status["failed"] += 1
        status["last_error"] = str(exc)
        logger.warning(
            "uclass ws failed",
            extra={"wsfunction": wsfunction, "error": str(exc)},
        )
        if _ws_required(required_ws, wsfunction):
            raise
        return None


def _payload_shape(payload: Any) -> dict[str, Any]:
    if isinstance(payload, dict):
        return {"type": "dict", "keys": sorted(payload.keys())[:8]}
    if isinstance(payload, list):
        return {"type": "list", "len": len(payload)}
    return {"type": type(payload).__name__}


def _payload_shape_fingerprint(payload: Any, depth: int = 0, max_depth: int = 2) -> str:
    if depth >= max_depth:
        return type(payload).__name__
    if isinstance(payload, dict):
        if not payload:
            return "dict{}"
        parts: list[str] = []
        keys = sorted(payload.keys())
        for key in keys[:10]:
            value = payload.get(key)
            if isinstance(value, dict):
                value_type = _payload_shape_fingerprint(value, depth + 1, max_depth)
            elif isinstance(value, list):
                value_type = _payload_shape_fingerprint(value, depth + 1, max_depth)
            else:
                value_type = type(value).__name__
            parts.append(f"{key}:{value_type}")
        suffix = ",..." if len(keys) > 10 else ""
        return "dict{" + ",".join(parts) + suffix + "}"
    if isinstance(payload, list):
        if not payload:
            return "list[len=0]"
        first = payload[0]
        first_type = _payload_shape_fingerprint(first, depth + 1, max_depth)
        return f"list[len={len(payload)};first={first_type}]"
    return type(payload).__name__


def _log_missing_semantic_fields(
    *,
    category: str,
    wsfunction: str,
    external_id: str,
    raw: dict[str, Any],
    missing_fields: list[str],
) -> int:
    if not missing_fields:
        return 0
    logger.warning(
        "uclass semantic fields missing",
        extra={
            "category": category,
            "wsfunction": wsfunction,
            "external_id": external_id,
            "missing_fields": missing_fields,
            "shape_fingerprint": _payload_shape_fingerprint(raw),
        },
    )
    return 1


def _warn_popup_notification_semantics(
    notifications: list[Any],
    wsfunction: str,
) -> int:
    warned = 0
    for item in notifications:
        raw = item.metadata.get("raw") if hasattr(item, "metadata") else None
        if not isinstance(raw, dict):
            continue
        missing: list[str] = []
        if not any(raw.get(key) not in (None, "") for key in ("id", "notificationid")):
            missing.append("id_or_notificationid")
        if not any(
            str(raw.get(key) or "").strip()
            for key in ("subject", "title", "name", "fullmessage", "smallmessage", "message", "text")
        ):
            missing.append("title_or_message")
        if not any(raw.get(key) not in (None, "") for key in ("timecreated", "createdat", "timemodified")):
            missing.append("timecreated_or_createdat_or_timemodified")
        warned += _log_missing_semantic_fields(
            category="popup_notifications",
            wsfunction=wsfunction,
            external_id=str(getattr(item, "external_id", "")),
            raw=raw,
            missing_fields=missing,
        )
    return warned


def _warn_assignment_semantics(tasks: list[Any], wsfunction: str) -> int:
    warned = 0
    for task in tasks:
        external_id = str(getattr(task, "external_id", ""))
        if not external_id.startswith("uclass:assign"):
            continue
        raw = task.metadata.get("raw") if hasattr(task, "metadata") else None
        if not isinstance(raw, dict):
            continue
        missing: list[str] = []
        if raw.get("id") in (None, "", 0, "0"):
            missing.append("id")
        if not str(raw.get("name") or "").strip():
            missing.append("name")
        if not any(
            raw.get(key) not in (None, "", 0, "0")
            for key in ("duedate", "cutoffdate", "allowsubmissionsfromdate")
        ):
            missing.append("due_date_family")
        warned += _log_missing_semantic_fields(
            category="assignments",
            wsfunction=wsfunction,
            external_id=external_id,
            raw=raw,
            missing_fields=missing,
        )
    return warned


def _warn_forum_semantics(notifications: list[Any], wsfunction: str) -> int:
    warned = 0
    for item in notifications:
        raw = item.metadata.get("raw") if hasattr(item, "metadata") else None
        if not isinstance(raw, dict):
            continue
        missing: list[str] = []
        if raw.get("discussion") in (None, "", 0, "0"):
            missing.append("discussion")
        if not any(
            str(raw.get(key) or "").strip() for key in ("name", "subject", "message")
        ):
            missing.append("name_or_subject_or_message")
        if not any(raw.get(key) not in (None, "") for key in ("timemodified", "created")):
            missing.append("timemodified_or_created")
        warned += _log_missing_semantic_fields(
            category="forum_discussions",
            wsfunction=wsfunction,
            external_id=str(getattr(item, "external_id", "")),
            raw=raw,
            missing_fields=missing,
        )
    return warned


def _extract_first_forum_id(payload: Any) -> int | None:
    if isinstance(payload, list):
        forums = payload
    elif isinstance(payload, dict):
        forums = payload.get("forums") if isinstance(payload.get("forums"), list) else []
    else:
        forums = []
    for forum in forums:
        if not isinstance(forum, dict):
            continue
        raw = forum.get("id")
        try:
            return int(raw)
        except (TypeError, ValueError):
            continue
    return None


def _parse_dt(value: str | None) -> datetime | None:
    if not value:
        return None
    try:
        parsed = dt_parser.isoparse(value)
    except Exception:
        return None
    if parsed.tzinfo is None:
        parsed = parsed.replace(tzinfo=timezone.utc)
    return parsed


def _parse_datetime_like(value: Any) -> datetime | None:
    if value in (None, ""):
        return None
    if isinstance(value, (int, float)) and not isinstance(value, bool):
        try:
            return datetime.fromtimestamp(float(value), tz=timezone.utc)
        except Exception:
            return None
    text = str(value).strip()
    if not text:
        return None
    if re.fullmatch(r"\d{10}(?:\.\d+)?", text) or re.fullmatch(r"\d{13}", text):
        try:
            numeric = float(text)
            if len(text.split(".", 1)[0]) >= 13:
                numeric /= 1000.0
            return datetime.fromtimestamp(numeric, tz=timezone.utc)
        except Exception:
            return None
    return _parse_dt(text)


def _safe_inbox_suffix(external_id: str) -> str:
    parts = [part for part in external_id.split(":") if part]
    raw = parts[-1] if parts else external_id
    cleaned = re.sub(r"[^a-zA-Z0-9_-]+", "", raw)
    if cleaned:
        return cleaned
    return sha1(external_id.encode("utf-8")).hexdigest()[:16]


def _normalize_token_text(value: str) -> str:
    text = str(value or "").strip().lower()
    return re.sub(r"[\s\-\(\)_]+", "", text)


def _parse_uos_location(location: str | None) -> dict[str, Any]:
    raw = str(location or "").strip()
    if not raw:
        return {"raw": "", "building_no": None, "room": None, "extra": None, "ok": False}
    match = LOCATION_BUILDING_RE.match(raw)
    if not match:
        return {
            "raw": raw,
            "building_no": None,
            "room": None,
            "extra": None,
            "ok": False,
        }
    return {
        "raw": raw,
        "building_no": str(match.group("building_no") or "").strip() or None,
        "room": str(match.group("room") or "").strip() or None,
        "extra": str(match.group("extra") or "").strip() or None,
        "ok": True,
    }


def _timetable_weekday_code(item: dict[str, Any]) -> str | None:
    metadata = dict(item.get("metadata") or {}) if isinstance(item.get("metadata"), dict) else {}
    raw_weekday = str(metadata.get("weekday_code") or "").strip().upper()
    if raw_weekday in {"MO", "TU", "WE", "TH", "FR", "SA", "SU"}:
        return raw_weekday
    rrule_text = str(item.get("rrule") or "").strip().upper()
    match = re.search(r"BYDAY=([A-Z]{2})", rrule_text)
    if match:
        return match.group(1)
    start_at = _parse_dt(str(item.get("start_at") or "").strip() or None)
    if start_at is None:
        return None
    return ("MO", "TU", "WE", "TH", "FR", "SA", "SU")[start_at.astimezone(ZoneInfo("Asia/Seoul")).weekday()]


def _timetable_match_key(item: dict[str, Any]) -> tuple[str, str | None, str | None, str | None, str | None]:
    title_key = normalize_course_alias(str(item.get("title") or "").strip())
    location_meta = _parse_uos_location(str(item.get("location") or "").strip() or None)
    building_no = str(location_meta.get("building_no") or "").strip() or None
    room = str(location_meta.get("room") or "").strip() or None
    start_at = _parse_dt(str(item.get("start_at") or "").strip() or None)
    end_at = _parse_dt(str(item.get("end_at") or "").strip() or None)
    start_hm = start_at.astimezone(ZoneInfo("Asia/Seoul")).strftime("%H:%M") if start_at else None
    end_hm = end_at.astimezone(ZoneInfo("Asia/Seoul")).strftime("%H:%M") if end_at else None
    return title_key, _timetable_weekday_code(item), start_hm, end_hm, (
        f"{building_no}-{room}" if building_no and room else None
    )


def _best_official_catalog_match(
    item: dict[str, Any],
    *,
    official_index: dict[tuple[str, str | None, str | None, str | None, str | None], list[dict[str, Any]]],
) -> dict[str, Any] | None:
    key = _timetable_match_key(item)
    matches = list(official_index.get(key) or [])
    if len(matches) == 1:
        return matches[0]
    if len(matches) > 1:
        return None
    relaxed_key = (key[0], key[1], key[2], key[3], None)
    relaxed_matches = list(official_index.get(relaxed_key) or [])
    if len(relaxed_matches) == 1:
        return relaxed_matches[0]
    return None


def _enrich_ku_portal_events_with_official_catalog(
    settings: Settings,
    *,
    events: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    api_key = str(getattr(settings, "ku_openapi_timetable_api_key", "") or "").strip()
    api_url = str(getattr(settings, "ku_openapi_timetable_url", "") or "").strip()
    if not api_key:
        return events
    if not ku_openapi_uses_official_catalog_mode(api_url):
        return events
    timeout_sec = max(int(getattr(settings, "ku_openapi_timetable_timeout_sec", 15) or 15), 1)
    try:
        fetched = fetch_ku_openapi_timetable(
            api_url=api_url or KU_OPENAPI_OFFICIAL_TIMETABLE_URL,
            api_key=api_key,
            academic_year=getattr(settings, "ku_openapi_year", None),
            term=getattr(settings, "ku_openapi_term", None),
            timezone_name=str(getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul"),
            timeout_sec=timeout_sec,
        )
    except (KUOpenAPITimetableUnsupported, KUOpenAPITimetableMalformedPayload, requests.RequestException):
        logger.warning("uos official catalog enrichment skipped", exc_info=True)
        return events

    official_events = list(fetched.get("events") or [])
    if not official_events:
        return events

    official_index: dict[tuple[str, str | None, str | None, str | None, str | None], list[dict[str, Any]]] = {}
    for official_item in official_events:
        official_index.setdefault(_timetable_match_key(official_item), []).append(official_item)

    enriched_events: list[dict[str, Any]] = []
    source_url = str(fetched.get("source_url") or "").strip() or None
    for item in events:
        official_match = _best_official_catalog_match(item, official_index=official_index)
        if official_match is None:
            enriched_events.append(item)
            continue
        metadata = dict(item.get("metadata")) if isinstance(item.get("metadata"), dict) else {}
        official_metadata = (
            dict(official_match.get("metadata"))
            if isinstance(official_match.get("metadata"), dict)
            else {}
        )
        for key, value in official_metadata.items():
            if value in (None, "", []):
                continue
            if key.startswith("official_") or key in {"instructor", "class_nm"}:
                metadata[key] = value
        if source_url:
            metadata["official_catalog_source_url"] = source_url
        metadata["official_catalog_match_external_id"] = str(
            official_match.get("external_id") or ""
        ).strip() or None
        enriched_events.append({**item, "metadata": metadata})
    return enriched_events


def _event_occurrences_on_date(
    event: Any,
    target_date_local: datetime,
    timezone_name: str,
) -> list[tuple[datetime, datetime]]:
    start_dt = _parse_dt(getattr(event, "start_at", None))
    end_dt = _parse_dt(getattr(event, "end_at", None))
    if start_dt is None or end_dt is None:
        return []
    if end_dt <= start_dt:
        end_dt = start_dt + timedelta(hours=1)
    tz = ZoneInfo(timezone_name)
    day_start_local = target_date_local.astimezone(tz).replace(
        hour=0,
        minute=0,
        second=0,
        microsecond=0,
    )
    day_end_local = day_start_local + timedelta(days=1)
    day_start_utc = day_start_local.astimezone(timezone.utc)
    day_end_utc = day_end_local.astimezone(timezone.utc)
    duration = end_dt - start_dt
    rrule_value = str(getattr(event, "rrule", "") or "").strip()

    if not rrule_value:
        if start_dt < day_end_utc and end_dt > day_start_utc:
            return [(start_dt.astimezone(tz), end_dt.astimezone(tz))]
        return []
    rule = rrule_value if rrule_value.upper().startswith("RRULE:") else f"RRULE:{rrule_value}"
    try:
        parsed_rule = rrulestr(rule, dtstart=start_dt)
        starts = parsed_rule.between(day_start_utc, day_end_utc, inc=True)
    except Exception:
        starts = []
    output: list[tuple[datetime, datetime]] = []
    for occurrence_start in starts:
        if occurrence_start.tzinfo is None:
            occurrence_start = occurrence_start.replace(tzinfo=start_dt.tzinfo or timezone.utc)
        occurrence_end = occurrence_start + duration
        output.append((occurrence_start.astimezone(tz), occurrence_end.astimezone(tz)))
    return output


def _strip_markdown_fence(text: str) -> str:
    body = str(text or "").strip()
    if body.startswith("```"):
        body = re.sub(r"^```[a-zA-Z0-9_-]*\s*", "", body)
        body = re.sub(r"\s*```$", "", body).strip()
    return body


def _parse_llm_json_payload(text: str) -> dict[str, Any] | None:
    stripped = _strip_markdown_fence(text)
    if not stripped:
        return None
    candidates = [stripped]
    start = stripped.find("{")
    end = stripped.rfind("}")
    if start != -1 and end != -1 and end > start:
        candidates.append(stripped[start : end + 1])
    for candidate in candidates:
        try:
            parsed = json.loads(candidate)
        except Exception:
            continue
        if isinstance(parsed, dict):
            return parsed
    return None


BRIEFING_LLM_TIMEOUT_SEC_CAP = 20
MATERIAL_LLM_TIMEOUT_SEC_CAP = 30
MATERIAL_BRIEF_LLM_TEXT_EXCERPT_CHAR_CAP = 6000
MATERIAL_DEADLINE_LLM_TEXT_EXCERPT_CHAR_CAP = 6000


def _llm_timeout_sec(settings: Settings, *, cap: int | None = None) -> int:
    timeout = max(int(getattr(settings, "llm_timeout_sec", 30) or 30), 1)
    if cap is not None:
        timeout = min(timeout, max(int(cap), 1))
    return timeout


def _material_llm_text_excerpt_limit(settings: Settings, *, cap: int) -> int:
    configured = max(int(getattr(settings, "material_extract_max_chars", 12000) or 12000), 1)
    return min(configured, max(int(cap), 1))


def _llm_client(
    settings: Settings,
    *,
    timeout_sec: int | None = None,
) -> LLMClient:
    return LLMClient(
        LLMConfig(
            provider=getattr(settings, "llm_provider", "local"),
            model=getattr(settings, "llm_model", "gemma4"),
            timeout_sec=(
                max(int(timeout_sec), 1)
                if timeout_sec is not None
                else getattr(settings, "llm_timeout_sec", 30)
            ),
            local_endpoint=getattr(
                settings,
                "llm_local_endpoint",
                "http://127.0.0.1:11434/api/chat",
            ),
        )
    )


def _llm_client_with_timeout(
    settings: Settings,
    *,
    timeout_sec: int | None = None,
) -> LLMClient:
    try:
        return _llm_client(settings, timeout_sec=timeout_sec)
    except TypeError:
        # Some tests monkeypatch `_llm_client` with a narrower signature.
        return _llm_client(settings)


def _contains_material_task_hints(*values: str) -> bool:
    haystack = " ".join(str(value or "") for value in values).lower()
    return any(keyword in haystack for keyword in MATERIAL_TASK_HINT_KEYWORDS)


def _contains_submission_action_hints(*values: str) -> bool:
    haystack = " ".join(str(value or "") for value in values).lower()
    title = str(values[0] if values else "")
    submit_signals = (
        "제출",
        "작성",
        "등록",
        "첨부",
        "내세요",
        "내기",
        "챙겨",
        "submit",
        "turn in",
        "hand in",
    )
    deliverable_signals = (
        "과제물",
        "보고서",
        "레포트",
        "리포트",
        "프로젝트",
        "project",
        "report",
        "assignment",
        "homework",
        "quiz",
        "퀴즈",
        "hw",
    )
    deadline_signals = ("마감", "기한", "까지", "due", "deadline")
    title_has_strong_signal = any(
        signal in title.lower()
        for signal in (*submit_signals, *deliverable_signals)
    )
    if any(noise in title for noise in ("우수자", "우수 자", "발표 안내")) and not title_has_strong_signal:
        return False
    if (
        any(noise in haystack for noise in ("강의 업로드", "강의를 업로드", "업로드했", "업로드 했"))
        and "제출" not in haystack
        and "과제" not in haystack
    ):
        return False
    if re.search(
        r"(?:업로드|upload).{0,24}(?:제출|submit)|(?:제출|submit).{0,24}(?:업로드|upload)",
        haystack,
    ):
        return True
    if any(signal in haystack for signal in submit_signals):
        return True
    if any(signal in haystack for signal in deliverable_signals):
        return True
    return "과제" in haystack and any(signal in haystack for signal in deadline_signals)


def _snippet_has_explicit_time(text: str) -> bool:
    snippet = str(text or "")
    return any(pattern.search(snippet) for pattern in MATERIAL_TIME_PATTERNS)


def _snippet_has_explicit_date(text: str) -> bool:
    snippet = str(text or "")
    return any(pattern.search(snippet) for pattern in MATERIAL_DATE_PATTERNS)


def _coerce_due_iso(
    value: str | None,
    *,
    timezone_name: str,
    reference_local: datetime,
    default_end_of_day: bool = False,
) -> str | None:
    text = str(value or "").strip()
    if not text:
        return None
    parsed = _parse_dt(text)
    if parsed is not None:
        if parsed.tzinfo is None:
            parsed = parsed.replace(tzinfo=ZoneInfo(timezone_name))
            if default_end_of_day and not _snippet_has_explicit_time(text):
                parsed = parsed.replace(hour=23, minute=59, second=0, microsecond=0)
        return parsed.astimezone(timezone.utc).replace(microsecond=0).isoformat()

    tz = ZoneInfo(timezone_name)
    date_match = None
    for pattern in MATERIAL_DATE_PATTERNS:
        date_match = pattern.search(text)
        if date_match:
            break
    if date_match is None:
        return None
    year_raw = date_match.groupdict().get("year")
    month = int(date_match.group("month"))
    day = int(date_match.group("day"))
    year = int(year_raw) if year_raw else reference_local.year
    hour = 23 if default_end_of_day else 0
    minute = 59 if default_end_of_day else 0
    for pattern in MATERIAL_TIME_PATTERNS:
        time_match = pattern.search(text)
        if not time_match:
            continue
        hour = int(time_match.group("hour"))
        minute_raw = time_match.groupdict().get("minute")
        minute = int(minute_raw) if minute_raw not in {None, ""} else 0
        ampm = str(time_match.groupdict().get("ampm") or "").strip().lower()
        if ampm in {"pm", "오후"} and hour < 12:
            hour += 12
        elif ampm in {"am", "오전"} and hour == 12:
            hour = 0
        break
    try:
        parsed = datetime(year, month, day, hour, minute, tzinfo=tz)
    except ValueError:
        return None
    if not year_raw and parsed < (reference_local - timedelta(days=90)):
        try:
            parsed = parsed.replace(year=parsed.year + 1)
        except ValueError:
            pass
    return parsed.astimezone(timezone.utc).replace(microsecond=0).isoformat()


def _normalize_task_title_key(value: str) -> str:
    lowered = str(value or "").strip().lower()
    return re.sub(r"[^0-9a-zA-Z가-힣]+", "", lowered)


def _int_or_none(value: Any) -> int | None:
    if value in (None, ""):
        return None
    try:
        return int(value)
    except (TypeError, ValueError):
        return None


def _uclass_course_namespace(settings: Settings, *, ws_base_url: str | None = None) -> str:
    ws_base = str(ws_base_url or getattr(settings, "uclass_ws_base", "") or "").strip()
    parsed = urlparse(ws_base)
    host = str(parsed.netloc or "").strip().lower()
    return re.sub(r"[^a-z0-9]+", "-", host).strip("-") or "uclass"


def _uclass_canonical_course_id(
    settings: Settings,
    course_id: Any,
    *,
    ws_base_url: str | None = None,
) -> str | None:
    numeric_id = _int_or_none(course_id)
    if numeric_id is None:
        return None
    return f"uclass:{_uclass_course_namespace(settings, ws_base_url=ws_base_url)}:{numeric_id}"


_UOS_UCLASS_SHORTNAME_SECTION_RE = re.compile(
    r"(?<!\d)(?P<subject_no>\d{5})[_-](?P<dvcl_no>\d{2})(?:[_-](?P<section>[A-Za-z0-9]+))?(?!\d)"
)


def _extract_uos_official_course_metadata_from_uclass(
    course_meta: dict[str, Any],
) -> dict[str, str]:
    for key in ("shortname", "idnumber", "fullname", "displayname"):
        text = str(course_meta.get(key) or "").strip()
        if not text:
            continue
        match = _UOS_UCLASS_SHORTNAME_SECTION_RE.search(text)
        if not match:
            continue
        subject_no = str(match.group("subject_no") or "").strip()
        dvcl_no = str(match.group("dvcl_no") or "").strip()
        course_section = str(match.group("section") or "").strip()
        output = {
            "official_subject_no": subject_no,
            "official_dvcl_no": dvcl_no,
        }
        if course_section:
            output["official_course_section"] = course_section
        return output
    return {}


def _course_alias_pairs(course_meta: dict[str, Any]) -> list[tuple[str, str]]:
    subject_no = str(course_meta.get("official_subject_no") or "").strip()
    dvcl_no = str(
        course_meta.get("official_dvcl_no")
        or course_meta.get("official_course_section")
        or ""
    ).strip()
    candidates = [
        ("fullname", course_meta.get("fullname")),
        ("displayname", course_meta.get("displayname")),
        ("shortname", course_meta.get("shortname")),
        ("course_name", course_meta.get("course_name")),
        ("idnumber", course_meta.get("idnumber")),
        ("official_course_name", course_meta.get("official_course_name")),
        ("official_course_code", course_meta.get("official_course_code")),
        ("official_subject_no", subject_no),
        (
            "official_subject_section",
            f"{subject_no}-{dvcl_no}" if subject_no and dvcl_no else "",
        ),
    ]
    output: list[tuple[str, str]] = []
    seen: set[str] = set()
    for alias_type, raw_value in candidates:
        alias = str(raw_value or "").strip()
        normalized = normalize_course_alias(alias)
        if not alias or not normalized or normalized in seen:
            continue
        seen.add(normalized)
        output.append((alias_type, alias))
    return output


def _register_uclass_courses(
    settings: Settings,
    db: Database,
    course_index: dict[int, dict[str, Any]],
    *,
    user_id: int | None = 0,
    ws_base_url: str | None = None,
) -> dict[int, str]:
    canonical_by_course_id: dict[int, str] = {}
    for course_id, course_meta in sorted(course_index.items()):
        canonical_course_id = _uclass_canonical_course_id(
            settings,
            course_id,
            ws_base_url=ws_base_url,
        )
        if not canonical_course_id:
            continue
        existing_by_external = db.find_course_by_external_id(
            source="uclass",
            external_course_id=str(course_id),
            user_id=user_id,
        )
        if existing_by_external is not None:
            canonical_course_id = existing_by_external.canonical_course_id
        existing_course = db.get_course(canonical_course_id, user_id=user_id)
        existing_metadata = (
            _json_load(existing_course.metadata_json)
            if existing_course is not None
            else {}
        )
        display_name = (
            str(course_meta.get("fullname") or "").strip()
            or str(course_meta.get("displayname") or "").strip()
            or str(course_meta.get("shortname") or "").strip()
            or f"course-{course_id}"
        )
        course_metadata = {**existing_metadata, **dict(course_meta)}
        extracted_official = _extract_uos_official_course_metadata_from_uclass(course_metadata)
        for key, value in extracted_official.items():
            if value and not str(course_metadata.get(key) or "").strip():
                course_metadata[key] = value
        course_metadata["provider"] = "uclass"
        course_metadata["provider_namespace"] = _uclass_course_namespace(
            settings,
            ws_base_url=ws_base_url,
        )
        db.upsert_course(
            canonical_course_id=canonical_course_id,
            source="uclass",
            external_course_id=course_id,
            display_name=display_name,
            metadata_json=course_metadata,
            user_id=user_id,
        )
        for alias_type, alias in _course_alias_pairs(course_metadata):
            db.upsert_course_alias(
                canonical_course_id=canonical_course_id,
                alias=alias,
                alias_type=alias_type,
                source="uclass",
                metadata_json={"external_course_id": course_id},
                user_id=user_id,
            )
        canonical_by_course_id[course_id] = canonical_course_id
    return canonical_by_course_id


def _course_alias_candidates(metadata: dict[str, Any], *fallback_aliases: str) -> list[str]:
    output: list[str] = []
    seen: set[str] = set()
    for key in (
        "course_name",
        "coursefullname",
        "course_name_raw",
        "coursename",
        "displayname",
        "fullname",
        "shortname",
    ):
        raw_value = metadata.get(key)
        alias = str(raw_value or "").strip()
        normalized = normalize_course_alias(alias)
        if not alias or not normalized or normalized in seen:
            continue
        seen.add(normalized)
        output.append(alias)
    for raw_value in fallback_aliases:
        alias = str(raw_value or "").strip()
        normalized = normalize_course_alias(alias)
        if not alias or not normalized or normalized in seen:
            continue
        seen.add(normalized)
        output.append(alias)
    return output


def _resolve_canonical_course_id(
    metadata_json: str | dict[str, Any] | None,
    *,
    alias_map: dict[str, tuple[str, ...]],
    fallback_aliases: list[str] | None = None,
) -> str:
    metadata = metadata_json if isinstance(metadata_json, dict) else _json_load(metadata_json)
    canonical_course_id = str(metadata.get("canonical_course_id") or "").strip()
    if canonical_course_id:
        return canonical_course_id
    for alias in _course_alias_candidates(metadata, *(fallback_aliases or [])):
        normalized = normalize_course_alias(alias)
        matches = alias_map.get(normalized, ())
        if len(matches) == 1:
            return matches[0]
    return ""


def _attach_canonical_course_id(
    metadata_json: dict[str, Any] | None,
    *,
    alias_map: dict[str, tuple[str, ...]] | None = None,
    canonical_by_course_id: dict[int, str] | None = None,
    fallback_aliases: list[str] | None = None,
) -> dict[str, Any]:
    metadata = dict(metadata_json or {})
    canonical_course_id = str(metadata.get("canonical_course_id") or "").strip()
    if canonical_course_id:
        return metadata

    course_id = _int_or_none(metadata.get("course_id") or metadata.get("courseid"))
    if course_id is not None and canonical_by_course_id and course_id in canonical_by_course_id:
        metadata["canonical_course_id"] = canonical_by_course_id[course_id]
        return metadata

    if alias_map:
        resolved = _resolve_canonical_course_id(
            metadata,
            alias_map=alias_map,
            fallback_aliases=fallback_aliases,
        )
        if resolved:
            metadata["canonical_course_id"] = resolved
    return metadata


def _course_display_tokens(value: str) -> list[str]:
    raw = re.sub(r"\([^)]*\)", " ", str(value or "").strip())
    return [
        token
        for token in re.findall(r"[A-Za-z0-9\u3131-\u318E\uAC00-\uD7A3]{2,}", raw)
        if not token.isdigit()
    ]


def _canonical_course_key(value: str) -> str:
    tokens = _course_display_tokens(value)
    if tokens:
        return "".join(token.lower() for token in tokens)
    return _normalize_task_title_key(re.sub(r"\([^)]*\)", " ", str(value or "").strip()))


def _display_course_name(value: str) -> str:
    tokens = _course_display_tokens(value)
    if tokens:
        return " ".join(tokens)
    return re.sub(r"\s+", " ", re.sub(r"\([^)]*\)", " ", str(value or "").strip())).strip()


def _material_task_titles_similar(left: str, right: str) -> bool:
    left_key = _normalize_task_title_key(left)
    right_key = _normalize_task_title_key(right)
    if not left_key or not right_key:
        return False
    if left_key == right_key:
        return True
    shorter, longer = sorted((left_key, right_key), key=len)
    return len(shorter) >= 4 and shorter in longer


def _is_equivalent_existing_task(
    db: Database,
    *,
    title: str,
    due_at: str | None,
    course_name: str,
    canonical_course_id: str | None,
    external_id: str,
) -> bool:
    due_dt = _parse_dt(due_at)
    course_key = _canonical_course_key(course_name)
    linked_course_id = str(canonical_course_id or "").strip()
    for existing in db.list_tasks(limit=1000):
        if str(existing.external_id) == external_id:
            return False
        if not _material_task_titles_similar(existing.title, title):
            continue
        existing_meta = _json_load(existing.metadata_json)
        existing_linked_course_id = str(existing_meta.get("canonical_course_id") or "").strip()
        if linked_course_id and existing_linked_course_id and linked_course_id != existing_linked_course_id:
            continue
        existing_course_key = _canonical_course_key(str(existing_meta.get("course_name") or ""))
        if course_key and existing_course_key and course_key != existing_course_key:
            continue
        existing_due_dt = _parse_dt(existing.due_at)
        if due_dt is None and existing_due_dt is None:
            return True
        if due_dt is None or existing_due_dt is None:
            continue
        if abs((existing_due_dt - due_dt).total_seconds()) <= 15 * 60:
            return True
    return False


def _clean_material_task_title(value: str) -> str:
    text = str(value or "").strip()
    if not text:
        return ""
    text = re.sub(r"\([^)]*(?:월|화|수|목|금|토|일)[^)]*\)", " ", text)
    for pattern in MATERIAL_DATE_PATTERNS:
        match = pattern.search(text)
        if match:
            text = text[: match.start()] + " " + text[match.end() :]
    for pattern in MATERIAL_TIME_PATTERNS:
        text = pattern.sub(" ", text)
    text = re.sub(
        r"(?i)\b(?:due|deadline|assignment|homework|project|report|quiz|hw)\b",
        " ",
        text,
    )
    text = re.sub(r"(제출\s*기한|제출일|마감일|마감|제출|기한|까지)", " ", text)
    text = re.sub(r"[\[\]\(\):\-|]", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    generic = {"", "제출", "마감", "기한", "due", "deadline", "assignment"}
    if text.lower() in generic:
        return ""
    return text[:120]


def _material_title_from_lines(
    lines: list[str],
    index: int,
    *,
    fallback_title: str,
) -> str:
    candidates: list[str] = []
    if 0 <= index - 1 < len(lines):
        candidates.append(lines[index - 1])
    if 0 <= index < len(lines):
        candidates.append(lines[index])
    if 0 <= index + 1 < len(lines):
        candidates.append(lines[index + 1])
    for candidate in candidates:
        cleaned = _clean_material_task_title(candidate)
        if cleaned and _contains_material_task_hints(candidate):
            return cleaned
    for candidate in candidates:
        cleaned = _clean_material_task_title(candidate)
        if cleaned:
            return cleaned
    return Path(fallback_title).stem or fallback_title


def _heuristic_material_deadline_tasks(
    *,
    title: str,
    course_name: str,
    extracted_text: str,
    timezone_name: str,
    reference_local: datetime,
) -> list[dict[str, Any]]:
    lines = [line.strip() for line in str(extracted_text or "").splitlines() if line.strip()]
    output: list[dict[str, Any]] = []
    seen: set[tuple[str, str | None]] = set()
    for index, line in enumerate(lines):
        prev_line = lines[index - 1] if index > 0 else ""
        line_has_hint = _contains_material_task_hints(line)
        prev_has_hint = _contains_material_task_hints(prev_line)
        line_has_date = _snippet_has_explicit_date(line)
        if not (line_has_hint or (line_has_date and prev_has_hint)):
            continue
        block = " ".join(
            item
            for item in [
                prev_line,
                line,
                lines[index + 1] if index + 1 < len(lines) else "",
            ]
            if item
        )
        if not _contains_material_task_hints(block):
            continue
        due_at = _coerce_due_iso(
            block,
            timezone_name=timezone_name,
            reference_local=reference_local,
            default_end_of_day=True,
        )
        if not due_at:
            continue
        task_title = _material_title_from_lines(lines, index, fallback_title=title)
        key = (_normalize_task_title_key(task_title), due_at)
        if key in seen:
            continue
        seen.add(key)
        output.append(
            {
                "title": task_title,
                "due_at": due_at,
                "evidence": block[:240],
                "course_name": course_name,
                "method": "heuristic",
            }
        )
    return output


def _llm_material_deadline_tasks(
    *,
    settings: Settings,
    db: Database,
    title: str,
    course_name: str,
    extracted_text: str,
    local_path: str | None,
    reference_local: datetime,
) -> tuple[list[dict[str, Any]], str | None]:
    if not bool(getattr(settings, "llm_enabled", False)):
        return [], None
    gate = _identity_warning_gate(
        settings=settings,
        db=db,
        step="material_deadline_extract",
        destination="llm",
    )
    if gate is not None:
        return [], str(gate.get("error") or "identity_ack_required")
    attachment_paths = _material_llm_attachment_paths(settings, local_path)
    payload = {
        "mode": "material_deadline_extract",
        "title": title,
        "course_name": course_name,
        "timezone": settings.timezone,
        "reference_local": reference_local.isoformat(),
        "file_attached": bool(attachment_paths),
        "text_excerpt": str(extracted_text or "")[
            : _material_llm_text_excerpt_limit(
                settings,
                cap=MATERIAL_DEADLINE_LLM_TEXT_EXCERPT_CHAR_CAP,
            )
        ],
        "instructions": (
            "Extract only actual student deliverables with deadlines. "
            "Return JSON only with tasks[]."
        ),
    }
    try:
        client = _llm_client_with_timeout(
            settings,
            timeout_sec=_llm_timeout_sec(
                settings,
                cap=MATERIAL_LLM_TIMEOUT_SEC_CAP,
            ),
        )
        raw = client.generate_text(
            system_prompt=MATERIAL_DEADLINE_SYSTEM_PROMPT,
            prompt=(
                "Extract assignment deadlines from the following JSON payload.\n\n"
                + json.dumps(payload, ensure_ascii=False, sort_keys=True)
            ),
            attachment_paths=attachment_paths,
        )
        parsed = _parse_llm_json_payload(raw) or {}
    except Exception as exc:
        logger.warning(
            "material deadline llm fallback",
            extra={"title": title, "error": str(exc)},
        )
        return [], str(exc)
    items = parsed.get("tasks")
    if not isinstance(items, list):
        return [], None
    output: list[dict[str, Any]] = []
    seen: set[tuple[str, str | None]] = set()
    for item in items:
        if not isinstance(item, dict):
            continue
        task_title = _clean_material_task_title(str(item.get("title") or "")) or Path(title).stem or title
        due_at = _coerce_due_iso(
            str(item.get("due_at") or ""),
            timezone_name=settings.timezone,
            reference_local=reference_local,
            default_end_of_day=True,
        )
        if not due_at:
            evidence_due = _coerce_due_iso(
                str(item.get("evidence") or ""),
                timezone_name=settings.timezone,
                reference_local=reference_local,
                default_end_of_day=True,
            )
            due_at = evidence_due
        if not due_at:
            continue
        key = (_normalize_task_title_key(task_title), due_at)
        if key in seen:
            continue
        seen.add(key)
        output.append(
            {
                "title": task_title,
                "due_at": due_at,
                "evidence": str(item.get("evidence") or "")[:240],
                "course_name": course_name,
                "method": "llm",
            }
        )
    return output, None


def _build_material_deadline_scan(
    *,
    settings: Settings,
    db: Database,
    artifact_external_id: str,
    title: str,
    course_name: str,
    canonical_course_id: str | None,
    extracted_text: str,
    local_path: str | None,
    reference_local: datetime,
    artifact_provenance_source: str,
    artifact_evidence_links: list[str] | None = None,
) -> dict[str, Any]:
    heuristic_tasks = _heuristic_material_deadline_tasks(
        title=title,
        course_name=course_name,
        extracted_text=extracted_text,
        timezone_name=settings.timezone,
        reference_local=reference_local,
    )
    llm_tasks: list[dict[str, Any]] = []
    llm_error: str | None = None
    llm_attempted = False
    should_attempt_llm = (
        (not heuristic_tasks and _contains_material_task_hints(title, extracted_text[:2000]))
        or (local_path and not str(extracted_text or "").strip())
    )
    if should_attempt_llm:
        llm_attempted = True
        llm_tasks, llm_error = _llm_material_deadline_tasks(
            settings=settings,
            db=db,
            title=title,
            course_name=course_name,
            extracted_text=extracted_text,
            local_path=local_path,
            reference_local=reference_local,
        )
    merged = list(llm_tasks or heuristic_tasks)
    if llm_tasks:
        existing_keys = {
            (_normalize_task_title_key(item.get("title", "")), item.get("due_at"))
            for item in merged
        }
        for item in heuristic_tasks:
            key = (_normalize_task_title_key(item.get("title", "")), item.get("due_at"))
            if key not in existing_keys:
                merged.append(item)
    upserted: list[dict[str, Any]] = []
    for item in merged:
        task_title = str(item.get("title") or "").strip() or Path(title).stem or title
        due_at = str(item.get("due_at") or "").strip() or None
        title_key = _normalize_task_title_key(task_title) or sha1(task_title.encode("utf-8")).hexdigest()[:12]
        external_id = f"uclass:material-task:{sha1(f'{artifact_external_id}|{title_key}'.encode('utf-8')).hexdigest()[:24]}"
        if _is_equivalent_existing_task(
            db,
            title=task_title,
            due_at=due_at,
            course_name=course_name,
            canonical_course_id=canonical_course_id,
            external_id=external_id,
        ):
            continue
        metadata = {
            "course_name": course_name,
            "canonical_course_id": canonical_course_id,
            "artifact_external_id": artifact_external_id,
            "artifact_title": title,
            "detected_via": "material_deadline",
            "detected_method": str(item.get("method") or "unknown"),
            "evidence": str(item.get("evidence") or "")[:240],
        }
        detected_method = str(item.get("method") or "unknown").strip().lower()
        provenance_source = "llm_inferred" if detected_method == "llm" else artifact_provenance_source
        provenance_confidence = "medium" if detected_method == "llm" else "low"
        metadata = attach_provenance(
            metadata,
            source=provenance_source,
            confidence=provenance_confidence,
            last_verified_at=now_utc_iso(),
            evidence_links=artifact_evidence_links,
            raw_source_ids=[artifact_external_id],
            derivation=f"{detected_method or 'unknown'}_deadline_scan",
        )
        db.upsert_task(
            external_id=external_id,
            source="uclass",
            due_at=due_at,
            title=task_title,
            status="open",
            metadata_json=metadata,
        )
        upserted.append(
            {
                "external_id": external_id,
                "title": task_title,
                "due_at": due_at,
                "method": metadata["detected_method"],
            }
        )
    return {
        "version": MATERIAL_DEADLINE_SCAN_VERSION,
        "ok": not (llm_attempted and llm_error and not heuristic_tasks and not upserted),
        "count": len(upserted),
        "items": upserted,
        "candidate_count": len(merged),
        "method": "llm" if llm_tasks else ("heuristic" if heuristic_tasks else "none"),
        "error": llm_error,
    }


def _buildings_for_event(
    db: Database,
    location: str | None,
    *,
    school_slug: str = "ku_online_class",
) -> tuple[str | None, str | None, str | None]:
    parsed = _parse_uos_location(location)
    building_no = parsed.get("building_no")
    room = parsed.get("room")
    if not building_no:
        return None, None, None
    name = db.get_building_name(str(building_no), school_slug=school_slug)
    return str(building_no), (str(name) if name else None), (str(room) if room else None)


def _identity_warning_gate(
    settings: Settings,
    db: Database,
    *,
    step: str,
    destination: str,
) -> dict[str, Any] | None:
    if not bool(getattr(settings, "include_identity", False)):
        return None
    active_ack = db.get_active_identity_ack()
    if active_ack is not None:
        return None
    token_suggestion = sha1(
        f"{step}:{destination}:{now_utc_iso()}".encode("utf-8")
    ).hexdigest()[:12]
    warning_gate = {
        "gate": "identity_ack",
        "step": step,
        "destination": destination,
        "include_identity": True,
        "ack_required": True,
        "ack_present": False,
        "message": "include_identity=true requires explicit ACK before external send",
        "token_suggestion": token_suggestion,
        "ack_command": f"kus ack identity --token {token_suggestion} --expires-hours 24",
    }
    return {
        "ok": False,
        "blocked": True,
        "error": "identity_ack_required",
        "warning_gate": warning_gate,
    }


def _heuristic_key_terms(text: str, limit: int = 3) -> list[str]:
    counts: Counter[str] = Counter()
    for token in WORD_RE.findall(text):
        lowered = token.lower()
        if lowered in STOPWORDS:
            continue
        counts[lowered] += 1
    terms = [term for term, _ in counts.most_common(limit)]
    while len(terms) < limit:
        terms.append("review")
    return terms[:limit]


def _build_heuristic_brief(title: str, extracted_text: str) -> dict[str, Any]:
    lines = [line.strip(" -\t") for line in extracted_text.splitlines() if line.strip()]
    headings: list[str] = []
    for line in lines:
        if len(line) > 90:
            continue
        if line.isupper() or line.endswith(":"):
            headings.append(line)
        elif len(line.split()) <= 8 and line[:1].isupper():
            headings.append(line)
        if len(headings) >= 3:
            break

    bullets: list[str] = []
    for heading in headings:
        bullets.append(f"{heading}")
        if len(bullets) >= 5:
            break
    for line in lines:
        if line in bullets:
            continue
        bullets.append(line[:140])
        if len(bullets) >= 5:
            break
    if not bullets:
        bullets = [f"No extractable summary lines for {title}."]
    while len(bullets) < 5:
        bullets.append("No additional material details detected.")

    key_terms = _heuristic_key_terms(extracted_text, limit=3)
    question = f"What is the core argument in {title} and how can I explain it from memory?"
    return {
        "mode": "heuristic",
        "bullets": bullets[:5],
        "key_terms": key_terms[:3],
        "question": question,
    }


def _material_llm_attachment_paths(
    settings: Settings,
    local_path: str | None,
) -> list[str] | None:
    # The local-only LLM path does not forward files directly.
    return None


def _build_material_brief(
    settings: Settings,
    db: Database,
    title: str,
    extracted_text: str,
    local_path: str | None = None,
    artifact_provenance_source: str | None = None,
) -> dict[str, Any]:
    heuristic = _build_heuristic_brief(title=title, extracted_text=extracted_text)
    source_text_hash = sha1(extracted_text.encode("utf-8")).hexdigest()
    heuristic["version"] = MATERIAL_BRIEF_VERSION
    heuristic["source_text_hash"] = source_text_hash
    heuristic["provenance"] = {
        "source": str(artifact_provenance_source or "uclass_html").strip().lower(),
        "confidence": "low",
        "last_verified_at": now_utc_iso(),
        "derivation": "heuristic_summary",
    }
    if not settings.llm_enabled:
        return heuristic
    gate = _identity_warning_gate(
        settings=settings,
        db=db,
        step="material_brief",
        destination="llm",
    )
    if gate is not None:
        return {
            "mode": "identity_gate_blocked",
            "version": MATERIAL_BRIEF_VERSION,
            "source_text_hash": source_text_hash,
            "bullets": heuristic["bullets"][:5],
            "key_terms": heuristic["key_terms"][:3],
            "question": heuristic["question"],
            "provenance": heuristic["provenance"],
            "warning_gate": gate["warning_gate"],
        }
    try:
        client = _llm_client_with_timeout(
            settings,
            timeout_sec=_llm_timeout_sec(
                settings,
                cap=MATERIAL_LLM_TIMEOUT_SEC_CAP,
            ),
        )
        attachment_paths = _material_llm_attachment_paths(settings, local_path)
        summary = client.summarize(
            {
                "mode": "material_brief",
                "title": title,
                "file_attached": bool(attachment_paths),
                "text_excerpt": extracted_text[
                    : _material_llm_text_excerpt_limit(
                        settings,
                        cap=MATERIAL_BRIEF_LLM_TEXT_EXCERPT_CHAR_CAP,
                    )
                ],
                "requirements": (
                    "Return concise learning takeaways for today's class and mention what to review next."
                ),
            },
            system_prompt=MATERIAL_BRIEF_SYSTEM_PROMPT,
            attachment_paths=attachment_paths,
        )
        bullets = [item.strip() for item in summary.bullets if item.strip()]
        for extra in heuristic["bullets"]:
            if len(bullets) >= 5:
                break
            if extra not in bullets:
                bullets.append(extra)
        while len(bullets) < 5:
            bullets.append("No additional material details detected.")
        return {
            "mode": "llm",
            "version": MATERIAL_BRIEF_VERSION,
            "source_text_hash": source_text_hash,
            "bullets": bullets[:5],
            "key_terms": heuristic["key_terms"][:3],
            "question": summary.action_item.strip() or heuristic["question"],
            "provenance": {
                "source": "llm_inferred",
                "confidence": "medium",
                "last_verified_at": now_utc_iso(),
                "derivation": "llm_summary",
            },
        }
    except Exception as exc:
        logger.warning(
            "material brief llm fallback",
            extra={"title": title, "error": str(exc)},
        )
        heuristic["mode"] = "heuristic_fallback"
        return heuristic


def _extract_text_from_pptx(path: Path, max_chars: int) -> tuple[str | None, str | None]:
    try:
        from pptx import Presentation
    except Exception:
        return None, "python-pptx not installed"
    try:
        presentation = Presentation(str(path))
        chunks: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", None)
                if not isinstance(text, str):
                    continue
                cleaned = text.strip()
                if cleaned:
                    chunks.append(cleaned)
        combined = "\n".join(chunks).strip()
        if not combined:
            return None, "no text extracted from pptx"
        return combined[:max_chars], None
    except Exception as exc:
        return None, str(exc)


def _extract_text_from_pdf(path: Path, max_chars: int) -> tuple[str | None, str | None]:
    try:
        from pypdf import PdfReader
    except Exception:
        PdfReader = None
    if PdfReader is not None:
        try:
            reader = PdfReader(str(path))
            pages: list[str] = []
            for page in reader.pages:
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(text.strip())
            combined = "\n".join(pages).strip()
            if combined:
                return combined[:max_chars], None
        except Exception as exc:
            logger.warning("pypdf extraction failed", extra={"path": str(path), "error": str(exc)})
    try:
        raw = path.read_bytes()
    except Exception as exc:
        return None, str(exc)
    snippets: list[str] = []
    for chunk in re.findall(rb"\((.{1,200}?)\)", raw):
        try:
            decoded = chunk.decode("utf-8", errors="ignore").strip()
        except Exception:
            continue
        if len(decoded) < 4:
            continue
        snippets.append(decoded)
        if sum(len(item) for item in snippets) > max_chars:
            break
    if not snippets:
        return None, "no text extracted from pdf"
    return "\n".join(snippets)[:max_chars], None


def _normalize_extracted_text(raw: str, *, max_chars: int) -> str:
    text = html.unescape(str(raw or ""))
    text = text.replace("\r\n", "\n").replace("\r", "\n").replace("\xa0", " ")
    text = MATERIAL_CONTROL_CHAR_RE.sub(" ", text)
    lines: list[str] = []
    previous = ""
    for line in text.split("\n"):
        cleaned = re.sub(r"[ \t\u200b]+", " ", line).strip()
        if not cleaned:
            if lines and lines[-1] != "":
                lines.append("")
            previous = ""
            continue
        if cleaned == previous:
            continue
        lines.append(cleaned)
        previous = cleaned
    normalized = "\n".join(lines).strip()
    normalized = re.sub(r"\n{3,}", "\n\n", normalized)
    if max_chars > 0 and len(normalized) > max_chars:
        normalized = normalized[:max_chars].rstrip()
    return normalized


def _looks_meaningful_extracted_text(text: str) -> bool:
    meaningful = len(re.findall(r"[가-힣A-Za-z0-9]", str(text or "")))
    if meaningful < 6:
        return False
    noise = len(re.findall(r"[^가-힣A-Za-z0-9\s\-\.,:;!?()<>/]", str(text or "")))
    return noise <= meaningful * 2


def _html_fragment_to_text(fragment: str, *, max_chars: int) -> str:
    text = MATERIAL_HTML_DROP_BLOCK_RE.sub(" ", fragment)
    text = MATERIAL_BREAK_TAG_RE.sub("\n", text)
    text = MATERIAL_BLOCK_TAG_RE.sub("\n", text)
    text = MATERIAL_HTML_TAG_RE.sub(" ", text)
    return _normalize_extracted_text(text, max_chars=max_chars)


def _extract_text_from_html(path: Path, max_chars: int) -> tuple[str | None, str | None]:
    try:
        raw = path.read_text(encoding="utf-8", errors="ignore")
    except Exception as exc:
        return None, str(exc)
    if not raw.strip():
        return None, "empty html file"
    title_match = MATERIAL_HTML_TITLE_RE.search(raw)
    title = (
        _html_fragment_to_text(title_match.group(1), max_chars=max_chars)
        if title_match
        else ""
    )
    best_body = ""
    for pattern in MATERIAL_HTML_SECTION_PATTERNS:
        for match in pattern.findall(raw):
            text = _html_fragment_to_text(str(match), max_chars=max_chars)
            if len(text) > len(best_body):
                best_body = text
    if not best_body:
        best_body = _html_fragment_to_text(raw, max_chars=max_chars)
    combined = _normalize_extracted_text(
        "\n\n".join(part for part in [title, best_body] if part),
        max_chars=max_chars,
    )
    if combined:
        return combined, None
    return None, "no text extracted from html"


def _extract_text_from_hwp_prvtext_bytes(data: bytes, *, max_chars: int) -> str:
    if not data:
        return ""
    try:
        decoded = data.decode("utf-16le", errors="ignore")
    except Exception:
        return ""
    return _normalize_extracted_text(decoded, max_chars=max_chars)


def _extract_text_from_hwp_body_bytes(
    data: bytes,
    *,
    compressed: bool,
    max_chars: int,
) -> str:
    try:
        payload = zlib.decompress(data, -15) if compressed else data
    except Exception:
        return ""
    chunks: list[str] = []
    position = 0
    budget = max(max_chars, 1) * 2
    while position + 4 <= len(payload):
        header = struct.unpack("<I", payload[position : position + 4])[0]
        tag_id = header & 0x3FF
        size = (header >> 20) & 0xFFF
        position += 4
        if size == 0xFFF:
            if position + 4 > len(payload):
                break
            size = struct.unpack("<I", payload[position : position + 4])[0]
            position += 4
        if size < 0 or position + size > len(payload):
            break
        record = payload[position : position + size]
        position += size
        if tag_id != HWP_PARA_TEXT_TAG or not record:
            continue
        text = _normalize_extracted_text(
            record.decode("utf-16le", errors="ignore"),
            max_chars=max_chars,
        )
        if not _looks_meaningful_extracted_text(text):
            continue
        chunks.append(text)
        if sum(len(item) for item in chunks) >= budget:
            break
    return _normalize_extracted_text("\n\n".join(chunks), max_chars=max_chars)


def _extract_text_from_hwp(path: Path, max_chars: int) -> tuple[str | None, str | None]:
    try:
        import olefile
    except Exception:
        return None, "olefile not installed"
    try:
        with olefile.OleFileIO(str(path)) as ole:
            if ole.exists("PrvText"):
                preview = _extract_text_from_hwp_prvtext_bytes(
                    ole.openstream("PrvText").read(),
                    max_chars=max_chars,
                )
                if preview:
                    return preview, None
            compressed = True
            if ole.exists("FileHeader"):
                header = ole.openstream("FileHeader").read()
                if len(header) >= 40:
                    flags = struct.unpack("<I", header[36:40])[0]
                    compressed = bool(flags & 1)
            section_names = [
                "/".join(parts)
                for parts in ole.listdir()
                if parts and parts[0] == "BodyText"
            ]
            section_names.sort()
            texts: list[str] = []
            for name in section_names:
                text = _extract_text_from_hwp_body_bytes(
                    ole.openstream(name).read(),
                    compressed=compressed,
                    max_chars=max_chars,
                )
                if text:
                    texts.append(text)
                if sum(len(item) for item in texts) >= max(max_chars, 1) * 2:
                    break
            merged = _normalize_extracted_text("\n\n".join(texts), max_chars=max_chars)
            if merged:
                return merged, None
    except Exception as exc:
        return None, str(exc)
    return None, "no text extracted from hwp"


def extract_material_text(path: Path, max_chars: int) -> tuple[str | None, str | None, str]:
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        text, error = _extract_text_from_pptx(path, max_chars=max_chars)
        return text, error, "pptx"
    if suffix == ".pdf":
        text, error = _extract_text_from_pdf(path, max_chars=max_chars)
        return text, error, "pdf"
    if suffix in {".php", ".html", ".htm"}:
        text, error = _extract_text_from_html(path, max_chars=max_chars)
        return text, error, "html"
    if suffix == ".hwp":
        text, error = _extract_text_from_hwp(path, max_chars=max_chars)
        return text, error, "hwp"
    return None, "unsupported file extension", suffix.lstrip(".") or "unknown"


def _resolve_uclass_token(
    settings: Settings,
    *,
    prefer_static: bool = True,
    username: str | None = None,
    password: str | None = None,
    ws_base_url: str | None = None,
) -> str:
    static_token = str(getattr(settings, "uclass_wstoken", "") or "").strip()
    if prefer_static and static_token:
        return static_token
    resolved_username = str(username or "").strip()
    resolved_password = str(password if password is not None else "")
    if not resolved_username or not resolved_password:
        raise ValueError(
            "UClass token unavailable: set UCLASS_WSTOKEN or provide explicit username/password"
        )
    ws_base = str(ws_base_url or getattr(settings, "uclass_ws_base", "") or "").strip()
    if not ws_base:
        raise ValueError("UCLASS_WS_BASE is required")
    service = str(getattr(settings, "uclass_token_service", "moodle_mobile_app") or "moodle_mobile_app").strip()
    token_endpoint = str(getattr(settings, "uclass_token_endpoint", "") or "").strip() or None
    request_method = str(getattr(settings, "uclass_request_method", "GET") or "GET").strip().upper()
    html_session = getattr(settings, "_uclass_html_session", None)
    token_error: Exception | None = None
    try:
        return request_moodle_ws_token(
            ws_base_url=ws_base,
            username=resolved_username,
            password=resolved_password,
            service=service,
            token_endpoint=token_endpoint,
        )
    except Exception as exc:
        token_error = exc
    try:
        return request_moodle_mobile_launch_token(
            ws_base_url=ws_base,
            username=resolved_username,
            password=resolved_password,
            service=service,
            timeout_sec=30,
            session=html_session,
            request_method=request_method,
        )
    except Exception as mobile_exc:
        if token_error is not None:
            raise RuntimeError(
                f"{token_error}; mobile launch token fallback failed: {mobile_exc}"
            ) from mobile_exc
        raise


def build_uclass_probe_report(
    settings: Settings,
    output_json_path: Path | None = None,
) -> dict[str, Any]:
    if not settings.uclass_ws_base:
        raise ValueError("UCLASS_WS_BASE is required")
    token = _resolve_uclass_token(settings)
    client = MoodleWSClient(
        base_url=settings.uclass_ws_base,
        token=token,
        request_method=settings.uclass_request_method,
    )
    rows: list[dict[str, Any]] = []

    site_info: dict[str, Any] = {}
    site_error: str | None = None
    try:
        payload = client.get_site_info(settings.uclass_func_site_info)
        if isinstance(payload, dict):
            site_info = payload
        rows.append(
            {
                "key": "site_info",
                "wsfunction": settings.uclass_func_site_info,
                "enabled": True,
                "status": "OK",
                "error": None,
                "detail": _payload_shape(payload),
                "shape_fingerprint": _payload_shape_fingerprint(payload),
            }
        )
    except Exception as exc:
        site_error = str(exc)
        rows.append(
            {
                "key": "site_info",
                "wsfunction": settings.uclass_func_site_info,
                "enabled": True,
                "status": "FAIL",
                "error": site_error,
                "detail": None,
                "shape_fingerprint": None,
            }
        )

    helper_courses: list[int] = []
    helper_forum_id: int | None = None

    def run_probe(
        key: str,
        wsfunction: str,
        enabled: bool,
        loader: Callable[[], Any] | None,
        skip_reason: str | None = None,
    ) -> Any:
        if not enabled:
            rows.append(
                {
                    "key": key,
                    "wsfunction": wsfunction,
                    "enabled": False,
                    "status": "SKIP",
                    "error": "disabled",
                    "detail": None,
                    "shape_fingerprint": None,
                }
            )
            return None
        if skip_reason or loader is None:
            rows.append(
                {
                    "key": key,
                    "wsfunction": wsfunction,
                    "enabled": True,
                    "status": "SKIP",
                    "error": skip_reason or "missing loader",
                    "detail": None,
                    "shape_fingerprint": None,
                }
            )
            return None
        try:
            payload = loader()
            rows.append(
                {
                    "key": key,
                    "wsfunction": wsfunction,
                    "enabled": True,
                    "status": "OK",
                    "error": None,
                    "detail": _payload_shape(payload),
                    "shape_fingerprint": _payload_shape_fingerprint(payload),
                }
            )
            return payload
        except Exception as exc:
            rows.append(
                {
                    "key": key,
                    "wsfunction": wsfunction,
                    "enabled": True,
                    "status": "FAIL",
                    "error": str(exc),
                    "detail": None,
                    "shape_fingerprint": None,
                }
            )
            return None

    courses_payload = run_probe(
        key="courses",
        wsfunction=settings.uclass_func_courses,
        enabled=settings.uclass_enable_courses,
        loader=lambda: client.get_users_courses(settings.uclass_func_courses),
    )
    helper_courses = sorted(extract_course_index(courses_payload).keys())
    if not helper_courses and (
        settings.uclass_enable_contents
        or settings.uclass_enable_assignments
        or settings.uclass_enable_forums
    ):
        try:
            fallback_courses = client.get_users_courses(settings.uclass_func_courses)
            helper_courses = sorted(extract_course_index(fallback_courses).keys())
        except Exception:
            helper_courses = []

    run_probe(
        key="popup_notifications",
        wsfunction=settings.uclass_func_popup_notifications,
        enabled=settings.uclass_enable_popup_notifications,
        loader=lambda: client.get_popup_notifications(
            settings.uclass_func_popup_notifications,
            limit=1,
        ),
    )
    run_probe(
        key="action_events",
        wsfunction=settings.uclass_func_action_events,
        enabled=settings.uclass_enable_action_events,
        loader=lambda: client.get_action_events(
            settings.uclass_func_action_events,
            limitnum=1,
        ),
    )
    run_probe(
        key="course_contents",
        wsfunction=settings.uclass_func_course_contents,
        enabled=settings.uclass_enable_contents,
        loader=(
            (lambda cid=helper_courses[0]: client.get_course_contents(settings.uclass_func_course_contents, course_id=cid))
            if helper_courses
            else None
        ),
        skip_reason=None if helper_courses else "no course id available",
    )
    run_probe(
        key="assignments",
        wsfunction=settings.uclass_func_assignments,
        enabled=settings.uclass_enable_assignments,
        loader=lambda: client.get_assignments(
            settings.uclass_func_assignments,
            course_ids=helper_courses[:1],
        ),
    )
    forums_payload = run_probe(
        key="forums",
        wsfunction=settings.uclass_func_forums,
        enabled=settings.uclass_enable_forums,
        loader=lambda: client.get_forums(
            settings.uclass_func_forums,
            course_ids=helper_courses[:1],
        ),
    )
    helper_forum_id = _extract_first_forum_id(forums_payload)
    run_probe(
        key="forum_discussions",
        wsfunction=settings.uclass_func_forum_discussions,
        enabled=settings.uclass_enable_forums,
        loader=(
            (lambda fid=helper_forum_id: client.get_forum_discussions(
                settings.uclass_func_forum_discussions,
                forum_id=fid,
                page=0,
                per_page=1,
            ))
            if helper_forum_id is not None
            else None
        ),
        skip_reason=None if helper_forum_id is not None else "no forum id available",
    )
    report = {
        "generated_at": now_utc_iso(),
        "site_info": {
            "userid": site_info.get("userid"),
            "username": site_info.get("username"),
            "fullname": site_info.get("fullname"),
            "sitename": site_info.get("sitename"),
            "siteurl": site_info.get("siteurl"),
            "release": site_info.get("release"),
            "version": site_info.get("version"),
        },
        "site_info_error": site_error,
        "rows": rows,
    }
    if output_json_path:
        output_json_path.parent.mkdir(parents=True, exist_ok=True)
        output_json_path.write_text(
            json.dumps(report, ensure_ascii=True, indent=2, sort_keys=True),
            encoding="utf-8",
        )
    return report


def _is_ku_portal_browser_session(item: dict[str, Any] | None) -> bool:
    payload = item or {}
    school_slug = str(payload.get("school_slug") or "").strip().lower()
    provider = str(payload.get("provider") or "").strip().lower()
    return school_slug == KU_PORTAL_SCHOOL_SLUG or provider == KU_PORTAL_PROVIDER


def _is_ku_portal_browser_target(item: dict[str, Any] | None) -> bool:
    payload = item or {}
    if not _is_ku_portal_browser_session(payload):
        return False
    if str(payload.get("profile_dir") or "").strip():
        return True
    secret_kind = str(payload.get("secret_kind") or "").strip()
    secret_ref = str(payload.get("secret_ref") or "").strip()
    return bool(secret_kind and secret_ref)


def _has_ku_portal_browser_session(
    db: Database,
    *,
    user_id: int | None = None,
    chat_id: str | None = None,
) -> bool:
    owner_id = _safe_int(user_id)
    chat = str(chat_id or "").strip()
    if owner_id is None and not chat:
        return False
    sessions = db.list_lms_browser_sessions(
        user_id=owner_id,
        chat_id=chat or None,
        status="active",
        limit=20,
    )
    return any(_is_ku_portal_browser_session(item) for item in sessions)


def _has_ready_ku_portal_timetable_session(
    db: Database,
    *,
    user_id: int | None = None,
    chat_id: str | None = None,
    settings: Settings | None = None,
) -> bool:
    owner_id = _safe_int(user_id)
    chat = str(chat_id or "").strip()
    if owner_id is None and not chat:
        return False
    sessions = db.list_lms_browser_sessions(
        user_id=owner_id,
        chat_id=chat or None,
        status="active",
        limit=20,
    )
    for item in sessions:
        if not _is_ku_portal_browser_session(item):
            continue
        metadata = (
            dict(item.get("metadata_json"))
            if isinstance(item.get("metadata_json"), dict)
            else {}
        )
        sync_meta = (
            dict(metadata.get("portal_timetable_sync"))
            if isinstance(metadata.get("portal_timetable_sync"), dict)
            else {}
        )
        status = str(sync_meta.get("status") or "").strip().lower()
        auth_required = bool(sync_meta.get("auth_required"))
        if status == "success" and not auth_required:
            return True
    if settings is not None:
        state = db.get_sync_state("sync_ku_portal_timetable", user_id=owner_id or 0)
        if state is not None:
            cursor = _json_load(state.last_cursor_json)
            dashboard = (
                dict(cursor.get("_sync_dashboard"))
                if isinstance(cursor.get("_sync_dashboard"), dict)
                else {}
            )
            status = str(
                cursor.get("status")
                or dashboard.get("status")
                or ""
            ).strip().lower()
            payload_source = str(cursor.get("payload_source") or "").strip()
            auth_required = bool(cursor.get("auth_required"))
            if (
                status == "success"
                and payload_source == KU_OPENAPI_TIMETABLE_SOURCE
                and not auth_required
            ):
                return True
    return False


def _resolve_ku_portal_timetable_targets(
    settings: Settings,
    db: Database,
) -> list[dict[str, Any]]:
    api_url = str(getattr(settings, "ku_openapi_timetable_url", "") or "").strip()
    api_key = str(getattr(settings, "ku_openapi_timetable_api_key", "") or "").strip() or None
    if ku_openapi_timetable_configured(api_url, api_key) and ku_openapi_uses_official_catalog_mode(api_url):
        targets: list[dict[str, Any]] = []
        seen: set[tuple[int, str, str]] = set()
        for item in _resolve_uclass_sync_targets(settings, db):
            school_slug = str(item.get("school_slug") or "").strip().lower()
            if school_slug and school_slug != UOS_ONLINE_CLASS_SCHOOL_SLUG:
                continue
            owner_id = _safe_int(item.get("user_id")) or 0
            chat_id = str(item.get("chat_id") or "").strip() or None
            ws_base_url = str(item.get("ws_base_url") or "").strip()
            key = (owner_id, chat_id or "", ws_base_url)
            if key in seen:
                continue
            seen.add(key)
            targets.append(
                {
                    "user_id": owner_id,
                    "chat_id": chat_id,
                    "school_slug": KU_PORTAL_SCHOOL_SLUG,
                    "provider": KU_PORTAL_PROVIDER,
                    "display_name": str(item.get("display_name") or "고려대학교 공식 시간표 API").strip()
                    or "고려대학교 공식 시간표 API",
                    "ws_base_url": ws_base_url,
                    "token": str(item.get("token") or "").strip(),
                    "token_error": str(item.get("token_error") or "").strip(),
                    "connection_id": int(item.get("connection_id") or 0),
                    "source_connection": "uclass",
                }
            )
        return targets
    targets: list[dict[str, Any]] = []
    for item in db.list_lms_browser_sessions(status="active", limit=500):
        if not _is_ku_portal_browser_session(item):
            continue
        owner_id = _safe_int(item.get("user_id")) or 0
        chat_id = str(item.get("chat_id") or "").strip()
        if owner_id <= 0 and chat_id:
            user = db.ensure_user_for_chat(
                chat_id=chat_id,
                timezone_name=settings.timezone,
                metadata_json={"source": "ku_portal_browser_session"},
            )
            owner_id = _safe_int(user.get("id")) or 0
        metadata = (
            dict(item.get("metadata_json"))
            if isinstance(item.get("metadata_json"), dict)
            else {}
        )
        browser_result = (
            dict(metadata.get("browser_result"))
            if isinstance(metadata.get("browser_result"), dict)
            else {}
        )
        targets.append(
            {
                **item,
                "user_id": owner_id,
                "chat_id": chat_id or None,
                "session_metadata": metadata,
                "browser_result": browser_result,
                "current_url": str(
                    browser_result.get("current_url")
                    or item.get("login_url")
                    or KU_WISE_INDEX_URL
                ).strip()
                or KU_WISE_INDEX_URL,
            }
        )
    return targets


def _user_has_timetable_events(
    db: Database,
    *,
    user_id: int | None = None,
) -> bool:
    owner_id = _safe_int(user_id)
    if owner_id is None or owner_id <= 0:
        return False
    for event in db.list_events(limit=3000, user_id=owner_id):
        if _is_timetable_event(event):
            return True
    return False


def prime_ku_portal_timetable_for_user(
    settings: Settings,
    db: Database,
    *,
    chat_id: str | None = None,
    user_id: int | None = None,
    force: bool = False,
) -> dict[str, Any]:
    owner_id = _safe_int(user_id)
    chat = str(chat_id or "").strip() or None
    if owner_id is None and chat:
        owner_id = int(
            _resolve_user_scope(
                settings,
                db,
                chat_id=chat,
                create_if_missing=False,
                metadata_source="portal_timetable_prime",
            )["user_id"]
            or 0
        ) or None
    if owner_id is None or owner_id <= 0:
        return {"skipped": True, "reason": "user_id missing"}
    if not force and _user_has_timetable_events(db, user_id=owner_id):
        return {"skipped": True, "reason": "timetable already present", "user_id": owner_id}
    targets = [
        item
        for item in _resolve_ku_portal_timetable_targets(settings, db)
        if int(item.get("user_id") or 0) == owner_id
    ]
    if chat:
        targets = [item for item in targets if str(item.get("chat_id") or "").strip() == chat]
    if not targets:
        reason = "no active KU timetable target"
        _record_sync_dashboard_state(
            db,
            "sync_ku_portal_timetable",
            status="skipped",
            last_error=reason,
            cursor_payload={"skipped": True, "reason": reason},
            user_id=owner_id,
        )
        return {"skipped": True, "reason": reason, "user_id": owner_id}
    target = targets[0]
    try:
        result = _sync_ku_portal_timetable_target(
            settings=settings,
            db=db,
            target=target,
        )
    except Exception as exc:
        logger.warning(
            "portal timetable prime failed",
            extra={"chat_id": chat, "user_id": owner_id, "error": str(exc)},
        )
        _record_sync_dashboard_state(
            db,
            "sync_ku_portal_timetable",
            status="error",
            action_required=1,
            last_error=str(exc),
            cursor_payload={"skipped": True, "reason": str(exc)},
            user_id=owner_id,
        )
        return {
            "ok": False,
            "error": str(exc),
            "chat_id": chat,
            "user_id": owner_id,
        }
    result_status = str(result.get("_sync_status") or "success").strip().lower() or "success"
    result_error = str(result.get("_sync_last_error") or result.get("reason") or "").strip() or None
    _record_sync_dashboard_state(
        db,
        "sync_ku_portal_timetable",
        status=result_status,
        new_items=int(result.get("upserted_events") or 0),
        action_required=int(result.get("_sync_action_required") or 0),
        last_error=result_error,
        cursor_payload=result.get("_sync_cursor") if isinstance(result.get("_sync_cursor"), dict) else None,
        user_id=owner_id,
    )
    return {
        "ok": result_status == "success",
        "skipped": result_status == "skipped",
        "status": result_status,
        "error": result_error if result_status == "error" else None,
        "reason": str(result.get("reason") or "").strip() or None,
        "chat_id": chat,
        "user_id": owner_id,
        "upserted_events": int(result.get("upserted_events") or 0),
    }


def _prime_ku_portal_timetable_for_day_if_missing(
    settings: Settings,
    db: Database,
    *,
    target_day_local: datetime,
    chat_id: str | None = None,
    user_id: int | None = None,
) -> dict[str, Any]:
    owner_id = _safe_int(user_id)
    chat = str(chat_id or "").strip() or None
    if owner_id is None and chat:
        owner_id = int(
            _resolve_user_scope(
                settings,
                db,
                chat_id=chat,
                create_if_missing=False,
                metadata_source="portal_timetable_day_prime",
            )["user_id"]
            or 0
        ) or None
    if owner_id is None or owner_id <= 0:
        return {"skipped": True, "reason": "user_id missing"}
    if _collect_class_occurrences(
        settings=settings,
        db=db,
        target_day_local=target_day_local,
        max_items=1,
        user_id=owner_id,
    ):
        return {"skipped": True, "reason": "day already present", "user_id": owner_id}
    targets = [
        item
        for item in _resolve_ku_portal_timetable_targets(settings, db)
        if int(item.get("user_id") or 0) == owner_id
    ]
    if chat:
        targets = [item for item in targets if str(item.get("chat_id") or "").strip() == chat]
    if not targets:
        reason = "no active KU timetable target"
        _record_sync_dashboard_state(
            db,
            "sync_ku_portal_timetable",
            status="skipped",
            last_error=reason,
            cursor_payload={"skipped": True, "reason": reason},
            user_id=owner_id,
        )
        return {"skipped": True, "reason": reason, "user_id": owner_id}
    return prime_ku_portal_timetable_for_user(
        settings=settings,
        db=db,
        chat_id=chat,
        user_id=owner_id,
        force=True,
    )


def _normalize_timetable_source_attempts(value: Any) -> list[dict[str, Any]]:
    attempts: list[dict[str, Any]] = []
    if not isinstance(value, list):
        return attempts
    for item in value:
        if not isinstance(item, dict):
            continue
        attempts.append(
            {
                "source": str(item.get("source") or "").strip() or None,
                "status": str(item.get("status") or "").strip() or None,
                "reason": str(item.get("reason") or "").strip() or None,
                "source_url": str(item.get("source_url") or "").strip() or None,
            }
        )
    return attempts


def _timetable_source_attempt(
    source: str,
    *,
    status: str,
    reason: str | None = None,
    source_url: str | None = None,
) -> dict[str, Any]:
    return {
        "source": str(source or "").strip() or None,
        "status": str(status or "").strip() or None,
        "reason": str(reason or "").strip() or None,
        "source_url": str(source_url or "").strip() or None,
    }


def _kupid_sso_timetable_enabled(settings: Settings) -> bool:
    """Opt-in flag for the KUPID SSO timetable adapter.

    Defaults to False so existing OpenAPI / browser-scraping tests and
    deployments are unaffected until the operator explicitly opts in.
    Source priority: settings attr → env var ``KUPID_SSO_TIMETABLE_ENABLED``.
    """
    raw = getattr(settings, "kupid_sso_timetable_enabled", None)
    if isinstance(raw, bool):
        return raw
    candidate = str(raw or "").strip().lower()
    if not candidate:
        candidate = os.environ.get("KUPID_SSO_TIMETABLE_ENABLED", "").strip().lower()
    return candidate in {"1", "true", "yes", "on"}


def _try_fetch_kupid_sso_timetable(
    settings: Settings,
    *,
    target: dict[str, Any],
) -> tuple[dict[str, Any] | None, list[dict[str, Any]]]:
    """Try the vendored ku-portal-mcp KUPID SSO timetable path.

    Returns ``(fetched, attempts)``. ``fetched`` is None when the adapter
    is not configured (no credentials) or fails — caller should fall through
    to the next strategy in the chain.
    """
    if not _kupid_sso_timetable_enabled(settings):
        return None, []
    timezone_name = str(getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul")
    try:
        fetched = fetch_kupid_sso_timetable(
            settings=settings,
            target=target,
            timezone_name=timezone_name,
        )
    except RuntimeError as exc:
        attempt = _timetable_source_attempt(
            KUPID_SSO_TIMETABLE_SOURCE,
            status="skipped",
            reason=str(exc),
        )
        logger.debug("kupid sso timetable skipped: %s", exc)
        return None, [attempt]
    except Exception as exc:  # noqa: BLE001 — log + fall through to next source
        attempt = _timetable_source_attempt(
            KUPID_SSO_TIMETABLE_SOURCE,
            status="fallback",
            reason=str(exc),
        )
        logger.warning("kupid sso timetable fetch failed; falling back: %s", exc)
        return None, [attempt]

    attempt = _timetable_source_attempt(
        KUPID_SSO_TIMETABLE_SOURCE,
        status="selected",
        source_url=str(fetched.get("source_url") or "").strip() or None,
    )
    fetched = dict(fetched)
    fetched["payload_source"] = KUPID_SSO_TIMETABLE_SOURCE
    fetched["source_attempts"] = [attempt]
    fetched["fallback_used"] = False
    return fetched, [attempt]


def _try_fetch_ku_openapi_timetable(
    settings: Settings,
    *,
    target: dict[str, Any],
) -> tuple[dict[str, Any] | None, list[dict[str, Any]]]:
    api_url = str(getattr(settings, "ku_openapi_timetable_url", "") or "").strip()
    api_key = str(getattr(settings, "ku_openapi_timetable_api_key", "") or "").strip() or None
    if not ku_openapi_timetable_configured(api_url, api_key):
        return None, []
    if ku_openapi_uses_official_catalog_mode(api_url):
        return None, []
    timeout_sec = max(int(getattr(settings, "ku_openapi_timetable_timeout_sec", 15) or 15), 1)
    try:
        fetched = fetch_ku_openapi_timetable(
            api_url=api_url,
            api_key=api_key,
            academic_year=getattr(settings, "ku_openapi_year", None),
            term=getattr(settings, "ku_openapi_term", None),
            timezone_name=str(getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul"),
            timeout_sec=timeout_sec,
            target=target,
        )
    except KUOpenAPITimetableUnsupported as exc:
        attempt = _timetable_source_attempt(
            KU_OPENAPI_TIMETABLE_SOURCE,
            status="unsupported",
            reason=str(exc),
            source_url=api_url,
        )
        logger.info("uos official timetable API unsupported for target", extra=attempt)
        return None, [attempt]
    except (KUOpenAPITimetableMalformedPayload, requests.RequestException) as exc:
        attempt = _timetable_source_attempt(
            KU_OPENAPI_TIMETABLE_SOURCE,
            status="fallback",
            reason=str(exc),
            source_url=api_url,
        )
        logger.warning("uos official timetable API failed; falling back to portal", extra=attempt)
        return None, [attempt]

    resolved_source_url = str(fetched.get("source_url") or api_url).strip() or None
    result = dict(fetched)
    result["payload_source"] = KU_OPENAPI_TIMETABLE_SOURCE
    result["source_url"] = resolved_source_url
    result["allow_empty_success"] = True
    result["source_attempts"] = [
        _timetable_source_attempt(
            KU_OPENAPI_TIMETABLE_SOURCE,
            status="selected",
            source_url=resolved_source_url,
        )
    ]
    result["fallback_used"] = False
    return result, list(result.get("source_attempts") or [])


def _collect_text_fragments(value: Any, *, limit: int = 200) -> list[str]:
    output: list[str] = []

    def _walk(item: Any) -> None:
        if len(output) >= limit:
            return
        if isinstance(item, dict):
            for inner in item.values():
                _walk(inner)
            return
        if isinstance(item, (list, tuple, set)):
            for inner in item:
                _walk(inner)
            return
        text = str(item or "").strip()
        if text:
            output.append(text)

    _walk(value)
    return output


def _official_catalog_section_key(item: dict[str, Any]) -> str:
    metadata = dict(item.get("metadata") or {}) if isinstance(item.get("metadata"), dict) else {}
    subject_no = str(metadata.get("official_subject_no") or "").strip()
    dvcl_no = str(metadata.get("official_dvcl_no") or "").strip()
    if subject_no or dvcl_no:
        return f"{subject_no}:{dvcl_no}"
    return "|".join(
        [
            normalize_course_alias(str(item.get("title") or "").strip()),
            str(metadata.get("instructor") or "").strip(),
            str(item.get("rrule") or "").strip(),
        ]
    )


def _build_uos_official_course_profiles(
    db: Database,
    *,
    user_id: int,
) -> list[dict[str, Any]]:
    aliases_by_course: dict[str, list[str]] = {}
    for alias in db.list_course_aliases(limit=5000, user_id=user_id):
        course_id = str(alias.canonical_course_id or "").strip()
        alias_text = str(alias.alias or "").strip()
        if not course_id or not alias_text:
            continue
        aliases_by_course.setdefault(course_id, []).append(alias_text)

    profiles: list[dict[str, Any]] = []
    for course in db.list_courses(limit=2000, user_id=user_id):
        canonical_course_id = str(course.canonical_course_id or "").strip()
        if not canonical_course_id:
            continue
        metadata = _json_load(course.metadata_json)
        aliases = list(aliases_by_course.get(canonical_course_id) or [])
        for candidate in (
            str(course.display_name or "").strip(),
            str(metadata.get("fullname") or "").strip(),
            str(metadata.get("displayname") or "").strip(),
            str(metadata.get("shortname") or "").strip(),
            str(metadata.get("course_name") or "").strip(),
            str(metadata.get("idnumber") or "").strip(),
        ):
            if candidate and candidate not in aliases:
                aliases.append(candidate)
        normalized_aliases = {
            normalize_course_alias(alias)
            for alias in aliases
            if normalize_course_alias(alias)
        }
        if not normalized_aliases:
            continue
        subject_no = str(metadata.get("official_subject_no") or "").strip()
        dvcl_no = str(metadata.get("official_dvcl_no") or "").strip()
        metadata_text = " ".join(_collect_text_fragments(metadata))
        profiles.append(
            {
                "canonical_course_id": canonical_course_id,
                "display_name": str(course.display_name or "").strip() or None,
                "aliases": aliases,
                "normalized_aliases": normalized_aliases,
                "metadata_text": metadata_text,
                "official_subject_no": subject_no or None,
                "official_dvcl_no": dvcl_no or None,
                "official_section_key": (
                    f"{subject_no}:{dvcl_no}"
                    if subject_no and dvcl_no
                    else None
                ),
            }
        )
    return profiles


def _persist_official_course_binding(
    db: Database,
    *,
    canonical_course_id: str,
    event_metadata: dict[str, Any],
    user_id: int,
) -> None:
    course = db.get_course(canonical_course_id, user_id=user_id)
    if course is None:
        return
    existing_metadata = _json_load(course.metadata_json)
    merged_metadata = dict(existing_metadata)
    changed = False
    for key in (
        "official_subject_no",
        "official_dvcl_no",
        "official_course_section",
        "official_course_name",
        "official_course_code",
        "official_syllabus_url",
        "official_syllabus_id",
        "official_term_code",
    ):
        value = event_metadata.get(key)
        if value in (None, ""):
            continue
        if merged_metadata.get(key) == value:
            continue
        merged_metadata[key] = value
        changed = True
    if changed:
        db.upsert_course(
            canonical_course_id=canonical_course_id,
            source=course.source,
            external_course_id=course.external_course_id,
            display_name=course.display_name,
            metadata_json=merged_metadata,
            user_id=user_id,
        )
    alias_candidates = [
        ("official_course_name", merged_metadata.get("official_course_name")),
        ("official_course_code", merged_metadata.get("official_course_code")),
        ("official_subject_no", merged_metadata.get("official_subject_no")),
    ]
    subject_no = str(merged_metadata.get("official_subject_no") or "").strip()
    dvcl_no = str(merged_metadata.get("official_dvcl_no") or "").strip()
    if subject_no and dvcl_no:
        alias_candidates.append(
            ("official_subject_section", f"{subject_no}-{dvcl_no}")
        )
    for alias_type, alias in alias_candidates:
        alias_text = str(alias or "").strip()
        if not alias_text:
            continue
        db.upsert_course_alias(
            canonical_course_id=canonical_course_id,
            alias=alias_text,
            alias_type=alias_type,
            source=KU_OPENAPI_TIMETABLE_SOURCE,
            metadata_json={
                "official_subject_no": subject_no or None,
                "official_dvcl_no": dvcl_no or None,
            },
            user_id=user_id,
        )


def _normalize_uos_room_lookup_key(
    *,
    building_no: str | None,
    room: str | None,
) -> str:
    building_key = str(building_no or "").strip()
    room_key = re.sub(
        r"[^0-9a-z]+",
        "",
        str(_normalize_room_label(room, building_no=building_key) or "").strip().lower(),
    )
    if not building_key or not room_key:
        return ""
    return f"{building_key}:{room_key}"


def _uos_building_catalog_lookup_keys(item: dict[str, Any]) -> set[str]:
    keys: set[str] = set()
    building_no = str(item.get("building_code") or "").strip()
    for raw_room in (
        item.get("room_code"),
        item.get("room_name"),
    ):
        lookup_key = _normalize_uos_room_lookup_key(
            building_no=building_no,
            room=str(raw_room or "").strip() or None,
        )
        if lookup_key:
            keys.add(lookup_key)
    space_name = str(item.get("space_name") or "").strip()
    parsed_space = _parse_uos_location(space_name or None)
    if parsed_space.get("ok"):
        lookup_key = _normalize_uos_room_lookup_key(
            building_no=str(parsed_space.get("building_no") or "").strip() or None,
            room=str(parsed_space.get("room") or "").strip() or None,
        )
        if lookup_key:
            keys.add(lookup_key)
    return keys


def _load_ku_openapi_building_catalog(
    settings: Settings,
    db: Database,
) -> dict[str, Any]:
    api_key = str(getattr(settings, "ku_openapi_timetable_api_key", "") or "").strip()
    cache_key = f"{KU_OPENAPI_OFFICIAL_BUILDING_URL}|{api_key}"
    cache = getattr(settings, "_ku_openapi_building_catalog_cache", None)
    if not isinstance(cache, dict):
        cache = {}
        setattr(settings, "_ku_openapi_building_catalog_cache", cache)
    cached = cache.get(cache_key)
    if isinstance(cached, dict):
        return cached
    if not api_key:
        payload = {
            "ok": False,
            "items": [],
            "building_names": {},
            "rooms_by_key": {},
            "reason": "KU official building API key is not configured",
            "source_url": None,
            "payload_source": KU_OPENAPI_BUILDING_SOURCE,
        }
        cache[cache_key] = payload
        return payload
    timeout_sec = max(int(getattr(settings, "ku_openapi_timetable_timeout_sec", 15) or 15), 1)
    try:
        fetched = fetch_ku_openapi_building_catalog(
            api_key=api_key,
            api_url=KU_OPENAPI_OFFICIAL_BUILDING_URL,
            timeout_sec=timeout_sec,
        )
    except (
        KUOpenAPIBuildingCatalogError,
        KUOpenAPIBuildingCatalogMalformedPayload,
        requests.RequestException,
    ) as exc:
        payload = {
            "ok": False,
            "items": [],
            "building_names": {},
            "rooms_by_key": {},
            "reason": str(exc).strip() or "KU official building API request failed",
            "source_url": KU_OPENAPI_OFFICIAL_BUILDING_URL,
            "payload_source": KU_OPENAPI_BUILDING_SOURCE,
        }
        cache[cache_key] = payload
        logger.warning("uos official building catalog load failed", exc_info=True)
        return payload

    building_names: dict[str, str] = {}
    rooms_by_key: dict[str, list[dict[str, Any]]] = {}
    for item in list(fetched.get("items") or []):
        if not isinstance(item, dict):
            continue
        building_no = str(item.get("building_code") or "").strip()
        building_name = str(item.get("building_name") or "").strip()
        if building_no and building_name:
            building_names[building_no] = building_name
            db.upsert_building(
                building_no=building_no,
                building_name=building_name,
                metadata_json={
                    "source": KU_OPENAPI_BUILDING_SOURCE,
                    "source_url": fetched.get("source_url"),
                },
                school_slug=KU_PORTAL_SCHOOL_SLUG,
            )
        for lookup_key in _uos_building_catalog_lookup_keys(item):
            rooms_by_key.setdefault(lookup_key, []).append(item)
    payload = {
        **dict(fetched),
        "ok": True,
        "building_names": building_names,
        "rooms_by_key": rooms_by_key,
    }
    cache[cache_key] = payload
    return payload


def _enrich_events_with_uos_building_catalog(
    settings: Settings,
    db: Database,
    *,
    events: list[dict[str, Any]],
) -> tuple[list[dict[str, Any]], dict[str, Any] | None]:
    catalog = _load_ku_openapi_building_catalog(settings, db)
    if not catalog.get("ok"):
        return events, None
    building_names = dict(catalog.get("building_names") or {})
    rooms_by_key = dict(catalog.get("rooms_by_key") or {})
    enriched_events: list[dict[str, Any]] = []
    resolved_locations = 0
    for item in events:
        payload = dict(item)
        metadata = (
            dict(payload.get("metadata"))
            if isinstance(payload.get("metadata"), dict)
            else {}
        )
        location_meta = _parse_uos_location(str(payload.get("location") or "").strip() or None)
        building_no = (
            str(metadata.get("official_building_no") or "").strip()
            or str(location_meta.get("building_no") or "").strip()
        )
        if building_no and not str(metadata.get("official_building_no") or "").strip():
            metadata["official_building_no"] = building_no
        if building_no:
            building_name = building_names.get(building_no)
            if building_name and not str(metadata.get("official_building_name") or "").strip():
                metadata["official_building_name"] = building_name
        room_hint = (
            str(metadata.get("official_room") or "").strip()
            or str(location_meta.get("room") or "").strip()
            or None
        )
        lookup_key = _normalize_uos_room_lookup_key(
            building_no=building_no or None,
            room=room_hint,
        )
        matches = list(rooms_by_key.get(lookup_key) or []) if lookup_key else []
        if len(matches) == 1:
            match = matches[0]
            match_building_no = str(match.get("building_code") or "").strip()
            match_building_name = str(match.get("building_name") or "").strip()
            resolved_room = _normalize_room_label(
                str(match.get("room_name") or match.get("room_code") or "").strip() or None,
                building_no=match_building_no or building_no or None,
            )
            if match_building_no:
                metadata["official_building_no"] = match_building_no
            if match_building_name:
                metadata["official_building_name"] = match_building_name
            if resolved_room:
                metadata["official_room"] = resolved_room
            room_name = str(match.get("room_name") or "").strip()
            if room_name:
                metadata["official_room_name"] = room_name
            space_name = str(match.get("space_name") or "").strip()
            if space_name:
                metadata["official_space_name"] = space_name
            resolved_locations += 1
        payload["metadata"] = metadata
        enriched_events.append(payload)
    return enriched_events, {
        "payload_source": KU_OPENAPI_BUILDING_SOURCE,
        "source_url": catalog.get("source_url"),
        "building_count": len(building_names),
        "room_count": len(list(catalog.get("items") or [])),
        "resolved_locations": resolved_locations,
    }


def _refresh_uclass_course_aliases_for_timetable_target(
    settings: Settings,
    db: Database,
    *,
    target: dict[str, Any],
    owner_id: int,
) -> None:
    ws_base_url = str(target.get("ws_base_url") or "").strip()
    token = str(target.get("token") or "").strip()
    if not ws_base_url or not token:
        return
    client = MoodleWSClient(
        base_url=ws_base_url,
        token=token,
        request_method=str(getattr(settings, "uclass_request_method", "GET") or "GET").strip().upper(),
    )
    try:
        site_info = client.get_site_info(settings.uclass_func_site_info)
        userid = _int_or_none(site_info.get("userid")) if isinstance(site_info, dict) else None
        if userid is not None:
            client.site_userid = userid
        courses_payload = client.get_users_courses(settings.uclass_func_courses)
        course_index = extract_course_index(courses_payload)
        if course_index:
            _register_uclass_courses(
                settings,
                db,
                course_index,
                user_id=owner_id,
                ws_base_url=ws_base_url,
            )
    except Exception:
        logger.warning(
            "failed to refresh uclass course aliases for official timetable sync",
            extra={"user_id": owner_id, "ws_base_url": ws_base_url},
            exc_info=True,
        )


def _match_uos_official_catalog_events_to_uclass_courses(
    settings: Settings,
    db: Database,
    *,
    target: dict[str, Any],
    fetched: dict[str, Any],
) -> dict[str, Any]:
    owner_id = _safe_int(target.get("user_id")) or 0
    profiles = _build_uos_official_course_profiles(db, user_id=owner_id)
    if not profiles and owner_id > 0:
        _refresh_uclass_course_aliases_for_timetable_target(
            settings,
            db,
            target=target,
            owner_id=owner_id,
        )
        profiles = _build_uos_official_course_profiles(db, user_id=owner_id)

    grouped_sections: dict[str, dict[str, Any]] = {}
    for item in list(fetched.get("events") or []):
        if not isinstance(item, dict):
            continue
        section_key = _official_catalog_section_key(item)
        metadata = dict(item.get("metadata") or {}) if isinstance(item.get("metadata"), dict) else {}
        bucket = grouped_sections.setdefault(
            section_key,
            {
                "section_key": section_key,
                "title": str(item.get("title") or "").strip(),
                "events": [],
                "metadata": metadata,
            },
        )
        bucket["events"].append(item)

    if not profiles:
        result = dict(fetched)
        result["events"] = []
        result["table_count"] = 0
        result["reason"] = "No active UClass course aliases available for official timetable matching"
        result["course_match_summary"] = {
            "selected_sections": 0,
            "ambiguous_courses": 0,
            "unmatched_courses": 0,
            "catalog_sections": len(grouped_sections),
        }
        return result

    selected_sections: list[dict[str, Any]] = []
    ambiguous_courses: list[dict[str, Any]] = []
    unmatched_courses: list[dict[str, Any]] = []
    used_section_keys: set[str] = set()
    grouped_items = list(grouped_sections.values())

    for profile in profiles:
        scored_sections: list[tuple[int, dict[str, Any]]] = []
        metadata_text_key = _normalize_token_text(profile.get("metadata_text") or "")
        profile_subject_no = str(profile.get("official_subject_no") or "").strip()
        profile_dvcl_no = str(profile.get("official_dvcl_no") or "").strip()
        profile_section_key = str(profile.get("official_section_key") or "").strip()
        for section in grouped_items:
            metadata = dict(section.get("metadata") or {})
            title = str(section.get("title") or "").strip()
            normalized_title = normalize_course_alias(title)
            score = _best_title_match_score(list(profile.get("aliases") or []), title)
            if normalized_title and normalized_title in set(profile.get("normalized_aliases") or set()):
                score = max(score, 160)
            subject_no = str(metadata.get("official_subject_no") or "").strip()
            course_code = str(metadata.get("official_course_code") or "").strip()
            instructor = str(metadata.get("instructor") or "").strip()
            dvcl_no = str(metadata.get("official_dvcl_no") or "").strip()
            if profile_section_key and section.get("section_key") == profile_section_key:
                score += 600
            if profile_subject_no and subject_no == profile_subject_no:
                score += 120
            if profile_dvcl_no and dvcl_no == profile_dvcl_no:
                score += 120
            if subject_no and subject_no in metadata_text_key:
                score += 90
            if course_code and _normalize_token_text(course_code) in metadata_text_key:
                score += 40
            if instructor and _normalize_token_text(instructor) in metadata_text_key:
                score += 30
            if dvcl_no and re.search(rf"(?<!\d){re.escape(dvcl_no)}(?!\d)", str(profile.get('metadata_text') or "")):
                score += 15
            if score > 0:
                scored_sections.append((score, section))
        scored_sections.sort(
            key=lambda row: (
                -int(row[0]),
                str(row[1].get("title") or "").lower(),
                str(row[1].get("section_key") or ""),
            )
        )
        if not scored_sections or scored_sections[0][0] < 40:
            unmatched_courses.append(
                {
                    "canonical_course_id": profile["canonical_course_id"],
                    "display_name": profile.get("display_name"),
                }
            )
            continue
        best_score, best_section = scored_sections[0]
        second_score = scored_sections[1][0] if len(scored_sections) > 1 else 0
        if second_score >= 40 and (best_score - second_score) < 15:
            ambiguous_courses.append(
                {
                    "canonical_course_id": profile["canonical_course_id"],
                    "display_name": profile.get("display_name"),
                    "candidate_sections": [
                        {
                            "section_key": str(section.get("section_key") or ""),
                            "title": str(section.get("title") or ""),
                            "score": int(score),
                        }
                        for score, section in scored_sections[:3]
                    ],
                }
            )
            continue
        section_key = str(best_section.get("section_key") or "")
        if section_key in used_section_keys:
            continue
        used_section_keys.add(section_key)
        selected_sections.append(
            {
                "canonical_course_id": profile["canonical_course_id"],
                "display_name": profile.get("display_name"),
                "score": best_score,
                "section": best_section,
            }
        )

    matched_events: list[dict[str, Any]] = []
    for item in selected_sections:
        canonical_course_id = str(item.get("canonical_course_id") or "").strip()
        display_name = str(item.get("display_name") or "").strip() or None
        score = int(item.get("score") or 0)
        section = item.get("section") if isinstance(item.get("section"), dict) else {}
        for event in list(section.get("events") or []):
            event_payload = dict(event)
            metadata = dict(event_payload.get("metadata") or {}) if isinstance(event_payload.get("metadata"), dict) else {}
            metadata["canonical_course_id"] = canonical_course_id
            metadata["official_catalog_match_score"] = score
            if display_name:
                metadata["course_display_name"] = display_name
            event_payload["metadata"] = metadata
            matched_events.append(event_payload)

    result = dict(fetched)
    result["events"] = matched_events
    result["table_count"] = len(selected_sections)
    result["catalog_section_count"] = len(grouped_sections)
    result["course_match_summary"] = {
        "selected_sections": len(selected_sections),
        "ambiguous_courses": len(ambiguous_courses),
        "unmatched_courses": len(unmatched_courses),
        "catalog_sections": len(grouped_sections),
    }
    if not matched_events:
        result["reason"] = "Official timetable API did not match active UClass courses"
    return result


def _sync_uos_official_catalog_timetable_target(
    settings: Settings,
    db: Database,
    *,
    target: dict[str, Any],
) -> dict[str, Any]:
    api_url = str(getattr(settings, "ku_openapi_timetable_url", "") or "").strip()
    api_key = str(getattr(settings, "ku_openapi_timetable_api_key", "") or "").strip() or None
    owner_id = _safe_int(target.get("user_id")) or 0
    chat_id = str(target.get("chat_id") or "").strip() or None
    if not ku_openapi_timetable_configured(api_url, api_key):
        reason = "KU official timetable API is not configured"
        return {
            "user_id": owner_id,
            "chat_id": chat_id,
            "school_slug": KU_PORTAL_SCHOOL_SLUG,
            "display_name": str(target.get("display_name") or "").strip() or None,
            "upserted_events": 0,
            "fetched_events": 0,
            "table_count": 0,
            "current_url": None,
            "title": None,
            "skipped": True,
            "reason": reason,
            "payload_source": KU_OPENAPI_TIMETABLE_SOURCE,
            "source_attempts": [],
            "fallback_used": False,
            "_sync_status": "skipped",
            "_sync_last_error": reason,
            "_sync_action_required": 0,
            "_sync_cursor": {
                "upserted": 0,
                "fetched_events": 0,
                "table_count": 0,
                "reason": reason,
                "auth_required": False,
                "network_sample_count": 0,
                "payload_source": KU_OPENAPI_TIMETABLE_SOURCE,
                "source_url": None,
                "source_attempts": [],
                "fallback_used": False,
            },
        }
    timeout_sec = max(int(getattr(settings, "ku_openapi_timetable_timeout_sec", 15) or 15), 1)
    source_url = str(api_url or KU_OPENAPI_OFFICIAL_TIMETABLE_URL).strip() or KU_OPENAPI_OFFICIAL_TIMETABLE_URL
    try:
        fetched = fetch_ku_openapi_timetable(
            api_url=api_url or KU_OPENAPI_OFFICIAL_TIMETABLE_URL,
            api_key=api_key,
            academic_year=getattr(settings, "ku_openapi_year", None),
            term=getattr(settings, "ku_openapi_term", None),
            timezone_name=str(getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul"),
            timeout_sec=timeout_sec,
            target=target,
        )
    except KUOpenAPITimetableUnsupported as exc:
        reason = str(exc).strip() or "KU official timetable API request is unsupported"
        return {
            "user_id": owner_id,
            "chat_id": chat_id,
            "school_slug": KU_PORTAL_SCHOOL_SLUG,
            "display_name": str(target.get("display_name") or "").strip() or None,
            "upserted_events": 0,
            "fetched_events": 0,
            "table_count": 0,
            "current_url": None,
            "title": None,
            "skipped": True,
            "reason": reason,
            "payload_source": KU_OPENAPI_TIMETABLE_SOURCE,
            "source_attempts": [
                _timetable_source_attempt(
                    KU_OPENAPI_TIMETABLE_SOURCE,
                    status="unsupported",
                    reason=reason,
                    source_url=source_url,
                )
            ],
            "fallback_used": False,
            "_sync_status": "skipped",
            "_sync_last_error": reason,
            "_sync_action_required": 0,
            "_sync_cursor": {
                "upserted": 0,
                "fetched_events": 0,
                "table_count": 0,
                "reason": reason,
                "auth_required": False,
                "network_sample_count": 0,
                "payload_source": KU_OPENAPI_TIMETABLE_SOURCE,
                "source_url": source_url,
                "source_attempts": [
                    _timetable_source_attempt(
                        KU_OPENAPI_TIMETABLE_SOURCE,
                        status="unsupported",
                        reason=reason,
                        source_url=source_url,
                    )
                ],
                "fallback_used": False,
            },
        }
    except (KUOpenAPITimetableMalformedPayload, requests.RequestException) as exc:
        reason = str(exc).strip() or "KU official timetable API request failed"
        return {
            "user_id": owner_id,
            "chat_id": chat_id,
            "school_slug": KU_PORTAL_SCHOOL_SLUG,
            "display_name": str(target.get("display_name") or "").strip() or None,
            "upserted_events": 0,
            "fetched_events": 0,
            "table_count": 0,
            "current_url": None,
            "title": None,
            "skipped": False,
            "reason": reason,
            "payload_source": KU_OPENAPI_TIMETABLE_SOURCE,
            "source_attempts": [
                _timetable_source_attempt(
                    KU_OPENAPI_TIMETABLE_SOURCE,
                    status="error",
                    reason=reason,
                    source_url=source_url,
                )
            ],
            "fallback_used": False,
            "_sync_status": "error",
            "_sync_last_error": reason,
            "_sync_action_required": 1,
            "_sync_cursor": {
                "upserted": 0,
                "fetched_events": 0,
                "table_count": 0,
                "reason": reason,
                "auth_required": False,
                "network_sample_count": 0,
                "payload_source": KU_OPENAPI_TIMETABLE_SOURCE,
                "source_url": source_url,
                "source_attempts": [
                    _timetable_source_attempt(
                        KU_OPENAPI_TIMETABLE_SOURCE,
                        status="error",
                        reason=reason,
                        source_url=source_url,
                    )
                ],
                "fallback_used": False,
            },
        }

    filtered = _match_uos_official_catalog_events_to_uclass_courses(
        settings,
        db,
        target=target,
        fetched=dict(fetched),
    )
    filtered["payload_source"] = KU_OPENAPI_TIMETABLE_SOURCE
    filtered["source_url"] = str(filtered.get("source_url") or fetched.get("source_url") or source_url).strip() or source_url
    filtered["source_attempts"] = [
        _timetable_source_attempt(
            KU_OPENAPI_TIMETABLE_SOURCE,
            status="selected",
            source_url=str(filtered.get("source_url") or "").strip() or None,
        )
    ]
    filtered["fallback_used"] = False
    filtered["allow_empty_success"] = False
    return _apply_ku_portal_timetable_fetch_result(
        settings=settings,
        db=db,
        target=target,
        fetched=filtered,
    )


def _sync_ku_portal_timetable_target(
    settings: Settings,
    db: Database,
    *,
    target: dict[str, Any],
) -> dict[str, Any]:
    owner_id = _safe_int(target.get("user_id")) or 0
    chat_id = str(target.get("chat_id") or "").strip() or None
    secret_kind = str(target.get("secret_kind") or "").strip()
    secret_ref = str(target.get("secret_ref") or "").strip()
    api_url = str(getattr(settings, "ku_openapi_timetable_url", "") or "").strip()
    api_key = str(getattr(settings, "ku_openapi_timetable_api_key", "") or "").strip() or None
    if ku_openapi_timetable_configured(api_url, api_key) and ku_openapi_uses_official_catalog_mode(api_url):
        return _sync_uos_official_catalog_timetable_target(
            settings=settings,
            db=db,
            target=target,
        )

    openapi_fetched, source_attempts = _try_fetch_ku_openapi_timetable(
        settings=settings,
        target=target,
    )
    if openapi_fetched is not None:
        return _apply_ku_portal_timetable_fetch_result(
            settings=settings,
            db=db,
            target=target,
            fetched=openapi_fetched,
        )

    # KUPID SSO path is opt-in. The legacy browser-scraping flow runs
    # untouched whenever the flag is off — keeps existing tests/deployments
    # working until the operator explicitly enables the new adapter.
    if _kupid_sso_timetable_enabled(settings):
        sso_fetched, sso_attempts = _try_fetch_kupid_sso_timetable(
            settings=settings,
            target=target,
        )
        source_attempts = source_attempts + sso_attempts
        if sso_fetched is not None:
            sso_fetched = dict(sso_fetched)
            sso_fetched["source_attempts"] = source_attempts
            sso_fetched["fallback_used"] = bool(
                any(
                    str(item.get("source") or "").strip() == KU_OPENAPI_TIMETABLE_SOURCE
                    for item in source_attempts
                    if str(item.get("status") or "").strip()
                    in {"unsupported", "fallback", "error"}
                )
            )
            return _apply_ku_portal_timetable_fetch_result(
                settings=settings,
                db=db,
                target=target,
                fetched=sso_fetched,
            )

    profile_dir = None
    if str(target.get("profile_dir") or "").strip():
        candidate_profile_dir = Path(str(target.get("profile_dir") or "")).expanduser()
        if candidate_profile_dir.exists():
            profile_dir = candidate_profile_dir
    storage_state_secret = None
    if profile_dir is None and secret_kind and secret_ref:
        try:
            storage_state_secret = default_secret_store(settings).read_secret(
                ref=StoredSecretRef(kind=secret_kind, ref=secret_ref)
            )
        except Exception as exc:
            reason = (
                PORTAL_SECURE_STORAGE_MISSING_REASON
                if _looks_like_secure_storage_missing(str(exc))
                else str(exc).strip() or "unknown error"
            )
            current_url = str(target.get("current_url") or "").strip() or None
            source_attempts = source_attempts + [
                _timetable_source_attempt(
                    KU_PORTAL_BROWSER_TIMETABLE_SOURCE,
                    status="error",
                    reason=reason,
                    source_url=current_url,
                )
            ]
            session_metadata = dict(target.get("session_metadata") or {})
            session_metadata["portal_timetable_sync"] = {
                "last_synced_at": now_utc_iso(),
                "event_count": 0,
                "table_count": 0,
                "status": "error",
                "auth_required": True,
                "reason": reason,
                "current_url": current_url,
                "title": None,
                "payload_source": None,
                "source_url": None,
                "source_attempts": source_attempts,
                "fallback_used": bool(
                    any(
                        str(item.get("source") or "").strip() == KU_OPENAPI_TIMETABLE_SOURCE
                        for item in source_attempts
                    )
                ),
            }
            db.upsert_lms_browser_session(
                chat_id=str(target.get("chat_id") or ""),
                school_slug=str(target.get("school_slug") or KU_PORTAL_SCHOOL_SLUG),
                provider=str(target.get("provider") or KU_PORTAL_PROVIDER),
                display_name=str(target.get("display_name") or "고려대학교 포털/대학행정"),
                login_url=str(target.get("login_url") or KU_PORTAL_LOGIN_URL),
                profile_dir=str(target.get("profile_dir") or ""),
                secret_kind=secret_kind or None,
                secret_ref=secret_ref or None,
                status=str(target.get("status") or "active"),
                last_opened_at=target.get("last_opened_at"),
                last_verified_at=target.get("last_verified_at"),
                metadata_json=session_metadata,
                user_id=owner_id,
            )
            return {
                "user_id": owner_id,
                "chat_id": chat_id,
                "school_slug": KU_PORTAL_SCHOOL_SLUG,
                "display_name": str(target.get("display_name") or "").strip() or None,
                "upserted_events": 0,
                "fetched_events": 0,
                "table_count": 0,
                "current_url": str(target.get("current_url") or "").strip() or None,
                "title": None,
                "skipped": False,
                "reason": reason,
                "_sync_status": "error",
                "_sync_last_error": reason,
                "_sync_action_required": 1,
                "_sync_cursor": {
                    "upserted": 0,
                    "fetched_events": 0,
                    "table_count": 0,
                    "current_url": current_url,
                    "title": None,
                    "reason": reason,
                    "auth_required": True,
                    "network_sample_count": 0,
                    "payload_source": None,
                    "source_url": None,
                    "source_attempts": source_attempts,
                    "fallback_used": bool(
                        any(
                            str(item.get("source") or "").strip() == KU_OPENAPI_TIMETABLE_SOURCE
                            for item in source_attempts
                        )
                    ),
                },
            }
    fetched = fetch_ku_portal_timetable(
        storage_state=storage_state_secret,
        profile_dir=profile_dir,
        current_url=str(target.get("current_url") or "").strip() or None,
        timezone_name=str(getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul"),
        browser_channel=str(getattr(settings, "onboarding_browser_channel", "") or ""),
        browser_executable_path=getattr(settings, "onboarding_browser_executable_path", None),
        # Sync runs should be non-interactive even if manual onboarding uses a visible browser.
        headless=True,
    )
    fetched = dict(fetched)
    portal_source_url = (
        str(fetched.get("source_url") or fetched.get("current_url") or target.get("current_url") or "").strip()
        or None
    )
    fetched["payload_source"] = str(
        fetched.get("payload_source") or KU_PORTAL_BROWSER_TIMETABLE_SOURCE
    ).strip() or KU_PORTAL_BROWSER_TIMETABLE_SOURCE
    fetched["source_url"] = portal_source_url
    fetched["source_attempts"] = source_attempts + [
        _timetable_source_attempt(
            str(fetched.get("payload_source") or KU_PORTAL_BROWSER_TIMETABLE_SOURCE),
            status="selected",
            source_url=portal_source_url,
        )
    ]
    fetched["fallback_used"] = bool(
        any(
            str(item.get("source") or "").strip() == KU_OPENAPI_TIMETABLE_SOURCE
            for item in source_attempts
        )
    )
    return _apply_ku_portal_timetable_fetch_result(
        settings=settings,
        db=db,
        target=target,
        fetched=fetched,
    )


def _apply_ku_portal_timetable_fetch_result(
    settings: Settings,
    db: Database,
    *,
    target: dict[str, Any],
    fetched: dict[str, Any],
) -> dict[str, Any]:
    owner_id = _safe_int(target.get("user_id")) or 0
    chat_id = str(target.get("chat_id") or "").strip() or None
    secret_kind = str(target.get("secret_kind") or "").strip()
    secret_ref = str(target.get("secret_ref") or "").strip()
    events = list(fetched.get("events") or [])
    payload_source = str(
        fetched.get("payload_source") or KU_PORTAL_BROWSER_TIMETABLE_SOURCE
    ).strip() or KU_PORTAL_BROWSER_TIMETABLE_SOURCE
    source_url = str(fetched.get("source_url") or fetched.get("current_url") or "").strip() or None
    source_attempts = _normalize_timetable_source_attempts(fetched.get("source_attempts"))
    if not source_attempts:
        source_attempts = [
            _timetable_source_attempt(
                payload_source,
                status="selected",
                source_url=source_url,
            )
        ]
    fallback_used = bool(
        fetched.get("fallback_used")
        or any(
            str(item.get("source") or "").strip() == KU_OPENAPI_TIMETABLE_SOURCE
            for item in source_attempts
            if str(item.get("source") or "").strip() != payload_source
        )
    )
    allow_empty_success = bool(fetched.get("allow_empty_success"))
    if payload_source != KU_OPENAPI_TIMETABLE_SOURCE and events:
        events = _enrich_ku_portal_events_with_official_catalog(
            settings,
            events=events,
        )
    building_catalog_summary: dict[str, Any] | None = None
    if events:
        events, building_catalog_summary = _enrich_events_with_uos_building_catalog(
            settings,
            db,
            events=events,
        )
    alias_map = db.course_alias_resolution_map(user_id=owner_id)
    upserted = 0
    for item in events:
        metadata = dict(item.get("metadata")) if isinstance(item.get("metadata"), dict) else {}
        metadata["school_slug"] = str(metadata.get("school_slug") or KU_PORTAL_SCHOOL_SLUG).strip() or KU_PORTAL_SCHOOL_SLUG
        metadata["timetable_source"] = KU_PORTAL_SCHOOL_SLUG
        metadata["timetable_payload_source"] = payload_source
        if source_url:
            metadata["timetable_payload_source_url"] = source_url
        if _is_ku_portal_browser_target(target):
            metadata["portal_session"] = {
                "chat_id": chat_id,
                "user_id": owner_id,
                "school_slug": KU_PORTAL_SCHOOL_SLUG,
                "display_name": str(target.get("display_name") or "").strip() or None,
                "current_url": str(fetched.get("current_url") or target.get("current_url") or "").strip() or None,
                "payload_source": payload_source,
            }
        elif payload_source == KU_OPENAPI_TIMETABLE_SOURCE:
            metadata["official_api_target"] = {
                "chat_id": chat_id,
                "user_id": owner_id,
                "display_name": str(target.get("display_name") or "").strip() or None,
                "connection_id": _safe_int(target.get("connection_id")),
                "source_url": source_url,
            }
        metadata = attach_provenance(
            metadata,
            source="portal_uos_timetable",
            confidence="high",
            last_verified_at=now_utc_iso(),
            evidence_links=[
                str(
                    source_url
                    or fetched.get("current_url")
                    or target.get("current_url")
                    or ""
                ).strip()
            ],
            raw_source_ids=[str(item.get("external_id") or "").strip()],
            derivation=(
                "official_api_timetable_import"
                if payload_source == KU_OPENAPI_TIMETABLE_SOURCE
                else "portal_timetable_import"
            ),
        )
        metadata = _attach_canonical_course_id(
            metadata,
            alias_map=alias_map,
            fallback_aliases=[str(item.get("title") or "").strip()],
        )
        location_meta = _parse_uos_location(str(item.get("location") or "").strip() or None)
        building_no = location_meta.get("building_no")
        if building_no:
            metadata["building_no"] = str(building_no)
            building_name = db.get_building_name(str(building_no), school_slug=str(metadata.get("school_slug") or KU_PORTAL_SCHOOL_SLUG))
            if building_name:
                metadata["building_name"] = str(building_name)
            if location_meta.get("room"):
                metadata["room"] = str(location_meta.get("room"))
            if location_meta.get("extra"):
                metadata["location_extra"] = str(location_meta.get("extra"))
        canonical_course_id = str(metadata.get("canonical_course_id") or "").strip()
        if canonical_course_id:
            _persist_official_course_binding(
                db,
                canonical_course_id=canonical_course_id,
                event_metadata=metadata,
                user_id=owner_id,
            )
            db.upsert_course_alias(
                canonical_course_id=canonical_course_id,
                alias=str(item.get("title") or "").strip(),
                alias_type="portal_title",
                source="portal",
                metadata_json={"event_external_id": str(item.get("external_id") or "").strip()},
                user_id=owner_id,
            )
        db.upsert_event(
            external_id=str(item.get("external_id") or "").strip(),
            source="portal",
            start=str(item.get("start_at") or "").strip(),
            end=str(item.get("end_at") or "").strip(),
            title=str(item.get("title") or "").strip(),
            location=str(item.get("location") or "").strip() or None,
            rrule=str(item.get("rrule") or "").strip() or None,
            metadata_json=metadata,
            user_id=owner_id,
        )
        upserted += 1

    title = str(fetched.get("title") or "").strip()
    auth_required = bool(fetched.get("auth_required"))
    has_timetable_surface = bool(fetched.get("has_timetable_surface"))
    fetch_reason = str(fetched.get("reason") or "").strip() or None
    has_timetable_surface = bool(
        has_timetable_surface or title == UOS_TIMETABLE_TITLE
    )
    if events:
        status = "success"
        reason = None
    elif auth_required:
        status = "error"
        reason = "KU portal session expired; reconnect required"
    elif fetch_reason:
        status = "skipped"
        reason = fetch_reason
    elif has_timetable_surface or (
        payload_source == KU_OPENAPI_TIMETABLE_SOURCE
        and bool(fetched.get("ok"))
        and allow_empty_success
    ):
        status = "success"
        reason = None
    else:
        status = "skipped"
        reason = "portal timetable not available"
    current_url = str(fetched.get("current_url") or target.get("current_url") or "").strip() or None
    session_metadata = dict(target.get("session_metadata") or {})
    browser_result = (
        dict(session_metadata.get("browser_result"))
        if isinstance(session_metadata.get("browser_result"), dict)
        else {}
    )
    if payload_source == KU_PORTAL_BROWSER_TIMETABLE_SOURCE:
        if current_url:
            browser_result["current_url"] = current_url
        if title:
            browser_result["title"] = title
        browser_result = sanitize_browser_session_result(browser_result)
        if browser_result:
            session_metadata["browser_result"] = browser_result
    session_metadata["portal_timetable_sync"] = {
        "last_synced_at": now_utc_iso(),
        "event_count": len(events),
        "table_count": int(fetched.get("table_count") or 0),
        "status": status,
        "auth_required": auth_required,
        "reason": reason,
        "current_url": current_url,
        "title": title or None,
        "payload_source": payload_source,
        "source_url": source_url,
        "source_attempts": source_attempts,
        "fallback_used": fallback_used,
    }
    if _is_ku_portal_browser_target(target):
        db.upsert_lms_browser_session(
            chat_id=str(target.get("chat_id") or ""),
            school_slug=str(target.get("school_slug") or KU_PORTAL_SCHOOL_SLUG),
            provider=str(target.get("provider") or KU_PORTAL_PROVIDER),
            display_name=str(target.get("display_name") or "고려대학교 포털/대학행정"),
            login_url=str(target.get("login_url") or KU_PORTAL_LOGIN_URL),
            profile_dir=str(target.get("profile_dir") or ""),
            secret_kind=secret_kind or None,
            secret_ref=secret_ref or None,
            status=str(target.get("status") or "active"),
            last_opened_at=target.get("last_opened_at"),
            last_verified_at=now_utc_iso(),
            metadata_json=session_metadata,
            user_id=owner_id,
        )
    cursor = {
        "upserted": upserted,
        "fetched_events": len(events),
        "table_count": int(fetched.get("table_count") or 0),
        "current_url": current_url,
        "title": title or None,
        "reason": reason,
        "auth_required": auth_required,
        "network_sample_count": len(list(fetched.get("network_samples") or [])),
        "payload_source": payload_source,
        "source_url": source_url,
        "source_attempts": source_attempts,
        "fallback_used": fallback_used,
    }
    if isinstance(fetched.get("course_match_summary"), dict):
        cursor["course_match_summary"] = dict(fetched.get("course_match_summary") or {})
    if fetched.get("catalog_section_count") is not None:
        cursor["catalog_section_count"] = int(fetched.get("catalog_section_count") or 0)
    if building_catalog_summary:
        cursor["building_catalog_summary"] = dict(building_catalog_summary)
    return {
        "user_id": owner_id,
        "chat_id": chat_id,
        "school_slug": KU_PORTAL_SCHOOL_SLUG,
        "display_name": str(target.get("display_name") or "").strip() or None,
        "upserted_events": upserted,
        "fetched_events": len(events),
        "table_count": int(fetched.get("table_count") or 0),
        "current_url": current_url,
        "title": title or None,
        "skipped": status == "skipped",
        "reason": reason,
        "payload_source": payload_source,
        "source_attempts": source_attempts,
        "fallback_used": fallback_used,
        "_sync_status": status,
        "_sync_last_error": reason,
        "_sync_action_required": 1 if auth_required or bool(fetch_reason) else 0,
        "_sync_cursor": cursor,
    }


def record_ku_portal_timetable_fetch_for_user(
    settings: Settings,
    db: Database,
    *,
    fetched: dict[str, Any],
    chat_id: str | None = None,
    user_id: int | None = None,
) -> dict[str, Any]:
    owner_id = _safe_int(user_id)
    chat = str(chat_id or "").strip() or None
    if owner_id is None and chat:
        owner_id = int(
            _resolve_user_scope(
                settings,
                db,
                chat_id=chat,
                create_if_missing=False,
                metadata_source="portal_timetable_prefetch",
            )["user_id"]
            or 0
        ) or None
    if owner_id is None or owner_id <= 0:
        return {"skipped": True, "reason": "user_id missing"}
    targets = [
        item
        for item in _resolve_ku_portal_timetable_targets(settings, db)
        if int(item.get("user_id") or 0) == owner_id
    ]
    if chat:
        targets = [item for item in targets if str(item.get("chat_id") or "").strip() == chat]
    if not targets:
        reason = "no active KU timetable target"
        _record_sync_dashboard_state(
            db,
            "sync_ku_portal_timetable",
            status="skipped",
            last_error=reason,
            cursor_payload={"skipped": True, "reason": reason},
            user_id=owner_id,
        )
        return {"skipped": True, "reason": reason, "user_id": owner_id}
    try:
        result = _apply_ku_portal_timetable_fetch_result(
            settings=settings,
            db=db,
            target=targets[0],
            fetched=fetched,
        )
    except Exception as exc:
        logger.warning(
            "portal timetable prefetch apply failed",
            extra={"chat_id": chat, "user_id": owner_id, "error": str(exc)},
        )
        _record_sync_dashboard_state(
            db,
            "sync_ku_portal_timetable",
            status="error",
            action_required=1,
            last_error=str(exc),
            cursor_payload={"skipped": True, "reason": str(exc)},
            user_id=owner_id,
        )
        return {
            "ok": False,
            "error": str(exc),
            "chat_id": chat,
            "user_id": owner_id,
        }
    result_status = str(result.get("_sync_status") or "success").strip().lower() or "success"
    result_error = str(result.get("_sync_last_error") or result.get("reason") or "").strip() or None
    _record_sync_dashboard_state(
        db,
        "sync_ku_portal_timetable",
        status=result_status,
        new_items=int(result.get("upserted_events") or 0),
        action_required=int(result.get("_sync_action_required") or 0),
        last_error=result_error,
        cursor_payload=result.get("_sync_cursor") if isinstance(result.get("_sync_cursor"), dict) else None,
        user_id=owner_id,
    )
    return {
        "ok": result_status == "success",
        "skipped": result_status == "skipped",
        "status": result_status,
        "error": result_error if result_status == "error" else None,
        "reason": str(result.get("reason") or "").strip() or None,
        "chat_id": chat,
        "user_id": owner_id,
        "upserted_events": int(result.get("upserted_events") or 0),
    }


def _portal_timetable_adapters() -> dict[str, PortalTimetableAdapter]:
    return {
        KU_PORTAL_PROVIDER: PortalTimetableAdapter(
            provider=KU_PORTAL_PROVIDER,
            school_slug=KU_PORTAL_SCHOOL_SLUG,
            job_name="sync_ku_portal_timetable",
            no_target_reason="No active KU timetable targets",
            resolve_targets=_resolve_ku_portal_timetable_targets,
            sync_target=_sync_ku_portal_timetable_target,
        )
    }


def _sync_portal_timetable_adapter(
    settings: Settings,
    db: Database,
    *,
    adapter: PortalTimetableAdapter,
) -> dict[str, Any]:
    targets = adapter.resolve_targets(settings, db)
    if not targets:
        reason = adapter.no_target_reason
        _record_sync_dashboard_state(
            db,
            adapter.job_name,
            status="skipped",
            last_error=reason,
            cursor_payload={"skipped": True, "reason": reason},
        )
        return {"skipped": True, "reason": reason}

    target_results: list[dict[str, Any]] = []
    skipped_targets: list[dict[str, Any]] = []
    error_targets: list[dict[str, Any]] = []
    failed_targets: list[dict[str, Any]] = []
    aggregate_upserted = 0
    for target in targets:
        owner_id = _safe_int(target.get("user_id")) or 0
        try:
            result = adapter.sync_target(settings=settings, db=db, target=target)
        except Exception as exc:
            failure = {
                "user_id": owner_id,
                "chat_id": str(target.get("chat_id") or "").strip() or None,
                "school_slug": adapter.school_slug,
                "provider": adapter.provider,
                "error": str(exc),
            }
            failed_targets.append(failure)
            logger.warning("portal timetable sync failed for target", extra=failure)
            _record_sync_dashboard_state(
                db,
                adapter.job_name,
                status="error",
                action_required=1,
                last_error=str(exc),
                cursor_payload={"skipped": True, "reason": str(exc), "target": failure},
                user_id=owner_id,
            )
            continue
        _record_sync_dashboard_state(
            db,
            adapter.job_name,
            status=str(result.get("_sync_status") or "success"),
            new_items=int(result.get("upserted_events") or 0),
            action_required=int(result.get("_sync_action_required") or 0),
            last_error=str(result.get("_sync_last_error") or "").strip() or None,
            cursor_payload=result.get("_sync_cursor") if isinstance(result.get("_sync_cursor"), dict) else None,
            user_id=owner_id,
        )
        aggregate_upserted += int(result.get("upserted_events") or 0)
        result_status = str(result.get("_sync_status") or "").strip().lower()
        if result_status == "error":
            error_targets.append(result)
        elif result.get("skipped"):
            skipped_targets.append(result)
        else:
            target_results.append(result)

    overall_status = (
        "error"
        if failed_targets or error_targets
        else "success"
        if target_results
        else "skipped"
    )
    overall_error = (
        str(failed_targets[0].get("error") or "").strip()
        if failed_targets
        else str(error_targets[0].get("reason") or error_targets[0].get("_sync_last_error") or "").strip() or None
        if error_targets
        else str(skipped_targets[0].get("reason") or "").strip() or None
        if skipped_targets and not target_results
        else None
    )
    all_results = target_results + error_targets + skipped_targets
    cursor_payload = {
        "upserted": aggregate_upserted,
        "target_count": len(targets),
        "synced_targets": len(target_results),
        "error_targets": [
            {
                "user_id": item.get("user_id"),
                "chat_id": item.get("chat_id"),
                "reason": item.get("reason"),
            }
            for item in error_targets
        ],
        "skipped_targets": [
            {
                "user_id": item.get("user_id"),
                "chat_id": item.get("chat_id"),
                "reason": item.get("reason"),
            }
            for item in skipped_targets
        ],
        "failed_targets": failed_targets,
        "payload_sources": [
            {
                "user_id": item.get("user_id"),
                "chat_id": item.get("chat_id"),
                "payload_source": item.get("payload_source"),
                "fallback_used": bool(item.get("fallback_used")),
            }
            for item in all_results
        ],
    }
    _record_sync_dashboard_state(
        db,
        adapter.job_name,
        status=overall_status,
        new_items=aggregate_upserted,
        action_required=len(failed_targets) + len(error_targets),
        last_error=overall_error,
        cursor_payload=cursor_payload,
    )
    return {
        "upserted_events": aggregate_upserted,
        "synced_targets": len(target_results),
        "error_targets": error_targets,
        "skipped_targets": skipped_targets,
        "failed_targets": failed_targets,
    }


def sync_ku_portal_timetable(settings: Settings, db: Database) -> dict[str, Any]:
    adapter = _portal_timetable_adapters()[KU_PORTAL_PROVIDER]
    return _sync_portal_timetable_adapter(settings=settings, db=db, adapter=adapter)

def _singleton_settings_owned_user_scope(
    settings: Settings,
    db: Database,
) -> dict[str, Any] | None:
    def _collect_candidates(items: list[dict[str, Any]]) -> list[dict[str, Any]]:
        candidates: dict[int, dict[str, Any]] = {}
        for item in items:
            owner_id = _safe_int(item.get("user_id") or item.get("id"))
            if owner_id is None or owner_id <= 0:
                continue
            chat_id = (
                str(item.get("chat_id") or item.get("telegram_chat_id") or "").strip()
                or None
            )
            existing = candidates.get(owner_id)
            if existing is None:
                candidates[owner_id] = {
                    "user_id": owner_id,
                    "chat_id": chat_id,
                }
                continue
            if not existing.get("chat_id") and chat_id:
                existing["chat_id"] = chat_id
        return list(candidates.values())

    preference_candidates = _collect_candidates(db.list_user_preferences(limit=1000))
    if len(preference_candidates) == 1:
        return preference_candidates[0]
    if len(preference_candidates) > 1:
        return None

    configured_candidates = _collect_candidates(
        [
            user
            for chat_id in _configured_telegram_chat_ids(settings)
            if (user := db.get_user_by_chat_id(chat_id)) is not None
        ]
    )
    if len(configured_candidates) == 1:
        return configured_candidates[0]
    if len(configured_candidates) > 1:
        return None

    active_candidates = _collect_candidates(db.list_users(status="active", limit=1000))
    if len(active_candidates) == 1:
        return active_candidates[0]
    return None


def _resolve_uclass_sync_targets(settings: Settings, db: Database) -> list[dict[str, Any]]:
    targets: list[dict[str, Any]] = []
    secret_store = default_secret_store(settings)
    for item in db.list_moodle_connections(status="active", limit=500):
        owner_id = _safe_int(item.get("user_id")) or 0
        chat_id = str(item.get("chat_id") or "").strip() or None
        user_scope = _resolve_user_scope(
            settings,
            db,
            chat_id=chat_id,
            user_id=owner_id if owner_id > 0 else None,
            create_if_missing=bool(chat_id),
            metadata_source="uclass_connection",
        )
        owner_id = _safe_int(user_scope.get("user_id")) or owner_id
        token = ""
        token_error = ""
        try:
            token = secret_store.read_secret(
                ref=StoredSecretRef(
                    kind=str(item.get("secret_kind") or "").strip(),
                    ref=str(item.get("secret_ref") or "").strip(),
                )
            )
        except Exception as exc:
            token_error = (
                UCLASS_SECURE_STORAGE_MISSING_REASON
                if _looks_like_secure_storage_missing(str(exc))
                else str(exc)
            )
        targets.append(
            {
                "user_id": owner_id,
                "chat_id": chat_id,
                "school_slug": str(item.get("school_slug") or "").strip() or None,
                "display_name": str(item.get("display_name") or item.get("school_slug") or "").strip() or "UClass",
                "ws_base_url": str(item.get("ws_base_url") or "").strip(),
                "timezone": str(
                    user_scope.get("timezone")
                    or getattr(settings, "timezone", "Asia/Seoul")
                    or "Asia/Seoul"
                ),
                "token": token,
                "token_error": token_error,
                "connection_id": int(item.get("id") or 0),
            }
        )
    if targets:
        return targets
    ws_base = str(getattr(settings, "uclass_ws_base", "") or "").strip()
    static_token = str(getattr(settings, "uclass_wstoken", "") or "").strip()
    if not ws_base or not static_token:
        return []
    default_scope = _singleton_settings_owned_user_scope(settings, db) or {}
    return [
        {
            "user_id": _safe_int(default_scope.get("user_id")) or 0,
            "chat_id": str(default_scope.get("chat_id") or "").strip() or None,
            "school_slug": None,
            "display_name": "UClass",
            "ws_base_url": ws_base,
            "timezone": str(default_scope.get("timezone") or getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul"),
            "token": static_token,
            "token_error": "",
            "connection_id": 0,
        }
    ]


@dataclass
class _UClassTargetAuthContext:
    owner_id: int
    ws_base_url: str
    token: str
    token_error: str
    html_material_candidates: list[Any]
    html_material_error: str
    client: MoodleWSClient
    ws_available: bool
    required_ws: set[str]
    ws_status: dict[str, dict[str, Any]]


@dataclass
class _UClassTargetWSResult:
    site_info: dict[str, Any]
    course_index: dict[int, dict[str, Any]]
    canonical_courses: dict[int, str]
    alias_map: dict[str, tuple[str, ...]]
    contents_payload_by_course: dict[int, Any]
    notifications: list[Any]
    tasks: list[Any]
    events: list[Any]
    semantic_warnings: int


@dataclass
class _UClassMaterialState:
    material: Any
    owner_id: int
    metadata: dict[str, Any]
    local_path: str | None
    content_hash: str | None
    previous_content_hash: str
    downloaded: bool
    resolved_filename: str
    previous_extract: Any
    previous_deadline_scan: Any
    previous_brief: Any
    material_provenance_source: str
    text_content: str | None = None
    text_error: str | None = None
    text_extract_type: str = ""

    @property
    def content_changed(self) -> bool:
        current_hash = str(self.content_hash or "").strip()
        return bool(current_hash and current_hash != str(self.previous_content_hash or "").strip())


@dataclass
class _UClassMaterialSyncResult:
    artifact_count: int = 0
    downloaded_count: int = 0
    reused_count: int = 0
    failed_downloads: int = 0
    extracted_count: int = 0
    extraction_failures: int = 0
    brief_count: int = 0
    material_task_count: int = 0
    generated_brief_items: list[dict[str, Any]] = field(default_factory=list)


def _prepare_uclass_target_auth(
    settings: Settings,
    *,
    target: dict[str, Any],
    owner_id: int,
) -> _UClassTargetAuthContext:
    ws_base_url = str(target.get("ws_base_url") or "").strip()
    if not ws_base_url:
        raise ValueError("UCLASS_WS_BASE is required")
    setattr(settings, "_uclass_resolved_token", "")
    html_material_candidates: list[Any] = []
    html_material_error = ""
    token = str(target.get("token") or "").strip()
    token_error = str(target.get("token_error") or "").strip()

    setattr(settings, "_uclass_resolved_token", token)
    if not token:
        raise RuntimeError(token_error or UCLASS_RECONNECT_REQUIRED_REASON)

    client = MoodleWSClient(
        base_url=ws_base_url,
        token=token,
        request_method=settings.uclass_request_method,
    )
    required_ws = {
        item.strip()
        for item in settings.uclass_required_wsfunctions
        if item.strip()
    }
    return _UClassTargetAuthContext(
        owner_id=owner_id,
        ws_base_url=ws_base_url,
        token=token,
        token_error=token_error,
        html_material_candidates=html_material_candidates,
        html_material_error=html_material_error,
        client=client,
        ws_available=bool(token),
        required_ws=required_ws,
        ws_status={},
    )


def _fetch_uclass_forum_notifications(
    settings: Settings,
    *,
    auth: _UClassTargetAuthContext,
    course_ids: list[int],
) -> list[Any]:
    forum_notifications: list[Any] = []
    forums_payload = _call_optional_ws(
        ws_status=auth.ws_status,
        wsfunction=settings.uclass_func_forums,
        enabled=auth.ws_available and settings.uclass_enable_forums,
        required_ws=auth.required_ws,
        loader=lambda: auth.client.get_forums(
            settings.uclass_func_forums,
            course_ids=course_ids,
        ),
        log_skip=auth.ws_available,
    )
    if isinstance(forums_payload, list):
        forum_items = forums_payload
    elif isinstance(forums_payload, dict):
        maybe_items = forums_payload.get("forums")
        forum_items = maybe_items if isinstance(maybe_items, list) else []
    else:
        forum_items = []
    for forum in forum_items:
        if not isinstance(forum, dict):
            continue
        forum_id = forum.get("id")
        if forum_id in (None, ""):
            continue
        discussions_payload = _call_optional_ws(
            ws_status=auth.ws_status,
            wsfunction=settings.uclass_func_forum_discussions,
            enabled=auth.ws_available and settings.uclass_enable_forums,
            required_ws=auth.required_ws,
            loader=lambda fid=forum_id: auth.client.get_forum_discussions(
                settings.uclass_func_forum_discussions,
                forum_id=int(fid),
                page=0,
                per_page=settings.uclass_page_limit,
            ),
            log_skip=auth.ws_available,
        )
        forum_notifications.extend(normalize_forum_notifications(discussions_payload, forum=forum))
    return forum_notifications


def _fetch_uclass_ws_stage(
    settings: Settings,
    db: Database,
    *,
    owner_id: int,
    auth: _UClassTargetAuthContext,
) -> _UClassTargetWSResult:
    site_info = _call_optional_ws(
        ws_status=auth.ws_status,
        wsfunction=settings.uclass_func_site_info,
        enabled=auth.ws_available,
        required_ws=auth.required_ws,
        loader=lambda: auth.client.get_site_info(settings.uclass_func_site_info),
        log_skip=auth.ws_available,
    ) or {}

    popup_payload = _call_optional_ws(
        ws_status=auth.ws_status,
        wsfunction=settings.uclass_func_popup_notifications,
        enabled=auth.ws_available and settings.uclass_enable_popup_notifications,
        required_ws=auth.required_ws,
        loader=lambda: auth.client.get_popup_notifications(
            settings.uclass_func_popup_notifications,
            limit=settings.uclass_page_limit,
        ),
        log_skip=auth.ws_available,
    )
    action_payload = _call_optional_ws(
        ws_status=auth.ws_status,
        wsfunction=settings.uclass_func_action_events,
        enabled=auth.ws_available and settings.uclass_enable_action_events,
        required_ws=auth.required_ws,
        loader=lambda: auth.client.get_action_events(
            settings.uclass_func_action_events,
            limitnum=settings.uclass_page_limit,
        ),
        log_skip=auth.ws_available,
    )

    need_courses = (
        settings.uclass_enable_courses
        or settings.uclass_enable_contents
        or settings.uclass_enable_assignments
        or settings.uclass_enable_forums
    )
    courses_payload = _call_optional_ws(
        ws_status=auth.ws_status,
        wsfunction=settings.uclass_func_courses,
        enabled=auth.ws_available and need_courses,
        required_ws=auth.required_ws,
        loader=lambda: auth.client.get_users_courses(settings.uclass_func_courses),
        log_skip=auth.ws_available,
    )
    course_index = extract_course_index(courses_payload)
    course_ids = sorted(course_index.keys())
    canonical_courses = _register_uclass_courses(
        settings,
        db,
        course_index,
        user_id=owner_id,
        ws_base_url=auth.ws_base_url,
    )
    alias_map = db.course_alias_resolution_map(user_id=owner_id)

    contents_payload_by_course: dict[int, Any] = {}
    if settings.uclass_enable_contents:
        for course_id in course_ids:
            payload = _call_optional_ws(
                ws_status=auth.ws_status,
                wsfunction=settings.uclass_func_course_contents,
                enabled=auth.ws_available,
                required_ws=auth.required_ws,
                loader=lambda cid=course_id: auth.client.get_course_contents(
                    settings.uclass_func_course_contents,
                    course_id=cid,
                ),
                log_skip=auth.ws_available,
            )
            if payload is not None:
                contents_payload_by_course[course_id] = payload

    assignments_payload = _call_optional_ws(
        ws_status=auth.ws_status,
        wsfunction=settings.uclass_func_assignments,
        enabled=auth.ws_available and settings.uclass_enable_assignments,
        required_ws=auth.required_ws,
        loader=lambda: auth.client.get_assignments(
            settings.uclass_func_assignments,
            course_ids=course_ids,
        ),
        log_skip=auth.ws_available,
    )
    forum_notifications = _fetch_uclass_forum_notifications(
        settings,
        auth=auth,
        course_ids=course_ids,
    )

    popup_notifications = normalize_notifications(popup_payload) if popup_payload is not None else []
    semantic_warnings = _warn_popup_notification_semantics(
        popup_notifications,
        settings.uclass_func_popup_notifications,
    )
    semantic_warnings += _warn_forum_semantics(
        forum_notifications,
        settings.uclass_func_forum_discussions,
    )
    notifications = list(
        {
            item.external_id: item
            for item in [*popup_notifications, *forum_notifications]
        }.values()
    )

    tasks: list[Any] = []
    events: list[Any] = []
    if action_payload is not None:
        base_tasks, base_events = normalize_action_events(action_payload)
        tasks.extend(base_tasks)
        events.extend(base_events)
    if assignments_payload is not None:
        tasks.extend(normalize_assignments(assignments_payload, course_index=course_index))
    semantic_warnings += _warn_assignment_semantics(tasks, settings.uclass_func_assignments)
    tasks = list({item.external_id: item for item in tasks}.values())
    events = list({item.external_id: item for item in events}.values())

    return _UClassTargetWSResult(
        site_info=site_info if isinstance(site_info, dict) else {},
        course_index=course_index,
        canonical_courses=canonical_courses,
        alias_map=alias_map,
        contents_payload_by_course=contents_payload_by_course,
        notifications=notifications,
        tasks=tasks,
        events=events,
        semantic_warnings=semantic_warnings,
    )


def _persist_uclass_ws_records(
    db: Database,
    *,
    owner_id: int,
    ws_data: _UClassTargetWSResult,
) -> tuple[int, int, int]:
    notif_count = 0
    for item in ws_data.notifications:
        metadata = _attach_canonical_course_id(
            item.metadata,
            alias_map=ws_data.alias_map,
            canonical_by_course_id=ws_data.canonical_courses,
        )
        db.upsert_notification(
            external_id=item.external_id,
            source="uclass",
            created_at=item.created_at,
            title=item.title,
            body=item.body,
            url=item.url,
            metadata_json=attach_provenance(
                metadata,
                source="uclass_ws",
                confidence="high",
                last_verified_at=item.created_at,
                evidence_links=[item.url] if item.url else None,
                derivation="uclass_notification",
            ),
            user_id=owner_id,
        )
        notif_count += 1

    task_count = 0
    for item in ws_data.tasks:
        metadata = _attach_canonical_course_id(
            item.metadata,
            alias_map=ws_data.alias_map,
            canonical_by_course_id=ws_data.canonical_courses,
        )
        db.upsert_task(
            external_id=item.external_id,
            source="uclass",
            due_at=item.due_at,
            title=item.title,
            status=item.status,
            metadata_json=attach_provenance(
                metadata,
                source="uclass_ws",
                confidence="high",
                last_verified_at=now_utc_iso(),
                derivation="uclass_task",
            ),
            user_id=owner_id,
        )
        task_count += 1

    event_count = 0
    for item in ws_data.events:
        metadata = _attach_canonical_course_id(
            item.metadata,
            alias_map=ws_data.alias_map,
            canonical_by_course_id=ws_data.canonical_courses,
        )
        db.upsert_event(
            external_id=item.external_id,
            source="uclass",
            start=item.start_at,
            end=item.end_at,
            title=item.title,
            location=item.location,
            rrule=item.rrule,
            metadata_json=attach_provenance(
                metadata,
                source="uclass_ws",
                confidence="high",
                last_verified_at=item.start_at,
                derivation="uclass_event",
            ),
            user_id=owner_id,
        )
        event_count += 1

    return notif_count, task_count, event_count


def _discover_uclass_material_candidates(
    ws_data: _UClassTargetWSResult,
    *,
    timezone_name: str,
    html_material_candidates: list[Any] | None = None,
) -> list[Any]:
    candidates = extract_material_candidates(
        ws_data.notifications,
        ws_data.tasks,
        ws_data.events,
        timezone_name=timezone_name,
    )
    candidates.extend(
        extract_material_candidates_from_course_contents(
            course_contents=ws_data.contents_payload_by_course,
            course_index=ws_data.course_index,
            timezone_name=timezone_name,
        )
    )
    if html_material_candidates:
        candidates.extend(html_material_candidates)
    return list({item.external_id: item for item in candidates}.values())


def _initialize_uclass_material_state(
    db: Database,
    *,
    owner_id: int,
    material: Any,
    alias_map: dict[str, tuple[str, ...]],
    canonical_courses: dict[int, str],
) -> _UClassMaterialState:
    resolved_filename = _safe_filename(material.filename)
    existing_artifact = db.get_artifact(
        external_id=material.external_id,
        source="uclass",
        user_id=owner_id,
    )
    previous_content_hash = (
        str(existing_artifact.content_hash or "").strip()
        if existing_artifact and existing_artifact.content_hash
        else ""
    )
    metadata = _json_load(existing_artifact.metadata_json if existing_artifact else None)
    metadata = metadata | material.metadata | {"url": material.url}
    metadata["original_url"] = material.url
    metadata = _attach_canonical_course_id(
        metadata,
        alias_map=alias_map,
        canonical_by_course_id=canonical_courses,
    )
    material_provenance_source = (
        "uclass_html"
        if str(material.metadata.get("source_kind") or "").strip()
        else "uclass_ws"
    )
    local_path = None
    content_hash = None
    if existing_artifact and existing_artifact.icloud_path and Path(existing_artifact.icloud_path).exists():
        local_path = existing_artifact.icloud_path
        content_hash = existing_artifact.content_hash
        resolved_filename = Path(existing_artifact.icloud_path).name
        metadata["resolved_filename"] = resolved_filename
    previous_extract = metadata.get("text_extract")
    previous_deadline_scan = metadata.get("deadline_scan")
    previous_brief = metadata.get("brief")
    return _UClassMaterialState(
        material=material,
        owner_id=owner_id,
        metadata=metadata,
        local_path=local_path,
        content_hash=content_hash,
        previous_content_hash=previous_content_hash,
        downloaded=False,
        resolved_filename=resolved_filename,
        previous_extract=previous_extract,
        previous_deadline_scan=previous_deadline_scan,
        previous_brief=previous_brief,
        material_provenance_source=material_provenance_source,
    )


def _uclass_materials_root(settings: Settings, *, owner_id: int) -> Path | None:
    storage_root = resolve_storage_root(settings)
    if storage_root is None:
        return None
    scope = f"user-{owner_id}" if owner_id > 0 else "global"
    return storage_materials_dir(storage_root) / scope


def _download_uclass_material_stage(
    settings: Settings,
    db: Database,
    *,
    state: _UClassMaterialState,
    materials_root: Path | None,
) -> tuple[int, int, int]:
    if not settings.uclass_download_materials or materials_root is None or not state.material.url:
        return 0, 0, 0
    filename = _safe_filename(state.material.filename)
    course = _safe_path_component(state.material.course, fallback="general")
    date_folder = _safe_path_component(state.material.date_folder, fallback="unknown-date")
    target_path = materials_root / course / date_folder / filename
    try:
        download_result = _download_material(
            db=db,
            settings=settings,
            external_id=state.material.external_id,
            url=state.material.url,
            target=target_path,
            owner_id=state.owner_id,
        )
        if len(download_result) == 3:
            local_path, content_hash, downloaded = download_result  # type: ignore[misc]
            resolved_filename = (
                Path(local_path).name if local_path else _safe_filename(state.material.filename)
            )
            download_meta: dict[str, Any] = {
                "original_url": state.material.url,
                "resolved_filename": resolved_filename,
            }
        else:
            local_path, content_hash, downloaded, download_meta = download_result  # type: ignore[misc]
            resolved_filename = str(download_meta.get("resolved_filename") or state.resolved_filename)
        state.local_path = local_path
        state.content_hash = content_hash
        state.downloaded = bool(downloaded)
        state.resolved_filename = resolved_filename
        state.metadata["original_url"] = download_meta.get("original_url") or state.material.url
        state.metadata["resolved_filename"] = resolved_filename
        if download_meta.get("content_type"):
            state.metadata["content_type"] = str(download_meta["content_type"])
        if state.downloaded:
            state.metadata["downloaded_at"] = now_utc_iso()
            return 1, 0, 0
        return 0, 1, 0
    except Exception as exc:
        logger.warning(
            "failed to download material",
            extra={
                "external_id": state.material.external_id,
                "error": str(exc),
                "user_id": state.owner_id,
            },
        )
        return 0, 0, 1


def _attach_uclass_material_provenance(state: _UClassMaterialState) -> None:
    state.metadata = attach_provenance(
        state.metadata,
        source=state.material_provenance_source,
        confidence="high",
        last_verified_at=state.metadata.get("downloaded_at") or now_utc_iso(),
        evidence_links=[state.material.url] if state.material.url else None,
        derivation="uclass_material",
    )


def _extract_uclass_material_stage(
    settings: Settings,
    *,
    state: _UClassMaterialState,
) -> tuple[int, int]:
    should_extract = (
        settings.material_extraction_enabled
        and bool(state.local_path)
        and (
            state.downloaded
            or state.content_changed
            or not isinstance(state.previous_extract, dict)
            or not state.previous_extract.get("ok")
        )
    )
    if not should_extract or not state.local_path:
        return 0, 0
    state.text_content, state.text_error, state.text_extract_type = extract_material_text(
        Path(state.local_path),
        max_chars=settings.material_extract_max_chars,
    )
    if state.text_content:
        text_hash = sha1(state.text_content.encode("utf-8")).hexdigest()
        state.metadata["text_extract"] = {
            "ok": True,
            "type": state.text_extract_type,
            "hash": text_hash,
            "chars": len(state.text_content),
            "excerpt": state.text_content[: min(2000, len(state.text_content))],
        }
        return 1, 0
    if state.text_error and not state.text_error.startswith("unsupported"):
        state.metadata["text_extract"] = {
            "ok": False,
            "type": state.text_extract_type,
            "error": state.text_error,
        }
        return 0, 1
    return 0, 0


def _hydrate_uclass_material_summary_input(
    settings: Settings,
    *,
    state: _UClassMaterialState,
) -> None:
    if (
        state.text_content is None
        and bool(getattr(settings, "material_briefing_enabled", False))
        and isinstance(state.previous_extract, dict)
        and bool(state.previous_extract.get("ok"))
    ):
        state.text_content = str(state.previous_extract.get("excerpt") or "").strip() or None
        state.text_extract_type = str(
            state.previous_extract.get("type") or state.text_extract_type or ""
        ).strip()


def _generate_uclass_material_summary_stage(
    settings: Settings,
    db: Database,
    *,
    state: _UClassMaterialState,
) -> tuple[int, dict[str, Any] | None]:
    should_generate_brief = (
        bool(getattr(settings, "material_briefing_enabled", False))
        and bool(state.local_path)
        and (
            state.downloaded
            or state.content_changed
            or not isinstance(state.previous_brief, dict)
        )
    )
    if not should_generate_brief or not state.text_content:
        return 0, None
    state.metadata["brief"] = _build_material_brief(
        settings=settings,
        db=db,
        title=state.material.filename,
        extracted_text=state.text_content,
        local_path=state.local_path,
        artifact_provenance_source=state.material_provenance_source,
    )
    brief_payload = state.metadata.get("brief")
    if not isinstance(brief_payload, dict):
        return 1, None
    return (
        1,
        {
            "external_id": state.material.external_id,
            "filename": str(state.resolved_filename or state.material.filename or "material"),
            "course_name": str(state.metadata.get("course_name") or "").strip(),
            "extract_type": str(state.text_extract_type or "").strip().lower(),
            "text_excerpt": str(state.text_content or "")[:500],
            "mode": str(brief_payload.get("mode") or "unknown"),
            "bullets": list(brief_payload.get("bullets") or []),
            "question": str(brief_payload.get("question") or "").strip(),
        },
    )


def _should_scan_uclass_material_tasks(state: _UClassMaterialState) -> bool:
    return bool(state.local_path) and (
        state.downloaded
        or not isinstance(state.previous_deadline_scan, dict)
        or not bool(state.previous_deadline_scan.get("ok", True))
        or int(state.previous_deadline_scan.get("version") or 0) < MATERIAL_DEADLINE_SCAN_VERSION
    )


def _ensure_uclass_material_text_for_task_stage(
    settings: Settings,
    *,
    state: _UClassMaterialState,
) -> tuple[int, int]:
    if state.text_content is not None or not bool(getattr(settings, "material_extraction_enabled", False)):
        return 0, 0
    if not state.local_path:
        return 0, 0
    state.text_content, state.text_error, state.text_extract_type = extract_material_text(
        Path(state.local_path),
        max_chars=settings.material_extract_max_chars,
    )
    if state.text_content and not isinstance(state.metadata.get("text_extract"), dict):
        text_hash = sha1(state.text_content.encode("utf-8")).hexdigest()
        state.metadata["text_extract"] = {
            "ok": True,
            "type": state.text_extract_type,
            "hash": text_hash,
            "chars": len(state.text_content),
            "excerpt": state.text_content[: min(2000, len(state.text_content))],
        }
        return 1, 0
    if (
        state.text_error
        and not state.text_error.startswith("unsupported")
        and not isinstance(state.metadata.get("text_extract"), dict)
    ):
        state.metadata["text_extract"] = {
            "ok": False,
            "type": state.text_extract_type,
            "error": state.text_error,
        }
        return 0, 1
    return 0, 0


def _extract_uclass_material_tasks_stage(
    settings: Settings,
    db: Database,
    *,
    state: _UClassMaterialState,
) -> int:
    scan = _build_material_deadline_scan(
        settings=settings,
        db=db,
        artifact_external_id=state.material.external_id,
        title=str(state.resolved_filename or state.material.filename or "material"),
        course_name=str(state.metadata.get("course_name") or "").strip(),
        canonical_course_id=str(state.metadata.get("canonical_course_id") or "").strip() or None,
        extracted_text=str(state.text_content or ""),
        local_path=state.local_path,
        reference_local=datetime.now(ZoneInfo(settings.timezone)),
        artifact_provenance_source=state.material_provenance_source,
        artifact_evidence_links=[state.material.url] if state.material.url else None,
    )
    state.metadata["deadline_scan"] = scan
    return int(scan.get("count") or 0)


def _record_uclass_material_state(db: Database, *, state: _UClassMaterialState) -> None:
    db.record_artifact(
        external_id=state.material.external_id,
        source="uclass",
        filename=state.resolved_filename,
        icloud_path=state.local_path,
        content_hash=state.content_hash,
        metadata_json=state.metadata,
        user_id=state.owner_id,
    )


def _sync_uclass_materials(
    settings: Settings,
    db: Database,
    *,
    owner_id: int,
    candidates: list[Any],
    alias_map: dict[str, tuple[str, ...]],
    canonical_courses: dict[int, str],
) -> _UClassMaterialSyncResult:
    result = _UClassMaterialSyncResult()
    materials_root = _uclass_materials_root(settings, owner_id=owner_id)
    for material in candidates:
        state = _initialize_uclass_material_state(
            db,
            owner_id=owner_id,
            material=material,
            alias_map=alias_map,
            canonical_courses=canonical_courses,
        )
        downloaded_delta, reused_delta, failed_delta = _download_uclass_material_stage(
            settings,
            db,
            state=state,
            materials_root=materials_root,
        )
        result.downloaded_count += downloaded_delta
        result.reused_count += reused_delta
        result.failed_downloads += failed_delta
        _attach_uclass_material_provenance(state)

        extracted_delta, extraction_failure_delta = _extract_uclass_material_stage(
            settings,
            state=state,
        )
        result.extracted_count += extracted_delta
        result.extraction_failures += extraction_failure_delta
        _hydrate_uclass_material_summary_input(settings, state=state)

        brief_delta, brief_item = _generate_uclass_material_summary_stage(
            settings,
            db,
            state=state,
        )
        result.brief_count += brief_delta
        if brief_item is not None:
            result.generated_brief_items.append(brief_item)

        if _should_scan_uclass_material_tasks(state):
            extracted_delta, extraction_failure_delta = _ensure_uclass_material_text_for_task_stage(
                settings,
                state=state,
            )
            result.extracted_count += extracted_delta
            result.extraction_failures += extraction_failure_delta
            result.material_task_count += _extract_uclass_material_tasks_stage(
                settings,
                db,
                state=state,
            )

        _record_uclass_material_state(db, state=state)
        result.artifact_count += 1
    return result


def _link_uclass_timetable_events(
    db: Database,
    *,
    owner_id: int,
    alias_map: dict[str, tuple[str, ...]],
) -> None:
    for event in db.list_events(limit=3000, user_id=owner_id):
        if not _is_timetable_event(event):
            continue
        existing_metadata = _json_load(event.metadata_json)
        linked_metadata = _attach_canonical_course_id(
            existing_metadata,
            alias_map=alias_map,
            fallback_aliases=[event.title],
        )
        canonical_course_id = str(linked_metadata.get("canonical_course_id") or "").strip()
        if not canonical_course_id or canonical_course_id == str(existing_metadata.get("canonical_course_id") or "").strip():
            continue
        db.upsert_course_alias(
            canonical_course_id=canonical_course_id,
            alias=event.title,
            alias_type="portal_title",
            source=str(event.source),
            metadata_json={"event_external_id": event.external_id},
            user_id=owner_id,
        )
        db.upsert_event(
            external_id=event.external_id,
            source=event.source,
            start=event.start_at,
            end=event.end_at,
            title=event.title,
            location=event.location,
            rrule=event.rrule,
            metadata_json=linked_metadata,
            user_id=owner_id,
        )


def _sync_uclass_target(
    settings: Settings,
    db: Database,
    *,
    target: dict[str, Any],
    send_material_brief_push_enabled: bool = True,
) -> dict[str, Any]:
    owner_id = _safe_int(target.get("user_id")) or 0
    timezone_name = str(
        target.get("timezone") or getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul"
    )
    auth = _prepare_uclass_target_auth(
        settings,
        target=target,
        owner_id=owner_id,
    )
    ws_data = _fetch_uclass_ws_stage(
        settings,
        db,
        owner_id=owner_id,
        auth=auth,
    )
    notif_count, task_count, event_count = _persist_uclass_ws_records(
        db,
        owner_id=owner_id,
        ws_data=ws_data,
    )
    candidates = _discover_uclass_material_candidates(
        ws_data,
        timezone_name=timezone_name,
        html_material_candidates=auth.html_material_candidates,
    )
    material_result = _sync_uclass_materials(
        settings,
        db,
        owner_id=owner_id,
        candidates=candidates,
        alias_map=ws_data.alias_map,
        canonical_courses=ws_data.canonical_courses,
    )

    if send_material_brief_push_enabled:
        brief_push_kwargs: dict[str, Any] = {
            "settings": settings,
            "db": db,
            "generated_brief_items": material_result.generated_brief_items,
        }
        target_chat_id = str(target.get("chat_id") or "").strip()
        if target_chat_id:
            brief_push_kwargs["chat_ids"] = [target_chat_id]
        brief_push_result = send_material_brief_push(**brief_push_kwargs)
    else:
        brief_push_result = {
            "skipped": True,
            "reason": "admin refresh disabled material brief push",
        }

    _link_uclass_timetable_events(
        db,
        owner_id=owner_id,
        alias_map=ws_data.alias_map,
    )

    semantic_warning_count = int(ws_data.semantic_warnings or 0)
    result = {
        "user_id": owner_id,
        "chat_id": str(target.get("chat_id") or "").strip() or None,
        "school_slug": str(target.get("school_slug") or "").strip() or None,
        "display_name": str(target.get("display_name") or "").strip() or None,
        "upserted_notifications": notif_count,
        "upserted_tasks": task_count,
        "upserted_events": event_count,
        "recorded_artifacts": material_result.artifact_count,
        "downloaded_artifacts": material_result.downloaded_count,
        "reused_artifacts": material_result.reused_count,
        "failed_artifact_downloads": material_result.failed_downloads,
        "extracted_artifacts": material_result.extracted_count,
        "failed_artifact_extractions": material_result.extraction_failures,
        "generated_material_briefs": material_result.brief_count,
        "detected_material_tasks": material_result.material_task_count,
        "material_brief_push": brief_push_result,
        "html_material_candidates": len(auth.html_material_candidates),
        "html_material_error": auth.html_material_error,
        "semantic_warnings": ws_data.semantic_warnings,
        "wsfunctions": auth.ws_status,
    }
    _record_sync_dashboard_state(
        db,
        "sync_uclass",
        status="success",
        new_items=notif_count + task_count + event_count + material_result.artifact_count,
        action_required=(
            material_result.failed_downloads
            + material_result.extraction_failures
            + semantic_warning_count
        ),
        cursor_payload={
            "site": ws_data.site_info.get("sitename", ""),
            "notifications": notif_count,
            "tasks": task_count,
            "events": event_count,
            "artifacts": material_result.artifact_count,
            "material_downloaded": material_result.downloaded_count,
            "material_reused": material_result.reused_count,
            "material_download_failures": material_result.failed_downloads,
            "material_extracted": material_result.extracted_count,
            "material_extract_failures": material_result.extraction_failures,
            "material_briefs": material_result.brief_count,
            "material_tasks": material_result.material_task_count,
            "material_brief_push": brief_push_result,
            "html_material_candidates": len(auth.html_material_candidates),
            "html_material_error": auth.html_material_error,
            "semantic_warnings": ws_data.semantic_warnings,
            "wsfunctions": auth.ws_status,
        },
        user_id=owner_id,
    )
    return result


def sync_uclass(settings: Settings, db: Database) -> dict[str, Any]:
    targets = _resolve_uclass_sync_targets(settings, db)
    if not targets:
        logger.info("skipping uclass sync; config missing")
        reason = (
            "No active moodle_connections and UCLASS_WSTOKEN missing"
            if str(getattr(settings, "uclass_ws_base", "") or "").strip()
            else "No active moodle_connections and UCLASS_WS_BASE missing"
        )
        _record_sync_dashboard_state(
            db,
            "sync_uclass",
            status="skipped",
            last_error=reason,
            cursor_payload={"skipped": True, "reason": reason},
        )
        return {"skipped": True, "reason": reason}

    aggregate = {
        "upserted_notifications": 0,
        "upserted_tasks": 0,
        "upserted_events": 0,
        "recorded_artifacts": 0,
        "html_material_candidates": 0,
        "downloaded_artifacts": 0,
        "reused_artifacts": 0,
        "failed_artifact_downloads": 0,
        "extracted_artifacts": 0,
        "failed_artifact_extractions": 0,
        "generated_material_briefs": 0,
        "detected_material_tasks": 0,
    }
    target_results: list[dict[str, Any]] = []
    failures: list[dict[str, Any]] = []
    for target in targets:
        owner_id = _safe_int(target.get("user_id")) or 0
        try:
            result = _sync_uclass_target(settings=settings, db=db, target=target)
        except Exception as exc:
            failure = {
                "user_id": owner_id,
                "chat_id": str(target.get("chat_id") or "").strip() or None,
                "school_slug": str(target.get("school_slug") or "").strip() or None,
                "error": str(exc),
            }
            failures.append(failure)
            logger.warning("uclass sync failed for target", extra=failure)
            _record_sync_dashboard_state(
                db,
                "sync_uclass",
                status="error",
                action_required=1,
                last_error=str(exc),
                cursor_payload={"skipped": True, "reason": str(exc), "target": failure},
                user_id=owner_id,
            )
            continue
        target_results.append(result)
        for key in aggregate:
            aggregate[key] += int(result.get(key) or 0)

    overall_status = "error" if failures else "success" if target_results else "skipped"
    overall_error = str(failures[0].get("error") or "") if failures else None
    semantic_warning_count = sum(
        len(item.get("semantic_warnings") or [])
        if isinstance(item.get("semantic_warnings"), list)
        else int(item.get("semantic_warnings") or 0)
        for item in target_results
    )
    _record_sync_dashboard_state(
        db,
        "sync_uclass",
        status=overall_status,
        new_items=(
            aggregate["upserted_notifications"]
            + aggregate["upserted_tasks"]
            + aggregate["upserted_events"]
            + aggregate["recorded_artifacts"]
        ),
        action_required=(
            aggregate["failed_artifact_downloads"]
            + aggregate["failed_artifact_extractions"]
            + semantic_warning_count
            + len(failures)
        ),
        last_error=overall_error,
        cursor_payload={
            **aggregate,
            "skipped": not bool(target_results),
            "reason": overall_error if not target_results else None,
            "targets": [
                {
                    "user_id": int(item.get("user_id") or 0),
                    "chat_id": item.get("chat_id"),
                    "school_slug": item.get("school_slug"),
                    "display_name": item.get("display_name"),
                    "semantic_warnings": item.get("semantic_warnings"),
                    "wsfunctions": item.get("wsfunctions"),
                }
                for item in target_results
            ],
            "target_count": len(targets),
            "synced_targets": len(target_results),
            "failed_targets": failures,
        },
    )
    if not target_results and failures:
        return {
            "ok": False,
            "error": overall_error or UCLASS_RECONNECT_REQUIRED_REASON,
            "failed_targets": failures,
        }
    if len(target_results) == 1 and not failures:
        return {
            **target_results[0],
            "targets": target_results,
            "failed_targets": failures,
        }
    return {
        **aggregate,
        "targets": target_results,
        "failed_targets": failures,
    }


def _weather_target_cache_key(target: dict[str, Any]) -> str:
    payload = {
        "lat": round(float(target.get("lat") or 0.0), 6),
        "lon": round(float(target.get("lon") or 0.0), 6),
        "air_quality_district_codes": [
            str(code).strip()
            for code in list(target.get("air_quality_district_codes") or [])
            if str(code).strip()
        ],
    }
    return sha1(json.dumps(payload, ensure_ascii=True, sort_keys=True).encode("utf-8")).hexdigest()


def _weather_target_payload(target: dict[str, Any]) -> dict[str, Any]:
    payload = {
        "location_label": str(target.get("location_label") or "").strip() or "기준 위치",
        "lat": float(target.get("lat") or 0.0),
        "lon": float(target.get("lon") or 0.0),
        "air_quality_district_codes": [
            str(code).strip()
            for code in list(target.get("air_quality_district_codes") or [])
            if str(code).strip()
        ],
        "source": str(target.get("source") or "").strip() or "unknown",
    }
    payload["cache_key"] = str(target.get("cache_key") or _weather_target_cache_key(payload))
    return payload


def _default_weather_target(settings: Settings) -> dict[str, Any]:
    target = {
        "location_label": str(getattr(settings, "weather_location_label", "") or "").strip() or "서울특별시",
        "lat": float(getattr(settings, "weather_lat", 37.583801) or 37.583801),
        "lon": float(getattr(settings, "weather_lon", 127.058701) or 127.058701),
        "air_quality_district_codes": [
            str(code).strip()
            for code in list(getattr(settings, "air_quality_district_codes", []) or [])
            if str(code).strip()
        ],
        "source": "global_default",
    }
    target["cache_key"] = _weather_target_cache_key(target)
    return target


def _user_weather_target(
    settings: Settings,
    db: Database,
    *,
    user_id: int | None = None,
    chat_id: str | None = None,
) -> dict[str, Any]:
    preferences = db.get_user_preferences(user_id=user_id, chat_id=chat_id)
    if not isinstance(preferences, dict):
        return _default_weather_target(settings)
    location_label = str(preferences.get("weather_location_label") or "").strip()
    lat = preferences.get("weather_lat")
    lon = preferences.get("weather_lon")
    if not location_label or lat is None or lon is None:
        return _default_weather_target(settings)
    district_code = str(preferences.get("weather_air_quality_district_code") or "").strip()
    target = {
        "location_label": location_label,
        "lat": float(lat),
        "lon": float(lon),
        "air_quality_district_codes": [district_code] if district_code else [],
        "source": "user_preference",
    }
    target["cache_key"] = _weather_target_cache_key(target)
    return target


def _weather_snapshot_matches_target(
    snapshot: dict[str, Any] | None,
    target: dict[str, Any],
) -> bool:
    if not isinstance(snapshot, dict):
        return False
    target_payload = _weather_target_payload(target)
    snapshot_target = snapshot.get("target") if isinstance(snapshot.get("target"), dict) else {}
    snapshot_cache_key = str(snapshot_target.get("cache_key") or "").strip()
    if snapshot_cache_key:
        return snapshot_cache_key == target_payload["cache_key"]
    if target_payload["source"] != "global_default":
        return False
    return True


def _latest_weather_snapshot(
    db: Database,
    *,
    user_id: int | None = None,
    allow_global_fallback: bool = True,
) -> dict[str, Any] | None:
    snapshot = db.latest_weather_snapshot(
        user_id=user_id,
        allow_global_fallback=allow_global_fallback,
    )
    if not isinstance(snapshot, dict):
        return None
    return snapshot


def _latest_matching_weather_snapshot(
    settings: Settings,
    db: Database,
    *,
    target: dict[str, Any],
    user_id: int | None = None,
) -> dict[str, Any] | None:
    owner_id = _safe_int(user_id) or 0
    if owner_id > 0:
        user_snapshot = _latest_weather_snapshot(
            db,
            user_id=owner_id,
            allow_global_fallback=False,
        )
        if _weather_snapshot_matches_target(user_snapshot, target):
            return user_snapshot
    if str(target.get("source") or "").strip() == "global_default":
        global_snapshot = _latest_weather_snapshot(
            db,
            user_id=0,
            allow_global_fallback=False,
        )
        if _weather_snapshot_matches_target(global_snapshot, target):
            return global_snapshot
    return None


def _weather_snapshot_reference_at(snapshot: dict[str, Any] | None) -> datetime | None:
    if not isinstance(snapshot, dict):
        return None
    for key in ("generated_at", "observed_at", "last_run_at"):
        parsed = _parse_dt(str(snapshot.get(key) or "").strip())
        if parsed is not None:
            return parsed
    return None


def _weather_snapshot_is_fresh(
    snapshot: dict[str, Any] | None,
    *,
    now_local: datetime,
    max_age_minutes: int = WEATHER_SNAPSHOT_MAX_AGE_MINUTES,
) -> bool:
    reference_at = _weather_snapshot_reference_at(snapshot)
    if reference_at is None:
        return False
    age = now_local.astimezone(timezone.utc) - reference_at.astimezone(timezone.utc)
    return age <= timedelta(minutes=max(int(max_age_minutes), 1))


def _weather_state_owner_id(target: dict[str, Any], user_id: int | None) -> int | None:
    owner_id = _safe_int(user_id)
    if owner_id is not None and owner_id > 0 and str(target.get("source") or "").strip() == "user_preference":
        return owner_id
    return 0


def _fetch_weather_snapshot_for_target(
    settings: Settings,
    *,
    target: dict[str, Any],
    timezone_name: str,
    weather_client: KMAWeatherClient,
    air_client: SeoulAirQualityClient | None = None,
    now_local: datetime | None = None,
) -> tuple[dict[str, Any], str]:
    snapshot = weather_client.fetch_snapshot(
        lat=float(target["lat"]),
        lon=float(target["lon"]),
        location_label=str(target.get("location_label") or "").strip() or "기준 위치",
        timezone_name=timezone_name,
        now_local=now_local,
    )
    snapshot["target"] = _weather_target_payload(target)

    air_error = ""
    district_codes = [
        str(code).strip()
        for code in list(target.get("air_quality_district_codes") or [])
        if str(code).strip()
    ]
    if bool(getattr(settings, "air_quality_enabled", False)) and district_codes:
        try:
            resolved_air_client = air_client or SeoulAirQualityClient(
                api_key=getattr(settings, "air_quality_seoul_api_key", None)
            )
            snapshot["air_quality"] = {
                "ok": True,
                **resolved_air_client.fetch_snapshot(
                    district_codes=district_codes,
                    timezone_name=timezone_name,
                ),
            }
        except Exception as exc:
            air_error = str(exc)
            snapshot["air_quality"] = {
                "ok": False,
                "error": air_error,
                "districts": [],
            }
    elif bool(getattr(settings, "air_quality_enabled", False)):
        snapshot["air_quality"] = {
            "ok": False,
            "skipped": True,
            "reason": "no district mapping for selected weather location",
            "districts": [],
        }
    else:
        snapshot["air_quality"] = {
            "ok": False,
            "skipped": True,
            "reason": "AIR_QUALITY_ENABLED is false",
            "districts": [],
        }
    return snapshot, air_error


def _get_or_refresh_weather_snapshot(
    settings: Settings,
    db: Database,
    *,
    user_id: int | None = None,
    chat_id: str | None = None,
    now_local: datetime | None = None,
    force_refresh: bool = False,
    snapshot_cache: dict[str, dict[str, Any] | None] | None = None,
    weather_client: KMAWeatherClient | None = None,
    air_client: SeoulAirQualityClient | None = None,
) -> dict[str, Any] | None:
    if not bool(getattr(settings, "weather_enabled", True)):
        return None
    timezone_name = str(getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul")
    tz = ZoneInfo(timezone_name)
    anchor_local = now_local.astimezone(tz) if now_local is not None else datetime.now(tz)
    current_local = datetime.now(tz)
    fetch_local = anchor_local if anchor_local <= current_local else current_local
    target = _user_weather_target(settings, db, user_id=user_id, chat_id=chat_id)
    cache_key = str(target.get("cache_key") or _weather_target_cache_key(target))
    if snapshot_cache is not None and cache_key in snapshot_cache:
        return snapshot_cache[cache_key]

    existing_snapshot = _latest_matching_weather_snapshot(
        settings,
        db,
        target=target,
        user_id=user_id,
    )
    if not force_refresh and existing_snapshot is not None:
        if snapshot_cache is not None:
            snapshot_cache[cache_key] = existing_snapshot
        return existing_snapshot

    resolved_weather_client = weather_client or KMAWeatherClient(
        auth_key=getattr(settings, "weather_kma_auth_key", None)
    )
    state_owner_id = _weather_state_owner_id(target, user_id)
    try:
        snapshot, air_error = _fetch_weather_snapshot_for_target(
            settings,
            target=target,
            timezone_name=timezone_name,
            weather_client=resolved_weather_client,
            air_client=air_client,
            now_local=fetch_local,
        )
    except Exception as exc:
        if existing_snapshot is None:
            if snapshot_cache is not None:
                snapshot_cache[cache_key] = None
            return None
        error_payload = dict(existing_snapshot)
        error_payload.update(
            {
                "error": str(exc),
                "location_label": str(target.get("location_label") or "").strip() or "기준 위치",
                "target": _weather_target_payload(target),
            }
        )
        _record_sync_dashboard_state(
            db,
            "sync_weather",
            status="error",
            action_required=1,
            last_error=str(exc),
            cursor_payload=error_payload,
            user_id=state_owner_id,
        )
        if snapshot_cache is not None:
            snapshot_cache[cache_key] = error_payload
        return error_payload

    _record_sync_dashboard_state(
        db,
        "sync_weather",
        status="success",
        new_items=1,
        action_required=1 if air_error else 0,
        cursor_payload=snapshot,
        user_id=state_owner_id,
    )
    if snapshot_cache is not None:
        snapshot_cache[cache_key] = snapshot
    return snapshot


def sync_weather(settings: Settings, db: Database) -> dict[str, Any]:
    if not bool(getattr(settings, "weather_enabled", True)):
        reason = "WEATHER_ENABLED is false"
        _record_sync_dashboard_state(
            db,
            "sync_weather",
            status="skipped",
            last_error=reason,
            cursor_payload={"skipped": True, "reason": reason},
        )
        return {"skipped": True, "reason": reason}

    timezone_name = str(getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul")
    weather_client = KMAWeatherClient(auth_key=getattr(settings, "weather_kma_auth_key", None))
    air_client = (
        SeoulAirQualityClient(api_key=getattr(settings, "air_quality_seoul_api_key", None))
        if bool(getattr(settings, "air_quality_enabled", False))
        else None
    )
    snapshot_cache: dict[str, dict[str, Any] | None] = {}
    default_target = _default_weather_target(settings)

    try:
        snapshot, air_error = _fetch_weather_snapshot_for_target(
            settings,
            target=default_target,
            timezone_name=timezone_name,
            weather_client=weather_client,
            air_client=air_client,
        )
    except Exception as exc:
        previous_snapshot = _latest_matching_weather_snapshot(
            settings,
            db,
            target=default_target,
            user_id=None,
        ) or {}
        error_payload = dict(previous_snapshot)
        error_payload.update(
            {
                "error": str(exc),
                "location_label": str(default_target.get("location_label") or "").strip() or "서울특별시",
                "target": _weather_target_payload(default_target),
            }
        )
        _record_sync_dashboard_state(
            db,
            "sync_weather",
            status="error",
            action_required=1,
            last_error=str(exc),
            cursor_payload=error_payload,
        )
        return {"ok": False, "error": str(exc), "location_label": default_target["location_label"]}

    snapshot_cache[default_target["cache_key"]] = snapshot
    _record_sync_dashboard_state(
        db,
        "sync_weather",
        status="success",
        new_items=1,
        action_required=1 if air_error else 0,
        cursor_payload=snapshot,
    )

    user_targets: dict[str, dict[str, Any]] = {}
    for row in db.list_user_weather_locations(limit=5000):
        owner_id = _safe_int(row.get("user_id"))
        if owner_id is None or owner_id <= 0:
            continue
        target = {
            "location_label": str(row.get("weather_location_label") or "").strip() or "기준 위치",
            "lat": float(row.get("weather_lat") or 0.0),
            "lon": float(row.get("weather_lon") or 0.0),
            "air_quality_district_codes": [
                str(row.get("weather_air_quality_district_code") or "").strip()
            ]
            if str(row.get("weather_air_quality_district_code") or "").strip()
            else [],
            "source": "user_preference",
        }
        target["cache_key"] = _weather_target_cache_key(target)
        bucket = user_targets.setdefault(
            target["cache_key"],
            {"target": target, "user_ids": []},
        )
        bucket["user_ids"].append(owner_id)

    warmed_targets = 0
    reused_targets = 0
    failed_targets: list[dict[str, Any]] = []
    for item in user_targets.values():
        target = item["target"] if isinstance(item.get("target"), dict) else {}
        owner_ids = [
            int(owner_id)
            for owner_id in list(item.get("user_ids") or [])
            if _safe_int(owner_id) is not None and int(owner_id) > 0
        ]
        if not owner_ids:
            continue
        cache_key = str(target.get("cache_key") or "")
        cached_snapshot = snapshot_cache.get(cache_key) if cache_key else None
        if cached_snapshot is not None:
            reused_targets += 1
            for owner_id in owner_ids:
                _record_sync_dashboard_state(
                    db,
                    "sync_weather",
                    status="success",
                    new_items=1,
                    action_required=1
                    if bool(
                        isinstance(cached_snapshot.get("air_quality"), dict)
                        and cached_snapshot["air_quality"].get("ok") is False
                        and cached_snapshot["air_quality"].get("error")
                    )
                    else 0,
                    cursor_payload=cached_snapshot,
                    user_id=owner_id,
                )
            continue
        try:
            user_snapshot, user_air_error = _fetch_weather_snapshot_for_target(
                settings,
                target=target,
                timezone_name=timezone_name,
                weather_client=weather_client,
                air_client=air_client,
            )
        except Exception as exc:
            failed_targets.append(
                {
                    "location_label": str(target.get("location_label") or "").strip() or "기준 위치",
                    "user_ids": owner_ids,
                    "error": str(exc),
                }
            )
            for owner_id in owner_ids:
                previous_snapshot = _latest_matching_weather_snapshot(
                    settings,
                    db,
                    target=target,
                    user_id=owner_id,
                ) or {}
                error_payload = dict(previous_snapshot)
                error_payload.update(
                    {
                        "error": str(exc),
                        "location_label": str(target.get("location_label") or "").strip() or "기준 위치",
                        "target": _weather_target_payload(target),
                    }
                )
                _record_sync_dashboard_state(
                    db,
                    "sync_weather",
                    status="error",
                    action_required=1,
                    last_error=str(exc),
                    cursor_payload=error_payload,
                    user_id=owner_id,
                )
            continue
        snapshot_cache[cache_key] = user_snapshot
        warmed_targets += 1
        for owner_id in owner_ids:
            _record_sync_dashboard_state(
                db,
                "sync_weather",
                status="success",
                new_items=1,
                action_required=1 if user_air_error else 0,
                cursor_payload=user_snapshot,
                user_id=owner_id,
            )

    current = snapshot.get("current") if isinstance(snapshot.get("current"), dict) else {}
    return {
        "ok": True,
        "location_label": default_target["location_label"],
        "observed_at": snapshot.get("observed_at"),
        "temperature_c": current.get("temperature_c"),
        "condition_text": current.get("condition_text"),
        "air_quality_ok": bool(
            isinstance(snapshot.get("air_quality"), dict)
            and snapshot["air_quality"].get("ok")
        ),
        "air_quality_error": air_error or None,
        "user_target_count": len(user_targets),
        "warmed_user_targets": warmed_targets,
        "reused_user_targets": reused_targets,
        "failed_user_targets": failed_targets,
    }


def _weather_day_for_date(
    snapshot: dict[str, Any] | None,
    *,
    target_day_local: datetime,
) -> dict[str, Any] | None:
    if not isinstance(snapshot, dict):
        return None
    for key in ("today", "tomorrow"):
        payload = snapshot.get(key)
        if isinstance(payload, dict) and str(payload.get("date") or "") == target_day_local.date().isoformat():
            return payload
    return None


def _safe_float(value: Any) -> float | None:
    if value in (None, ""):
        return None
    try:
        return float(value)
    except Exception:
        return None


def _format_weather_temperature(value: Any) -> str:
    number = _safe_float(value)
    if number is None:
        return "?"
    if abs(number - round(number)) < 0.05:
        return str(int(round(number)))
    return f"{number:.1f}".rstrip("0").rstrip(".")


def _format_weather_range(min_value: Any, max_value: Any) -> str:
    low = _safe_float(min_value)
    high = _safe_float(max_value)
    if low is None and high is None:
        return "?C"
    if low is None:
        return f"{_format_weather_temperature(high)}C"
    if high is None:
        return f"{_format_weather_temperature(low)}C"
    if abs(low - high) < 0.05:
        return f"{_format_weather_temperature(low)}C"
    return f"{_format_weather_temperature(low)}~{_format_weather_temperature(high)}C"


def _format_air_quality_district_line(row: dict[str, Any] | None) -> str | None:
    if not isinstance(row, dict):
        return None
    district_name = str(row.get("district_name") or row.get("district_code") or "").strip()
    if not district_name:
        return None
    pieces = [district_name]
    has_measurement = False
    cai_grade = str(row.get("cai_grade") or "").strip()
    if cai_grade:
        pieces.append(cai_grade)
        has_measurement = True
    cai = _safe_int(row.get("cai"))
    if cai is not None:
        pieces.append(f"CAI {cai}")
        has_measurement = True
    pm10 = _safe_int(row.get("pm10"))
    pm25 = _safe_int(row.get("pm25"))
    if pm10 is not None:
        pieces.append(f"PM10 {pm10}")
        has_measurement = True
    if pm25 is not None:
        pieces.append(f"PM2.5 {pm25}")
        has_measurement = True
    dominant = str(row.get("dominant_pollutant") or "").strip()
    if dominant:
        pieces.append(f"주오염 {dominant}")
        has_measurement = True
    if not has_measurement:
        pieces.append("측정값 미수신")
    return " ".join(pieces)


def _pm25_grade_label(value: Any) -> str | None:
    pm25 = _safe_int(value)
    if pm25 is None:
        return None
    if pm25 <= 15:
        return "좋음"
    if pm25 <= 35:
        return "보통"
    if pm25 <= 75:
        return "나쁨"
    return "매우나쁨"


def _air_quality_grade_severity(value: Any) -> int:
    label = str(value or "").strip()
    return {
        "좋음": 0,
        "보통": 1,
        "나쁨": 2,
        "매우나쁨": 3,
    }.get(label, -1)


def _worst_air_quality_grade(row: dict[str, Any] | None) -> str | None:
    if not isinstance(row, dict):
        return None
    candidates = [
        str(row.get("cai_grade") or "").strip() or None,
        _pm25_grade_label(row.get("pm25")),
    ]
    best_label = None
    best_severity = -1
    for label in candidates:
        severity = _air_quality_grade_severity(label)
        if severity > best_severity:
            best_label = label
            best_severity = severity
    return best_label


def _build_briefing_weather_lines(
    settings: Settings,
    db: Database,
    *,
    now_local: datetime,
    user_id: int | None = None,
    snapshot_cache: dict[str, dict[str, Any] | None] | None = None,
) -> list[str]:
    snapshot = _get_or_refresh_weather_snapshot(
        settings,
        db,
        user_id=user_id,
        now_local=now_local,
        snapshot_cache=snapshot_cache,
    )
    return _build_weather_message_lines(snapshot if isinstance(snapshot, dict) else None, now_local=now_local)


def _weather_precip_probability_max(day: dict[str, Any] | None) -> int | None:
    if not isinstance(day, dict):
        return None
    values: list[int] = []
    direct_value = _safe_int(day.get("precip_probability_max"))
    if direct_value is not None:
        values.append(direct_value)
    for key in ("morning", "afternoon"):
        daypart = day.get(key)
        if not isinstance(daypart, dict):
            continue
        value = _safe_int(daypart.get("precip_probability_max"))
        if value is not None:
            values.append(value)
    return max(values) if values else None


def _format_weather_command_current_line(current: dict[str, Any] | None) -> str | None:
    if not isinstance(current, dict):
        return None
    temperature_c = current.get("temperature_c")
    condition_text = str(current.get("condition_text") or "").strip()
    if temperature_c is None and not condition_text:
        return None
    pieces = []
    if temperature_c is not None:
        pieces.append(f"현재 {_format_weather_temperature(temperature_c)}C")
    else:
        pieces.append("현재 정보 없음")
    if condition_text:
        pieces.append(condition_text)
    return "- " + " / ".join(pieces)


def _format_weather_command_daypart_line(label: str, daypart: dict[str, Any] | None) -> str | None:
    if not isinstance(daypart, dict):
        return None
    parts = [
        f"{label} : {_format_weather_range(daypart.get('temperature_min_c'), daypart.get('temperature_max_c'))}",
        str(daypart.get("condition_text") or "정보 없음").strip() or "정보 없음",
    ]
    precip_probability = _safe_int(daypart.get("precip_probability_max"))
    if precip_probability is not None:
        parts.append(f"강수확률 {precip_probability}%")
    return "- " + ", ".join(parts)


def _format_weather_command_air_quality_line(row: dict[str, Any] | None) -> str | None:
    if not isinstance(row, dict):
        return None
    district_name = str(row.get("district_name") or row.get("district_code") or "").strip()
    if not district_name:
        return None
    grade = _worst_air_quality_grade(row)
    if not grade:
        return f"- {district_name} 측정값 미수신"
    return f"- {district_name} {grade}"


def _latest_air_quality_measured_at(air_quality: dict[str, Any] | None) -> str | None:
    if not isinstance(air_quality, dict):
        return None
    measured_at = str(air_quality.get("measured_at") or "").strip()
    if measured_at:
        return measured_at
    districts = air_quality.get("districts") if isinstance(air_quality.get("districts"), list) else []
    measured_candidates = [
        str(item.get("measured_at") or "").strip()
        for item in districts
        if isinstance(item, dict) and str(item.get("measured_at") or "").strip()
    ]
    return max(measured_candidates) if measured_candidates else None


def _format_weather_command_air_quality_title(
    air_quality: dict[str, Any] | None,
    *,
    now_local: datetime,
) -> str:
    parsed = _parse_dt(_latest_air_quality_measured_at(air_quality))
    if parsed is None:
        return "미세먼지"
    measured_time = parsed.astimezone(now_local.tzinfo or timezone.utc).strftime("%H:%M")
    return f"미세먼지 ({measured_time} 기준)"


def _format_weather_command_today_title(
    snapshot: dict[str, Any] | None,
    *,
    now_local: datetime,
) -> str:
    if not isinstance(snapshot, dict):
        return "오늘 날씨"
    for key in ("observed_at", "generated_at"):
        parsed = _parse_dt(str(snapshot.get(key) or "").strip())
        if parsed is None:
            continue
        measured_time = parsed.astimezone(now_local.tzinfo or timezone.utc).strftime("%H:%M")
        return f"오늘 날씨 ({measured_time} 기준)"
    return "오늘 날씨"


def _format_weather_command_tomorrow_title(
    snapshot: dict[str, Any] | None,
    *,
    now_local: datetime,
) -> str:
    if not isinstance(snapshot, dict):
        return "내일 날씨"
    for key in ("generated_at", "observed_at"):
        parsed = _parse_dt(str(snapshot.get(key) or "").strip())
        if parsed is None:
            continue
        measured_time = parsed.astimezone(now_local.tzinfo or timezone.utc).strftime("%H:%M")
        return f"내일 날씨 ({measured_time} 기준)"
    return "내일 날씨"


def _format_weather_command_location_line(snapshot: dict[str, Any] | None) -> str | None:
    if not isinstance(snapshot, dict):
        return None
    label = str(snapshot.get("location_label") or "").strip()
    if not label:
        return None
    return f"- 지역 {label}"


def _format_weather_command_tomorrow_temperature_line(day: dict[str, Any] | None) -> str | None:
    if not isinstance(day, dict):
        return None
    low = day.get("temperature_min_c")
    high = day.get("temperature_max_c")
    if low is None and high is None:
        return None
    return (
        f"- 최저 {_format_weather_temperature(low)}C / 최고 {_format_weather_temperature(high)}C"
    )


def _format_weather_command_tomorrow_precip_line(day: dict[str, Any] | None) -> str | None:
    precip_probability = _weather_precip_probability_max(day)
    if precip_probability is None:
        return None
    return f"- 강수확률 {precip_probability}%"


def _build_weather_message_lines(
    snapshot: dict[str, Any] | None,
    *,
    now_local: datetime,
) -> list[str]:
    if not isinstance(snapshot, dict):
        return []
    current = snapshot.get("current") if isinstance(snapshot.get("current"), dict) else {}
    today = _weather_day_for_date(snapshot, target_day_local=now_local)
    tomorrow = _weather_day_for_date(snapshot, target_day_local=now_local + timedelta(days=1))
    if not isinstance(today, dict):
        today = snapshot.get("today") if isinstance(snapshot.get("today"), dict) else None
    if not isinstance(tomorrow, dict):
        tomorrow = snapshot.get("tomorrow") if isinstance(snapshot.get("tomorrow"), dict) else None

    lines = [_format_weather_command_today_title(snapshot, now_local=now_local)]
    location_line = _format_weather_command_location_line(snapshot)
    if location_line:
        lines.append(location_line)
    current_line = _format_weather_command_current_line(current)
    if current_line:
        lines.append(current_line)
    today_morning = _format_weather_command_daypart_line(
        "오전",
        today.get("morning") if isinstance(today, dict) else None,
    )
    today_afternoon = _format_weather_command_daypart_line(
        "오후",
        today.get("afternoon") if isinstance(today, dict) else None,
    )
    if today_morning:
        lines.append(today_morning)
    if today_afternoon:
        lines.append(today_afternoon)

    air_quality = snapshot.get("air_quality") if isinstance(snapshot.get("air_quality"), dict) else {}
    if bool(air_quality.get("ok")):
        districts = air_quality.get("districts") if isinstance(air_quality.get("districts"), list) else []
        district_lines = [
            _format_weather_command_air_quality_line(item if isinstance(item, dict) else None)
            for item in districts[:2]
        ]
        district_lines = [item for item in district_lines if item]
        if district_lines:
            lines.append("")
            lines.append(_format_weather_command_air_quality_title(
                air_quality,
                now_local=now_local,
            ))
            lines.extend(district_lines)

    tomorrow_lines = [
        _format_weather_command_tomorrow_temperature_line(tomorrow),
        _format_weather_command_tomorrow_precip_line(tomorrow),
    ]
    tomorrow_lines = [item for item in tomorrow_lines if item]
    if tomorrow_lines:
        lines.append("")
        lines.append("내일 날씨")
        lines.extend(tomorrow_lines)
    return lines


def _format_telegram_todayweather(
    settings: Settings,
    db: Database,
    *,
    user_id: int | None = None,
    chat_id: str | None = None,
    snapshot_cache: dict[str, dict[str, Any] | None] | None = None,
) -> str:
    snapshot = _get_or_refresh_weather_snapshot(
        settings,
        db,
        user_id=user_id,
        chat_id=chat_id,
        snapshot_cache=snapshot_cache,
    )
    if snapshot is None:
        preferences = db.get_user_preferences(user_id=user_id, chat_id=chat_id)
        custom_label = (
            str(preferences.get("weather_location_label") or "").strip()
            if isinstance(preferences, dict)
            else ""
        )
        default_target = _default_weather_target(settings)
        location_label = custom_label or str(default_target.get("location_label") or "서울특별시")
        return "\n".join(
            [
                "[KU] 오늘 날씨",
                "",
                "- 아직 날씨 데이터가 준비되지 않았습니다.",
                f"- 현재 지역: {location_label}",
                "- 첫 동기화 전이거나 데이터가 아직 갱신되지 않았습니다.",
                "- 지역 변경: `/region 고려대`",
                "- 기본값 복귀: `/region reset`",
            ]
        )
    timezone_name = str(getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul")
    now_local = datetime.now(ZoneInfo(timezone_name))
    return "\n".join(_build_weather_message_lines(snapshot, now_local=now_local))


def _format_telegram_tomorrowweather(
    settings: Settings,
    db: Database,
    *,
    user_id: int | None = None,
    chat_id: str | None = None,
    snapshot_cache: dict[str, dict[str, Any] | None] | None = None,
) -> str:
    snapshot = _get_or_refresh_weather_snapshot(
        settings,
        db,
        user_id=user_id,
        chat_id=chat_id,
        snapshot_cache=snapshot_cache,
    )
    timezone_name = str(getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul")
    now_local = datetime.now(ZoneInfo(timezone_name))
    if snapshot is None:
        preferences = db.get_user_preferences(user_id=user_id, chat_id=chat_id)
        custom_label = (
            str(preferences.get("weather_location_label") or "").strip()
            if isinstance(preferences, dict)
            else ""
        )
        default_target = _default_weather_target(settings)
        location_label = custom_label or str(default_target.get("location_label") or "서울특별시")
        return "\n".join(
            [
                "[KU] 내일 날씨",
                "",
                "- 아직 내일 날씨 데이터가 준비되지 않았습니다.",
                f"- 현재 지역: {location_label}",
                "- 첫 동기화 전이거나 데이터가 아직 갱신되지 않았습니다.",
                "- 지역 변경: `/region 고려대`",
                "- 기본값 복귀: `/region reset`",
            ]
        )

    tomorrow = _weather_day_for_date(snapshot, target_day_local=now_local + timedelta(days=1))
    if not isinstance(tomorrow, dict):
        tomorrow = snapshot.get("tomorrow") if isinstance(snapshot.get("tomorrow"), dict) else None
    if not isinstance(tomorrow, dict):
        location_line = _format_weather_command_location_line(snapshot)
        lines = [
            _format_weather_command_tomorrow_title(snapshot, now_local=now_local),
        ]
        if location_line:
            lines.append(location_line)
        lines.extend(
            [
                "- 아직 내일 예보가 준비되지 않았습니다.",
                "- 잠시 후 다시 시도해 주세요.",
            ]
        )
        return "\n".join(lines)

    lines = [_format_weather_command_tomorrow_title(snapshot, now_local=now_local)]
    location_line = _format_weather_command_location_line(snapshot)
    if location_line:
        lines.append(location_line)
    temperature_line = _format_weather_command_tomorrow_temperature_line(tomorrow)
    if temperature_line:
        lines.append(temperature_line)
    precip_line = _format_weather_command_tomorrow_precip_line(tomorrow)
    if precip_line:
        lines.append(precip_line)
    morning_line = _format_weather_command_daypart_line(
        "오전",
        tomorrow.get("morning") if isinstance(tomorrow, dict) else None,
    )
    afternoon_line = _format_weather_command_daypart_line(
        "오후",
        tomorrow.get("afternoon") if isinstance(tomorrow, dict) else None,
    )
    if morning_line:
        lines.append(morning_line)
    if afternoon_line:
        lines.append(afternoon_line)
    if len(lines) == 1:
        lines.append("- 아직 내일 예보가 준비되지 않았습니다.")
    return "\n".join(lines)


def _format_weather_region_status(
    settings: Settings,
    db: Database,
    *,
    user_id: int | None = None,
    chat_id: str | None = None,
    error_message: str | None = None,
) -> str:
    preferences = db.get_user_preferences(user_id=user_id, chat_id=chat_id)
    custom_label = (
        str(preferences.get("weather_location_label") or "").strip()
        if isinstance(preferences, dict)
        else ""
    )
    default_target = _default_weather_target(settings)
    lines = ["[KU] 날씨 지역", ""]
    if error_message:
        lines.append(f"- {error_message}")
    if custom_label:
        lines.append(f"- 현재 설정: {custom_label}")
    else:
        lines.append(
            f"- 현재 설정: 기본값 ({str(default_target.get('location_label') or '서울특별시')})"
        )
    lines.extend(
        [
            "- 반영 대상: /weather, 아침/저녁 브리핑",
            "- 변경 예: /region 고려대",
            "- 변경 예: /region 동대문구",
            "- 직접 지정: /region 37.5665,126.9780 서울시청",
            "- 기본값 복귀: /region reset",
        ]
    )
    return "\n".join(lines)


def _handle_telegram_region_command(
    settings: Settings,
    db: Database,
    *,
    query: str | None = None,
    user_id: int | None = None,
    chat_id: str | None = None,
) -> dict[str, Any]:
    owner_id = _safe_int(user_id)
    chat = str(chat_id or "").strip() or None
    if owner_id is None and not chat:
        return {"ok": False, "error": "user scope missing for /region"}

    raw_query = str(query or "").strip()
    if not raw_query:
        return {
            "ok": True,
            "message": _format_weather_region_status(
                settings,
                db,
                user_id=owner_id,
                chat_id=chat,
            ),
        }

    if raw_query.lower() in WEATHER_REGION_RESET_TOKENS or raw_query in WEATHER_REGION_RESET_TOKENS:
        db.upsert_user_preferences(
            user_id=owner_id,
            chat_id=chat,
            weather_location_label=None,
            weather_lat=None,
            weather_lon=None,
            weather_air_quality_district_code=None,
            metadata_json={"weather_region_source": "reset"},
        )
        default_target = _default_weather_target(settings)
        return {
            "ok": True,
            "message": "\n".join(
                [
                    "[KU] 날씨 지역",
                    "",
                    f"- 기본값으로 되돌렸습니다: {default_target['location_label']}",
                    "- 이제 /weather 와 아침/저녁 브리핑에 기본 지역을 사용합니다.",
                ]
            ),
        }

    try:
        resolved = resolve_weather_location_query(raw_query)
    except Exception as exc:
        return {
            "ok": False,
            "message": _format_weather_region_status(
                settings,
                db,
                user_id=owner_id,
                chat_id=chat,
                error_message=f"지역 확인 중 오류가 발생했습니다: {exc}",
            ),
        }
    if not isinstance(resolved, dict):
        return {
            "ok": False,
            "message": _format_weather_region_status(
                settings,
                db,
                user_id=owner_id,
                chat_id=chat,
                error_message=f"'{raw_query}' 위치를 찾지 못했습니다.",
            ),
        }

    label = str(resolved.get("label") or "").strip()
    lat = resolved.get("lat")
    lon = resolved.get("lon")
    if not label or lat is None or lon is None:
        return {
            "ok": False,
            "message": _format_weather_region_status(
                settings,
                db,
                user_id=owner_id,
                chat_id=chat,
                error_message=f"'{raw_query}' 위치를 저장할 수 없습니다.",
            ),
        }

    db.upsert_user_preferences(
        user_id=owner_id,
        chat_id=chat,
        weather_location_label=label,
        weather_lat=float(lat),
        weather_lon=float(lon),
        weather_air_quality_district_code=resolved.get("air_quality_district_code"),
        metadata_json={"weather_region_source": str(resolved.get("source") or "unknown")},
    )
    return {
        "ok": True,
        "message": "\n".join(
            [
                "[KU] 날씨 지역",
                "",
                f"- 저장됨: {label}",
                "- 이제 /weather 와 아침/저녁 브리핑에 이 지역을 사용합니다.",
                "- 기본값으로 돌아가려면 /region reset",
            ]
        ),
    }


def _task_due_today(task_due_at: str | None, timezone_name: str, today_local: datetime) -> bool:
    due_dt = _parse_dt(task_due_at)
    if due_dt is None:
        return False
    return due_dt.astimezone(ZoneInfo(timezone_name)).date() == today_local.date()


def _task_due_on_day(task_due_at: str | None, timezone_name: str, target_day_local: datetime) -> bool:
    return _task_due_today(task_due_at, timezone_name, target_day_local)


def _truncate_telegram_text(value: str, limit: int = 120) -> str:
    text = str(value or "").strip()
    capped = max(int(limit), 1)
    if len(text) <= capped:
        return text
    return text[: max(capped - 1, 1)].rstrip() + "…"


def _extract_material_week_markers(*values: Any) -> tuple[int, ...]:
    output: list[int] = []
    seen: set[int] = set()
    for value in values:
        text = str(value or "")
        if not text:
            continue
        for pattern in MATERIAL_WEEK_MARKER_PATTERNS:
            for match in pattern.finditer(text):
                try:
                    week_no = int(match.group(1))
                except Exception:
                    continue
                if week_no <= 0 or week_no in seen:
                    continue
                output.append(week_no)
                seen.add(week_no)
    return tuple(output)


def _class_occurrence_week_index(
    event: Any,
    *,
    occurrence_start_local: datetime,
    timezone_name: str,
    anchor_start_local: datetime | None = None,
) -> int | None:
    start_dt = _parse_dt(getattr(event, "start_at", None))
    if start_dt is None:
        return None
    tz = ZoneInfo(timezone_name)
    start_local = (
        anchor_start_local.astimezone(tz)
        if isinstance(anchor_start_local, datetime)
        else start_dt.astimezone(tz)
    )
    if occurrence_start_local.date() < start_local.date():
        return None
    rrule_value = str(getattr(event, "rrule", "") or "").strip().upper()
    if rrule_value and "FREQ=WEEKLY" not in rrule_value:
        return None
    delta_days = (occurrence_start_local.date() - start_local.date()).days
    return max(delta_days // 7 + 1, 1)


def _course_week_anchor_start_local(
    db: Database,
    *,
    canonical_course_id: str,
    timezone_name: str,
    user_id: int | None = None,
) -> datetime | None:
    canonical = str(canonical_course_id or "").strip()
    if not canonical:
        return None
    course = db.get_course(canonical, user_id=user_id)
    if course is None:
        return None
    metadata = _json_load(course.metadata_json)
    for key in ("startdate", "course_start_at", "start_at"):
        start_dt = _parse_datetime_like(metadata.get(key))
        if start_dt is not None:
            return start_dt.astimezone(ZoneInfo(timezone_name))
    return None


def _material_brief_candidate_item(
    *,
    filename: str,
    extract_type: str = "",
    text_excerpt: str = "",
    bullets: list[str] | None = None,
    question: str = "",
) -> dict[str, Any]:
    return {
        "filename": str(filename or "").strip(),
        "extract_type": str(extract_type or "").strip().lower(),
        "text_excerpt": str(text_excerpt or "").strip(),
        "bullets": list(bullets or []),
        "question": str(question or "").strip(),
    }


def _material_brief_candidate_item_from_artifact(
    filename: str,
    metadata: dict[str, Any],
) -> dict[str, Any]:
    brief = metadata.get("brief") if isinstance(metadata.get("brief"), dict) else {}
    text_extract = metadata.get("text_extract") if isinstance(metadata.get("text_extract"), dict) else {}
    extract_type = str(text_extract.get("type") or "").strip().lower()
    if not extract_type:
        content_type = str(metadata.get("content_type") or "").strip().lower()
        if "html" in content_type:
            extract_type = "html"
        elif Path(filename).suffix.lower() in {".html", ".htm", ".php"}:
            extract_type = "html"
    return _material_brief_candidate_item(
        filename=filename,
        extract_type=extract_type,
        text_excerpt=str(text_extract.get("excerpt") or "").strip(),
        bullets=list(brief.get("bullets") or []),
        question=str(brief.get("question") or "").strip(),
    )


def _is_invalid_material_brief_candidate_item(item: dict[str, Any]) -> bool:
    filename = str(item.get("filename") or "").strip()
    extract_type = str(item.get("extract_type") or "").strip().lower()
    text_excerpt = str(item.get("text_excerpt") or "").strip()
    bullets = [str(value).strip() for value in list(item.get("bullets") or []) if str(value).strip()]
    question = str(item.get("question") or "").strip()
    combined = "\n".join(
        str(value).strip()
        for value in [
            filename,
            text_excerpt,
            *bullets,
            question,
        ]
        if str(value).strip()
    )
    if any(pattern.search(combined) for pattern in MATERIAL_BRIEF_INVALID_ITEM_PATTERNS):
        return True
    if (
        extract_type == "html"
        and MATERIAL_GENERIC_HTML_FILENAME_RE.match(filename)
        and not text_excerpt
        and not bullets
        and not question
    ):
        return True
    if (
        extract_type == "html"
        and MATERIAL_GENERIC_HTML_FILENAME_RE.match(filename)
        and any(hint in combined.lower() for hint in MATERIAL_HTML_LOGIN_TEXT_HINTS)
    ):
        return True
    return False


def _class_preparation_line_from_matched_artifacts(
    matched_artifacts: list[dict[str, Any]],
) -> str | None:
    if not matched_artifacts:
        return None
    row = matched_artifacts[0]
    bullets = list(row.get("clean_bullets") or []) if isinstance(row, dict) else []
    if bullets:
        return _truncate_telegram_text(str(bullets[0]), limit=110)
    filename = str(row.get("filename") or "").strip()
    if not filename:
        return None
    if Path(filename).suffix.lower() in {".php", ".html", ".htm"}:
        return None
    return _truncate_telegram_text(f"자료: {filename}", limit=110)


def _preferred_material_rows(
    matched_artifacts: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    if not matched_artifacts:
        return []
    top_week_rank = max((int(row.get("week_rank") or 0) for row in matched_artifacts), default=0)
    return [
        row for row in matched_artifacts if int(row.get("week_rank") or 0) == top_week_rank
    ]


def _dedupe_preserve_order(
    values: list[str],
    *,
    key_fn: Callable[[str], str] | None = None,
    limit: int | None = None,
) -> list[str]:
    output: list[str] = []
    seen: set[str] = set()
    max_items = int(limit) if limit is not None else 0
    for value in values:
        text = str(value or "").strip()
        if not text:
            continue
        key = key_fn(text) if callable(key_fn) else text
        normalized_key = str(key or "").strip()
        if not normalized_key or normalized_key in seen:
            continue
        seen.add(normalized_key)
        output.append(text)
        if max_items > 0 and len(output) >= max_items:
            break
    return output


def _clean_course_summary_bullets(raw_items: Any, limit: int = 3) -> list[str]:
    if isinstance(raw_items, list):
        source_items = list(raw_items)
    elif isinstance(raw_items, str):
        source_items = [line.strip(" -*\t") for line in raw_items.splitlines() if line.strip()]
    else:
        source_items = []
    return _clean_telegram_brief_bullets(source_items, limit=limit)


def _clean_course_summary_sentence(value: Any, *, limit: int = 120) -> str:
    text = html.unescape(str(value or "").strip())
    text = re.sub(r"^[\-\*\d\.\)\s]+", "", text)
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return ""
    if any(pattern.match(text) for pattern in LOW_SIGNAL_BRIEF_PATTERNS):
        return ""
    return _truncate_telegram_text(text, limit=max(int(limit), 1))


def _compose_course_short_summary(bullets: list[str], *, limit: int = 120) -> str:
    parts: list[str] = []
    for bullet in bullets[:2]:
        text = str(bullet or "").strip()
        if not text:
            continue
        if text[-1] not in ".!?":
            text += "."
        parts.append(text)
    if not parts:
        return ""
    return _clean_course_summary_sentence(" ".join(parts), limit=limit)


def _course_day_summary_fallback_payload(
    *,
    aggregated_bullets: list[str],
    fallback_question: str,
    fallback_provenance: dict[str, Any] | None,
    filenames: list[str],
) -> dict[str, Any]:
    long_bullets = aggregated_bullets[:3]
    review = _clean_telegram_brief_question(fallback_question)
    short_summary = _compose_course_short_summary(long_bullets, limit=120)
    provenance = (
        dict(fallback_provenance)
        if isinstance(fallback_provenance, dict)
        else {
            "source": "uclass_html",
            "confidence": "low",
            "last_verified_at": now_utc_iso(),
            "derivation": "aggregated_material_brief",
        }
    )
    best_brief = None
    if long_bullets:
        best_brief = {
            "mode": "aggregated_material_brief",
            "version": COURSE_DAY_SUMMARY_VERSION,
            "bullets": long_bullets,
            "question": review,
            "provenance": provenance,
        }
    return {
        "best_brief": best_brief,
        "material_summary_tag": _provenance_brief_tag(provenance, fallback_source="uclass_html"),
        "material_summary_bits": tuple(long_bullets or filenames[:3]),
        "preparation": short_summary,
    }


def _llm_course_day_summary_payload(
    *,
    class_item: dict[str, Any],
    preferred_artifacts: list[dict[str, Any]],
    aggregated_bullets: list[str],
) -> dict[str, Any]:
    materials: list[dict[str, Any]] = []
    for row in preferred_artifacts[:6]:
        metadata = row.get("metadata") if isinstance(row.get("metadata"), dict) else {}
        materials.append(
            {
                "filename": str(row.get("filename") or "").strip(),
                "section_name": str(metadata.get("section_name") or "").strip(),
                "module_name": str(metadata.get("module_name") or "").strip(),
                "bullets": list(row.get("clean_bullets") or [])[:3],
                "question": str(row.get("clean_question") or "").strip(),
            }
        )
    return {
        "mode": "course_day_material_summary",
        "course_title": str(class_item.get("title") or "").strip(),
        "course_display_name": str(class_item.get("course_display_name") or "").strip(),
        "date": (
            class_item.get("start_local").date().isoformat()
            if isinstance(class_item.get("start_local"), datetime)
            else None
        ),
        "materials": materials,
        "aggregated_bullets": aggregated_bullets[:6],
        "requirements": {
            "short_summary": "1-2 short sentences for /today and scheduled briefings",
            "long_bullets": "2-3 combined bullets for /todaysummary",
            "review": "one short review action sentence",
        },
    }


def _llm_course_day_summary(
    settings: Settings,
    db: Database,
    *,
    class_item: dict[str, Any],
    preferred_artifacts: list[dict[str, Any]],
    aggregated_bullets: list[str],
    fallback_question: str,
) -> dict[str, Any] | None:
    if not bool(getattr(settings, "llm_enabled", False)):
        return None
    gate = _identity_warning_gate(
        settings=settings,
        db=db,
        step="course_day_summary",
        destination="llm",
    )
    if gate is not None:
        return None
    payload = _llm_course_day_summary_payload(
        class_item=class_item,
        preferred_artifacts=preferred_artifacts,
        aggregated_bullets=aggregated_bullets,
    )
    try:
        raw = _llm_client(settings).generate_text(
            system_prompt=COURSE_DAY_SUMMARY_SYSTEM_PROMPT,
            prompt=json.dumps(payload, ensure_ascii=False),
        )
        parsed = _parse_llm_json_payload(raw)
        if not isinstance(parsed, dict):
            return None
        long_bullets = _clean_course_summary_bullets(
            parsed.get("long_bullets") or parsed.get("bullets"),
            limit=3,
        )
        short_summary = _clean_course_summary_sentence(parsed.get("short_summary"), limit=120)
        review = _clean_telegram_brief_question(
            parsed.get("review") or parsed.get("question") or fallback_question
        )
        if not long_bullets:
            return None
        if not short_summary:
            short_summary = _compose_course_short_summary(long_bullets, limit=120)
        if not short_summary:
            return None
        return {
            "best_brief": {
                "mode": "llm_course_day_summary",
                "version": COURSE_DAY_SUMMARY_VERSION,
                "bullets": long_bullets,
                "question": review,
                "provenance": {
                    "source": "llm_inferred",
                    "confidence": "medium",
                    "last_verified_at": now_utc_iso(),
                    "derivation": "course_day_material_summary",
                },
            },
            "material_summary_tag": "[AI]",
            "material_summary_bits": tuple(long_bullets),
            "preparation": short_summary,
        }
    except Exception as exc:
        logger.warning(
            "course day summary llm fallback",
            extra={
                "course_title": str(class_item.get("title") or "").strip(),
                "error": str(exc),
            },
        )
        return None


def _aggregate_course_material_brief(
    settings: Settings,
    db: Database,
    *,
    class_item: dict[str, Any],
    matched_artifacts: list[dict[str, Any]],
) -> dict[str, Any]:
    preferred_artifacts = _preferred_material_rows(matched_artifacts)
    filenames = _dedupe_preserve_order(
        [
            str(row.get("filename") or "").strip()
            for row in preferred_artifacts
        ],
        key_fn=lambda text: text.lower(),
        limit=6,
    )
    if not preferred_artifacts:
        return {
            "best_brief": None,
            "best_brief_filename": "",
            "material_filenames": tuple(),
            "material_summary_tag": "",
            "material_summary_bits": tuple(),
            "preparation": None,
        }

    aggregated_bullets = _dedupe_preserve_order(
        [
            bullet
            for row in preferred_artifacts
            for bullet in list(row.get("clean_bullets") or [])
        ],
        key_fn=_normalize_task_title_key,
        limit=6,
    )
    first_question = ""
    first_provenance: dict[str, Any] | None = None
    material_summary_tag = ""
    for row in preferred_artifacts:
        brief = row.get("brief") if isinstance(row.get("brief"), dict) else None
        clean_question = str(row.get("clean_question") or "").strip()
        if not first_question and clean_question:
            first_question = clean_question
        if isinstance(brief, dict) and first_provenance is None:
            provenance = brief.get("provenance") if isinstance(brief.get("provenance"), dict) else None
            if isinstance(provenance, dict):
                first_provenance = dict(provenance)
        if not material_summary_tag:
            brief_tag = _provenance_brief_tag(
                brief.get("provenance") if isinstance(brief, dict) and isinstance(brief.get("provenance"), dict) else None,
                fallback_source="llm_inferred",
            )
            if brief_tag:
                material_summary_tag = brief_tag
    fallback_payload = _course_day_summary_fallback_payload(
        aggregated_bullets=aggregated_bullets,
        fallback_question=first_question,
        fallback_provenance=first_provenance,
        filenames=filenames,
    )
    llm_payload = _llm_course_day_summary(
        settings,
        db,
        class_item=class_item,
        preferred_artifacts=preferred_artifacts,
        aggregated_bullets=aggregated_bullets,
        fallback_question=first_question,
    )
    effective_payload = llm_payload or fallback_payload
    return {
        "best_brief": effective_payload["best_brief"],
        "best_brief_filename": filenames[0] if filenames else "",
        "material_filenames": tuple(filenames),
        "material_summary_tag": (
            str(effective_payload["material_summary_tag"] or "").strip()
            or material_summary_tag
        ),
        "material_summary_bits": tuple(effective_payload["material_summary_bits"]),
        "preparation": str(effective_payload["preparation"] or "").strip() or None,
    }


def _format_course_material_filenames(
    filenames: tuple[str, ...],
    *,
    limit: int = 3,
) -> str:
    items = [str(name or "").strip() for name in filenames if str(name or "").strip()]
    if not items:
        return ""
    capped = max(int(limit), 1)
    selected = items[:capped]
    if len(items) <= capped:
        return ", ".join(selected)
    return f"{', '.join(selected)} 외 {len(items) - capped}건"


def _start_of_local_day(value: datetime) -> datetime:
    aware = value if isinstance(value, datetime) else datetime.now(timezone.utc)
    if aware.tzinfo is None:
        aware = aware.replace(tzinfo=timezone.utc)
    local = aware.astimezone(aware.tzinfo or timezone.utc)
    return local.replace(hour=0, minute=0, second=0, microsecond=0)


def _local_day_window_start_utc_iso(value: datetime, *, lookback_days: int = 0) -> str:
    start_local = _start_of_local_day(value) - timedelta(days=max(int(lookback_days), 0))
    return start_local.astimezone(timezone.utc).replace(microsecond=0).isoformat()


def _effective_day_brief_lookahead_now_iso(
    lookahead_now_iso: str | None,
    *,
    lookahead_days: int,
) -> str | None:
    if int(lookahead_days) <= 0:
        return None
    normalized = normalize_datetime(lookahead_now_iso) if lookahead_now_iso else None
    parsed = _parse_dt(normalized or now_utc_iso())
    if parsed is None:
        parsed = datetime.now(timezone.utc)
    parsed = parsed.astimezone(timezone.utc).replace(microsecond=0)
    if not lookahead_now_iso:
        parsed = parsed.replace(second=0)
    return parsed.isoformat()


def _day_brief_cache_key(
    *,
    cache_scope: str,
    settings: Settings,
    user_id: int | None,
    target_day_local: datetime,
    reference_day_local: datetime,
    max_classes: int,
    artifact_limit: int,
    notification_limit: int,
    open_task_limit: int,
    lookahead_days: int,
    lookahead_limit: int,
    lookahead_now_iso: str | None,
    refresh_task_merge_cache: bool,
) -> str:
    payload = {
        "cache_scope": str(cache_scope or "").strip() or None,
        "timezone": str(getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul"),
        "user_id": _safe_int(user_id),
        "target_day": _start_of_local_day(target_day_local).date().isoformat(),
        "reference_day": _start_of_local_day(reference_day_local).date().isoformat(),
        "max_classes": int(max_classes),
        "artifact_limit": int(artifact_limit),
        "notification_limit": int(notification_limit),
        "open_task_limit": int(open_task_limit),
        "lookahead_days": int(lookahead_days),
        "lookahead_limit": int(lookahead_limit),
        "lookahead_now_iso": str(lookahead_now_iso or "").strip() or None,
        "refresh_task_merge_cache": bool(refresh_task_merge_cache),
    }
    return sha1(json.dumps(payload, ensure_ascii=True, sort_keys=True).encode("utf-8")).hexdigest()


def _day_brief_cache_fingerprint(snapshot: dict[str, Any]) -> str:
    return sha1(json.dumps(snapshot, ensure_ascii=True, sort_keys=True).encode("utf-8")).hexdigest()


def _get_cached_day_brief(*, cache_key: str, fingerprint: str) -> DayBrief | None:
    with _DAY_BRIEF_CACHE_LOCK:
        entry = _DAY_BRIEF_CACHE.get(cache_key)
        if not isinstance(entry, dict) or str(entry.get("fingerprint") or "") != fingerprint:
            return None
        cached_value = entry.get("value")
        if not isinstance(cached_value, DayBrief):
            return None
        _DAY_BRIEF_CACHE.pop(cache_key, None)
        _DAY_BRIEF_CACHE[cache_key] = entry
        return deepcopy(cached_value)


def _store_day_brief_cache(
    *,
    cache_key: str,
    fingerprint: str,
    day_brief: DayBrief,
) -> None:
    with _DAY_BRIEF_CACHE_LOCK:
        _DAY_BRIEF_CACHE.pop(cache_key, None)
        _DAY_BRIEF_CACHE[cache_key] = {
            "fingerprint": fingerprint,
            "value": deepcopy(day_brief),
        }
        while len(_DAY_BRIEF_CACHE) > DAY_BRIEF_CACHE_MAX_ENTRIES:
            oldest_key = next(iter(_DAY_BRIEF_CACHE))
            _DAY_BRIEF_CACHE.pop(oldest_key, None)


@dataclass(frozen=True)
class CourseDayBrief:
    class_item: dict[str, Any]
    matched_artifacts: tuple[dict[str, Any], ...]
    best_brief: dict[str, Any] | None
    best_brief_filename: str
    latest_material_filename: str
    material_filenames: tuple[str, ...]
    material_summary_tag: str
    material_summary_bits: tuple[str, ...]
    preparation: str | None
    notice_titles: tuple[str, ...]
    task_lines: tuple[str, ...]
    file_task_lines: tuple[str, ...]


@dataclass(frozen=True)
class DayBrief:
    target_day_local: datetime
    meetings_result: dict[str, Any]
    meeting_items: tuple[dict[str, Any], ...]
    course_briefs: tuple[CourseDayBrief, ...]
    tasks_due_on_day: tuple[Any, ...]
    tasks_due_within_window: tuple[Any, ...]


@dataclass(frozen=True)
class _DayBriefArtifactCandidate:
    artifact: Any
    metadata: dict[str, Any]
    canonical_course_id: str
    course_name: str
    module_name: str
    section_name: str
    filename: str
    brief: dict[str, Any] | None
    clean_bullets: tuple[str, ...]
    clean_question: str
    material_week_markers: tuple[int, ...]
    reference_dt: datetime | None
    updated_ts: float
    is_attachment_like: int
    is_non_html_file: int


@dataclass(frozen=True)
class _DayBriefNotificationCandidate:
    notification: Any
    metadata: dict[str, Any]
    canonical_course_id: str
    title: str
    body: str


@dataclass(frozen=True)
class _DayBriefTaskCandidate:
    task: Any
    metadata: dict[str, Any]
    canonical_course_id: str
    source: str
    due_dt: datetime | None
    title: str
    course_name: str
    detected_via: str


@dataclass(frozen=True)
class _DayBriefMatchContext:
    alias_map: dict[str, tuple[str, ...]]
    all_artifact_candidates: tuple[_DayBriefArtifactCandidate, ...]
    artifact_candidates_by_course_id: dict[str, tuple[_DayBriefArtifactCandidate, ...]]
    fallback_artifact_candidates: tuple[_DayBriefArtifactCandidate, ...]
    all_notification_candidates: tuple[_DayBriefNotificationCandidate, ...]
    notification_candidates_by_course_id: dict[str, tuple[_DayBriefNotificationCandidate, ...]]
    fallback_notification_candidates: tuple[_DayBriefNotificationCandidate, ...]
    all_task_candidates: tuple[_DayBriefTaskCandidate, ...]
    task_candidates_by_course_id: dict[str, tuple[_DayBriefTaskCandidate, ...]]
    fallback_task_candidates: tuple[_DayBriefTaskCandidate, ...]

    def artifact_candidates_for_course(
        self,
        canonical_course_id: str,
    ) -> tuple[_DayBriefArtifactCandidate, ...]:
        key = str(canonical_course_id or "").strip()
        if not key:
            return self.all_artifact_candidates
        return self.artifact_candidates_by_course_id.get(key, ()) + self.fallback_artifact_candidates

    def notification_candidates_for_course(
        self,
        canonical_course_id: str,
    ) -> tuple[_DayBriefNotificationCandidate, ...]:
        key = str(canonical_course_id or "").strip()
        if not key:
            return self.all_notification_candidates
        return self.notification_candidates_by_course_id.get(key, ()) + self.fallback_notification_candidates

    def task_candidates_for_course(
        self,
        canonical_course_id: str,
    ) -> tuple[_DayBriefTaskCandidate, ...]:
        key = str(canonical_course_id or "").strip()
        if not key:
            return self.all_task_candidates
        return self.task_candidates_by_course_id.get(key, ()) + self.fallback_task_candidates


def _build_day_brief_match_context(
    db: Database,
    *,
    artifacts: list[Any],
    notifications: list[Any],
    open_tasks: list[Any],
    alias_map: dict[str, tuple[str, ...]],
) -> _DayBriefMatchContext:
    artifact_candidates: list[_DayBriefArtifactCandidate] = []
    artifact_by_course_id: dict[str, list[_DayBriefArtifactCandidate]] = {}
    fallback_artifacts: list[_DayBriefArtifactCandidate] = []
    for artifact in artifacts:
        if str(getattr(artifact, "source", "")) != "uclass":
            continue
        metadata = _json_load(getattr(artifact, "metadata_json", None))
        canonical_course_id = _resolve_canonical_course_id(metadata, alias_map=alias_map)
        course_name = str(metadata.get("course_name") or "").strip()
        module_name = str(metadata.get("module_name") or "").strip()
        section_name = str(metadata.get("section_name") or "").strip()
        filename = str(getattr(artifact, "filename", "material") or "material")
        candidate_item = _material_brief_candidate_item_from_artifact(filename, metadata)
        if _is_invalid_material_brief_candidate_item(candidate_item):
            continue
        brief = metadata.get("brief") if isinstance(metadata.get("brief"), dict) else None
        updated_dt = _parse_dt(getattr(artifact, "updated_at", None))
        content_type = str(metadata.get("content_type") or "").strip().lower()
        source_kind = str(metadata.get("source_kind") or "").strip().lower()
        candidate = _DayBriefArtifactCandidate(
            artifact=artifact,
            metadata=metadata,
            canonical_course_id=canonical_course_id,
            course_name=course_name,
            module_name=module_name,
            section_name=section_name,
            filename=filename,
            brief=dict(brief) if isinstance(brief, dict) else None,
            clean_bullets=tuple(_clean_telegram_brief_bullets(candidate_item.get("bullets"), limit=2)),
            clean_question=_clean_telegram_brief_question(candidate_item.get("question")),
            material_week_markers=_extract_material_week_markers(
                filename,
                course_name,
                module_name,
                section_name,
            ),
            reference_dt=_material_reference_dt(metadata, artifact=artifact) or updated_dt,
            updated_ts=updated_dt.timestamp() if updated_dt else 0.0,
            is_attachment_like=1 if "attachment" in source_kind else 0,
            is_non_html_file=(
                1
                if Path(filename).suffix.lower() not in {".php", ".html", ".htm"}
                and "text/html" not in content_type
                else 0
            ),
        )
        artifact_candidates.append(candidate)
        if canonical_course_id:
            artifact_by_course_id.setdefault(canonical_course_id, []).append(candidate)
        else:
            fallback_artifacts.append(candidate)

    notification_candidates: list[_DayBriefNotificationCandidate] = []
    notification_by_course_id: dict[str, list[_DayBriefNotificationCandidate]] = {}
    fallback_notifications: list[_DayBriefNotificationCandidate] = []
    for item in notifications:
        source = str(getattr(item, "source", "")).strip().lower()
        if source not in {"uclass", "conflict"}:
            continue
        metadata = _json_load(getattr(item, "metadata_json", None))
        candidate = _DayBriefNotificationCandidate(
            notification=item,
            metadata=metadata,
            canonical_course_id=_resolve_canonical_course_id(metadata, alias_map=alias_map),
            title=str(getattr(item, "title", "")).strip(),
            body=str(getattr(item, "body", "") or "").strip(),
        )
        notification_candidates.append(candidate)
        if candidate.canonical_course_id:
            notification_by_course_id.setdefault(candidate.canonical_course_id, []).append(candidate)
        else:
            fallback_notifications.append(candidate)

    task_candidates: list[_DayBriefTaskCandidate] = []
    task_by_course_id: dict[str, list[_DayBriefTaskCandidate]] = {}
    fallback_tasks: list[_DayBriefTaskCandidate] = []
    for task in open_tasks:
        source = str(getattr(task, "source", "")).strip().lower()
        if source not in {"uclass", "inbox"}:
            continue
        due_dt = _parse_dt(getattr(task, "due_at", None))
        if due_dt is None:
            continue
        metadata = _json_load(getattr(task, "metadata_json", None))
        candidate = _DayBriefTaskCandidate(
            task=task,
            metadata=metadata,
            canonical_course_id=_resolve_canonical_course_id(metadata, alias_map=alias_map),
            source=source,
            due_dt=due_dt,
            title=str(getattr(task, "title", "Task")).strip() or "Task",
            course_name=str(metadata.get("course_name") or "").strip(),
            detected_via=str(metadata.get("detected_via") or "").strip(),
        )
        task_candidates.append(candidate)
        if candidate.canonical_course_id:
            task_by_course_id.setdefault(candidate.canonical_course_id, []).append(candidate)
        else:
            fallback_tasks.append(candidate)

    return _DayBriefMatchContext(
        alias_map=alias_map,
        all_artifact_candidates=tuple(artifact_candidates),
        artifact_candidates_by_course_id={
            key: tuple(values) for key, values in artifact_by_course_id.items()
        },
        fallback_artifact_candidates=tuple(fallback_artifacts),
        all_notification_candidates=tuple(notification_candidates),
        notification_candidates_by_course_id={
            key: tuple(values) for key, values in notification_by_course_id.items()
        },
        fallback_notification_candidates=tuple(fallback_notifications),
        all_task_candidates=tuple(task_candidates),
        task_candidates_by_course_id={key: tuple(values) for key, values in task_by_course_id.items()},
        fallback_task_candidates=tuple(fallback_tasks),
    )


class DayBriefService:
    def __init__(
        self,
        settings: Settings,
        db: Database,
        *,
        user_id: int | None = None,
    ) -> None:
        self.settings = settings
        self.db = db
        self.user_id = user_id

    def build_day_brief(
        self,
        *,
        target_day_local: datetime,
        reference_day_local: datetime | None = None,
        max_classes: int,
        artifact_limit: int = 500,
        notification_limit: int = 200,
        open_task_limit: int = 500,
        lookahead_days: int = 0,
        lookahead_limit: int = 20,
        lookahead_now_iso: str | None = None,
        refresh_task_merge_cache: bool = False,
    ) -> DayBrief:
        reference_day = reference_day_local or target_day_local
        effective_lookahead_now_iso = _effective_day_brief_lookahead_now_iso(
            lookahead_now_iso,
            lookahead_days=lookahead_days,
        )
        artifact_since_iso = _local_day_window_start_utc_iso(
            target_day_local,
            lookback_days=DAY_BRIEF_ARTIFACT_LOOKBACK_DAYS,
        )
        notification_since_iso = _local_day_window_start_utc_iso(
            reference_day,
            lookback_days=DAY_BRIEF_NOTIFICATION_LOOKBACK_DAYS,
        )
        open_task_since_iso = _local_day_window_start_utc_iso(reference_day)
        open_task_rows: list[Any] | None = None
        if refresh_task_merge_cache:
            open_task_rows = self.db.list_open_tasks_due_from(
                open_task_since_iso,
                limit=open_task_limit,
                user_id=self.user_id,
            )
            _precompute_task_merge_cache(
                self.settings,
                self.db,
                user_id=self.user_id,
                tasks=open_task_rows,
            )
        cache_key = _day_brief_cache_key(
            cache_scope=str(self.db.db_path.expanduser().resolve()),
            settings=self.settings,
            user_id=self.user_id,
            target_day_local=target_day_local,
            reference_day_local=reference_day,
            max_classes=max_classes,
            artifact_limit=artifact_limit,
            notification_limit=notification_limit,
            open_task_limit=open_task_limit,
            lookahead_days=lookahead_days,
            lookahead_limit=lookahead_limit,
            lookahead_now_iso=effective_lookahead_now_iso,
            refresh_task_merge_cache=refresh_task_merge_cache,
        )
        cache_fingerprint = _day_brief_cache_fingerprint(
            self.db.day_brief_cache_snapshot(user_id=self.user_id)
        )
        cached_day_brief = _get_cached_day_brief(
            cache_key=cache_key,
            fingerprint=cache_fingerprint,
        )
        if cached_day_brief is not None:
            return cached_day_brief
        alias_map = self.db.course_alias_resolution_map(user_id=self.user_id)
        meetings_result = _collect_primary_meetings_scoped(
            settings=self.settings,
            db=self.db,
            target_day_local=target_day_local,
            user_id=self.user_id,
        )
        meeting_items = tuple(
            list(meetings_result.get("events") or [])
            if isinstance(meetings_result, dict) and bool(meetings_result.get("ok"))
            else []
        )
        class_items = _collect_class_occurrences(
            settings=self.settings,
            db=self.db,
            target_day_local=target_day_local,
            max_items=max_classes,
            user_id=self.user_id,
            alias_map=alias_map,
        )
        artifacts = self.db.list_artifacts_since(
            artifact_since_iso,
            limit=artifact_limit,
            user_id=self.user_id,
        )
        notifications = self.db.list_notifications_since(
            notification_since_iso,
            limit=notification_limit,
            user_id=self.user_id,
        )
        if open_task_rows is None:
            open_task_rows = self.db.list_open_tasks_due_from(
                open_task_since_iso,
                limit=open_task_limit,
                user_id=self.user_id,
            )
        open_tasks = _apply_task_merge_cache(
            self.settings,
            self.db,
            tasks=open_task_rows,
            user_id=self.user_id,
        )
        match_context = _build_day_brief_match_context(
            self.db,
            artifacts=artifacts,
            notifications=notifications,
            open_tasks=open_tasks,
            alias_map=alias_map,
        )
        tasks_due_on_day = tuple(
            task
            for task in open_tasks
            if _task_due_on_day(task.due_at, self.settings.timezone, target_day_local)
        )
        tasks_due_within_window: tuple[Any, ...] = ()
        if lookahead_days > 0:
            due_window_rows = self.db.list_tasks_due_within(
                days=lookahead_days,
                now_iso=effective_lookahead_now_iso or now_utc_iso(),
                limit=lookahead_limit,
                user_id=self.user_id,
            )
            tasks_due_within_window = tuple(
                _apply_task_merge_cache(
                    self.settings,
                    self.db,
                    tasks=due_window_rows,
                    user_id=self.user_id,
                )
            )
        day_brief = DayBrief(
            target_day_local=target_day_local,
            meetings_result=meetings_result,
            meeting_items=meeting_items,
            course_briefs=tuple(
                self._build_course_day_brief(
                    class_item,
                    artifacts=artifacts,
                    notifications=notifications,
                    open_tasks=open_tasks,
                    reference_day_local=reference_day,
                    match_context=match_context,
                )
                for class_item in class_items
            ),
            tasks_due_on_day=tasks_due_on_day,
            tasks_due_within_window=tasks_due_within_window,
        )
        _store_day_brief_cache(
            cache_key=cache_key,
            fingerprint=cache_fingerprint,
            day_brief=day_brief,
        )
        return day_brief

    def _build_course_day_brief(
        self,
        class_item: dict[str, Any],
        *,
        artifacts: list[Any],
        notifications: list[Any],
        open_tasks: list[Any],
        reference_day_local: datetime,
        match_context: _DayBriefMatchContext | None = None,
    ) -> CourseDayBrief:
        matched_artifacts = _matched_artifacts_for_class(
            self.db,
            class_item,
            artifacts,
            limit=20,
            user_id=self.user_id,
            match_context=match_context,
        )
        aggregated_material = _aggregate_course_material_brief(
            self.settings,
            self.db,
            class_item=class_item,
            matched_artifacts=matched_artifacts,
        )
        latest_material_filename = (
            str(matched_artifacts[0].get("filename") or "").strip()
            if matched_artifacts
            else ""
        )
        task_lines, file_task_lines = _matched_tasks_for_class(
            self.db,
            class_item,
            open_tasks,
            reference_day_local=reference_day_local,
            limit=2,
            user_id=self.user_id,
            match_context=match_context,
        )
        return CourseDayBrief(
            class_item=class_item,
            matched_artifacts=tuple(matched_artifacts),
            best_brief=aggregated_material["best_brief"],
            best_brief_filename=str(aggregated_material["best_brief_filename"] or "").strip(),
            latest_material_filename=latest_material_filename,
            material_filenames=tuple(aggregated_material["material_filenames"]),
            material_summary_tag=str(aggregated_material["material_summary_tag"] or "").strip(),
            material_summary_bits=tuple(aggregated_material["material_summary_bits"]),
            preparation=str(aggregated_material["preparation"] or "").strip() or None,
            notice_titles=tuple(
                _matched_notifications_for_class(
                    self.db,
                    class_item,
                    notifications,
                    limit=2,
                    user_id=self.user_id,
                    match_context=match_context,
                )
            ),
            task_lines=tuple(task_lines),
            file_task_lines=tuple(file_task_lines),
        )


def _format_status_time(value: str | None, timezone_name: str) -> str:
    parsed = _parse_dt(value)
    if parsed is None:
        return "기록 없음"
    return parsed.astimezone(ZoneInfo(timezone_name)).strftime("%Y-%m-%d %H:%M %Z")


def _append_message_section(lines: list[str], title: str) -> None:
    if lines and lines[-1] != "":
        lines.append("")
    lines.append(title)


def _inbox_item_type_label(item_type: str) -> str:
    return {
        "event_draft": "일정",
        "task_draft": "과제",
        "note": "메모",
        "command": "명령",
    }.get(str(item_type or "").strip(), "항목")


def _format_telegram_access_denied(chat_id: str | None) -> str:
    lines = ["[KU] 접근 제한", ""]
    lines.append("- 이 채팅은 아직 사용할 수 있도록 활성화되지 않았습니다.")
    lines.append("- `/setup`으로 연결 상태를 확인하세요.")
    return "\n".join(lines)


def _format_telegram_apply_result(result: dict[str, Any]) -> str:
    processed = int(result.get("processed") or 0)
    lines = ["[KU] Inbox 반영", ""]
    if processed <= 0:
        lines.append("- 반영할 draft가 없습니다.")
        lines.append("- 새 메모를 보내면 `/inbox`에서 다시 확인할 수 있습니다.")
        return "\n".join(lines)
    lines.extend(
        [
            "처리 결과",
            f"- 처리 {processed}건",
            (
                "- 생성: "
                f"일정 {int(result.get('created_events') or 0)}건 / "
                f"과제 {int(result.get('created_tasks') or 0)}건 / "
                f"메모 {int(result.get('notes') or 0)}건"
            ),
            "",
            "다음으로 해볼 것",
            "- `/today`로 오늘 일정과 과제를 확인하세요.",
            "- 남은 초안이 있으면 `/inbox`를 다시 확인하세요.",
        ]
    )
    return "\n".join(lines)


def _format_telegram_done_result(
    result: dict[str, Any],
    *,
    timezone_name: str,
) -> str:
    if not bool(result.get("ok")):
        reason = str(result.get("reason") or result.get("error") or "처리 실패").strip()
        return "\n".join(
            [
                "[KU] 과제 상태 변경",
                "",
                f"- 실패: {reason}",
            ]
        )
    task = result.get("task") if isinstance(result.get("task"), dict) else {}
    lines = ["[KU] 과제 상태 변경", "", "처리 결과"]
    title = str(task.get("title") or "과제").strip() or "과제"
    lines.append(f"- 과제: {title}")
    lines.append(f"- 상태: {str(task.get('status') or '').strip() or 'unknown'}")
    external_id = str(task.get("external_id") or "").strip()
    if external_id:
        lines.append(f"- ID: {external_id}")
    lines.extend(
        [
            "",
            "다음으로 해볼 것",
            "- `/today`로 남은 마감 과제를 확인하세요.",
        ]
    )
    return "\n".join(lines)


def _humanize_telegram_plan_error(reason: str) -> str:
    normalized = str(reason or "").strip()
    if not normalized:
        return "리마인더 예약에 실패했습니다."
    if normalized == "TELEGRAM_SMART_COMMANDS_ENABLED is false":
        return "리마인더 예약 기능이 현재 꺼져 있습니다."
    if normalized == "LLM is disabled and instruction could not be parsed":
        return "문장에서 예약 시각을 찾지 못했습니다. 예: `/plan 내일 오전 8시에 과제 제출 알림`"
    if normalized == "invalid run_at_iso":
        return "예약 시각을 해석하지 못했습니다. 조금 더 구체적으로 다시 적어 주세요."
    if normalized == "empty reminder message":
        return "리마인더 내용이 비어 있습니다. 다시 적어 주세요."
    return normalized


def _format_telegram_plan_result(
    result: dict[str, Any],
    *,
    timezone_name: str,
) -> str:
    if not bool(result.get("ok")):
        reason = _humanize_telegram_plan_error(
            str(result.get("error") or result.get("reason") or "예약 실패").strip()
        )
        return "\n".join(
            [
                "[KU] 리마인더 예약",
                "",
                f"- 실패: {reason}",
            ]
        )
    if not bool(result.get("scheduled")):
        reason = str(result.get("reason") or "예약할 리마인더를 찾지 못했습니다.").strip()
        return "\n".join(
            [
                "[KU] 리마인더 예약",
                "",
                f"- 예약 안 함: {reason}",
            ]
        )
    reminder = result.get("reminder") if isinstance(result.get("reminder"), dict) else {}
    run_at = _format_status_time(str(reminder.get("run_at") or None), timezone_name)
    mode = str(result.get("mode") or "").strip() or "unknown"
    message = str(reminder.get("message") or "").strip() or "리마인더"
    return "\n".join(
        [
            "[KU] 리마인더 예약",
            "",
            "예약 결과",
            f"- 시각: {run_at}",
            f"- 내용: {message}",
            f"- 방식: {mode}",
            "",
            "다음으로 해볼 것",
            "- 다른 일정도 같은 방식으로 `/plan`에 자연어로 보내세요.",
        ]
    )


def _uclass_status_summary(cursor: dict[str, Any]) -> str:
    if bool(cursor.get("skipped")):
        reason = str(cursor.get("reason") or "이유 없음").strip()
        return f"건너뜀 - {reason}"
    wsfunctions = cursor.get("wsfunctions") if isinstance(cursor.get("wsfunctions"), dict) else {}
    failed = sum(
        int(item.get("failed") or 0)
        for item in wsfunctions.values()
        if isinstance(item, dict)
    )
    if failed > 0:
        return f"일부 실패 - ws 실패 {failed}건"
    if wsfunctions:
        course_calls = int(
            (
                wsfunctions.get("core_course_get_contents")
                if isinstance(wsfunctions.get("core_course_get_contents"), dict)
                else {}
            ).get("ok")
            or 0
        )
        return f"정상 - 강의 내용 WS {course_calls}건 확인"
    if int(cursor.get("html_material_candidates") or 0) > 0:
        return "HTML fallback만 사용"
    return "최근 결과 정보 없음"


def _format_telegram_status(
    settings: Settings,
    db: Database,
    *,
    user_id: int | None = None,
) -> str:
    counts = db.counts(user_id=user_id)
    state_map = {row.job_name: row for row in db.list_sync_states(user_id=user_id)}
    telegram_state = state_map.get("sync_telegram")
    if telegram_state is None or (user_id and not telegram_state.last_run_at and not telegram_state.last_cursor_json):
        telegram_state = db.get_sync_state("sync_telegram")
    telegram_cursor = _json_load(telegram_state.last_cursor_json if telegram_state else None)
    uclass_state = state_map.get("sync_uclass")
    if uclass_state is None or (user_id and not uclass_state.last_run_at and not uclass_state.last_cursor_json):
        uclass_state = db.get_sync_state("sync_uclass")
    uclass_cursor = _json_load(uclass_state.last_cursor_json if uclass_state else None)
    sync_dashboard = db.sync_dashboard_snapshot(user_id=user_id)
    display_last_success = sync_dashboard.get("last_successful_sync_at")
    dashboard_sources = list(sync_dashboard.get("sources") or [])
    if user_id:
        global_sync_dashboard = db.sync_dashboard_snapshot()
        global_sources = {
            str(item.get("key") or ""): item
            for item in list(global_sync_dashboard.get("sources") or [])
        }
        merged_sources: list[dict[str, Any]] = []
        for item in dashboard_sources:
            fallback = global_sources.get(str(item.get("key") or ""))
            if (
                str(item.get("status") or "").strip().lower() == "never"
                and fallback
                and (
                    fallback.get("last_run_at")
                    or fallback.get("last_success_at")
                    or fallback.get("last_error")
                )
            ):
                merged_sources.append(dict(fallback))
                continue
            merged_sources.append(item)
        dashboard_sources = merged_sources
        if not display_last_success:
            display_last_success = global_sync_dashboard.get("last_successful_sync_at")
    tracked_sources = [
        item
        for item in dashboard_sources
        if item.get("key") in {"portal", "uclass", "weather", "telegram"}
    ]

    def _status_label(value: str) -> str:
        mapping = {
            "success": "준비됨",
            "skipped": "확인 필요",
            "error": "문제 있음",
            "never": "아직 없음",
        }
        return mapping.get(str(value or "").strip().lower(), "아직 없음")

    lines = ["[KU] 상태 요약", "", "전체"]
    lines.append(
        f"- 마지막 성공 동기화: {_format_status_time(display_last_success, settings.timezone)}"
    )
    lines.append(
        "- "
        + ", ".join(
            [
                f"일정 {int(counts.get('events', 0))}개",
                f"과제 {int(counts.get('tasks', 0))}개",
                f"자료 {int(counts.get('artifacts', 0))}개",
            ]
        )
    )
    action_required = int(sync_dashboard.get("action_required_count") or 0)
    if action_required > 0:
        lines.append(f"- 지금 확인할 항목: {action_required}건")
    else:
        lines.append("- 지금 확인할 항목은 없습니다.")

    _append_message_section(lines, "서비스별 상태")
    for item in tracked_sources:
        pieces = [
            f"{item.get('label')}: {_status_label(str(item.get('status') or 'never'))}",
        ]
        last_run_at = str(item.get("last_run_at") or "").strip()
        if last_run_at:
            pieces.append(f"마지막 확인 {_format_status_time(last_run_at, settings.timezone)}")
        if int(item.get("action_required") or 0) > 0:
            pieces.append("확인 필요")
        lines.append("- " + " | ".join(pieces))

    _append_message_section(lines, "추천")
    if int(sync_dashboard.get("pending_inbox_count") or 0) > 0:
        lines.append("- `/inbox`로 새 초안을 확인하세요.")
    if int(counts.get("events", 0)) or int(counts.get("tasks", 0)):
        lines.append("- `/today`로 오늘 일정과 과제를 확인하세요.")
    else:
        lines.append("- `/setup`으로 연결 상태를 다시 점검하세요.")
    if bool(getattr(settings, "telegram_smart_commands_enabled", False)):
        lines.append("- `/plan 내일 오전 8시에 과제 제출 알림`으로 리마인더를 추가할 수 있습니다.")
    return "\n".join(lines)


def _format_telegram_day_empty_message(state: DayAgendaState) -> str:
    lines = [f"[KU] {state.day_label} 보기 ({state.target_day_local.date().isoformat()})", ""]
    if state.empty_reason == "no_connection":
        lines.extend(
            [
                "- 아직 학교 계정이 연결되지 않았습니다.",
                "- `/connect`로 학교 계정을 연결하세요.",
                "- 연결 상태 확인: `/setup`",
            ]
        )
        return "\n".join(lines)
    if state.empty_reason == "first_sync_pending":
        lines.extend(
            [
                "- 학교 계정은 연결됐지만 아직 첫 동기화가 끝나지 않았습니다.",
                "- 잠시 후 다시 확인하거나 `/setup`으로 연결 상태를 점검하세요.",
            ]
        )
        return "\n".join(lines)
    lines.append(f"{state.day_label}은 등록된 일정, 수업, 마감 과제가 없습니다.")
    return "\n".join(lines)


def _render_telegram_day_message(state: DayAgendaState) -> str:
    if state.is_empty:
        return _format_telegram_day_empty_message(state)

    lines = [f"[KU] {state.day_label} 보기 ({state.target_day_local.date().isoformat()})", ""]
    if state.show_meeting_section:
        lines.append(f"일정: {len(state.meeting_items)}" if not state.meetings_failed else "일정: 확인 불가")
        if state.meeting_items:
            for item in state.meeting_items:
                if item.location:
                    lines.append(f"- {item.when} {item.title} @ {item.location}")
                else:
                    lines.append(f"- {item.when} {item.title}")
        elif state.meetings_failed:
            lines.append("- 일정 정보를 불러오지 못했습니다.")
        else:
            lines.append("- 일정 없음")
        if state.skipped_reason:
            lines.append(f"- 참고: {state.skipped_reason}")

    _append_message_section(lines, f"수업: {len(state.course_items)}")
    if not state.course_items:
        lines.append("- 수업 없음")
    for item in state.course_items:
        if item.location_text and item.location_text != "TBD":
            lines.append(f"- {item.when} {item.title} @ {item.location_text}")
        else:
            lines.append(f"- {item.when} {item.title}")
        if item.preparation:
            lines.append(f"  준비: {item.preparation}")
        if item.notice_titles:
            lines.append(f"  공지: {'; '.join(item.notice_titles)}")
        if item.task_lines:
            lines.append(f"  수업 과제: {'; '.join(item.task_lines)}")
        if item.file_task_lines:
            lines.append(f"  파일 감지 과제: {'; '.join(item.file_task_lines)}")

    _append_message_section(lines, f"마감 과제: {len(state.task_lines)}")
    if not state.task_lines:
        lines.append("- 마감 과제 없음")
    for item in state.task_lines:
        lines.append(f"- {item}")

    if state.upcoming_task_lines:
        _append_message_section(lines, f"다음 마감: {len(state.upcoming_task_lines)}")
        for item in state.upcoming_task_lines:
            lines.append(f"- {item}")

    if state.course_items:
        _append_message_section(lines, "바로가기")
        lines.append(f"- 자세한 수업 자료 요약: {state.summary_hint_command}")
    return "\n".join(lines)


def _format_telegram_day(
    settings: Settings,
    db: Database,
    *,
    target_day_local: datetime,
    day_label: str,
    summary_hint_command: str,
    include_upcoming_tasks: bool,
    user_id: int | None = None,
) -> str:
    day_brief = DayBriefService(settings, db, user_id=user_id).build_day_brief(
        target_day_local=target_day_local,
        reference_day_local=target_day_local,
        max_classes=6,
        artifact_limit=TELEGRAM_DAY_ARTIFACT_LIMIT,
        notification_limit=TELEGRAM_DAY_NOTIFICATION_LIMIT,
        open_task_limit=TELEGRAM_DAY_OPEN_TASK_LIMIT,
        lookahead_days=7 if include_upcoming_tasks else 0,
        lookahead_limit=20,
    )
    state = build_day_agenda_state(
        settings,
        db,
        day_brief=day_brief,
        target_day_local=target_day_local,
        day_label=day_label,
        summary_hint_command=summary_hint_command,
        include_upcoming_tasks=include_upcoming_tasks,
        user_id=user_id,
        format_time_range_local=_format_time_range_local,
        format_task_line=lambda task, reference_day_local: _format_briefing_task_line(
            task,
            reference_day_local=reference_day_local,
            include_course=True,
        ),
        is_task_due_on_target_day=lambda task: _task_due_on_day(
            task.due_at,
            settings.timezone,
            target_day_local,
        ),
    )
    return _render_telegram_day_message(state)


def _format_telegram_today(settings: Settings, db: Database, *, user_id: int | None = None) -> str:
    return _format_telegram_day(
        settings=settings,
        db=db,
        target_day_local=datetime.now(ZoneInfo(settings.timezone)),
        day_label="오늘",
        summary_hint_command="/todaysummary",
        include_upcoming_tasks=True,
        user_id=user_id,
    )


def _format_telegram_tomorrow(settings: Settings, db: Database, *, user_id: int | None = None) -> str:
    now_local = datetime.now(ZoneInfo(settings.timezone))
    return _format_telegram_day(
        settings=settings,
        db=db,
        target_day_local=now_local + timedelta(days=1),
        day_label="내일",
        summary_hint_command="/tomorrowsummary",
        include_upcoming_tasks=False,
        user_id=user_id,
    )

def _clean_telegram_brief_bullets(raw_items: Any, limit: int = 2) -> list[str]:
    output: list[str] = []
    for text in _clean_brief_bullets(raw_items, limit=max(int(limit), 1) * 3):
        stripped = str(text or "").strip()
        if not stripped:
            continue
        if re.fullmatch(r"\d+[\.\)]?", stripped):
            continue
        if len(stripped) < 5:
            continue
        if any(pattern.match(stripped) for pattern in LOW_SIGNAL_BRIEF_PATTERNS):
            continue
        output.append(stripped)
        if len(output) >= max(int(limit), 1):
            break
    return output


def _clean_telegram_brief_question(value: Any) -> str:
    text = str(value or "").strip()
    if not text:
        return ""
    if any(pattern.match(text) for pattern in LOW_SIGNAL_BRIEF_PATTERNS):
        return ""
    return text


def _reject_material_brief_push_item(item: dict[str, Any]) -> bool:
    return _is_invalid_material_brief_candidate_item(item)


def _sanitize_material_brief_push_items(items: list[dict[str, Any]], limit: int) -> list[dict[str, Any]]:
    sanitized: list[dict[str, Any]] = []
    for item in items:
        if _reject_material_brief_push_item(item):
            continue
        bullets = _clean_telegram_brief_bullets(item.get("bullets"), limit=3)
        if not bullets:
            continue
        sanitized.append(
            {
                **item,
                "bullets": bullets,
                "question": _clean_telegram_brief_question(item.get("question")),
            }
        )
        if len(sanitized) >= max(int(limit), 1):
            break
    return sanitized


def _telegram_notice_view_state_job(kind: str) -> str:
    normalized = str(kind or "").strip().lower()
    return f"telegram_notice_{normalized or 'unknown'}"


def _portal_notice_snapshot_job(kind: str) -> str:
    normalized = str(kind or "").strip().lower()
    return f"uos_notice_snapshot_{normalized or 'unknown'}"


def _notice_item_value(item: Any, key: str) -> Any:
    if isinstance(item, dict):
        return item.get(key)
    return getattr(item, key, None)


def _normalize_portal_notice_row(item: Any) -> dict[str, Any] | None:
    title = str(_notice_item_value(item, "title") or "").strip()
    if not title:
        return None
    return {
        "seq": str(_notice_item_value(item, "seq") or "").strip() or None,
        "sort": str(_notice_item_value(item, "sort") or "").strip() or None,
        "title": title,
        "department": str(_notice_item_value(item, "department") or "").strip() or None,
        "posted_on": str(_notice_item_value(item, "posted_on") or "").strip() or None,
        "list_id": str(_notice_item_value(item, "list_id") or "").strip() or None,
        "menuid": str(_notice_item_value(item, "menuid") or "").strip() or None,
        "source_url": str(_notice_item_value(item, "source_url") or "").strip() or None,
        "article_url": str(_notice_item_value(item, "article_url") or "").strip() or None,
    }


def _portal_notice_rows(value: Any) -> list[dict[str, Any]]:
    if not isinstance(value, list):
        return []
    rows: list[dict[str, Any]] = []
    seen: set[str] = set()
    for item in value:
        row = _normalize_portal_notice_row(item)
        if not row:
            continue
        dedupe_key = str(row.get("seq") or row.get("title") or "").strip()
        if not dedupe_key or dedupe_key in seen:
            continue
        seen.add(dedupe_key)
        rows.append(row)
        if len(rows) >= TELEGRAM_UOS_NOTICE_LIMIT:
            break
    return rows


def _portal_notice_snapshot_payload(
    notices: list[Any],
    *,
    fetched_at: str | None,
    empty: bool,
) -> dict[str, Any]:
    rows = _portal_notice_rows(notices)
    return {
        "fetched_at": str(fetched_at or "").strip() or None,
        "notice_count": len(rows),
        "notices": rows,
        "empty": bool(empty),
    }


def _portal_notice_snapshot_from_cursor(cursor: dict[str, Any]) -> dict[str, Any]:
    raw_snapshot = cursor.get("snapshot") if isinstance(cursor.get("snapshot"), dict) else {}
    notices = _portal_notice_rows(raw_snapshot.get("notices"))
    return {
        "fetched_at": str(raw_snapshot.get("fetched_at") or "").strip() or None,
        "notice_count": len(notices),
        "notices": notices,
        "empty": bool(raw_snapshot.get("empty")),
    }


def _portal_notice_source_payload(
    config: dict[str, Any],
    *,
    previous_cursor: dict[str, Any] | None = None,
    metadata: dict[str, Any] | None = None,
) -> dict[str, Any]:
    previous_source = (
        previous_cursor.get("source")
        if isinstance(previous_cursor, dict) and isinstance(previous_cursor.get("source"), dict)
        else {}
    )
    current = metadata if isinstance(metadata, dict) else {}
    list_url = (
        str(current.get("source_url") or "").strip()
        or str(previous_source.get("list_url") or "").strip()
        or None
    )
    resolved_url = (
        str(current.get("resolved_url") or "").strip()
        or str(previous_source.get("resolved_url") or "").strip()
        or list_url
    )
    return {
        "label": "학교 포털",
        "list_id": str(config.get("list_id") or "").strip() or None,
        "menuid": str(config.get("menuid") or "").strip() or None,
        "list_url": list_url,
        "resolved_url": resolved_url,
    }


def _portal_notice_attempt_payload(
    *,
    ok: bool,
    attempted_at: str | None,
    metadata: dict[str, Any] | None = None,
    error: str | None = None,
) -> dict[str, Any]:
    details = metadata if isinstance(metadata, dict) else {}
    return {
        "ok": bool(ok),
        "attempted_at": str(
            attempted_at or details.get("requested_at") or details.get("fetched_at") or now_utc_iso()
        ).strip()
        or now_utc_iso(),
        "error": str(error or "").strip() or None,
        "http_status": _safe_int(details.get("http_status")),
        "page_title": str(details.get("page_title") or "").strip() or None,
        "parser": str(details.get("parser") or "").strip() or None,
        "parsed_count": int(details.get("parsed_count") or 0),
        "empty_detected": bool(details.get("empty_detected")),
    }


def _portal_notice_error_details(exc: Exception) -> tuple[str, dict[str, Any]]:
    metadata: dict[str, Any] = {}
    if isinstance(exc, (PortalNoticeFetchError, KuNoticeFetchError)):
        metadata = dict(exc.metadata)
    message = str(exc).strip() or "unknown error"
    status_code = _safe_int(metadata.get("http_status"))
    if status_code is not None and f"{status_code}" not in message:
        message = f"HTTP {status_code}: {message}"
    return message, metadata


def _refresh_portal_notice_snapshot(db: Database, kind: str) -> dict[str, Any]:
    config = UOS_NOTICE_FEEDS.get(str(kind).strip().lower())
    if not isinstance(config, dict):
        return {}
    state = db.get_sync_state(_portal_notice_snapshot_job(kind))
    previous_cursor = _json_load(state.last_cursor_json)
    previous_snapshot = _portal_notice_snapshot_from_cursor(previous_cursor)
    try:
        result = fetch_ku_notice_feed(
            board_id=str(config["board_id"]),
            limit=TELEGRAM_UOS_NOTICE_LIMIT,
        )
    except Exception as exc:
        error_message, error_metadata = _portal_notice_error_details(exc)
        cursor = {
            "kind": str(kind).strip().lower(),
            "label": str(config["label"]),
            "requested_limit": TELEGRAM_UOS_NOTICE_LIMIT,
            "source": _portal_notice_source_payload(
                config,
                previous_cursor=previous_cursor,
                metadata=error_metadata,
            ),
            "snapshot": previous_snapshot,
            "last_attempt": _portal_notice_attempt_payload(
                ok=False,
                attempted_at=(
                    str(error_metadata.get("requested_at") or error_metadata.get("fetched_at") or "").strip()
                    or now_utc_iso()
                ),
                metadata=error_metadata,
                error=error_message,
            ),
        }
        db.update_sync_state(
            _portal_notice_snapshot_job(kind),
            last_run_at=str(cursor["last_attempt"]["attempted_at"]),
            last_cursor_json=cursor,
        )
        return cursor

    metadata = asdict(result.metadata)
    snapshot_fetched_at = (
        str(metadata.get("fetched_at") or metadata.get("requested_at") or "").strip()
        or now_utc_iso()
    )
    cursor = {
        "kind": str(kind).strip().lower(),
        "label": str(config["label"]),
        "requested_limit": TELEGRAM_UOS_NOTICE_LIMIT,
        "source": _portal_notice_source_payload(config, metadata=metadata),
        "snapshot": _portal_notice_snapshot_payload(
            result.notices,
            fetched_at=snapshot_fetched_at,
            empty=bool(metadata.get("empty_detected")),
        ),
        "last_attempt": _portal_notice_attempt_payload(
            ok=True,
            attempted_at=snapshot_fetched_at,
            metadata=metadata,
        ),
    }
    db.update_sync_state(
        _portal_notice_snapshot_job(kind),
        last_run_at=snapshot_fetched_at,
        last_cursor_json=cursor,
    )
    return cursor


def _portal_notice_source_line(cursor: dict[str, Any], *, timezone_name: str) -> str:
    source = cursor.get("source") if isinstance(cursor.get("source"), dict) else {}
    snapshot = _portal_notice_snapshot_from_cursor(cursor)
    attempt = cursor.get("last_attempt") if isinstance(cursor.get("last_attempt"), dict) else {}
    label = str(source.get("label") or "학교 포털").strip() or "학교 포털"
    if attempt.get("ok") is False and snapshot.get("fetched_at"):
        label = f"{label} 캐시"
    reference_at = (
        str(snapshot.get("fetched_at") or "").strip()
        or str(attempt.get("attempted_at") or "").strip()
        or None
    )
    if reference_at:
        return f"- 출처: {label} ({_format_status_time(reference_at, timezone_name)})"
    return f"- 출처: {label}"


def _portal_notice_status_lines(cursor: dict[str, Any], *, timezone_name: str) -> list[str]:
    attempt = cursor.get("last_attempt") if isinstance(cursor.get("last_attempt"), dict) else {}
    if attempt.get("ok") is not False:
        return []
    snapshot = _portal_notice_snapshot_from_cursor(cursor)
    error = str(attempt.get("error") or "").strip() or "unknown error"
    snapshot_fetched_at = str(snapshot.get("fetched_at") or "").strip() or None
    if snapshot_fetched_at:
        return [
            f"- 학교 포털 응답을 확인하지 못해 저장된 최근 공지를 보여줍니다: {error}",
            f"- 마지막 정상 반영: {_format_status_time(snapshot_fetched_at, timezone_name)}",
        ]
    return []


def _portal_notice_new_seq_set(notices: list[Any], *, previous_top_seq: str | None) -> set[str]:
    anchor = str(previous_top_seq or "").strip()
    if not anchor:
        return set()
    output: set[str] = set()
    for item in notices:
        seq = str(_notice_item_value(item, "seq") or "").strip()
        if not seq:
            continue
        if seq == anchor:
            break
        output.add(seq)
    return output


def _remember_portal_notice_view(
    db: Database,
    *,
    kind: str,
    notices: list[Any],
    user_id: int | None = None,
) -> None:
    top_notice = notices[0] if notices else None
    db.update_sync_state(
        _telegram_notice_view_state_job(kind),
        last_run_at=now_utc_iso(),
        last_cursor_json={
            "top_seq": str(_notice_item_value(top_notice, "seq") or "").strip() or None,
            "top_title": str(_notice_item_value(top_notice, "title") or "").strip() or None,
            "notice_count": len(notices),
        },
        user_id=user_id,
    )


def _format_telegram_uos_notice(
    db: Database,
    kind: str,
    *,
    timezone_name: str = "Asia/Seoul",
    user_id: int | None = None,
) -> str:
    config = UOS_NOTICE_FEEDS.get(str(kind).strip().lower())
    if not isinstance(config, dict):
        return "[KU] 학교 공지\n- 지원하지 않는 공지 종류입니다."
    previous_state = db.get_sync_state(_telegram_notice_view_state_job(kind), user_id=user_id)
    previous_cursor = _json_load(previous_state.last_cursor_json)
    previous_top_seq = str(previous_cursor.get("top_seq") or "").strip() or None
    cursor = _refresh_portal_notice_snapshot(db, kind)
    snapshot = _portal_notice_snapshot_from_cursor(cursor)
    notices = list(snapshot.get("notices") or [])
    new_seqs = _portal_notice_new_seq_set(notices, previous_top_seq=previous_top_seq)
    lines = [
        f"[KU] {config['label']}",
        "",
        f"최근 공지 {TELEGRAM_UOS_NOTICE_LIMIT}개",
    ]
    if previous_top_seq and new_seqs:
        lines.append(f"- 새 공지: {len(new_seqs)}건")
    lines.append("")
    if not notices:
        attempt = cursor.get("last_attempt") if isinstance(cursor.get("last_attempt"), dict) else {}
        if attempt.get("ok") is False and not snapshot.get("fetched_at"):
            lines.append("- 최근 공지 목록을 불러오지 못했습니다.")
            lines.append(f"- 오류: {str(attempt.get('error') or 'unknown error').strip() or 'unknown error'}")
        else:
            lines.append("- 표시할 공지가 없습니다.")
    else:
        for item in notices:
            posted_on = str(_notice_item_value(item, "posted_on") or "").strip()
            prefix = "[NEW] " if str(_notice_item_value(item, "seq") or "").strip() in new_seqs else ""
            title = str(_notice_item_value(item, "title") or "").strip() or "공지"
            if posted_on:
                lines.append(f"- {prefix}{posted_on} | {title}")
            else:
                lines.append(f"- {prefix}{title}")
    if notices or snapshot.get("fetched_at"):
        _remember_portal_notice_view(db, kind=kind, notices=notices, user_id=user_id)
    status_lines = _portal_notice_status_lines(cursor, timezone_name=timezone_name)
    if status_lines:
        lines.extend(["", "상태"])
        lines.extend(status_lines)
    lines.extend(["", _portal_notice_source_line(cursor, timezone_name=timezone_name)])
    return "\n".join(lines)


def _format_telegram_class_summary_for_day(
    settings: Settings,
    db: Database,
    *,
    target_day_local: datetime,
    day_label: str,
    day_command: str,
    user_id: int | None = None,
) -> str:
    day_brief = DayBriefService(settings, db, user_id=user_id).build_day_brief(
        target_day_local=target_day_local,
        reference_day_local=target_day_local,
        max_classes=8,
        artifact_limit=TELEGRAM_DAY_SUMMARY_ARTIFACT_LIMIT,
        notification_limit=TELEGRAM_DAY_SUMMARY_NOTIFICATION_LIMIT,
        open_task_limit=TELEGRAM_DAY_SUMMARY_OPEN_TASK_LIMIT,
    )
    course_briefs = list(day_brief.course_briefs)
    lines = [f"[KU] {day_label} 수업 자료 요약 ({target_day_local.date().isoformat()})"]
    if not course_briefs:
        lines.extend(
            [
                "",
                "수업: 0",
                "- 시간표 기준 등록된 수업이 없습니다.",
                f"- 확인: {day_command}",
            ]
        )
        return "\n".join(lines)

    lines.append("")
    lines.append(f"수업: {len(course_briefs)}")
    for idx, course_brief in enumerate(course_briefs, start=1):
        item = course_brief.class_item
        start_time = item["start_local"].strftime("%H:%M")
        title = str(item["title"] or "").strip() or "Untitled class"
        location_text = str(item.get("location_text") or "").strip()
        if location_text and location_text != "TBD":
            lines.append(f"{idx}. {start_time} {title} @ {location_text}")
        else:
            lines.append(f"{idx}. {start_time} {title}")

        if course_brief.best_brief is None:
            if course_brief.latest_material_filename:
                lines.append(
                    f"- 자료는 있지만 요약은 아직 없습니다: {course_brief.latest_material_filename}"
                )
            else:
                lines.append("- 연결된 강의자료 요약이 아직 없습니다.")
            if course_brief.notice_titles:
                lines.append(f"- 공지: {'; '.join(course_brief.notice_titles)}")
            if course_brief.task_lines:
                lines.append(f"- 수업 과제: {'; '.join(course_brief.task_lines)}")
            if course_brief.file_task_lines:
                lines.append(f"- 파일 감지 과제: {'; '.join(course_brief.file_task_lines)}")
            continue

        brief_tag = _provenance_brief_tag(
            (
                course_brief.best_brief.get("provenance")
                if isinstance(course_brief.best_brief.get("provenance"), dict)
                else None
            ),
            fallback_source="llm_inferred",
        )
        if brief_tag:
            lines.append(f"- 출처: {brief_tag[1:-1]}")
        bullets = _clean_telegram_brief_bullets(course_brief.best_brief.get("bullets"), limit=3)
        if not bullets:
            material_files = _format_course_material_filenames(course_brief.material_filenames, limit=3)
            if material_files:
                lines.append(f"- 자료: {material_files}")
            elif course_brief.best_brief_filename:
                lines.append(f"- 최신 자료: {course_brief.best_brief_filename}")
        for bullet in bullets:
            lines.append(f"- {bullet}")
        question = _clean_telegram_brief_question(course_brief.best_brief.get("question"))
        if question:
            lines.append(f"복습: {question}")
        material_files = _format_course_material_filenames(course_brief.material_filenames, limit=3)
        if material_files:
            lines.append(f"자료: {material_files}")
        elif course_brief.best_brief_filename:
            lines.append(f"자료: {course_brief.best_brief_filename}")
        if course_brief.notice_titles:
            lines.append(f"- 공지: {'; '.join(course_brief.notice_titles)}")
        if course_brief.task_lines:
            lines.append(f"- 수업 과제: {'; '.join(course_brief.task_lines)}")
        if course_brief.file_task_lines:
            lines.append(f"- 파일 감지 과제: {'; '.join(course_brief.file_task_lines)}")
    return "\n".join(lines)


def _format_telegram_today_summary(settings: Settings, db: Database, *, user_id: int | None = None) -> str:
    return _format_telegram_class_summary_for_day(
        settings=settings,
        db=db,
        target_day_local=datetime.now(ZoneInfo(settings.timezone)),
        day_label="오늘",
        day_command="/today",
        user_id=user_id,
    )


def _format_telegram_tomorrow_summary(
    settings: Settings,
    db: Database,
    *,
    user_id: int | None = None,
) -> str:
    now_local = datetime.now(ZoneInfo(settings.timezone))
    return _format_telegram_class_summary_for_day(
        settings=settings,
        db=db,
        target_day_local=now_local + timedelta(days=1),
        day_label="내일",
        day_command="/tomorrow",
        user_id=user_id,
    )


def _resolve_telegram_lms_credentials(
    *,
    settings: Settings | None = None,
    db: Database | None = None,
    user_id: int | None = None,
    chat_id: str | int | None = None,
) -> tuple[str, str] | None:
    if db is not None and (user_id is not None or chat_id is not None):
        try:
            connections = db.list_moodle_connections(
                user_id=user_id,
                chat_id=str(chat_id).strip() if chat_id is not None else None,
                status="active",
                limit=10,
            )
        except Exception:
            connections = []
        for connection in connections:
            login_kind = str(connection.get("login_secret_kind") or "").strip()
            login_ref = str(connection.get("login_secret_ref") or "").strip()
            login_id = str(connection.get("username") or "").strip()
            if not login_id or not login_kind or not login_ref:
                continue
            try:
                password = default_secret_store(settings).read_secret(
                    ref=StoredSecretRef(kind=login_kind, ref=login_ref)
                )
            except SecretStoreError:
                continue
            if password:
                return login_id, password

    env_user_id = os.environ.get("KU_PORTAL_ID", "").strip()
    env_password = os.environ.get("KU_PORTAL_PW", "").strip()
    if env_user_id and env_password:
        return env_user_id, env_password

    dotenv_credentials = _resolve_telegram_lms_credentials_from_dotenv(settings)
    if dotenv_credentials is not None:
        return dotenv_credentials
    return None


def _resolve_telegram_lms_credentials_from_dotenv(
    settings: Settings | None,
) -> tuple[str, str] | None:
    if settings is None:
        return None

    candidates: list[Path] = []
    seen: set[Path] = set()

    def add_candidate(path: Path) -> None:
        resolved = path.expanduser().resolve()
        if resolved in seen:
            return
        seen.add(resolved)
        candidates.append(resolved)

    storage_root_dir = getattr(settings, "storage_root_dir", None)
    if storage_root_dir:
        add_candidate(Path(storage_root_dir) / ".env")

    database_path = getattr(settings, "database_path", None)
    if database_path:
        db_parent = Path(database_path).expanduser().resolve().parent
        add_candidate(db_parent / ".env")
        add_candidate(db_parent.parent / ".env")

    for env_path in candidates:
        try:
            if not env_path.is_file():
                continue
            values = dotenv_values(env_path)
        except OSError:
            continue
        env_user_id = str(values.get("KU_PORTAL_ID") or "").strip()
        env_password = str(values.get("KU_PORTAL_PW") or "").strip()
        if env_user_id and env_password:
            return env_user_id, env_password
    return None


@dataclass(frozen=True)
class _TelegramAssignmentHint:
    course_name: str
    source_label: str
    title: str
    due_at: str
    evidence: str


def _telegram_assignments_cache_job_name(
    *,
    login_id: str,
    user_id: int | None,
    chat_id: str | int | None,
) -> str:
    payload = {
        "login_id_hash": sha1(str(login_id or "").encode("utf-8")).hexdigest(),
        "user_id": _safe_int(user_id),
        "chat_id": str(chat_id or "").strip() or None,
    }
    digest = sha1(json.dumps(payload, ensure_ascii=True, sort_keys=True).encode("utf-8")).hexdigest()
    return f"{TELEGRAM_ASSIGNMENTS_CACHE_JOB_PREFIX}:{digest}"


def _get_cached_telegram_assignments(
    db: Database | None,
    *,
    job_name: str,
    user_id: int | None,
    now: datetime,
) -> str | None:
    if db is None:
        return None
    try:
        state = db.get_sync_state(job_name, user_id=user_id)
    except AttributeError:
        return None
    payload = _json_load(state.last_cursor_json)
    message = str(payload.get("message") or "").strip()
    generated_at = _parse_dt(str(payload.get("generated_at") or ""))
    if not message or generated_at is None:
        return None
    age = (now.astimezone(timezone.utc) - generated_at.astimezone(timezone.utc)).total_seconds()
    if age < 0 or age > TELEGRAM_ASSIGNMENTS_CACHE_TTL_SECONDS:
        return None
    return message


def _store_cached_telegram_assignments(
    db: Database | None,
    *,
    job_name: str,
    user_id: int | None,
    message: str,
    generated_at: datetime,
) -> None:
    if db is None or not str(message or "").strip():
        return
    try:
        db.update_sync_state(
            job_name,
            last_run_at=generated_at.astimezone(timezone.utc).replace(microsecond=0).isoformat(),
            last_cursor_json={
                "generated_at": generated_at.astimezone(timezone.utc).replace(microsecond=0).isoformat(),
                "ttl_seconds": TELEGRAM_ASSIGNMENTS_CACHE_TTL_SECONDS,
                "message": message,
            },
            user_id=user_id,
        )
    except Exception:
        logger.debug("failed to store telegram assignments cache", exc_info=True)


def _format_telegram_assignments(
    *,
    settings: Settings | None = None,
    db: Database | None = None,
    user_id: int | None = None,
    chat_id: str | int | None = None,
) -> str:
    """Render Canvas LMS todo + upcoming events for /assignments replies.

    Live hit on mylms.korea.ac.kr (cached 25 min by ku-portal-mcp lms session).
    """
    from ku_secretary.connectors import ku_lms

    credentials = _resolve_telegram_lms_credentials(
        settings=settings,
        db=db,
        user_id=user_id,
        chat_id=chat_id,
    )
    if credentials is None:
        return "KU_PORTAL_ID / KU_PORTAL_PW 환경변수가 비어 있습니다."
    login_id, password = credentials
    cache_job_name = _telegram_assignments_cache_job_name(
        login_id=login_id,
        user_id=user_id,
        chat_id=chat_id,
    )
    cache_now = datetime.now(timezone.utc)
    cached_message = _get_cached_telegram_assignments(
        db,
        job_name=cache_job_name,
        user_id=user_id,
        now=cache_now,
    )
    if cached_message is not None:
        return cached_message

    try:
        session = ku_lms.login(user_id=login_id, password=password)
    except Exception as exc:  # noqa: BLE001
        return f"LMS 로그인 실패: {exc}"

    try:
        todos = ku_lms.get_todo(session)
    except Exception as exc:  # noqa: BLE001
        return f"LMS 할 일 조회 실패: {exc}"

    try:
        events = ku_lms.get_upcoming_events(session)
    except Exception:
        events = []
    timezone_name = str(getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul")
    reference_local = datetime.now(ZoneInfo(timezone_name))

    def _pretty_dt(value: str | None) -> str:
        if not value:
            return ""
        v = value.replace("T", " ")
        return v[:16] if len(v) > 16 else v

    def _assignment_key(item: dict[str, Any]) -> str:
        assignment = item.get("assignment") if isinstance(item.get("assignment"), dict) else item
        raw_id = assignment.get("id") or assignment.get("assignment_id") or item.get("assignment_id")
        if raw_id is not None:
            return f"id:{raw_id}"
        url = str(assignment.get("html_url") or assignment.get("url") or item.get("html_url") or "").strip()
        if url:
            return f"url:{url}"
        return "|".join(
            [
                str(assignment.get("name") or item.get("title") or "").strip().lower(),
                str(assignment.get("due_at") or item.get("due_at") or "").strip(),
            ]
        )

    def _assignment_event_key(title: Any, value: Any) -> tuple[str, str] | None:
        title_key = _normalize_task_title_key(str(title or ""))
        if not title_key:
            return None
        parsed = _parse_dt(str(value or ""))
        if parsed is not None:
            try:
                local = parsed.astimezone(ZoneInfo(timezone_name))
            except Exception:
                local = parsed
            return title_key, local.replace(second=0, microsecond=0).isoformat()
        raw = str(value or "").strip()
        return (title_key, raw) if raw else None

    seen_assignment_keys = {
        _assignment_key(item)
        for item in todos
        if isinstance(item, dict) and _assignment_key(item)
    }
    assignment_event_keys: set[tuple[str, str]] = set()
    for item in todos:
        if not isinstance(item, dict):
            continue
        assignment = item.get("assignment") if isinstance(item.get("assignment"), dict) else {}
        title = assignment.get("name") or item.get("title")
        due_at = assignment.get("due_at") or item.get("due_at")
        key = _assignment_event_key(title, due_at)
        if key is not None:
            assignment_event_keys.add(key)
    course_assignment_rows: list[dict[str, Any]] = []
    source_hint_rows: list[_TelegramAssignmentHint] = []
    seen_source_hint_keys: set[tuple[str, str, str, str]] = set()
    scanned_courses = 0
    assignment_scan_failures = 0
    source_scan_failures = 0
    try:
        courses = ku_lms.get_courses(session) or []
    except Exception:
        courses = []
        assignment_scan_failures += 1
    course_ids: list[int] = []
    course_name_by_id: dict[int, str] = {}
    for course in _lms_scannable_courses(courses)[:TELEGRAM_LMS_COURSE_SCAN_LIMIT]:
        if not isinstance(course, dict):
            continue
        cid, course_name = _lms_course_id_and_name(course)
        if cid is None:
            continue
        course_ids.append(cid)
        course_name_by_id[cid] = course_name
        scanned_courses += 1
        try:
            assignments = ku_lms.get_assignments(session, cid, upcoming_only=True) or []
        except Exception:
            assignment_scan_failures += 1
            continue
        for assignment in assignments:
            if not isinstance(assignment, dict):
                continue
            key = _assignment_key(assignment)
            if key and key in seen_assignment_keys:
                continue
            if key:
                seen_assignment_keys.add(key)
            name = str(assignment.get("name") or assignment.get("title") or "(제목 없음)").strip()
            due_at = _pretty_dt(str(assignment.get("due_at") or ""))
            event_key = _assignment_event_key(name, assignment.get("due_at") or due_at)
            if event_key is not None:
                assignment_event_keys.add(event_key)
            course_assignment_rows.append(
                {
                    "course_id": cid,
                    "course_name": course_name,
                    "title": name,
                    "due_at": due_at,
                }
            )

    def add_source_hints(hints: list[_TelegramAssignmentHint]) -> None:
        for hint in hints:
            if not _contains_submission_action_hints(hint.title, hint.evidence):
                continue
            due_dt = _parse_dt(str(hint.due_at or ""))
            if due_dt is None:
                continue
            local_tz = ZoneInfo(timezone_name)
            due_local = due_dt.astimezone(local_tz)
            if due_local < reference_local and due_local.year < reference_local.year:
                corrected_dt = None
                if hint.evidence:
                    corrected_text = re.sub(
                        rf"\b{due_local.year}\b",
                        str(reference_local.year),
                        hint.evidence,
                        count=1,
                    )
                    corrected_due = _coerce_due_iso(
                        corrected_text,
                        timezone_name=timezone_name,
                        reference_local=reference_local,
                        default_end_of_day=True,
                    )
                    corrected_dt = _parse_dt(corrected_due)
                corrected = (
                    corrected_dt.astimezone(local_tz)
                    if corrected_dt is not None
                    else due_local.replace(year=reference_local.year)
                )
                if reference_local <= corrected <= reference_local + timedelta(days=180):
                    hint = _TelegramAssignmentHint(
                        course_name=hint.course_name,
                        source_label=hint.source_label,
                        title=hint.title,
                        due_at=corrected.isoformat(),
                        evidence=hint.evidence,
                    )
                    due_local = corrected
            if due_local < reference_local:
                continue
            key = (
                _normalize_task_title_key(hint.course_name),
                _normalize_task_title_key(hint.title),
                str(hint.due_at or "").strip(),
                str(hint.source_label or "").strip(),
            )
            if key in seen_source_hint_keys:
                continue
            seen_source_hint_keys.add(key)
            source_hint_rows.append(hint)
            event_key = _assignment_event_key(hint.title, hint.due_at)
            if event_key is not None:
                assignment_event_keys.add(event_key)

    try:
        announcements = ku_lms.get_announcements(session, course_ids) if course_ids else []
    except Exception:
        announcements = []
        source_scan_failures += 1
    for announcement in announcements or []:
        if not isinstance(announcement, dict):
            continue
        cid = _lms_course_id_of_item(announcement)
        course_name = course_name_by_id.get(cid or -1, "강의")
        title = str(announcement.get("title") or announcement.get("subject") or "공지").strip()
        text_lines = [
            title,
            *_lms_text_lines_from_dict(
                announcement,
                (
                    "message",
                    "body",
                    "description",
                    "content",
                ),
            ),
        ]
        add_source_hints(
            _telegram_assignment_hints_from_text(
                title=title,
                course_name=course_name,
                source_label="공지",
                text_lines=text_lines,
                timezone_name=timezone_name,
                reference_local=reference_local,
            )
        )

    for cid in course_ids:
        course_name = course_name_by_id.get(cid, f"course {cid}")
        try:
            modules = ku_lms.get_modules(session, cid, include_items=True) or []
        except Exception:
            modules = []
            source_scan_failures += 1
        for module in modules:
            if not isinstance(module, dict):
                continue
            module_name = str(module.get("name") or module.get("title") or "").strip()
            module_lines = _lms_text_lines_from_dict(
                module,
                ("name", "title", "description", "unlock_at", "due_at"),
            )
            items = module.get("items")
            if not isinstance(items, list):
                continue
            for item in items:
                if not isinstance(item, dict):
                    continue
                title = str(item.get("title") or item.get("name") or module_name or "자료").strip()
                text_lines = [
                    *module_lines,
                    title,
                    *_lms_text_lines_from_dict(
                        item,
                        (
                            "title",
                            "name",
                            "type",
                            "content_type",
                            "html_url",
                            "url",
                            "due_at",
                            "unlock_at",
                            "completion_requirement",
                        ),
                    ),
                ]
                add_source_hints(
                    _telegram_assignment_hints_from_text(
                        title=title,
                        course_name=course_name,
                        source_label="모듈/자료",
                        text_lines=text_lines,
                        timezone_name=timezone_name,
                        reference_local=reference_local,
                    )
                )

        try:
            boards = ku_lms.list_boards(session, cid) or []
        except Exception:
            boards = []
            source_scan_failures += 1
        for board in boards[:TELEGRAM_LMS_BOARD_SCAN_LIMIT_PER_COURSE]:
            if not isinstance(board, dict):
                continue
            bid_raw = board.get("id") or board.get("board_id")
            try:
                bid = int(bid_raw)
            except (TypeError, ValueError):
                continue
            board_name = str(board.get("name") or board.get("title") or "게시판").strip()
            try:
                resp = ku_lms.list_board_posts(session, cid, bid)
            except Exception:
                source_scan_failures += 1
                continue
            detail_count = 0
            for post in _lms_board_posts_from_response(resp)[:TELEGRAM_LMS_BOARD_POST_LIMIT_PER_BOARD]:
                title = str(post.get("title") or post.get("subject") or "게시글").strip()
                post_text_lines = [
                    board_name,
                    title,
                    *_lms_text_lines_from_dict(
                        post,
                        ("title", "subject", "message", "body", "content"),
                    ),
                ]
                post_id = _lms_post_id(post)
                if post_id is not None and detail_count < TELEGRAM_LMS_ASSIGNMENT_BOARD_DETAIL_LIMIT_PER_BOARD:
                    detail_count += 1
                    try:
                        detail = ku_lms.get_board_post(session, cid, bid, post_id) or {}
                    except Exception:
                        detail = {}
                        source_scan_failures += 1
                    if isinstance(detail, dict):
                        post_text_lines.extend(
                            _lms_text_lines_from_dict(
                                detail,
                                (
                                    "title",
                                    "subject",
                                    "message",
                                    "body",
                                    "content",
                                    "description",
                                    "attachments",
                                    "files",
                                ),
                            )
                        )
                add_source_hints(
                    _telegram_assignment_hints_from_text(
                        title=title,
                        course_name=course_name,
                        source_label=f"게시판 {board_name}",
                        text_lines=post_text_lines,
                        timezone_name=timezone_name,
                        reference_local=reference_local,
                    )
                )

    if not todos and not events and not course_assignment_rows and not source_hint_rows:
        lines = ["[KU] 내야 할 과제", "- 마감 임박한 과제가 없습니다."]
        if scanned_courses:
            lines.append(f"- 확인: {scanned_courses}개 과목의 과제 목록과 공지/자료/게시판 제출 항목을 직접 확인했습니다.")
        if assignment_scan_failures or source_scan_failures:
            lines.append(
                f"- 참고: 일부 조회 실패 과제 {assignment_scan_failures}건 / "
                f"공지·자료·게시판 {source_scan_failures}건"
            )
        message = "\n".join(lines)
        _store_cached_telegram_assignments(
            db,
            job_name=cache_job_name,
            user_id=user_id,
            message=message,
            generated_at=cache_now,
        )
        return message

    lines = ["[KU] 내야 할 과제"]
    known_course_order = list(course_ids)
    course_groups: dict[int | str, dict[str, Any]] = {}

    def course_group_key(course_id: int | None, fallback_name: str = "기타") -> int | str:
        if course_id is not None:
            return course_id
        return fallback_name

    def ensure_course_group(course_id: int | None, course_name: str = "기타") -> dict[str, Any]:
        key = course_group_key(course_id, course_name)
        if key not in course_groups:
            course_groups[key] = {
                "name": course_name_by_id.get(course_id or -1, course_name),
                "todos": [],
                "course_assignments": [],
                "source_hints": [],
                "events": [],
            }
        return course_groups[key]

    if todos:
        for t in todos:
            if not isinstance(t, dict):
                continue
            a = t.get("assignment") if isinstance(t.get("assignment"), dict) else {}
            cid = _lms_course_id_of_item(a) or _lms_course_id_of_item(t)
            course_name = course_name_by_id.get(cid or -1, "Canvas 할 일")
            name = (a.get("name") or t.get("title") or "(제목 없음)").strip()
            due_at = _pretty_dt(a.get("due_at"))
            ensure_course_group(cid, course_name)["todos"].append({"title": name, "due_at": due_at})

    for row in course_assignment_rows:
        cid = row.get("course_id") if isinstance(row.get("course_id"), int) else None
        course_name = str(row.get("course_name") or course_name_by_id.get(cid or -1, "강의"))
        ensure_course_group(cid, course_name)["course_assignments"].append(row)

    for hint in source_hint_rows:
        matched_cid = None
        for cid, course_name in course_name_by_id.items():
            if course_name == hint.course_name:
                matched_cid = cid
                break
        ensure_course_group(matched_cid, hint.course_name)["source_hints"].append(hint)

    if events:
        for ev in events[:8]:
            if not isinstance(ev, dict):
                continue
            cid = _lms_course_id_of_item(ev)
            course_name = course_name_by_id.get(cid or -1, "다가오는 이벤트")
            title = (ev.get("title") or "(제목 없음)").strip()
            start = _pretty_dt(ev.get("start_at"))
            event_key = _assignment_event_key(title, ev.get("start_at") or start)
            if event_key is not None and event_key in assignment_event_keys:
                continue
            ensure_course_group(cid, course_name)["events"].append({"title": title, "start_at": start})

    ordered_keys: list[int | str] = [cid for cid in known_course_order if cid in course_groups]
    ordered_keys.extend(key for key in course_groups if key not in ordered_keys)
    rendered_course_items = 0
    rendered_source_hints = 0
    for key in ordered_keys:
        group = course_groups[key]
        lines.append("")
        lines.append(f"[{_compact_lms_course_name(group['name'])}]")
        if group["todos"] or group["course_assignments"]:
            lines.append("과제")
        for row in group["todos"]:
            details = []
            if row.get("due_at"):
                details.append(f"마감 {_format_lms_list_dt(row['due_at'], timezone_name=timezone_name)}")
            _append_lms_list_item(lines, title=row["title"], details=details)
        for row in group["course_assignments"]:
            if rendered_course_items >= TELEGRAM_LMS_ASSIGNMENT_DISPLAY_LIMIT:
                continue
            details = []
            if row.get("due_at"):
                details.append(f"마감 {_format_lms_list_dt(row['due_at'], timezone_name=timezone_name)}")
            _append_lms_list_item(lines, title=row["title"], details=details)
            rendered_course_items += 1
        if group["source_hints"]:
            lines.append("공지/자료/게시판 제출 항목")
        for hint in group["source_hints"]:
            if rendered_source_hints >= TELEGRAM_LMS_ASSIGNMENT_HINT_DISPLAY_LIMIT:
                continue
            details = [str(hint.source_label or "").strip()]
            if hint.due_at:
                details.append(f"마감 {_format_lms_list_dt(hint.due_at, timezone_name=timezone_name)}")
            _append_lms_list_item(lines, title=hint.title, details=details)
            rendered_source_hints += 1
        if group["events"]:
            lines.append("이벤트")
        for event in group["events"]:
            details = []
            if event.get("start_at"):
                details.append(f"일정 {_format_lms_list_dt(event['start_at'], timezone_name=timezone_name)}")
            _append_lms_list_item(lines, title=event["title"], details=details)
    remaining_course_items = len(course_assignment_rows) - rendered_course_items
    remaining_source_hints = len(source_hint_rows) - rendered_source_hints
    if remaining_course_items > 0 or remaining_source_hints > 0:
        lines.append("")
        omitted = []
        if remaining_course_items > 0:
            omitted.append(f"과제 {remaining_course_items}건")
        if remaining_source_hints > 0:
            omitted.append(f"공지/자료/게시판 제출 항목 {remaining_source_hints}건")
        lines.append(f"- 외 {' / '.join(omitted)}")

    if scanned_courses:
        lines.append("")
        lines.append(f"확인: {scanned_courses}개 과목의 과제 목록과 공지/자료/게시판 제출 항목을 직접 확인했습니다.")

    if assignment_scan_failures or source_scan_failures:
        lines.append("")
        lines.append(
            f"참고: 일부 조회 실패 과제 {assignment_scan_failures}건 / "
            f"공지·자료·게시판 {source_scan_failures}건"
        )

    message = "\n".join(lines)
    _store_cached_telegram_assignments(
        db,
        job_name=cache_job_name,
        user_id=user_id,
        message=message,
        generated_at=cache_now,
    )
    return message


def _format_lms_list_dt(value: str | None, *, timezone_name: str = "Asia/Seoul") -> str:
    text = str(value or "").strip()
    if not text:
        return ""
    parsed = _parse_dt(text)
    if parsed is not None:
        try:
            local = parsed.astimezone(ZoneInfo(timezone_name))
        except Exception:
            local = parsed
        return local.strftime("%m/%d %H:%M")
    text = text.replace("T", " ").replace("Z", "")
    match = re.match(r"(?P<year>\d{4})-(?P<month>\d{2})-(?P<day>\d{2})\s+(?P<hm>\d{2}:\d{2})", text)
    if match:
        return f"{match.group('month')}/{match.group('day')} {match.group('hm')}"
    return text[:16] if len(text) > 16 else text


def _truncate_lms_text(value: Any, limit: int = 64) -> str:
    text = re.sub(r"\s+", " ", str(value or "").strip())
    if len(text) <= max(int(limit), 1):
        return text
    return text[: max(int(limit), 1) - 1].rstrip() + "…"


def _compact_lms_course_name(value: Any) -> str:
    text = re.sub(r"\s+", " ", str(value or "").strip())
    if not text:
        return "강의"
    text = re.sub(r"^\s*[A-Za-z0-9][A-Za-z0-9_-]*\s+", "", text)
    text = re.sub(r"^\s*\([^)]*\)\s*", "", text)

    def strip_ascii_parenthetical(match: re.Match[str]) -> str:
        inner = match.group(1)
        return "" if re.search(r"[A-Za-z]", inner) else match.group(0)

    text = re.sub(r"\(([^)]*)\)", strip_ascii_parenthetical, text)
    text = re.sub(r"[-_\s]*(\d{2})\s*분반\s*$", r" \1분반", text)
    text = re.sub(r"\s+", " ", text).strip(" -_/")
    return _truncate_lms_text(text or str(value or "").strip() or "강의", 36)


def _lms_submission_status_label(value: Any) -> str:
    raw = str(value or "").strip().lower()
    return {
        "submitted": "제출됨",
        "graded": "채점됨",
        "pending_review": "검토중",
    }.get(raw, raw or "상태 확인")


def _append_lms_list_item(
    lines: list[str],
    *,
    title: Any,
    details: list[str],
    title_limit: int = 62,
) -> None:
    lines.append(f"- {_truncate_lms_text(title, title_limit)}")
    compact_details = [str(item or "").strip() for item in details if str(item or "").strip()]
    if compact_details:
        lines.append(f"  {' | '.join(compact_details)}")


def _format_telegram_submitted_assignments(
    *,
    settings: Settings | None = None,
    db: Database | None = None,
    user_id: int | None = None,
    chat_id: str | int | None = None,
) -> str:
    """Render submitted/graded Canvas assignments for /submitted replies."""
    from ku_secretary.connectors import ku_lms

    credentials = _resolve_telegram_lms_credentials(
        settings=settings,
        db=db,
        user_id=user_id,
        chat_id=chat_id,
    )
    if credentials is None:
        return "KU_PORTAL_ID / KU_PORTAL_PW 환경변수가 비어 있습니다."
    login_id, password = credentials

    try:
        session = ku_lms.login(user_id=login_id, password=password)
    except Exception as exc:  # noqa: BLE001
        return f"LMS 로그인 실패: {exc}"

    try:
        courses = ku_lms.get_courses(session) or []
    except Exception as exc:  # noqa: BLE001
        return f"강의 목록 조회 실패: {exc}"
    if not courses:
        return "[KU] 제출 완료 과제\n- 등록된 강의가 없습니다."
    timezone_name = str(getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul")

    def _pretty_dt(value: str | None) -> str:
        if not value:
            return ""
        v = value.replace("T", " ")
        return v[:16] if len(v) > 16 else v

    rows: list[dict[str, Any]] = []
    scanned_courses = 0
    failures = 0
    for course in _lms_scannable_courses(courses)[:TELEGRAM_LMS_COURSE_SCAN_LIMIT]:
        if not isinstance(course, dict):
            continue
        cid, course_name = _lms_course_id_and_name(course)
        if cid is None:
            continue
        scanned_courses += 1
        try:
            submissions = ku_lms.get_submissions(session, cid) or []
        except Exception:
            failures += 1
            continue
        for submission in submissions:
            if not isinstance(submission, dict):
                continue
            assignment = submission.get("assignment") if isinstance(submission.get("assignment"), dict) else {}
            submitted_at = str(submission.get("submitted_at") or "").strip()
            workflow_state = str(submission.get("workflow_state") or "").strip().lower()
            grade = str(submission.get("grade") or submission.get("entered_grade") or "").strip()
            score = submission.get("score")
            if not submitted_at and workflow_state not in {"submitted", "graded", "pending_review"}:
                continue
            if bool(submission.get("missing")) and not submitted_at:
                continue
            title = str(
                assignment.get("name")
                or submission.get("assignment_name")
                or submission.get("title")
                or "(제목 없음)"
            ).strip()
            rows.append(
                {
                    "course_name": course_name,
                    "course_id": cid,
                    "title": title,
                    "submitted_at": submitted_at,
                    "workflow_state": workflow_state or "submitted",
                    "due_at": str(assignment.get("due_at") or submission.get("cached_due_date") or "").strip(),
                    "grade": grade,
                    "score": score,
                    "late": bool(submission.get("late")),
                }
            )

    rows.sort(key=lambda item: str(item.get("submitted_at") or ""), reverse=True)
    lines = ["[KU] 제출 완료 과제"]
    if not rows:
        lines.append("- 제출 완료로 확인된 과제가 없습니다.")
        if scanned_courses:
            lines.append(f"- 확인: {scanned_courses}개 과목의 제출 상태를 직접 확인했습니다.")
        if failures:
            lines.append(f"- 참고: 일부 과목 제출 상태 조회 실패 {failures}건")
        return "\n".join(lines)

    lines.append("")
    lines.append(f"제출 완료 ({len(rows)}건)")
    grouped_rows: dict[int, list[dict[str, Any]]] = {}
    course_names: dict[int, str] = {}
    for row in rows:
        cid = int(row["course_id"])
        grouped_rows.setdefault(cid, []).append(row)
        course_names[cid] = str(row.get("course_name") or f"course {cid}")
    rendered = 0
    for course in _lms_scannable_courses(courses)[:TELEGRAM_LMS_COURSE_SCAN_LIMIT]:
        cid, fallback_name = _lms_course_id_and_name(course)
        if cid is None or cid not in grouped_rows:
            continue
        lines.append("")
        lines.append(f"[{_compact_lms_course_name(course_names.get(cid, fallback_name))}]")
        for row in grouped_rows[cid]:
            if rendered >= TELEGRAM_LMS_SUBMITTED_DISPLAY_LIMIT:
                continue
            lines.append(f"- {_truncate_lms_text(row['title'], 62)}")
            pieces = []
            submitted_at = _pretty_dt(str(row.get("submitted_at") or ""))
            if submitted_at:
                pieces.append(f"제출 {_format_lms_list_dt(submitted_at, timezone_name=timezone_name)}")
            workflow_state = str(row.get("workflow_state") or "").strip()
            if workflow_state:
                pieces.append(_lms_submission_status_label(workflow_state))
            due_at = _pretty_dt(str(row.get("due_at") or ""))
            if due_at:
                pieces.append(f"마감 {_format_lms_list_dt(due_at, timezone_name=timezone_name)}")
            grade = str(row.get("grade") or "").strip()
            score = row.get("score")
            if grade:
                pieces.append(f"성적 {grade}")
            elif score not in (None, ""):
                pieces.append(f"점수 {score}")
            if bool(row.get("late")):
                pieces.append("지각 제출")
            lines.append(f"  {' | '.join(item for item in pieces if item)}")
            rendered += 1
    remaining = len(rows) - rendered
    if remaining > 0:
        lines.append(f"- 외 {remaining}건")
    lines.append("")
    lines.append(f"확인: {scanned_courses}개 과목의 제출 상태를 직접 확인했습니다.")
    if failures:
        lines.append(f"참고: 일부 과목 제출 상태 조회 실패 {failures}건")
    return "\n".join(lines)


def _lms_course_id_and_name(course: dict[str, Any]) -> tuple[int | None, str]:
    cid_raw = course.get("id") or course.get("course_id")
    try:
        cid = int(cid_raw)
    except (TypeError, ValueError):
        cid = None
    name = (
        str(course.get("name") or course.get("course_name") or course.get("course_code") or "")
        .strip()
        or (f"course {cid}" if cid is not None else "강의")
    )
    return cid, name


def _lms_course_is_scannable(course: dict[str, Any]) -> bool:
    if bool(course.get("access_restricted_by_date")):
        return False
    cid, name = _lms_course_id_and_name(course)
    if cid is None:
        return False
    # Canvas may return old/restricted shell courses with only an id. Those
    # 401 on course detail endpoints, so exclude them from live scans.
    if str(name).strip() == str(cid):
        return False
    workflow_state = str(course.get("workflow_state") or "").strip().lower()
    if workflow_state and workflow_state not in {"available", "completed", "created"}:
        return False
    return True


def _lms_scannable_courses(courses: Any) -> list[dict[str, Any]]:
    if not isinstance(courses, list):
        return []
    return [
        course
        for course in courses
        if isinstance(course, dict) and _lms_course_is_scannable(course)
    ]


def _lms_board_posts_from_response(response: Any) -> list[dict[str, Any]]:
    if isinstance(response, list):
        return [item for item in response if isinstance(item, dict)]
    if not isinstance(response, dict):
        return []
    for key in ("posts", "items", "results", "data"):
        value = response.get(key)
        if isinstance(value, list):
            return [item for item in value if isinstance(item, dict)]
    return []


def _lms_first_present(item: dict[str, Any], keys: tuple[str, ...]) -> Any:
    for key in keys:
        value = item.get(key)
        if value not in (None, ""):
            return value
    return None


def _lms_course_id_of_item(item: dict[str, Any]) -> int | None:
    for key in ("course_id", "courseId", "context_id"):
        value = item.get(key)
        if value is None:
            continue
        try:
            return int(value)
        except (TypeError, ValueError):
            pass
    context_code = str(item.get("context_code") or "").strip()
    if context_code.startswith("course_"):
        try:
            return int(context_code.split("_", 1)[1])
        except ValueError:
            return None
    return None


def _lms_post_id(item: dict[str, Any]) -> int | None:
    for key in ("id", "post_id", "postId", "article_id", "articleId"):
        value = item.get(key)
        if value is None:
            continue
        try:
            return int(value)
        except (TypeError, ValueError):
            continue
    return None


def _lms_plain_text(value: Any) -> str:
    text = str(value or "").strip()
    if not text:
        return ""
    if "<" in text and ">" in text:
        text = MATERIAL_HTML_DROP_BLOCK_RE.sub(" ", text)
        text = MATERIAL_BREAK_TAG_RE.sub("\n", text)
        text = MATERIAL_BLOCK_TAG_RE.sub("\n", text)
        text = MATERIAL_HTML_TAG_RE.sub(" ", text)
    return re.sub(r"\s+", " ", html.unescape(text)).strip()


def _lms_text_lines_from_dict(item: dict[str, Any], keys: tuple[str, ...]) -> list[str]:
    lines: list[str] = []
    for key in keys:
        value = item.get(key)
        if isinstance(value, list):
            for subvalue in value:
                if isinstance(subvalue, dict):
                    nested = _lms_plain_text(
                        _lms_first_present(
                            subvalue,
                            ("display_name", "filename", "name", "title", "content", "body"),
                        )
                    )
                    if nested:
                        lines.append(nested)
                else:
                    nested = _lms_plain_text(subvalue)
                    if nested:
                        lines.append(nested)
            continue
        if isinstance(value, dict):
            nested = _lms_plain_text(
                _lms_first_present(
                    value,
                    ("display_name", "filename", "name", "title", "content", "body"),
                )
            )
            if nested:
                lines.append(nested)
            continue
        text = _lms_plain_text(value)
        if text:
            lines.append(text)
    return lines


def _telegram_assignment_hints_from_text(
    *,
    title: str,
    course_name: str,
    source_label: str,
    text_lines: list[str],
    timezone_name: str,
    reference_local: datetime,
) -> list[_TelegramAssignmentHint]:
    extracted_text = "\n".join(line for line in text_lines if str(line or "").strip())
    if not _contains_material_task_hints(title, extracted_text):
        return []
    rows: list[_TelegramAssignmentHint] = []
    for task in _heuristic_material_deadline_tasks(
        title=title,
        course_name=course_name,
        extracted_text=extracted_text,
        timezone_name=timezone_name,
        reference_local=reference_local,
    ):
        due_at = str(task.get("due_at") or "").strip()
        task_title = str(task.get("title") or title or "과제").strip()
        source_title = _clean_material_task_title(title)
        if re.search(r"(?i)\bhw\s*#?\s*\d+", str(title or "")):
            source_title = _truncate_lms_text(title, 80)
        if source_title and (
            len(task_title) > 80
            or task_title.startswith("안녕하세요")
            or task_title.startswith("#")
        ):
            task_title = source_title
        if not due_at or not task_title:
            continue
        rows.append(
            _TelegramAssignmentHint(
                course_name=course_name,
                source_label=source_label,
                title=task_title,
                due_at=due_at,
                evidence=str(task.get("evidence") or "").strip(),
            )
        )
    return rows


def _append_limited_telegram_line(lines: list[str], line: str) -> bool:
    candidate = "\n".join([*lines, line])
    if len(candidate) > TELEGRAM_LMS_MESSAGE_SOFT_LIMIT:
        return False
    lines.append(line)
    return True


def _format_telegram_lms_board(
    *,
    settings: Settings | None = None,
    db: Database | None = None,
    user_id: int | None = None,
    chat_id: str | int | None = None,
) -> str:
    """Render recent course-level announcements + board posts.

    Hits Canvas (mylms.korea.ac.kr) per-course. Conservative caps so the
    response stays under Telegram's 4 KB limit and we do not hammer LMS.
    """
    from ku_secretary.connectors import ku_lms

    credentials = _resolve_telegram_lms_credentials(
        settings=settings,
        db=db,
        user_id=user_id,
        chat_id=chat_id,
    )
    if credentials is None:
        return "KU_PORTAL_ID / KU_PORTAL_PW 환경변수가 비어 있습니다."
    login_id, password = credentials

    try:
        session = ku_lms.login(user_id=login_id, password=password)
    except Exception as exc:  # noqa: BLE001
        return f"LMS 로그인 실패: {exc}"

    try:
        courses = ku_lms.get_courses(session) or []
    except Exception as exc:  # noqa: BLE001
        return f"강의 목록 조회 실패: {exc}"
    if not courses:
        return "[KU] 과목 게시판\n- 등록된 강의가 없습니다."

    course_ids: list[int] = []
    course_name: dict[int, str] = {}
    for c in _lms_scannable_courses(courses):
        cid, nm = _lms_course_id_and_name(c)
        if cid is None:
            continue
        course_ids.append(cid)
        course_name[cid] = nm

    course_ids = course_ids[:TELEGRAM_LMS_COURSE_SCAN_LIMIT]
    timezone_name = str(getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul")

    try:
        announcements = ku_lms.get_announcements(session, course_ids) or []
    except Exception:
        announcements = []

    bucket: dict[int, list[tuple[str, str, str]]] = {cid: [] for cid in course_ids}

    for ann in announcements:
        cid = _lms_course_id_of_item(ann)
        if cid is None or cid not in bucket:
            continue
        title = (ann.get("title") or ann.get("subject") or "").strip() or "(제목 없음)"
        when = _format_lms_list_dt(
            ann.get("posted_at") or ann.get("created_at") or ann.get("delayed_post_at")
            or "",
            timezone_name=timezone_name,
        )
        bucket[cid].append(("공지", when, title))

    for cid in course_ids:
        try:
            boards = ku_lms.list_boards(session, cid) or []
        except Exception:
            boards = []
        for b in boards[:TELEGRAM_LMS_BOARD_SCAN_LIMIT_PER_COURSE]:
            bid_raw = b.get("id") or b.get("board_id")
            try:
                bid = int(bid_raw)
            except (TypeError, ValueError):
                continue
            try:
                resp = ku_lms.list_board_posts(session, cid, bid)
            except Exception:
                continue
            posts = _lms_board_posts_from_response(resp)
            for p in posts[:TELEGRAM_LMS_BOARD_POST_LIMIT_PER_BOARD]:
                title = (p.get("title") or p.get("subject") or "").strip() or "(제목 없음)"
                when = _format_lms_list_dt(
                    p.get("posted_at") or p.get("created_at") or p.get("date") or "",
                    timezone_name=timezone_name,
                )
                board_name = (b.get("name") or b.get("title") or "보드").strip()
                bucket[cid].append((board_name, when, title))

    lines = ["[KU] 과목별 게시판/공지"]
    rendered_any = False
    for cid in course_ids:
        items = bucket[cid]
        if not items:
            continue
        rendered_any = True
        lines.append("")
        lines.append(f"[{_compact_lms_course_name(course_name.get(cid, cid))}]")
        for kind, when, title in items[:6]:
            before_len = len(lines)
            _append_lms_list_item(
                lines,
                title=title,
                details=[str(kind or "").strip(), str(when or "").strip()],
            )
            candidate = "\n".join(lines)
            if len(candidate) > TELEGRAM_LMS_MESSAGE_SOFT_LIMIT:
                del lines[before_len:]
                lines.append("- 외 항목은 길이 제한 때문에 생략했습니다.")
                return "\n".join(lines)

    if not rendered_any:
        lines.append("- 최근 글이 없습니다.")
    lines.append("")
    lines.append(
        f"- 확인: {len(course_ids)}개 과목, 과목당 최대 "
        f"{TELEGRAM_LMS_BOARD_SCAN_LIMIT_PER_COURSE}개 게시판을 직접 확인했습니다."
    )
    return "\n".join(lines)


def _format_telegram_lms_materials(
    *,
    settings: Settings | None = None,
    db: Database | None = None,
    user_id: int | None = None,
    chat_id: str | int | None = None,
) -> str:
    """Render likely LMS material locations by scanning each Canvas course."""
    from ku_secretary.connectors import ku_lms

    credentials = _resolve_telegram_lms_credentials(
        settings=settings,
        db=db,
        user_id=user_id,
        chat_id=chat_id,
    )
    if credentials is None:
        return "KU_PORTAL_ID / KU_PORTAL_PW 환경변수가 비어 있습니다."
    login_id, password = credentials

    try:
        session = ku_lms.login(user_id=login_id, password=password)
    except Exception as exc:  # noqa: BLE001
        return f"LMS 로그인 실패: {exc}"

    try:
        courses = ku_lms.get_courses(session) or []
    except Exception as exc:  # noqa: BLE001
        return f"강의 목록 조회 실패: {exc}"
    if not courses:
        return "[KU] 강의자료 위치\n- 등록된 강의가 없습니다."

    lines = ["[KU] 강의자료 위치", ""]
    scanned_courses = 0
    module_failures = 0
    board_failures = 0
    rendered_any = False

    for course in _lms_scannable_courses(courses)[:TELEGRAM_LMS_COURSE_SCAN_LIMIT]:
        if not isinstance(course, dict):
            continue
        cid, course_name = _lms_course_id_and_name(course)
        if cid is None:
            continue
        scanned_courses += 1
        course_lines: list[str] = []
        course_item_count = 0
        seen_titles: set[str] = set()

        try:
            modules = ku_lms.get_modules(session, cid, include_items=True) or []
        except Exception:
            modules = []
            module_failures += 1
        for module in modules:
            if not isinstance(module, dict):
                continue
            module_name = str(module.get("name") or module.get("title") or "").strip()
            items = module.get("items")
            if not isinstance(items, list):
                continue
            for item in items:
                if not isinstance(item, dict):
                    continue
                title = str(item.get("title") or item.get("name") or "").strip()
                if not title:
                    continue
                item_type = str(item.get("type") or item.get("content_type") or "자료").strip()
                key = f"{item_type}|{title}".lower()
                if key in seen_titles:
                    continue
                seen_titles.add(key)
                prefix = f"모듈 {item_type}"
                if module_name:
                    prefix += f" ({module_name})"
                course_lines.extend([f"- {_truncate_lms_text(title, 62)}", f"  {prefix}"])
                course_item_count += 1
                if course_item_count >= TELEGRAM_LMS_MATERIAL_DISPLAY_LIMIT_PER_COURSE:
                    break
            if course_item_count >= TELEGRAM_LMS_MATERIAL_DISPLAY_LIMIT_PER_COURSE:
                break

        try:
            boards = ku_lms.list_boards(session, cid) or []
        except Exception:
            boards = []
            board_failures += 1
        board_material_lines = 0
        for board in boards[:TELEGRAM_LMS_BOARD_SCAN_LIMIT_PER_COURSE]:
            bid_raw = board.get("id") or board.get("board_id")
            try:
                bid = int(bid_raw)
            except (TypeError, ValueError):
                continue
            board_name = str(board.get("name") or board.get("title") or "게시판").strip()
            try:
                resp = ku_lms.list_board_posts(session, cid, bid)
            except Exception:
                board_failures += 1
                continue
            for post in _lms_board_posts_from_response(resp)[:2]:
                title = str(post.get("title") or post.get("subject") or "").strip()
                if not title:
                    continue
                key = f"{board_name}|{title}".lower()
                if key in seen_titles:
                    continue
                seen_titles.add(key)
                course_lines.extend([f"- {_truncate_lms_text(title, 62)}", f"  게시판 {board_name}"])
                course_item_count += 1
                board_material_lines += 1
                if board_material_lines >= 3 or course_item_count >= TELEGRAM_LMS_MATERIAL_DISPLAY_LIMIT_PER_COURSE:
                    break
            if board_material_lines >= 3 or course_item_count >= TELEGRAM_LMS_MATERIAL_DISPLAY_LIMIT_PER_COURSE:
                break

        if not course_lines:
            continue
        rendered_any = True
        if not _append_limited_telegram_line(lines, f"[{_compact_lms_course_name(course_name)}]"):
            break
        for row in course_lines:
            if not _append_limited_telegram_line(lines, row):
                lines.append("- 외 항목은 길이 제한 때문에 생략했습니다.")
                return "\n".join(lines)
        _append_limited_telegram_line(lines, "")

    if not rendered_any:
        lines.append("- 최근 강의자료 위치를 찾지 못했습니다.")
    lines.append(
        f"- 확인: {scanned_courses}개 과목의 모듈과 게시판을 직접 확인했습니다."
    )
    if module_failures or board_failures:
        lines.append(f"- 참고: 모듈 실패 {module_failures}건 / 게시판 실패 {board_failures}건")
    return "\n".join(lines).rstrip()


def _format_telegram_library(query: str | None) -> str:
    """Render KU library seat availability for /library replies.

    Hits the public HODI API (no auth). Real-time, no caching layer.
    """
    target = (query or "").strip() or None
    try:
        result = get_library_seats(target) if target else get_library_seats()
    except ValueError:
        names = ", ".join(list_known_libraries())
        return f"도서관을 찾을 수 없습니다: {target}\n사용 가능: {names}"
    except Exception as exc:  # noqa: BLE001 — surface upstream message verbatim
        return f"도서관 좌석 조회 실패: {exc}"

    libraries = result.get("libraries") or {}
    summary = result.get("summary") or {}

    lines: list[str] = []
    header = "KU 도서관 좌석"
    if target and libraries:
        header = f"KU 도서관 좌석 — {next(iter(libraries))}"
    lines.append(header)
    total_seats = int(summary.get("total_seats") or 0)
    total_available = int(summary.get("total_available") or 0)
    occupancy_rate = str(summary.get("occupancy_rate") or "0%")
    lines.append(
        f"- 합계: {total_available:,}/{total_seats:,}석 가용 (점유율 {occupancy_rate})"
    )

    for lib_name, rooms in libraries.items():
        if not rooms:
            continue
        lib_total = sum(int(r.get("total_seats") or 0) for r in rooms)
        lib_avail = sum(int(r.get("available") or 0) for r in rooms)
        lines.append("")
        lines.append(f"{lib_name} — {lib_avail}/{lib_total}석 가용")
        for room in rooms:
            note = " (노트북)" if room.get("is_notebook_allowed") else ""
            lines.append(
                f"- {room.get('room_name', '')}{note}: "
                f"{int(room.get('available') or 0)}/{int(room.get('total_seats') or 0)}석 가용 "
                f"(이용 {int(room.get('in_use') or 0)})"
            )

    return "\n".join(lines).rstrip()


def _format_telegram_uclass_notice(
    settings: Settings,
    db: Database,
    *,
    user_id: int | None = None,
) -> str:
    timezone_name = str(getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul")
    uclass_state = db.get_sync_state("sync_uclass", user_id=user_id)
    if user_id and not uclass_state.last_run_at and not uclass_state.last_cursor_json:
        uclass_state = db.get_sync_state("sync_uclass")
    uclass_card = _sync_dashboard_source_card(db, "uclass", user_id=user_id)
    notifications = [
        item
        for item in db.list_notifications(limit=200, user_id=user_id)
        if str(getattr(item, "source", "")).strip().lower() == "uclass"
    ]
    status = str(uclass_card.get("status") or "").strip().lower()
    last_error = str(uclass_card.get("last_error") or "").strip()
    last_success_at = str(uclass_card.get("last_success_at") or "").strip() or None
    last_run_at = str(uclass_card.get("last_run_at") or uclass_state.last_run_at or "").strip() or None
    last_success_dt = _parse_dt(last_success_at or last_run_at)
    is_stale = bool(
        last_success_dt is not None
        and (datetime.now(timezone.utc) - last_success_dt.astimezone(timezone.utc))
        > timedelta(hours=TELEGRAM_UCLASS_NOTICE_STALE_HOURS)
    )
    status_lines: list[str] = []
    if status == "error":
        status_lines.append(f"- 최근 UClass 동기화가 실패했습니다: {last_error or 'unknown error'}")
        if last_success_at:
            status_lines.append(f"- 마지막 정상 반영: {_format_status_time(last_success_at, timezone_name)}")
        if _looks_like_auth_or_session_issue(last_error):
            status_lines.append("- `/connect`로 학교 계정을 다시 연결해 주세요.")
    elif status == "never":
        status_lines.append("- 아직 UClass 동기화 기록이 없습니다.")
    elif status == "skipped" and not notifications:
        status_lines.append(f"- 최근 UClass 동기화가 건너뛰어졌습니다: {last_error or '이유 없음'}")
    if is_stale and status != "error":
        status_lines.append(f"- 마지막 UClass 동기화가 오래됐습니다: {_format_status_time(last_success_at or last_run_at, timezone_name)}")
        status_lines.append("- 최근 알림이 아직 반영되지 않았을 수 있습니다.")
    lines = [
        "[KU] 온라인강의실 알림",
        "",
        f"알림: {min(len(notifications), 10)}",
        "",
    ]
    if not notifications:
        lines.append("- 표시할 알림이 없습니다.")
    else:
        for idx, item in enumerate(notifications[:10], start=1):
            created_local = _parse_dt(getattr(item, "created_at", None))
            created_text = (
                created_local.astimezone(ZoneInfo(timezone_name)).strftime("%m-%d %H:%M")
                if created_local is not None
                else "시각 없음"
            )
            title = str(getattr(item, "title", "") or "알림").strip() or "알림"
            lines.append(f"- {created_text} | {title}")
    if status_lines:
        lines.extend(["", "상태"])
        lines.extend(status_lines)
    lines.extend(
        [
            "",
            f"- 출처: UClass ({_format_status_time(last_run_at, timezone_name)})",
        ]
    )
    return "\n".join(lines)


def _format_telegram_inbox(db: Database, *, user_id: int | None = None) -> str:
    drafts = [
        item
        for item in db.list_unprocessed_inbox(limit=100, user_id=user_id)
        if item.item_type in {"event_draft", "task_draft", "note"}
    ]
    lines = ["[KU] Inbox", "", f"대기 draft: {len(drafts)}"]
    if not drafts:
        lines.append("- 없음")
        lines.append("- 메모를 그냥 보내면 다시 초안으로 저장됩니다.")
        return "\n".join(lines)
    for item in drafts[:10]:
        item_type = _inbox_item_type_label(str(item.item_type))
        title = str(item.title or "untitled").strip() or "untitled"
        lines.append(f"- #{item.id} [{item_type}] {title}")
    lines.extend(
        [
            "",
            "바로 반영",
            "- `/apply <id>` 예: `/apply 12`",
            "- 여러 개를 한 번에 반영: `/apply all`",
        ]
    )
    return "\n".join(lines)


USER_PREFERENCE_TELEGRAM_CHAT_ALLOWED = "telegram_chat_allowed"
USER_PREFERENCE_MATERIAL_BRIEF_PUSH_ENABLED = "material_brief_push_enabled"
USER_PREFERENCE_SCHEDULED_BRIEFINGS_ENABLED = "scheduled_briefings_enabled"
USER_PREFERENCE_DAILY_DIGEST_ENABLED = "daily_digest_enabled"

NOTIFICATION_POLICY_KIND_BRIEFING_MORNING = "briefing_morning"
NOTIFICATION_POLICY_KIND_BRIEFING_EVENING = "briefing_evening"
NOTIFICATION_POLICY_KIND_DAILY_DIGEST = "daily_digest"
NOTIFICATION_POLICY_KIND_MATERIAL_BRIEF_PUSH = "material_brief_push"

_NOTIFICATION_POLICY_KIND_ALIASES: dict[str, tuple[str, ...]] = {
    NOTIFICATION_POLICY_KIND_BRIEFING_MORNING: (
        NOTIFICATION_POLICY_KIND_BRIEFING_MORNING,
        "morning_briefing",
    ),
    NOTIFICATION_POLICY_KIND_BRIEFING_EVENING: (
        NOTIFICATION_POLICY_KIND_BRIEFING_EVENING,
        "evening_briefing",
    ),
    NOTIFICATION_POLICY_KIND_DAILY_DIGEST: (
        NOTIFICATION_POLICY_KIND_DAILY_DIGEST,
    ),
    NOTIFICATION_POLICY_KIND_MATERIAL_BRIEF_PUSH: (
        NOTIFICATION_POLICY_KIND_MATERIAL_BRIEF_PUSH,
    ),
}

_NOTIFICATION_POLICY_WEEKDAY_NAMES = (
    "mon",
    "tue",
    "wed",
    "thu",
    "fri",
    "sat",
    "sun",
)

_NOTIFICATION_POLICY_WEEKDAY_ALIASES = {
    "mo": "mon",
    "mon": "mon",
    "monday": "mon",
    "tu": "tue",
    "tue": "tue",
    "tues": "tue",
    "tuesday": "tue",
    "we": "wed",
    "wed": "wed",
    "wednesday": "wed",
    "th": "thu",
    "thu": "thu",
    "thur": "thu",
    "thur.": "thu",
    "thurs": "thu",
    "thursday": "thu",
    "fr": "fri",
    "fri": "fri",
    "friday": "fri",
    "sa": "sat",
    "sat": "sat",
    "saturday": "sat",
    "su": "sun",
    "sun": "sun",
    "sunday": "sun",
}


def _configured_telegram_chat_ids(settings: Settings) -> list[str]:
    return [
        str(item).strip()
        for item in list(getattr(settings, "telegram_allowed_chat_ids", []) or [])
        if str(item).strip()
    ]


def _default_chat_ids_for_user_preference(settings: Settings, db: Database) -> set[str]:
    configured = set(_configured_telegram_chat_ids(settings))
    connected = {
        str(chat).strip()
        for chat in db.list_chat_ids_with_active_school_connections(limit=max(len(configured) * 10, 1000))
        if str(chat).strip()
    }
    return configured | connected


def _chat_ids_by_preference_override(
    db: Database,
    preference: str,
) -> tuple[set[str], set[str]]:
    enabled = set(db.list_chat_ids_by_preference(preference, enabled=True, limit=1000))
    disabled = set(db.list_chat_ids_by_preference(preference, enabled=False, limit=1000))
    return enabled, disabled


def _chat_ids_for_user_preference(
    settings: Settings,
    db: Database,
    preference: str,
) -> list[str]:
    default_chat_ids = _default_chat_ids_for_user_preference(settings, db)
    enabled, disabled = _chat_ids_by_preference_override(db, preference)
    resolved = (default_chat_ids | enabled) - disabled
    return sorted(chat for chat in resolved if chat)


def _effective_telegram_allowed_chat_ids(settings: Settings, db: Database) -> list[str]:
    default_chat_ids = _default_chat_ids_for_user_preference(settings, db)
    enabled, disabled = _chat_ids_by_preference_override(db, USER_PREFERENCE_TELEGRAM_CHAT_ALLOWED)
    allowed = default_chat_ids | enabled
    return sorted(chat for chat in allowed if chat and chat not in disabled)


def _notification_policy_kind_aliases(*policy_kinds: str) -> tuple[str, ...]:
    output: list[str] = []
    seen: set[str] = set()
    for raw_kind in policy_kinds:
        kind = str(raw_kind or "").strip().lower()
        if not kind:
            continue
        for candidate in _NOTIFICATION_POLICY_KIND_ALIASES.get(kind, (kind,)):
            if candidate in seen:
                continue
            seen.add(candidate)
            output.append(candidate)
    return tuple(output)


def _notification_policy_reference_timezone(
    settings: Settings,
    db: Database,
    policy: dict[str, Any],
) -> ZoneInfo:
    default_timezone_name = str(getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul")
    timezone_name = str(policy.get("timezone") or "").strip()
    if not timezone_name:
        user_scope = _resolve_user_scope(
            settings,
            db,
            user_id=int(policy.get("user_id") or 0) or None,
            chat_id=str(policy.get("chat_id") or "").strip() or None,
            create_if_missing=False,
            metadata_source="notification_policy_eval",
        )
        timezone_name = str(user_scope.get("timezone") or default_timezone_name or "Asia/Seoul")
    try:
        return ZoneInfo(timezone_name)
    except Exception:
        return ZoneInfo(default_timezone_name)


def _normalize_notification_policy_weekday(value: Any) -> str | None:
    if isinstance(value, bool):
        return None
    if isinstance(value, int):
        if value == 0 or value == 7:
            return "sun"
        if 1 <= value <= 6:
            return _NOTIFICATION_POLICY_WEEKDAY_NAMES[value - 1]
        return None
    text = str(value or "").strip().lower()
    if not text:
        return None
    if text.isdigit():
        return _normalize_notification_policy_weekday(int(text))
    return _NOTIFICATION_POLICY_WEEKDAY_ALIASES.get(text)


def _notification_policy_matches_reference_time(
    settings: Settings,
    db: Database,
    policy: dict[str, Any],
    *,
    reference_local: datetime | None = None,
) -> bool:
    if not bool(policy.get("enabled")):
        return False
    tz = _notification_policy_reference_timezone(settings, db, policy)
    if reference_local is None:
        current_local = datetime.now(tz)
    elif reference_local.tzinfo is None:
        current_local = reference_local.replace(tzinfo=tz)
    else:
        current_local = reference_local.astimezone(tz)

    normalized_days = {
        day
        for day in (
            _normalize_notification_policy_weekday(item)
            for item in list(policy.get("days_of_week_json") or [])
        )
        if day
    }
    if policy.get("days_of_week_json") and not normalized_days:
        return False
    if normalized_days:
        weekday_name = _NOTIFICATION_POLICY_WEEKDAY_NAMES[current_local.weekday()]
        if weekday_name not in normalized_days:
            return False

    time_local = str(policy.get("time_local") or "").strip()
    if not time_local:
        return True
    try:
        hour, minute = _parse_clock_time(time_local, "notification_policy.time_local")
    except ValueError:
        return False
    target_local = current_local.replace(hour=hour, minute=minute, second=0, microsecond=0)
    return current_local >= target_local


def _notification_policies_for_dispatch(
    db: Database,
    *,
    policy_kinds: tuple[str, ...],
) -> list[dict[str, Any]]:
    if not policy_kinds:
        return []
    rows: list[dict[str, Any]] = []
    seen: set[int] = set()
    for kind in policy_kinds:
        for row in db.list_notification_policies(policy_kind=kind, limit=1000):
            row_id = int(row.get("id") or 0)
            if row_id > 0 and row_id in seen:
                continue
            if row_id > 0:
                seen.add(row_id)
            rows.append(row)
    rows.sort(
        key=lambda item: (
            str(item.get("updated_at") or ""),
            int(item.get("id") or 0),
        ),
        reverse=True,
    )
    latest_by_user: dict[int, dict[str, Any]] = {}
    latest_by_chat: dict[str, dict[str, Any]] = {}
    for row in rows:
        owner_id = int(row.get("user_id") or 0)
        chat_id = str(row.get("chat_id") or "").strip()
        if owner_id > 0 and owner_id in latest_by_user:
            continue
        if owner_id <= 0 and chat_id and chat_id in latest_by_chat:
            continue
        if owner_id > 0:
            latest_by_user[owner_id] = row
            continue
        if chat_id:
            latest_by_chat[chat_id] = row
    return list(latest_by_user.values()) + list(latest_by_chat.values())


def _chat_ids_for_notification_dispatch(
    settings: Settings,
    db: Database,
    *,
    preference: str,
    policy_kinds: tuple[str, ...],
    reference_local: datetime | None = None,
) -> list[str]:
    legacy_chat_ids = _chat_ids_for_user_preference(settings, db, preference)
    effective_policy_kinds = _notification_policy_kind_aliases(*policy_kinds)
    if not effective_policy_kinds:
        return legacy_chat_ids

    policy_rows = _notification_policies_for_dispatch(db, policy_kinds=effective_policy_kinds)
    if not policy_rows:
        return legacy_chat_ids

    overridden_user_ids: set[int] = set()
    overridden_chat_ids: set[str] = set()
    policy_chat_ids: set[str] = set()
    for row in policy_rows:
        owner_id = int(row.get("user_id") or 0)
        chat_id = str(row.get("chat_id") or "").strip()
        if owner_id > 0:
            overridden_user_ids.add(owner_id)
        if chat_id:
            overridden_chat_ids.add(chat_id)
        if not chat_id:
            continue
        if _notification_policy_matches_reference_time(
            settings,
            db,
            row,
            reference_local=reference_local,
        ):
            policy_chat_ids.add(chat_id)

    fallback_chat_ids: set[str] = set()
    for raw_chat_id in legacy_chat_ids:
        chat_id = str(raw_chat_id or "").strip()
        if not chat_id:
            continue
        if chat_id in overridden_chat_ids:
            continue
        user = db.get_user_by_chat_id(chat_id)
        if user is not None and int(user.get("id") or 0) in overridden_user_ids:
            continue
        fallback_chat_ids.add(chat_id)

    return sorted(fallback_chat_ids | policy_chat_ids)


def _is_telegram_chat_allowed(settings: Settings, db: Database, chat_id: str | None) -> bool:
    chat = str(chat_id or "").strip()
    allowed = _effective_telegram_allowed_chat_ids(settings, db)
    if not allowed:
        return True
    if not chat:
        return False
    return chat in allowed


def _format_chat_lms_connections(db: Database, chat_id: str | None) -> list[str]:
    snapshot = _chat_lms_connection_snapshot(db, chat_id)
    return list(snapshot.all_labels)


def _format_setup_mark(value: str) -> str:
    normalized = str(value or "").strip().upper()
    mapping = {
        "OK": "준비됨",
        "WARN": "확인 필요",
        "TODO": "미연결",
    }
    return mapping.get(normalized, "미연결")


def _format_optional_setup_mark(ok: bool) -> str:
    return "준비됨" if ok else "선택"


def _render_telegram_setup_message(
    state: TelegramSetupState,
    *,
    smart_commands_enabled: bool,
) -> str:
    lines = ["[KU] 연결 상태", "", "핵심 상태"]
    lines.extend(
        [
            f"- Telegram 채팅: {_format_setup_mark('OK' if state.allowed else 'TODO')}",
            f"- 온라인강의실 연결: {_format_setup_mark(state.online_connection_level)}",
            f"- UClass 계정: {_format_setup_mark(state.uclass_account_level)}",
            f"- 시간표 소스: {_format_setup_mark(state.portal_level)}",
            f"- 로컬 LLM 요약(선택): {_format_optional_setup_mark(state.local_llm_ready)}",
        ]
    )

    _append_message_section(lines, "지금 할 일")
    if not state.allowed:
        lines.append("- 이 채팅을 사용할 수 있도록 먼저 활성화해야 합니다.")
    if not (state.uclass_labels or state.legacy_uclass_account):
        lines.append("- `/connect`로 학교 계정을 연결하세요.")
    elif not state.core_ready:
        lines.append("- `/setup` 결과를 확인한 뒤 필요하면 `/connect`로 다시 연결하세요.")
    else:
        lines.append("- 핵심 연결은 준비됐습니다. `/today`와 `/weather`를 사용해 보세요.")

    _append_message_section(lines, "현재 연결")
    if state.uclass_labels:
        lines.append(f"- 학교 계정: {', '.join(state.uclass_labels)}")
    else:
        lines.append("- 학교 계정: 아직 연결되지 않았습니다.")
    if state.portal_labels:
        lines.append(f"- 시간표 연결: {', '.join(state.portal_labels)}")
    elif state.show_official_api_connection:
        lines.append("- 시간표 연결: 고려대 학교 공식 API 자동 동기화")
    else:
        lines.append("- 시간표 연결: 아직 준비되지 않았습니다.")
    if not state.local_llm_ready:
        lines.append("- 로컬 LLM 요약은 선택 사항입니다. 필요하면 나중에 추가로 설정하세요.")
    if state.core_ready and not state.local_llm_ready:
        lines.append("- 핵심 연결은 준비됐습니다. 로컬 LLM은 선택 사항입니다.")

    if state.uclass_notes or state.portal_notes:
        _append_message_section(lines, "추가 안내")
        lines.extend(state.uclass_notes)
        lines.extend(state.portal_notes)

    _append_message_section(lines, "추천 명령")
    lines.extend(
        [
            "- /connect",
            "- /today",
            "- /tomorrow",
            "- /weather",
            "- /region 고려대",
            "- /todaysummary",
            "- /tomorrowsummary",
            "- /notice_uclass",
            "- /notice_general",
            "- /notice_academic",
            "- /status",
        ]
    )
    if smart_commands_enabled:
        lines.append("- /plan 내일 오전 8시에 과제 제출 알림")
    return "\n".join(lines)


def _directory_entry_by_slug(db: Database, school_slug: str) -> dict[str, Any] | None:
    target = str(school_slug or "").strip().lower()
    if not target:
        return None
    for entry in db.list_moodle_school_directory(limit=2000):
        if str(entry.get("school_slug") or "").strip().lower() == target:
            return entry
    return None


def _school_connect_scope_lines(entry: dict[str, Any] | None) -> list[str]:
    summary = school_support_summary(entry)
    capabilities = dict(summary.get("capabilities") or {})
    lines: list[str] = []
    if capabilities.get("portal_timetable_sync"):
        lines.extend(
            [
                "- 이 학교는 온라인강의실 계정을 연결하면 시간표는 학교 공식 API로 자동 동기화합니다.",
                "- 포털 브라우저 세션은 연결하지 않습니다.",
                "- 현재 사용자-facing 공식 지원 학교입니다.",
            ]
        )
        return lines
    if capabilities.get("portal_shared_account_hint"):
        lines.extend(
            [
                "- 이 학교는 온라인강의실 연결을 먼저 확인합니다.",
                "- 포털은 같은 학교 계정을 쓰는 것으로 등록돼 있지만, 자동 시간표 연동은 아직 학교별 구현이 필요합니다.",
            ]
        )
    elif capabilities.get("lms_credential_onboarding"):
        lines.extend(
            [
                "- 이 학교는 온라인강의실 연결을 확인합니다.",
                "- 포털/시간표 자동 연동 범위는 학교별 후속 구현이 필요할 수 있습니다.",
            ]
        )
    else:
        lines.append("- 학교별로 포털/대학행정 자동 연동 범위는 다를 수 있습니다.")
    if not bool(summary.get("official_user_support")):
        lines.append("- 현재 사용자-facing 공식 지원 학교는 고려대학교입니다.")
    return lines


def _issue_moodle_connect_link(
    settings: Settings,
    db: Database,
    *,
    chat_id: str | None,
    school_query: str | None = None,
) -> dict[str, Any]:
    chat = str(chat_id or "").strip()
    if not chat:
        return {"ok": False, "error": "chat_id is missing"}
    public_base_url = str(getattr(settings, "onboarding_public_base_url", "") or "").strip()
    if not public_base_url:
        return {
            "ok": False,
            "error": "ONBOARDING_PUBLIC_BASE_URL is missing",
            "message": (
                "[KU] 학교 계정 연결\n\n"
                "- 지금은 연결 링크를 만들 수 없습니다.\n"
                "- 운영자가 ONBOARDING_PUBLIC_BASE_URL 과 `kus onboarding serve`를 먼저 준비해야 합니다."
            ),
        }
    try:
        normalized_public_base_url = normalize_public_moodle_connect_base_url(public_base_url)
    except Exception as exc:
        return {
            "ok": False,
            "error": str(exc),
            "message": (
                "[KU] 학교 계정 연결\n\n"
                "- ONBOARDING_PUBLIC_BASE_URL 설정이 올바르지 않습니다.\n"
                "- 공개 링크는 반드시 HTTPS 주소여야 합니다. 예: https://connect.example.invalid"
            ),
        }
    ttl_minutes = max(int(getattr(settings, "onboarding_session_ttl_minutes", 15) or 15), 1)
    expires_at = datetime.now(timezone.utc) + timedelta(minutes=ttl_minutes)
    school_query_text = str(school_query or "").strip()
    allowed_school_slugs = onboarding_allowed_school_slugs(settings)
    allowed_visible_entries = visible_onboarding_school_entries(
        db.list_moodle_school_directory(limit=2000),
        settings=settings,
    )
    matched_school = None
    if school_query_text:
        matches = db.find_moodle_school_directory(school_query_text, limit=5)
        allowed_matches = [
            entry
            for entry in matches
            if school_entry_allowed_for_onboarding(entry, settings=settings)
        ]
        if allowed_matches:
            matched_school = allowed_matches[0]
        elif allowed_school_slugs:
            allowed_names = [
                str(entry.get("display_name") or "").strip()
                for entry in allowed_visible_entries
                if str(entry.get("display_name") or "").strip()
            ]
            supported_label = ", ".join(allowed_names) if allowed_names else "운영 대상 학교"
            return {
                "ok": False,
                "error": "school_not_allowed",
                "message": (
                    "[KU] 학교 계정 연결\n\n"
                    f"- 현재 이 인스턴스는 {supported_label}만 지원합니다.\n"
                    "- `/connect`를 다시 실행해 지원 대상 학교 연결 링크를 열어 주세요."
                ),
            }
    elif len(allowed_visible_entries) == 1:
        matched_school = allowed_visible_entries[0]
    prefilled_school_name = str(
        (matched_school or {}).get("display_name") or school_query_text or ""
    ).strip()
    session = db.create_onboarding_session(
        session_kind=MOODLE_ONBOARDING_SESSION_KIND,
        chat_id=chat,
        expires_at=expires_at,
        metadata_json={
            "source": "telegram_school_account_connect",
            "school_query": prefilled_school_name or None,
            "directory_school_slug": (
                str((matched_school or {}).get("school_slug") or "").strip() or None
            ),
        },
    )
    link = build_public_moodle_connect_url(normalized_public_base_url, session["token"])
    expires_local = expires_at.astimezone(ZoneInfo(settings.timezone)).strftime("%m-%d %H:%M")
    lines = [
        "[KU] 학교 계정 연결",
        "",
    ]
    if prefilled_school_name:
        lines.append(f"- 대상: {prefilled_school_name}")
        lines.extend(_school_connect_scope_lines(matched_school))
    else:
        lines.append("- 링크를 열어 학교를 선택한 뒤 학교 계정으로 로그인하세요.")
        lines.append("- 학교별로 포털/대학행정 자동 연동 범위는 다를 수 있습니다.")
        lines.append("- 현재 사용자-facing 공식 지원 학교는 고려대학교입니다.")
    lines.extend(
        [
            "- 비밀번호는 이 사용자의 온라인강의실 재인증이 필요할 때만 보안 저장소에 저장합니다.",
            f"- 링크 만료: {expires_local}",
            "",
            link,
        ]
    )
    return {
        "ok": True,
        "link": link,
        "expires_at": session["expires_at"],
        "message": "\n".join(lines),
    }


def _telegram_bot_menu_commands(settings: Settings) -> list[dict[str, str]]:
    commands = [
        {"command": "start", "description": "시작 안내"},
        {"command": "help", "description": "사용 가능한 명령 보기"},
        {"command": "bot", "description": "자연어 비서"},
        {"command": "setup", "description": "연결 상태 점검"},
        {"command": "connect", "description": "학교 계정 연결"},
        {"command": "status", "description": "동기화 상태 보기"},
        {"command": "today", "description": "오늘 일정과 마감 보기"},
        {"command": "tomorrow", "description": "내일 일정과 마감 보기"},
        {"command": "weather", "description": "오늘/내일 날씨 보기"},
        {"command": "region", "description": "날씨 지역 설정"},
        {"command": "todaysummary", "description": "오늘 수업 자료 요약 보기"},
        {"command": "tomorrowsummary", "description": "내일 수업 자료 요약 보기"},
        {"command": "notice_uclass", "description": "온라인강의실 최근 알림 보기"},
        {"command": "notice_general", "description": "학교 일반공지 10개 보기"},
        {"command": "notice_academic", "description": "학교 학사공지 10개 보기"},
        {"command": "library", "description": "도서관 좌석 현황 보기"},
        {"command": "assignments", "description": "내야 할 과제 목록 보기"},
        {"command": "submitted", "description": "제출 완료 과제 보기"},
        {"command": "board", "description": "과목별 게시판/공지 모음 보기"},
        {"command": "materials", "description": "과목별 강의자료 위치 보기"},
        {"command": "inbox", "description": "임시 draft 목록 보기"},
        {"command": "apply", "description": "inbox draft 반영하기"},
        {"command": "done", "description": "과제 완료 처리하기"},
    ]
    if not bool(getattr(settings, "telegram_assistant_enabled", False)):
        commands = [item for item in commands if item.get("command") != "bot"]
    if bool(getattr(settings, "telegram_smart_commands_enabled", False)):
        commands.append({"command": "plan", "description": "자연어 리마인더 예약"})
    return commands


def _register_telegram_bot_menu(
    settings: Settings,
    client: TelegramBotClient,
    *,
    existing_hash: str | None = None,
    pending_hash: str | None = None,
    retry_after_iso: str | None = None,
) -> dict[str, Any]:
    if not getattr(settings, "telegram_commands_enabled", False):
        return {"skipped": True, "reason": "TELEGRAM_COMMANDS_ENABLED is false"}
    commands = _telegram_bot_menu_commands(settings)
    command_payload = json.dumps(commands, ensure_ascii=False, sort_keys=True)
    commands_hash = sha1(command_payload.encode("utf-8")).hexdigest()
    if existing_hash and existing_hash == commands_hash:
        return {
            "ok": True,
            "updated": False,
            "command_count": len(commands),
            "hash": commands_hash,
            "target_hash": commands_hash,
        }
    retry_after_dt = _parse_dt(retry_after_iso)
    if (
        pending_hash
        and pending_hash == commands_hash
        and retry_after_dt is not None
        and retry_after_dt > datetime.now(timezone.utc)
    ):
        return {
            "skipped": True,
            "reason": "telegram menu retry cooldown",
            "command_count": len(commands),
            "hash": existing_hash,
            "target_hash": commands_hash,
            "retry_after": retry_after_dt.isoformat(),
        }
    register = getattr(client, "set_my_commands", None)
    if not callable(register):
        return {
            "skipped": True,
            "reason": "Telegram client does not support set_my_commands",
            "command_count": len(commands),
        }
    try:
        ok = bool(register(commands))
    except Exception as exc:
        logger.warning("failed to register telegram bot menu", extra={"error": str(exc)})
        return {
            "ok": False,
            "error": str(exc),
            "command_count": len(commands),
            "hash": existing_hash,
            "target_hash": commands_hash,
            "retry_after": (
                datetime.now(timezone.utc)
                + timedelta(seconds=TELEGRAM_MENU_RETRY_COOLDOWN_SECONDS)
            ).replace(microsecond=0).isoformat(),
        }
    if not ok:
        return {
            "ok": False,
            "error": "setMyCommands returned false",
            "command_count": len(commands),
            "hash": existing_hash,
            "target_hash": commands_hash,
            "retry_after": (
                datetime.now(timezone.utc)
                + timedelta(seconds=TELEGRAM_MENU_RETRY_COOLDOWN_SECONDS)
            ).replace(microsecond=0).isoformat(),
        }
    return {
        "ok": True,
        "updated": True,
        "command_count": len(commands),
        "hash": commands_hash,
        "target_hash": commands_hash,
    }


def _format_telegram_help(settings: Settings) -> str:
    assistant_lines: list[str] = []
    if bool(getattr(settings, "telegram_assistant_enabled", False)):
        assistant_lines = [
            "- /bot <자연어> : 자연어 비서",
            "- 예: /bot 오늘 일정이랑 날씨 알려줘",
        ]
    smart_lines: list[str] = []
    if bool(getattr(settings, "telegram_smart_commands_enabled", False)):
        smart_lines = [
            "- /plan <자연어> : 자연어 리마인더 예약",
            "- 예: /plan 내일 오후 10시에 과제 제출 알림",
        ]
    lines = [
        "[KU] 도움말",
        "",
        "기본 명령",
        "- /connect : 학교 계정 연결",
        "- /today : 오늘 일정과 마감 과제",
        "- /tomorrow : 내일 일정과 마감 과제",
        "- /weather : 오늘/내일 날씨",
        "- /region <지역명> : 날씨 지역 설정",
        "- /todaysummary : 오늘 수업 자료 요약",
        "- /tomorrowsummary : 내일 수업 자료 요약",
        "- /notice_uclass : 온라인강의실 최근 알림",
        "- /notice_general : 학교 일반공지",
        "- /notice_academic : 학교 학사공지",
        "- /library [도서관명] : 도서관 좌석 현황 (예: /library 중앙도서관)",
        "- /assignments : 제출해야 할 LMS 과제와 공지/자료/게시판 제출 항목 (별칭: /due /todo /과제)",
        "- /submitted : 제출 완료 LMS 과제 (별칭: /submissions /제출완료)",
        "- /board : 과목별 LMS 공지/게시판 최근 글 (별칭: /announcements /공지)",
        "- /materials : 과목별 LMS 모듈/게시판 강의자료 위치 (별칭: /자료 /강의자료)",
        "- /status : 현재 동기화 상태",
        *assistant_lines,
        *smart_lines,
        "- /done task <id|external_id> : 과제 완료 처리",
        "",
        "관리 명령",
        "- /inbox : 임시 draft 목록",
        "- /apply <id|all> : inbox draft 반영",
        "",
        "시작",
        "- 처음 사용이면 /connect 또는 /setup",
    ]
    lines.extend(
        [
            "",
            "빠른 흐름",
            "- 메모를 그냥 보내면 `/inbox`에 초안으로 저장됩니다.",
            "- `/inbox`를 본 뒤 `/apply <id>` 또는 `/apply all`로 반영할 수 있습니다.",
        ]
    )
    return "\n".join(lines)


def _format_telegram_start(
    settings: Settings,
    *,
    db: Database,
    chat_id: str | None = None,
    user_id: int | None = None,
) -> str:
    allowed = _is_telegram_chat_allowed(settings, db, chat_id)
    chat_connections = _format_chat_lms_connections(db, chat_id)
    lines = [
        "[KU] 시작 안내",
        "",
        "할 수 있는 일",
        "- 학교 계정 연결",
        "- 자연어 비서",
        "- 고려대 공식 시간표 동기화",
        "- 아침/저녁 브리핑",
        "- 오늘/내일 날씨와 미세먼지",
        "- 강의자료 요약과 과제 마감 확인",
        "- 온라인강의실/KU 공지 확인",
    ]
    if not bool(getattr(settings, "telegram_assistant_enabled", False)):
        lines = [item for item in lines if item != "- 자연어 비서"]
    if bool(getattr(settings, "telegram_smart_commands_enabled", False)):
        lines.append("- 자연어 리마인더 예약")
    lines.extend(
        [
            "",
            "현재 공식 지원",
            "- 고려대학교",
            "- 다른 학교는 일부 연결 정보만 확인될 수 있습니다.",
            "",
            "바로 시작",
            "- /connect",
            "- 링크를 열어 학교를 선택한 뒤 학교 계정으로 로그인하세요.",
        ]
    )
    if bool(getattr(settings, "telegram_assistant_enabled", False)):
        lines.append("- /bot 오늘 일정이랑 날씨 알려줘")
    if bool(getattr(settings, "telegram_smart_commands_enabled", False)):
        lines.append("- /plan 내일 오전 8시에 과제 제출 알림")
    if chat_connections:
        lines.extend(["", "현재 연결됨", f"- {', '.join(chat_connections)}"])
    if not allowed and chat_id:
        lines.extend(
            [
                "",
                "확인할 것",
                "- 이 채팅은 아직 사용할 수 있도록 활성화되지 않았습니다.",
                "- `/setup`으로 연결 상태를 확인하세요.",
            ]
        )
    return "\n".join(lines)


def _format_telegram_setup(
    settings: Settings,
    *,
    db: Database,
    chat_id: str | None = None,
    user_id: int | None = None,
) -> str:
    allowed = _is_telegram_chat_allowed(settings, db, chat_id)

    def _read_secret_ref(secret_kind: str | None, secret_ref: str | None) -> tuple[str, str]:
        kind = str(secret_kind or "").strip()
        ref = str(secret_ref or "").strip()
        if not kind or not ref:
            return "", ""
        try:
            secret = default_secret_store(settings).read_secret(
                ref=StoredSecretRef(kind=kind, ref=ref)
            )
        except Exception as exc:
            return "", str(exc).strip()
        return str(secret or ""), ""

    state = build_telegram_setup_state(
        settings,
        db=db,
        allowed=allowed,
        chat_id=chat_id,
        user_id=user_id,
        read_secret_ref=_read_secret_ref,
    )
    return _render_telegram_setup_message(
        state,
        smart_commands_enabled=bool(getattr(settings, "telegram_smart_commands_enabled", False)),
    )


def _safe_int(value: Any) -> int | None:
    try:
        return int(str(value).strip())
    except Exception:
        return None


def _is_timetable_event(event: Any) -> bool:
    source = str(getattr(event, "source", "") or "").strip().lower()
    if source != "portal":
        return False
    metadata = _json_load(getattr(event, "metadata_json", None))
    if str(metadata.get("timetable_source") or "").strip().lower() == KU_PORTAL_SCHOOL_SLUG:
        return True
    provenance = normalize_provenance(metadata, fallback_source=source)
    return str(provenance.get("source") or "").strip().lower() == "portal_uos_timetable"


def _resolve_user_scope(
    settings: Settings,
    db: Database,
    *,
    chat_id: str | None = None,
    user_id: int | None = None,
    create_if_missing: bool = True,
    metadata_source: str = "pipeline",
) -> dict[str, Any]:
    owner_id = _safe_int(user_id)
    user: dict[str, Any] | None = None
    if owner_id is not None and owner_id > 0:
        user = db.get_user(owner_id)
    chat = str(chat_id or "").strip()
    if user is None and chat:
        if create_if_missing:
            user = db.ensure_user_for_chat(
                chat_id=chat,
                timezone_name=str(getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul"),
                metadata_json={"source": metadata_source},
            )
        else:
            user = db.get_user_by_chat_id(chat)
    if user is None:
        return {
            "user_id": 0,
            "chat_id": chat,
            "timezone": str(getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul"),
            "status": "unknown",
            "metadata_json": {},
        }
    return {
        "user_id": int(user.get("id") or user.get("user_id") or 0),
        "chat_id": str(user.get("telegram_chat_id") or chat),
        "timezone": str(user.get("timezone") or getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul"),
        "status": str(user.get("status") or "active"),
        "metadata_json": dict(user.get("metadata_json") or {}),
    }


def _clean_brief_bullets(raw_items: Any, limit: int = 3) -> list[str]:
    output: list[str] = []
    if not isinstance(raw_items, list):
        return output
    for item in raw_items:
        text = str(item or "").strip()
        if not text:
            continue
        output.append(text)
        if len(output) >= max(int(limit), 1):
            break
    return output


def _material_brief_push_message(items: list[dict[str, Any]]) -> str:
    lines: list[str] = [
        "[KU] 새 강의자료 요약",
        "",
        f"자료: {len(items)}",
        "",
    ]
    for idx, item in enumerate(items, start=1):
        filename = str(item.get("filename") or "material")
        course_name = _display_course_name(str(item.get("course_name") or ""))
        source_label = "AI" if str(item.get("mode") or "").strip().lower() == "llm" else "UClass"
        if course_name:
            lines.append(f"{idx}. [{course_name}] {filename}")
        else:
            lines.append(f"{idx}. {filename}")
        lines.append(f"- 출처: {source_label}")
        bullets = _clean_telegram_brief_bullets(item.get("bullets"), limit=3)
        if bullets:
            for bullet in bullets:
                lines.append(f"- {bullet}")
        else:
            lines.append("- 요약 항목을 만들지 못했습니다.")
        question = _clean_telegram_brief_question(item.get("question"))
        if question:
            lines.append(f"- 할 일: {question}")
        lines.append("")
    return "\n".join(lines).strip()


def send_material_brief_push(
    settings: Settings,
    db: Database,
    generated_brief_items: list[dict[str, Any]],
    *,
    chat_ids: list[str] | None = None,
) -> dict[str, Any]:
    if not bool(getattr(settings, "material_brief_push_enabled", False)):
        return {"skipped": True, "reason": "MATERIAL_BRIEF_PUSH_ENABLED is false"}
    if not generated_brief_items:
        return {"skipped": True, "reason": "No generated material briefs"}
    if not bool(getattr(settings, "telegram_enabled", False)):
        return {"skipped": True, "reason": "TELEGRAM_ENABLED is false"}
    token = str(getattr(settings, "telegram_bot_token", "") or "").strip()
    if not token:
        return {"skipped": True, "reason": "TELEGRAM_BOT_TOKEN missing"}
    timezone_name = str(getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul")
    now_local = datetime.now(ZoneInfo(timezone_name))
    target_chat_ids = (
        [str(chat).strip() for chat in list(chat_ids or []) if str(chat).strip()]
        if chat_ids is not None
        else _chat_ids_for_notification_dispatch(
            settings,
            db,
            preference=USER_PREFERENCE_MATERIAL_BRIEF_PUSH_ENABLED,
            policy_kinds=(NOTIFICATION_POLICY_KIND_MATERIAL_BRIEF_PUSH,),
            reference_local=now_local,
        )
    )
    if not target_chat_ids:
        return {"skipped": True, "reason": "No eligible chat_ids"}

    gate = _identity_warning_gate(
        settings=settings,
        db=db,
        step="material_brief_push",
        destination="telegram",
    )
    if gate is not None:
        return gate

    max_items = max(int(getattr(settings, "material_brief_push_max_items", 3) or 3), 1)
    selected = _sanitize_material_brief_push_items(generated_brief_items, limit=max_items)
    if not selected:
        return {"skipped": True, "reason": "No high-quality material briefs"}
    message = _material_brief_push_message(selected)
    client = TelegramBotClient(token)
    sent_to: list[str] = []
    for chat_id in target_chat_ids:
        try:
            if client.send_message(chat_id=chat_id, text=message[:3500]):
                sent_to.append(chat_id)
        except Exception as exc:
            logger.warning(
                "failed to send material brief push",
                extra={"chat_id": chat_id, "error": str(exc)},
            )
    if not sent_to:
        return {"ok": False, "error": "material_brief_push_send_failed"}
    return {
        "ok": True,
        "sent_to": sent_to,
        "brief_count": len(selected),
        "total_generated": len(generated_brief_items),
    }


def _fallback_instruction_schedule(
    *,
    text: str,
    timezone_name: str,
) -> dict[str, Any] | None:
    normalized = str(text or "").strip()
    if not normalized:
        return None
    now_local = datetime.now(ZoneInfo(timezone_name))
    default_dt = now_local.replace(second=0, microsecond=0)
    try:
        parsed = dt_parser.parse(normalized, fuzzy=True, default=default_dt)
    except Exception:
        return None
    if parsed.tzinfo is None:
        parsed = parsed.replace(tzinfo=ZoneInfo(timezone_name))
    run_at_local = parsed.astimezone(ZoneInfo(timezone_name))
    if run_at_local <= now_local:
        run_at_local = run_at_local + timedelta(days=1)
    return {
        "action": "schedule",
        "run_at_iso": run_at_local.isoformat(),
        "message": normalized[:180],
    }


def _plan_instruction_with_llm(
    settings: Settings,
    db: Database,
    text: str,
) -> dict[str, Any]:
    if not bool(getattr(settings, "llm_enabled", False)):
        fallback = _fallback_instruction_schedule(text=text, timezone_name=settings.timezone)
        if fallback is None:
            return {"ok": False, "error": "LLM is disabled and instruction could not be parsed"}
        return {"ok": True, "plan": fallback, "mode": "heuristic"}

    gate = _identity_warning_gate(
        settings=settings,
        db=db,
        step="telegram_instruction_parse",
        destination="llm",
    )
    if gate is not None:
        return {
            "ok": False,
            "error": gate["error"],
            "blocked": True,
            "warning_gate": gate["warning_gate"],
        }

    now_local = datetime.now(ZoneInfo(settings.timezone)).replace(microsecond=0)
    system_prompt = (
        "You convert Korean reminder instructions into JSON only. "
        "Return one JSON object with keys: action, run_at_iso, message. "
        "action must be 'schedule' or 'none'. "
        "run_at_iso must be timezone-aware ISO8601 in user's timezone. "
        "If no actionable reminder exists, return action='none'."
    )
    user_prompt = json.dumps(
        {
            "timezone": settings.timezone,
            "now_local_iso": now_local.isoformat(),
            "instruction": text,
        },
        ensure_ascii=False,
    )
    try:
        client = _llm_client(settings)
        raw = client.generate_text(system_prompt=system_prompt, prompt=user_prompt)
    except Exception as exc:
        fallback = _fallback_instruction_schedule(text=text, timezone_name=settings.timezone)
        if fallback is None:
            return {"ok": False, "error": str(exc)}
        return {"ok": True, "plan": fallback, "mode": "heuristic_fallback", "llm_error": str(exc)}
    parsed = _parse_llm_json_payload(raw)
    if not isinstance(parsed, dict):
        fallback = _fallback_instruction_schedule(text=text, timezone_name=settings.timezone)
        if fallback is None:
            return {"ok": False, "error": "LLM output is not valid JSON"}
        return {"ok": True, "plan": fallback, "mode": "heuristic_fallback"}
    action = str(parsed.get("action") or "none").strip().lower()
    if action != "schedule":
        return {"ok": True, "plan": {"action": "none"}, "mode": "llm"}
    run_at_iso = str(parsed.get("run_at_iso") or "").strip()
    run_at = _parse_dt(run_at_iso)
    if run_at is None:
        fallback = _fallback_instruction_schedule(text=text, timezone_name=settings.timezone)
        if fallback is None:
            return {"ok": False, "error": "LLM did not return a valid run_at_iso"}
        return {"ok": True, "plan": fallback, "mode": "heuristic_fallback"}
    message = str(parsed.get("message") or text).strip()
    if not message:
        message = text.strip()[:180]
    return {
        "ok": True,
        "plan": {
            "action": "schedule",
            "run_at_iso": run_at.isoformat(),
            "message": message[:500],
        },
        "mode": "llm",
    }


def _schedule_telegram_instruction(
    settings: Settings,
    db: Database,
    *,
    instruction: str,
    chat_id: str,
    user_id: int | None = None,
) -> dict[str, Any]:
    if not bool(getattr(settings, "telegram_smart_commands_enabled", False)):
        return {"ok": False, "error": "TELEGRAM_SMART_COMMANDS_ENABLED is false"}
    planned = _plan_instruction_with_llm(settings=settings, db=db, text=instruction)
    if not planned.get("ok"):
        return planned
    plan = planned.get("plan") if isinstance(planned.get("plan"), dict) else {}
    action = str(plan.get("action") or "none").strip().lower()
    if action != "schedule":
        return {"ok": True, "scheduled": False, "reason": "no actionable reminder found"}
    run_at = str(plan.get("run_at_iso") or "").strip()
    run_at_dt = _parse_dt(run_at)
    if run_at_dt is None:
        return {"ok": False, "error": "invalid run_at_iso"}
    message = str(plan.get("message") or instruction).strip()[:500]
    if not message:
        return {"ok": False, "error": "empty reminder message"}
    seed = f"{chat_id}|{run_at_dt.isoformat()}|{message}"
    external_id = f"tg-reminder:{sha1(seed.encode('utf-8')).hexdigest()[:24]}"
    reminder = db.upsert_telegram_reminder(
        external_id=external_id,
        chat_id=chat_id,
        run_at=run_at_dt.isoformat(),
        message=message,
        metadata_json={
            "source": "telegram_plan",
            "instruction": instruction,
            "mode": planned.get("mode") or "unknown",
        },
        user_id=user_id,
    )
    return {
        "ok": True,
        "scheduled": True,
        "mode": planned.get("mode") or "unknown",
        "reminder": {
            "id": reminder["id"],
            "external_id": reminder["external_id"],
            "chat_id": reminder["chat_id"],
            "run_at": reminder["run_at"],
            "message": reminder["message"],
        },
    }


def _format_telegram_assistant_usage() -> str:
    return "\n".join(
        [
            "[KU] 자연어 비서",
            "",
            "- `/bot <요청>` 형식으로 보내 주세요.",
            "- 예: `/bot 오늘 일정이랑 날씨 알려줘`",
            "- 예: `/bot 내일 오전 8시에 과제 제출 알림해줘`",
        ]
    )


def _format_telegram_assistant_disabled() -> str:
    return "\n".join(
        [
            "[KU] 자연어 비서",
            "",
            "- 현재 이 기능은 비활성화되어 있습니다.",
        ]
    )


def _format_telegram_assistant_write_disabled() -> str:
    return "\n".join(
        [
            "[KU] 자연어 비서",
            "",
            "- 현재는 읽기 전용 모드입니다.",
            "- 조회 요청은 가능하지만 리마인더 생성이나 설정 변경은 비활성화되어 있습니다.",
        ]
    )


def _assistant_plan_requests_write(plan: dict[str, Any] | None) -> bool:
    actions = plan.get("actions") if isinstance(plan, dict) else None
    if not isinstance(actions, list):
        return False
    for item in actions:
        if not isinstance(item, dict):
            continue
        capability = get_capability(str(item.get("capability") or "").strip())
        if capability is not None and capability.side_effect:
            return True
    return False


def _start_telegram_chat_action_loop(
    client: TelegramBotClient,
    *,
    chat_id: str,
    action: str,
    interval_sec: float,
) -> Callable[[], None] | None:
    send_action = getattr(client, "send_chat_action", None)
    if not callable(send_action):
        return None

    safe_chat_id = str(chat_id or "").strip()
    safe_action = str(action or "").strip() or "typing"
    if not safe_chat_id:
        return None

    def _send_once() -> None:
        try:
            send_action(chat_id=safe_chat_id, action=safe_action)
        except Exception as exc:
            logger.warning(
                "failed to send telegram chat action",
                extra={"chat_id": safe_chat_id, "action": safe_action, "error": str(exc)},
            )

    stop_event = threading.Event()

    def _worker() -> None:
        wait_seconds = max(float(interval_sec), 0.01)
        while not stop_event.wait(wait_seconds):
            _send_once()

    _send_once()
    thread = threading.Thread(target=_worker, daemon=True)
    thread.start()

    def _stop() -> None:
        stop_event.set()
        thread.join(timeout=1.0)

    return _stop


def _handle_telegram_assistant_command(
    settings: Settings,
    db: Database,
    *,
    request_text: str | None = None,
    chat_id: str | None = None,
    user_id: int | None = None,
) -> dict[str, Any]:
    if not bool(getattr(settings, "telegram_assistant_enabled", False)):
        return {"ok": False, "message": _format_telegram_assistant_disabled()}

    request = str(request_text or "").strip()
    if not request:
        return {"ok": False, "message": _format_telegram_assistant_usage()}

    assistant_run: dict[str, Any] | None = None
    try:
        assistant_run = db.create_assistant_run(
            request_raw=request,
            user_id=user_id,
            chat_id=chat_id,
            context_json={
                "source": "telegram_command",
                "command": "assistant",
                "telegram_assistant_enabled": bool(
                    getattr(settings, "telegram_assistant_enabled", False)
                ),
                "telegram_assistant_write_enabled": bool(
                    getattr(settings, "telegram_assistant_write_enabled", False)
                ),
            },
            status="pending",
        )
    except Exception:
        assistant_run = None

    planned = plan_assistant_request(settings, text=request)
    if assistant_run is not None:
        db.update_assistant_run(
            int(assistant_run["id"]),
            planner_output_json=planned,
            status="planned" if not bool(planned.get("needs_clarification")) else "clarification",
        )

    if (
        _assistant_plan_requests_write(planned)
        and not bool(getattr(settings, "telegram_assistant_write_enabled", False))
    ):
        blocked = {
            "ok": False,
            "error": "assistant_write_disabled",
            "reply": _format_telegram_assistant_write_disabled(),
        }
        if assistant_run is not None:
            db.update_assistant_run(
                int(assistant_run["id"]),
                executor_result_json=blocked,
                final_reply=blocked["reply"],
                status="blocked",
            )
        return {"ok": False, "message": blocked["reply"]}

    executed = execute_assistant_plan(
        settings,
        db,
        plan=planned,
        user_id=user_id,
        chat_id=chat_id,
    )
    reply = str(executed.get("reply") or planned.get("reply") or "").strip()
    if not reply:
        reply = _format_telegram_assistant_usage()
    if assistant_run is not None:
        db.update_assistant_run(
            int(assistant_run["id"]),
            executor_result_json=executed,
            final_reply=reply,
            status=(
                "completed"
                if bool(executed.get("ok")) and not bool(executed.get("needs_clarification"))
                else "clarification"
                if bool(executed.get("ok"))
                else "failed"
            ),
        )
    return {
        "ok": bool(executed.get("ok")),
        "message": reply,
    }


def _dispatch_due_telegram_reminders(
    settings: Settings,
    db: Database,
    client: TelegramBotClient,
) -> dict[str, Any]:
    due = db.list_due_telegram_reminders(now_iso=now_utc_iso(), limit=100)
    sent = 0
    failed = 0
    for item in due:
        reminder_id = int(item["id"])
        chat_id = str(item.get("chat_id") or "").strip()
        message = str(item.get("message") or "").strip()
        if not chat_id or not message:
            db.mark_telegram_reminder_status(reminder_id, status="failed")
            failed += 1
            continue
        try:
            ok = client.send_message(chat_id=chat_id, text=f"[Reminder] {message}"[:3500])
        except Exception as exc:
            logger.warning(
                "failed to send scheduled telegram reminder",
                extra={"chat_id": chat_id, "error": str(exc)},
            )
            ok = False
        if ok:
            db.mark_telegram_reminder_status(reminder_id, status="sent")
            sent += 1
        else:
            db.mark_telegram_reminder_status(reminder_id, status="failed")
            failed += 1
    return {
        "due": len(due),
        "sent": sent,
        "failed": failed,
    }


def _execute_telegram_command(
    settings: Settings,
    db: Database,
    command_payload: dict[str, Any],
    *,
    chat_id: str | None = None,
    user_id: int | None = None,
) -> dict[str, Any]:
    command = str(command_payload.get("command") or "").strip().lower()
    if user_id is None and chat_id:
        user_scope = _resolve_user_scope(
            settings,
            db,
            chat_id=chat_id,
            create_if_missing=True,
            metadata_source="telegram_command_execute",
        )
        user_id = int(user_scope["user_id"] or 0)
    if command == "assistant" and not bool(command_payload.get("ok")):
        return {"ok": False, "message": _format_telegram_assistant_usage()}
    if not bool(command_payload.get("ok")):
        return {"ok": False, "error": command_payload.get("error") or "invalid command"}
    if command == "start":
        return {
            "ok": True,
            "message": _format_telegram_start(settings, db=db, chat_id=chat_id, user_id=user_id),
        }
    if command == "help":
        return {"ok": True, "message": _format_telegram_help(settings)}
    if command == "setup":
        return {
            "ok": True,
            "message": _format_telegram_setup(settings, db=db, chat_id=chat_id, user_id=user_id),
        }
    if command == "connect_moodle":
        return _issue_moodle_connect_link(
            settings=settings,
            db=db,
            chat_id=chat_id,
            school_query=command_payload.get("school_query"),
        )
    if not _is_telegram_chat_allowed(settings, db, chat_id):
        return {"ok": False, "message": _format_telegram_access_denied(chat_id)}
    if command == "status":
        return {"ok": True, "message": _format_telegram_status(settings, db, user_id=user_id)}
    if command == "today":
        return {"ok": True, "message": _format_telegram_today(settings, db, user_id=user_id)}
    if command == "tomorrow":
        return {"ok": True, "message": _format_telegram_tomorrow(settings, db, user_id=user_id)}
    if command == "region":
        return _handle_telegram_region_command(
            settings,
            db,
            query=command_payload.get("query"),
            chat_id=chat_id,
            user_id=user_id,
        )
    if command in {"weather", "todayweather"}:
        return {
            "ok": True,
            "message": _format_telegram_todayweather(
                settings,
                db,
                user_id=user_id,
                chat_id=chat_id,
            ),
        }
    if command == "today_summary":
        return {"ok": True, "message": _format_telegram_today_summary(settings, db, user_id=user_id)}
    if command == "tomorrow_summary":
        return {"ok": True, "message": _format_telegram_tomorrow_summary(settings, db, user_id=user_id)}
    if command == "notice_uclass":
        return {
            "ok": True,
            "message": _format_telegram_uclass_notice(settings, db, user_id=user_id),
        }
    if command == "library":
        return {
            "ok": True,
            "message": _format_telegram_library(command_payload.get("library")),
        }
    if command == "assignments":
        return {
            "ok": True,
            "message": _format_telegram_assignments(
                settings=settings,
                db=db,
                user_id=user_id,
                chat_id=chat_id,
            ),
        }
    if command == "submitted_assignments":
        return {
            "ok": True,
            "message": _format_telegram_submitted_assignments(
                settings=settings,
                db=db,
                user_id=user_id,
                chat_id=chat_id,
            ),
        }
    if command == "lms_board":
        return {
            "ok": True,
            "message": _format_telegram_lms_board(
                settings=settings,
                db=db,
                user_id=user_id,
                chat_id=chat_id,
            ),
        }
    if command == "lms_materials":
        return {
            "ok": True,
            "message": _format_telegram_lms_materials(
                settings=settings,
                db=db,
                user_id=user_id,
                chat_id=chat_id,
            ),
        }
    if command == "notice_general":
        return {
            "ok": True,
            "message": _format_telegram_uos_notice(
                db,
                "general",
                timezone_name=settings.timezone,
                user_id=user_id,
            ),
        }
    if command == "notice_academic":
        return {
            "ok": True,
            "message": _format_telegram_uos_notice(
                db,
                "academic",
                timezone_name=settings.timezone,
                user_id=user_id,
            ),
        }
    if command == "inbox":
        return {"ok": True, "message": _format_telegram_inbox(db, user_id=user_id)}
    if command == "apply":
        scope = str(command_payload.get("scope") or "")
        if scope == "all":
            result = apply_inbox_items(settings=settings, db=db, apply_all=True, user_id=user_id)
            return {"ok": True, "message": _format_telegram_apply_result(result)}
        selected_id = _safe_int(command_payload.get("id"))
        if selected_id is None:
            return {"ok": False, "error": "apply id must be numeric inbox row id"}
        result = apply_inbox_items(settings=settings, db=db, item_id=selected_id, user_id=user_id)
        return {"ok": True, "message": _format_telegram_apply_result(result)}
    if command == "done":
        target = str(command_payload.get("target") or "").strip().lower()
        selector = str(command_payload.get("id") or "").strip()
        if target == "task":
            result = mark_task_status(
                settings=settings,
                db=db,
                selector=selector,
                status="done",
                user_id=user_id,
            )
            return {
                "ok": bool(result.get("ok")),
                "message": _format_telegram_done_result(result, timezone_name=settings.timezone),
            }
        return {"ok": False, "error": "done target must be task"}
    if command == "plan":
        if not chat_id:
            return {"ok": False, "error": "chat_id is missing for /plan"}
        instruction = str(command_payload.get("instruction") or "").strip()
        if not instruction:
            return {"ok": False, "error": "missing instruction text"}
        result = _schedule_telegram_instruction(
            settings=settings,
            db=db,
            instruction=instruction,
            chat_id=chat_id,
            user_id=user_id,
        )
        return {
            "ok": bool(result.get("ok")),
            "message": _format_telegram_plan_result(result, timezone_name=settings.timezone),
        }
    if command == "assistant":
        return _handle_telegram_assistant_command(
            settings,
            db,
            request_text=command_payload.get("request"),
            chat_id=chat_id,
            user_id=user_id,
        )
    return {"ok": False, "error": f"unsupported command: {command}"}


def _process_telegram_commands(
    settings: Settings,
    db: Database,
    client: TelegramBotClient,
) -> dict[str, Any]:
    if not getattr(settings, "telegram_commands_enabled", False):
        return {"skipped": True, "reason": "TELEGRAM_COMMANDS_ENABLED is false"}
    gate = _identity_warning_gate(
        settings=settings,
        db=db,
        step="telegram_command_reply",
        destination="telegram",
    )
    processed = 0
    failed = 0
    blocked_sends = 0
    for item in db.list_unprocessed_inbox_commands(limit=200):
        draft = _json_load(item.draft_json)
        metadata = _json_load(item.metadata_json)
        chat_id = str(metadata.get("chat_id") or "").strip() or None
        command = str(draft.get("command") or "").strip().lower()
        owner_id = int(item.user_id or 0) or _resolve_user_scope(
            settings,
            db,
            chat_id=chat_id,
            create_if_missing=True,
            metadata_source="telegram_command",
        )["user_id"]
        stop_chat_action = None
        if (
            gate is None
            and command == "assistant"
            and bool(draft.get("ok"))
            and bool(getattr(settings, "telegram_assistant_enabled", False))
            and bool(chat_id)
            and _is_telegram_chat_allowed(settings, db, chat_id)
        ):
            stop_chat_action = _start_telegram_chat_action_loop(
                client,
                chat_id=str(chat_id),
                action=TELEGRAM_ASSISTANT_CHAT_ACTION,
                interval_sec=TELEGRAM_ASSISTANT_CHAT_ACTION_INTERVAL_SEC,
            )
        try:
            result = _execute_telegram_command(
                settings=settings,
                db=db,
                command_payload=draft,
                chat_id=chat_id,
                user_id=owner_id,
            )
        finally:
            if callable(stop_chat_action):
                stop_chat_action()
        metadata["command_result"] = result
        db.upsert_inbox_item(
            external_id=item.external_id,
            source=item.source,
            received_at=item.received_at,
            title=item.title,
            body=item.body,
            item_type=item.item_type,
            draft_json=draft,
            processed=True,
            metadata_json=metadata,
            user_id=owner_id,
        )
        if not result.get("ok"):
            failed += 1
        processed += 1

        chat_id = metadata.get("chat_id")
        message = str(result.get("message") or result.get("error") or "command processed")
        if chat_id:
            if gate is not None:
                blocked_sends += 1
                continue
            try:
                client.send_message(chat_id=str(chat_id), text=message[:3500])
            except Exception as exc:
                logger.warning(
                    "failed to send command reply",
                    extra={"chat_id": str(chat_id), "error": str(exc)},
                )
    payload: dict[str, Any] = {"processed": processed, "failed": failed}
    if gate is not None:
        payload["blocked_sends"] = blocked_sends
        payload["warning_gate"] = gate["warning_gate"]
    return payload


def _sync_telegram_once(
    settings: Settings,
    db: Database,
    *,
    client: TelegramBotClient | None = None,
    poll_timeout: int = 10,
) -> dict[str, Any]:
    state = db.get_sync_state("sync_telegram")
    cursor = _json_load(state.last_cursor_json)
    next_offset = cursor.get("next_offset")
    if not isinstance(next_offset, int):
        next_offset = None
    previous_menu_hash = str(cursor.get("telegram_menu_hash") or "").strip() or None
    previous_menu_target_hash = str(cursor.get("telegram_menu_target_hash") or "").strip() or None
    previous_menu_retry_after = str(cursor.get("telegram_menu_retry_after") or "").strip() or None

    active_client = client or TelegramBotClient(settings.telegram_bot_token)
    menu_result = _register_telegram_bot_menu(
        settings=settings,
        client=active_client,
        existing_hash=previous_menu_hash,
        pending_hash=previous_menu_target_hash,
        retry_after_iso=previous_menu_retry_after,
    )
    updates = active_client.get_updates(
        offset=next_offset,
        limit=settings.telegram_poll_limit,
        timeout=max(int(poll_timeout), 0),
    )
    allowed_chat_ids = _effective_telegram_allowed_chat_ids(settings, db)
    items = normalize_updates(
        updates,
        timezone_name=settings.timezone,
        allowed_chat_ids=allowed_chat_ids,
    )
    stored = 0
    touched_owner_ids: set[int] = set()
    for item in items:
        metadata = dict(item.metadata or {})
        chat_id = str(metadata.get("chat_id") or "").strip()
        owner = _resolve_user_scope(
            settings,
            db,
            chat_id=chat_id,
            create_if_missing=bool(chat_id),
            metadata_source="telegram_update",
        )
        metadata["user_id"] = int(owner["user_id"] or 0)
        if int(owner["user_id"] or 0) > 0:
            touched_owner_ids.add(int(owner["user_id"] or 0))
        db.upsert_inbox_item(
            external_id=item.external_id,
            source="telegram",
            received_at=item.received_at,
            title=item.title,
            body=item.body,
            item_type=item.item_type,
            draft_json=item.draft,
            processed=False,
            metadata_json=metadata,
            user_id=int(owner["user_id"] or 0),
        )
        stored += 1

    for owner_id in sorted(touched_owner_ids):
        _precompute_task_merge_cache(
            settings,
            db,
            user_id=owner_id,
        )

    command_result = _process_telegram_commands(
        settings=settings,
        db=db,
        client=active_client,
    )
    reminder_result = _dispatch_due_telegram_reminders(
        settings=settings,
        db=db,
        client=active_client,
    )

    update_ids = [
        int(update.get("update_id"))
        for update in updates
        if isinstance(update.get("update_id"), int)
    ]
    next_cursor = (max(update_ids) + 1) if update_ids else next_offset
    menu_hash = previous_menu_hash
    if isinstance(menu_result, dict) and str(menu_result.get("hash") or "").strip():
        menu_hash = str(menu_result.get("hash")).strip()
    menu_target_hash = previous_menu_target_hash
    if isinstance(menu_result, dict) and str(menu_result.get("target_hash") or "").strip():
        menu_target_hash = str(menu_result.get("target_hash")).strip()
    menu_retry_after = None
    if isinstance(menu_result, dict) and str(menu_result.get("retry_after") or "").strip():
        menu_retry_after = str(menu_result.get("retry_after")).strip()
    if menu_hash and menu_target_hash and menu_hash == menu_target_hash:
        menu_retry_after = None
    failed_commands = int(
        (command_result.get("failed") if isinstance(command_result, dict) else 0) or 0
    )
    failed_reminders = int(
        (reminder_result.get("failed") if isinstance(reminder_result, dict) else 0) or 0
    )
    blocked_sends = int(
        (command_result.get("blocked_sends") if isinstance(command_result, dict) else 0) or 0
    )
    _record_sync_dashboard_state(
        db,
        "sync_telegram",
        status="success",
        new_items=stored,
        action_required=failed_commands + failed_reminders + blocked_sends,
        cursor_payload={
            "next_offset": next_cursor,
            "telegram_menu_hash": menu_hash,
            "telegram_menu_target_hash": menu_target_hash,
            "telegram_menu_retry_after": menu_retry_after,
            "fetched": len(updates),
            "stored": stored,
            "menu": menu_result,
            "commands": command_result,
            "reminders": reminder_result,
        },
    )
    return {
        "fetched_updates": len(updates),
        "stored_messages": stored,
        "menu": menu_result,
        "commands": command_result,
        "reminders": reminder_result,
    }


def sync_telegram(
    settings: Settings,
    db: Database,
    *,
    client: TelegramBotClient | None = None,
    poll_timeout: int = 10,
) -> dict[str, Any]:
    if not settings.telegram_enabled:
        reason = "TELEGRAM_ENABLED is false"
        _record_sync_dashboard_state(
            db,
            "sync_telegram",
            status="skipped",
            last_error=reason,
            cursor_payload={"skipped": True, "reason": reason},
        )
        return {"skipped": True, "reason": reason}
    if not settings.telegram_bot_token:
        reason = "TELEGRAM_BOT_TOKEN missing"
        _record_sync_dashboard_state(
            db,
            "sync_telegram",
            status="skipped",
            last_error=reason,
            cursor_payload={"skipped": True, "reason": reason},
        )
        return {"skipped": True, "reason": reason}
    try:
        return _sync_telegram_once(
            settings=settings,
            db=db,
            client=client,
            poll_timeout=poll_timeout,
        )
    except Exception as exc:
        logger.warning("telegram sync failed", extra={"error": str(exc)})
        _record_sync_dashboard_state(
            db,
            "sync_telegram",
            status="error",
            action_required=1,
            last_error=str(exc),
            cursor_payload={"error": str(exc)},
        )
        return {"error": str(exc)}


def run_uclass_probe(
    settings: Settings,
    db: Database,
    output_json_path: Path | None = None,
) -> dict[str, Any]:
    report = build_uclass_probe_report(
        settings=settings,
        output_json_path=output_json_path,
    )
    ok = all(row.get("status") in {"OK", "SKIP"} for row in report.get("rows", []))
    db.update_sync_state(
        "uclass_probe",
        last_run_at=now_utc_iso(),
        last_cursor_json={
            "ok": ok,
            "output_json_path": str(output_json_path) if output_json_path else None,
        },
    )
    return report


def apply_inbox_items(
    settings: Settings,
    db: Database,
    item_id: int | None = None,
    apply_all: bool = False,
    *,
    user_id: int | None = None,
) -> dict[str, Any]:
    if item_id is None and not apply_all:
        raise ValueError("provide --id or --all")
    if item_id is not None and apply_all:
        raise ValueError("choose either --id or --all, not both")

    items = (
        db.list_unprocessed_inbox(limit=500, user_id=user_id)
        if apply_all
        else [db.get_inbox_item_by_id(item_id, user_id=user_id)] if item_id is not None else []
    )
    selected = [
        item for item in items if item is not None and item.item_type != "command"
    ]
    if not selected:
        return {"processed": 0, "created_events": 0, "created_tasks": 0, "notes": 0}

    created_events = 0
    created_tasks = 0
    note_count = 0
    processed = 0
    applied_ids: list[int] = []
    requested_user_id = _safe_int(user_id) or None
    per_owner_counts: dict[int, dict[str, Any]] = {}
    for item in selected:
        if item.processed:
            continue
        owner_id = requested_user_id if requested_user_id is not None else int(item.user_id or 0) or None
        owner_key = int(owner_id or 0)
        owner_stats = per_owner_counts.setdefault(
            owner_key,
            {
                "processed": 0,
                "created_events": 0,
                "created_tasks": 0,
                "notes": 0,
                "item_ids": [],
            },
        )
        draft = _json_load(item.draft_json)
        external_id = f"inbox:{_safe_inbox_suffix(item.external_id)}"
        if item.item_type == "event_draft":
            start_dt = _parse_dt(str(draft.get("start_at") or item.received_at))
            if start_dt is None:
                start_dt = datetime.now(timezone.utc).replace(microsecond=0)
            end_dt = _parse_dt(str(draft.get("end_at") or "")) or (start_dt + timedelta(hours=1))
            if end_dt <= start_dt:
                end_dt = start_dt + timedelta(hours=1)
            metadata = attach_provenance(
                {
                    "inbox_id": item.id,
                    "inbox_external_id": item.external_id,
                    "draft": draft,
                },
                source="telegram_draft",
                confidence="low",
                last_verified_at=item.received_at,
                raw_source_ids=[item.external_id, str(item.id or "")],
                derivation="telegram_event_parse",
            )
            db.upsert_event(
                external_id=external_id,
                source="inbox",
                start=start_dt.isoformat(),
                end=end_dt.isoformat(),
                title=str(draft.get("title") or item.title or "Inbox Event"),
                location=str(draft.get("location") or "").strip() or None,
                rrule=None,
                metadata_json=metadata,
                user_id=owner_id,
            )
            created_events += 1
            owner_stats["created_events"] = int(owner_stats.get("created_events") or 0) + 1
        elif item.item_type == "task_draft":
            due_at_raw = draft.get("due_at")
            due_dt = _parse_dt(str(due_at_raw)) if due_at_raw else None
            metadata = attach_provenance(
                {
                    "inbox_id": item.id,
                    "inbox_external_id": item.external_id,
                    "draft": draft,
                },
                source="telegram_draft",
                confidence="low",
                last_verified_at=item.received_at,
                raw_source_ids=[item.external_id, str(item.id or "")],
                derivation="telegram_task_parse",
            )
            db.upsert_task(
                external_id=external_id,
                source="inbox",
                due_at=due_dt.isoformat() if due_dt else None,
                title=str(draft.get("title") or item.title or "Inbox Task"),
                status=str(draft.get("status") or "open"),
                metadata_json=metadata,
                user_id=owner_id,
            )
            created_tasks += 1
            owner_stats["created_tasks"] = int(owner_stats.get("created_tasks") or 0) + 1
        else:
            note_count += 1
            owner_stats["notes"] = int(owner_stats.get("notes") or 0) + 1
        if item.id is not None:
            db.mark_inbox_processed_by_id(item.id, user_id=owner_id)
            applied_ids.append(item.id)
            owner_stats["item_ids"].append(item.id)
        else:
            db.mark_inbox_processed(item.external_id, item.source, user_id=owner_id)
        processed += 1
        owner_stats["processed"] = int(owner_stats.get("processed") or 0) + 1

    sync_cursor = {
        "processed": processed,
        "created_events": created_events,
        "created_tasks": created_tasks,
        "notes": note_count,
        "item_ids": applied_ids,
    }
    if requested_user_id is not None:
        db.update_sync_state(
            "apply_inbox",
            last_run_at=now_utc_iso(),
            last_cursor_json=sync_cursor,
            user_id=requested_user_id,
        )
    else:
        for owner_key, owner_stats in per_owner_counts.items():
            if owner_key <= 0:
                continue
            db.update_sync_state(
                "apply_inbox",
                last_run_at=now_utc_iso(),
                last_cursor_json={
                    "processed": int(owner_stats.get("processed") or 0),
                    "created_events": int(owner_stats.get("created_events") or 0),
                    "created_tasks": int(owner_stats.get("created_tasks") or 0),
                    "notes": int(owner_stats.get("notes") or 0),
                    "item_ids": list(owner_stats.get("item_ids") or []),
                },
                user_id=owner_key,
            )
        db.update_sync_state(
            "apply_inbox",
            last_run_at=now_utc_iso(),
            last_cursor_json=sync_cursor,
        )
    return {
        "processed": processed,
        "created_events": created_events,
        "created_tasks": created_tasks,
        "notes": note_count,
    }


def ignore_inbox_item(db: Database, item_id: int) -> dict[str, Any]:
    changed = db.mark_inbox_ignored_by_id(item_id)
    return {"ok": changed, "id": item_id}


def mark_task_status(
    settings: Settings,
    db: Database,
    selector: str,
    status: str,
    *,
    user_id: int | None = None,
) -> dict[str, Any]:
    updated = db.update_task_status(selector=selector, status=status, user_id=user_id)
    if not updated:
        return {"ok": False, "reason": "task not found", "selector": selector}

    return {
        "ok": True,
        "task": {
            "id": updated["id"],
            "external_id": updated["external_id"],
            "source": updated["source"],
            "status": updated["status"],
            "title": updated["title"],
        },
    }


def mark_review_status(
    settings: Settings,
    db: Database,
    selector: str,
    review_status: str,
) -> dict[str, Any]:
    updated = db.update_review_status(selector=selector, review_status=review_status)
    if not updated:
        return {"ok": False, "reason": "review not found", "selector": selector}

    return {
        "ok": True,
        "review": {
            "id": updated["id"],
            "external_id": updated["external_id"],
            "source": updated["source"],
            "title": updated["title"],
            "review_status": updated["review_status"],
        },
    }


def schedule_review_events(settings: Settings, db: Database) -> dict[str, Any]:
    if not settings.review_enabled:
        return {"skipped": True, "reason": "REVIEW_ENABLED is false"}
    intervals = sorted({day for day in settings.review_intervals_days if int(day) > 0})
    if not intervals:
        return {"skipped": True, "reason": "REVIEW_INTERVALS_DAYS is empty"}
    duration_min = max(5, settings.review_duration_min)
    tz = ZoneInfo(settings.timezone)
    morning_hour = min(max(int(settings.review_morning_hour), 0), 23)
    created = 0
    source_count = 0
    generated_ids: set[str] = set()
    existing_review_events = {
        item.external_id: item
        for item in db.list_events(limit=5000, include_inactive_reviews=True)
        if item.source == "review" or item.external_id.startswith("review:")
    }
    existing_review_ids = set(existing_review_events.keys())

    for event in db.list_events(limit=2000):
        if event.source == "review" or event.external_id.startswith("review:"):
            continue
        end_dt = _parse_dt(event.end_at)
        if end_dt is None:
            continue
        source_count += 1
        for day in intervals:
            review_external_id = f"review:{event.external_id}:D+{day}"
            if review_external_id in generated_ids:
                continue
            generated_ids.add(review_external_id)
            start_dt = end_dt + timedelta(days=day)
            end_review = start_dt + timedelta(minutes=duration_min)
            existing_meta = _json_load(
                existing_review_events.get(review_external_id).metadata_json
                if review_external_id in existing_review_events
                else None
            )
            review_status = str(existing_meta.get("review_status") or "scheduled")
            db.upsert_event(
                external_id=review_external_id,
                source="review",
                start=start_dt.isoformat(),
                end=end_review.isoformat(),
                title=f"Review: {event.title} (D+{day})",
                location=event.location,
                rrule=None,
                metadata_json={
                    "review_of_external_id": event.external_id,
                    "review_of_source": event.source,
                    "interval_days": day,
                    "kind": "event",
                    "review_status": review_status,
                },
            )
            if review_external_id not in existing_review_ids:
                created += 1
                existing_review_ids.add(review_external_id)

    for artifact in db.list_artifacts(limit=2000):
        if artifact.external_id.startswith("review:"):
            continue
        if artifact.updated_at is None:
            continue
        updated_dt = _parse_dt(artifact.updated_at)
        if updated_dt is None:
            continue
        local_updated = updated_dt.astimezone(tz)
        base_local = local_updated.replace(
            hour=morning_hour, minute=0, second=0, microsecond=0
        )
        if local_updated >= base_local:
            base_local = base_local + timedelta(days=1)
        meta = _json_load(artifact.metadata_json)
        source_title = (
            str(meta.get("course_name") or "").strip()
            or str(meta.get("module_name") or "").strip()
            or artifact.filename
        )
        source_count += 1
        for day in intervals:
            review_external_id = f"review:{artifact.external_id}:D+{day}"
            if review_external_id in generated_ids:
                continue
            generated_ids.add(review_external_id)
            start_local = base_local + timedelta(days=day)
            end_local = start_local + timedelta(minutes=duration_min)
            existing_meta = _json_load(
                existing_review_events.get(review_external_id).metadata_json
                if review_external_id in existing_review_events
                else None
            )
            review_status = str(existing_meta.get("review_status") or "scheduled")
            db.upsert_event(
                external_id=review_external_id,
                source="review",
                start=start_local.isoformat(),
                end=end_local.isoformat(),
                title=f"Review: {source_title} (D+{day})",
                location=None,
                rrule=None,
                metadata_json={
                    "review_of_external_id": artifact.external_id,
                    "review_of_source": artifact.source,
                    "interval_days": day,
                    "kind": "artifact",
                    "review_status": review_status,
                },
            )
            if review_external_id not in existing_review_ids:
                created += 1
                existing_review_ids.add(review_external_id)
    db.update_sync_state(
        "schedule_reviews",
        last_run_at=now_utc_iso(),
        last_cursor_json={
            "generated": created,
            "intervals": intervals,
            "sources": source_count,
        },
    )
    return {"generated": created, "intervals": intervals, "sources": source_count}


def _parse_clock_time(value: str, field_name: str) -> tuple[int, int]:
    match = re.match(r"^(\d{1,2}):(\d{2})$", str(value).strip())
    if not match:
        raise ValueError(f"{field_name} must be HH:MM")
    hour = int(match.group(1))
    minute = int(match.group(2))
    if hour < 0 or hour > 23 or minute < 0 or minute > 59:
        raise ValueError(f"{field_name} is out of range")
    return hour, minute


def _title_tokens(value: str) -> list[str]:
    return [item.lower() for item in re.findall(r"[A-Za-z0-9\u3131-\u318E\uAC00-\uD7A3]{2,}", str(value or ""))]


def _title_match_score(reference: str, candidate: str) -> int:
    ref = str(reference or "").strip()
    cand = str(candidate or "").strip()
    if not ref or not cand:
        return 0
    ref_norm = _normalize_token_text(ref)
    cand_norm = _normalize_token_text(cand)
    score = 0
    if ref_norm and cand_norm:
        if ref_norm == cand_norm:
            score = max(score, 100)
        elif ref_norm in cand_norm or cand_norm in ref_norm:
            score = max(score, 70)
    ref_tokens = set(_title_tokens(ref))
    cand_tokens = set(_title_tokens(cand))
    overlap = ref_tokens.intersection(cand_tokens)
    if overlap:
        score = max(score, len(overlap) * 12)
        if len(overlap) >= 2:
            score = max(score, 30)
        elif len(ref_tokens) <= 1 or len(cand_tokens) <= 1:
            score = max(score, 20)
    return score


def _format_time_range_local(start_dt: datetime, end_dt: datetime) -> str:
    return f"{start_dt.strftime('%H:%M')}-{end_dt.strftime('%H:%M')}"


def _render_class_location(
    *,
    raw_location: str | None,
    building_no: str | None,
    building_name: str | None,
    room: str | None,
) -> str:
    if building_name:
        room_part = f" {room}호" if room else ""
        return f"{building_name}{room_part}"
    if building_no and room:
        return f"{building_no}번 건물 {room}호"
    if building_no:
        return f"{building_no}번 건물"
    if str(raw_location or "").strip():
        return str(raw_location).strip()
    return "TBD"


def _metadata_text_value(metadata: dict[str, Any] | None, *keys: str) -> str | None:
    payload = metadata if isinstance(metadata, dict) else {}
    for key in keys:
        if not key:
            continue
        value = payload.get(key)
        if isinstance(value, (list, tuple)):
            value = " ".join(str(item).strip() for item in value if str(item).strip())
        text = str(value or "").strip()
        if text:
            return text
    return None


def _normalize_room_label(room: str | None, *, building_no: str | None = None) -> str | None:
    text = re.sub(r"\s*호$", "", str(room or "").strip())
    if not text:
        return None
    building_key = str(building_no or "").strip()
    if building_key:
        match = re.match(
            rf"^\s*{re.escape(building_key)}\s*-\s*(?P<room>[A-Za-z]?\d+(?:[.-]\d+)?)\s*$",
            text,
        )
        if match:
            return str(match.group("room") or "").strip() or text
    return text


def _class_name_candidates(class_item: dict[str, Any]) -> list[str]:
    output: list[str] = []
    seen: set[str] = set()
    for key in ("title", "official_course_name", "course_display_name"):
        value = str(class_item.get(key) or "").strip()
        normalized = normalize_course_alias(value)
        if not value or not normalized or normalized in seen:
            continue
        seen.add(normalized)
        output.append(value)
    return output


def _best_title_match_score(references: list[str], candidate: str) -> int:
    return max((_title_match_score(reference, candidate) for reference in references if reference), default=0)


def _parse_material_reference_dt(value: Any) -> datetime | None:
    return _parse_datetime_like(value)


def _material_reference_dt(metadata: dict[str, Any], *, artifact: Any) -> datetime | None:
    raw = metadata.get("raw") if isinstance(metadata.get("raw"), dict) else {}
    for value in (
        metadata.get("downloaded_at"),
        metadata.get("last_verified_at"),
        raw.get("timemodified"),
        raw.get("timecreated"),
        raw.get("created"),
        raw.get("createdat"),
        getattr(artifact, "updated_at", None),
    ):
        parsed = _parse_material_reference_dt(value)
        if parsed is not None:
            return parsed
    return None


def _enrich_class_item_metadata(
    db: Database,
    *,
    title: str,
    raw_location: str | None,
    metadata_json: dict[str, Any] | None,
    alias_map: dict[str, tuple[str, ...]] | None = None,
    school_slug: str = "ku_online_class",
    user_id: int | None = None,
) -> dict[str, Any]:
    metadata = dict(metadata_json or {})
    canonical_course_id = str(metadata.get("canonical_course_id") or "").strip()
    if not canonical_course_id and alias_map:
        canonical_course_id = _resolve_canonical_course_id(
            metadata,
            alias_map=alias_map,
            fallback_aliases=[title],
        )

    course = db.get_course(canonical_course_id, user_id=user_id) if canonical_course_id else None
    course_metadata = _json_load(getattr(course, "metadata_json", None)) if course is not None else {}

    official_location_present = any(
        _metadata_text_value(metadata, key)
        for key in ("official_building_no", "official_building_name", "official_room")
    )
    stored_location_present = any(
        _metadata_text_value(metadata, key)
        for key in ("building_no", "building_name", "room")
    )
    parsed_location = _parse_uos_location(raw_location)
    parsed_building_no = str(parsed_location.get("building_no") or "").strip() or None
    parsed_room = str(parsed_location.get("room") or "").strip() or None
    parsed_extra = str(parsed_location.get("extra") or "").strip() or None

    building_no = _metadata_text_value(metadata, "official_building_no") or _metadata_text_value(
        metadata, "building_no"
    )
    building_name = _metadata_text_value(metadata, "official_building_name") or _metadata_text_value(
        metadata, "building_name"
    )
    room = _normalize_room_label(
        _metadata_text_value(metadata, "official_room"),
        building_no=building_no,
    ) or _normalize_room_label(
        _metadata_text_value(metadata, "room"),
        building_no=building_no,
    )
    location_extra = _metadata_text_value(metadata, "location_extra") or parsed_extra

    if not building_no and parsed_building_no:
        building_no = parsed_building_no
    if not room and parsed_room and (not building_no or not parsed_building_no or building_no == parsed_building_no):
        room = parsed_room
    if building_no and not building_name:
        building_name = db.get_building_name(str(building_no), school_slug=school_slug)
    if not building_no and building_name and parsed_building_no:
        building_no = parsed_building_no

    location_source = "missing"
    location_confidence = "none"
    if official_location_present:
        location_source = "official"
        location_confidence = "high"
    elif stored_location_present:
        location_source = "metadata"
        location_confidence = "medium"
    elif parsed_location.get("ok"):
        location_source = "parsed"
        location_confidence = "low"
    elif str(raw_location or "").strip():
        location_source = "raw"
        location_confidence = "low"

    return {
        "canonical_course_id": canonical_course_id or None,
        "course_display_name": (
            str(getattr(course, "display_name", "") or "").strip() if course is not None else None
        )
        or None,
        "official_course_name": _metadata_text_value(metadata, "official_course_name")
        or _metadata_text_value(course_metadata, "official_course_name"),
        "official_course_code": _metadata_text_value(metadata, "official_course_code")
        or _metadata_text_value(
            course_metadata,
            "official_course_code",
            "course_code",
            "subject_code",
            "lecture_code",
        ),
        "syllabus_url": _metadata_text_value(metadata, "official_syllabus_url", "syllabus_url")
        or _metadata_text_value(
            course_metadata,
            "official_syllabus_url",
            "syllabus_url",
            "lecture_plan_url",
            "course_plan_url",
        ),
        "syllabus_id": _metadata_text_value(metadata, "official_syllabus_id")
        or _metadata_text_value(course_metadata, "official_syllabus_id", "syllabus_id", "plan_id"),
        "building_no": building_no,
        "building_name": building_name,
        "room": room,
        "location_extra": location_extra,
        "location_source": location_source,
        "location_confidence": location_confidence,
        "location_text": _render_class_location(
            raw_location=raw_location,
            building_no=building_no,
            building_name=building_name,
            room=room,
        ),
    }


def _collect_primary_meetings_scoped(
    settings: Settings,
    db: Database,
    *,
    target_day_local: datetime,
    user_id: int | None = None,
) -> dict[str, Any]:
    del settings, db, target_day_local, user_id
    return {"ok": True, "events": []}


def _collect_class_occurrences(
    settings: Settings,
    db: Database,
    *,
    target_day_local: datetime,
    max_items: int,
    user_id: int | None = None,
    alias_map: dict[str, tuple[str, ...]] | None = None,
) -> list[dict[str, Any]]:
    output: list[dict[str, Any]] = []
    seen: set[str] = set()
    resolved_alias_map = alias_map if alias_map is not None else db.course_alias_resolution_map(user_id=user_id)
    timezone_name = str(getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul")
    course_week_anchor_cache: dict[str, datetime | None] = {}
    for event in db.list_events(limit=3000, user_id=user_id):
        if not _is_timetable_event(event):
            continue
        occurrences = _event_occurrences_on_date(
            event=event,
            target_date_local=target_day_local,
            timezone_name=timezone_name,
        )
        for start_local, end_local in occurrences:
            dedupe_key = f"{event.external_id}|{start_local.isoformat()}"
            if dedupe_key in seen:
                continue
            seen.add(dedupe_key)
            metadata = _json_load(event.metadata_json)
            school_slug = str(metadata.get("school_slug") or "ku_online_class").strip() or "ku_online_class"
            enriched = _enrich_class_item_metadata(
                db,
                title=event.title,
                raw_location=event.location,
                metadata_json=metadata,
                alias_map=resolved_alias_map,
                school_slug=school_slug,
                user_id=user_id,
            )
            canonical_course_id = str(enriched["canonical_course_id"] or "").strip()
            if canonical_course_id not in course_week_anchor_cache:
                course_week_anchor_cache[canonical_course_id] = _course_week_anchor_start_local(
                    db,
                    canonical_course_id=canonical_course_id,
                    timezone_name=timezone_name,
                    user_id=user_id,
                )
            output.append(
                {
                    "external_id": event.external_id,
                    "title": event.title,
                    "canonical_course_id": canonical_course_id,
                    "start_local": start_local,
                    "end_local": end_local,
                    "course_display_name": enriched["course_display_name"],
                    "official_course_name": enriched["official_course_name"],
                    "official_course_code": enriched["official_course_code"],
                    "syllabus_url": enriched["syllabus_url"],
                    "syllabus_id": enriched["syllabus_id"],
                    "building_no": enriched["building_no"],
                    "building_name": enriched["building_name"],
                    "room": enriched["room"],
                    "location_extra": enriched["location_extra"],
                    "location_source": enriched["location_source"],
                    "location_confidence": enriched["location_confidence"],
                    "location_text": enriched["location_text"],
                    "occurrence_week_index": _class_occurrence_week_index(
                        event,
                        occurrence_start_local=start_local,
                        timezone_name=timezone_name,
                        anchor_start_local=course_week_anchor_cache.get(canonical_course_id),
                    ),
                }
            )
    output.sort(key=lambda row: (row["start_local"], str(row["title"]).lower()))
    capped = max(int(max_items), 1)
    return output[:capped]


def _matched_artifacts_for_class(
    db: Database,
    class_item: dict[str, Any],
    artifacts: list[Any],
    limit: int = 2,
    *,
    user_id: int | None = None,
    match_context: _DayBriefMatchContext | None = None,
) -> list[dict[str, Any]]:
    class_titles = _class_name_candidates(class_item)
    class_course_id = str(class_item.get("canonical_course_id") or "").strip()
    occurrence_week_index = _safe_int(class_item.get("occurrence_week_index"))
    class_week_markers = (
        (occurrence_week_index,)
        if occurrence_week_index is not None and occurrence_week_index > 0
        else ()
    )
    class_start_local = class_item.get("start_local")
    class_start_dt = class_start_local if isinstance(class_start_local, datetime) else None
    alias_map = (
        match_context.alias_map
        if match_context is not None
        else db.course_alias_resolution_map(user_id=user_id)
    )
    matches: list[dict[str, Any]] = []
    artifact_candidates = (
        match_context.artifact_candidates_for_course(class_course_id)
        if match_context is not None
        else ()
    )
    for artifact_candidate in artifact_candidates:
        artifact_course_id = artifact_candidate.canonical_course_id
        score = 200 if class_course_id and artifact_course_id and class_course_id == artifact_course_id else 0
        if score <= 0:
            score = _best_title_match_score(class_titles, artifact_candidate.course_name)
        if score <= 0 and artifact_candidate.module_name:
            score = _best_title_match_score(class_titles, artifact_candidate.module_name)
        if score <= 0 and artifact_candidate.section_name:
            score = _best_title_match_score(class_titles, artifact_candidate.section_name)
        if score <= 0:
            continue
        week_rank = 1
        if class_week_markers and artifact_candidate.material_week_markers:
            week_rank = 2 if set(class_week_markers) & set(artifact_candidate.material_week_markers) else 0
        date_distance_sec = float("inf")
        if class_start_dt is not None and artifact_candidate.reference_dt is not None:
            date_distance_sec = abs(
                (
                    artifact_candidate.reference_dt.astimezone(class_start_dt.tzinfo or timezone.utc)
                    - class_start_dt
                ).total_seconds()
            )
        matches.append(
            {
                "week_rank": week_rank,
                "score": score,
                "has_brief": 1 if artifact_candidate.clean_bullets else 0,
                "is_attachment_like": int(artifact_candidate.is_attachment_like),
                "is_non_html_file": int(artifact_candidate.is_non_html_file),
                "date_distance_sec": date_distance_sec,
                "updated_ts": float(artifact_candidate.updated_ts),
                "filename": artifact_candidate.filename,
                "metadata": artifact_candidate.metadata,
                "brief": artifact_candidate.brief,
                "clean_bullets": list(artifact_candidate.clean_bullets),
                "clean_question": artifact_candidate.clean_question,
            }
        )
    if match_context is None:
        for artifact in artifacts:
            if str(getattr(artifact, "source", "")) != "uclass":
                continue
            metadata = _json_load(getattr(artifact, "metadata_json", None))
            artifact_course_id = _resolve_canonical_course_id(metadata, alias_map=alias_map)
            if class_course_id and artifact_course_id and class_course_id != artifact_course_id:
                continue
            course_name = str(metadata.get("course_name") or "").strip()
            module_name = str(metadata.get("module_name") or "").strip()
            section_name = str(metadata.get("section_name") or "").strip()
            score = 200 if class_course_id and artifact_course_id and class_course_id == artifact_course_id else 0
            if score <= 0:
                score = _best_title_match_score(class_titles, course_name)
            if score <= 0 and module_name:
                score = _best_title_match_score(class_titles, module_name)
            if score <= 0 and section_name:
                score = _best_title_match_score(class_titles, section_name)
            if score <= 0:
                continue
            filename = str(getattr(artifact, "filename", "material") or "material")
            candidate_item = _material_brief_candidate_item_from_artifact(filename, metadata)
            if _is_invalid_material_brief_candidate_item(candidate_item):
                continue
            clean_bullets = _clean_telegram_brief_bullets(candidate_item.get("bullets"), limit=2)
            clean_question = _clean_telegram_brief_question(candidate_item.get("question"))
            updated_dt = _parse_dt(getattr(artifact, "updated_at", None))
            reference_dt = _material_reference_dt(metadata, artifact=artifact) or updated_dt
            brief = metadata.get("brief") if isinstance(metadata.get("brief"), dict) else None
            material_week_markers = _extract_material_week_markers(
                filename,
                course_name,
                module_name,
                section_name,
            )
            week_rank = 1
            if class_week_markers and material_week_markers:
                week_rank = 2 if set(class_week_markers) & set(material_week_markers) else 0
            date_distance_sec = float("inf")
            if class_start_dt is not None and reference_dt is not None:
                date_distance_sec = abs(
                    (reference_dt.astimezone(class_start_dt.tzinfo or timezone.utc) - class_start_dt).total_seconds()
                )
            suffix = Path(filename).suffix.lower()
            content_type = str(metadata.get("content_type") or "").strip().lower()
            source_kind = str(metadata.get("source_kind") or "").strip().lower()
            matches.append(
                {
                    "week_rank": week_rank,
                    "score": score,
                    "has_brief": 1 if clean_bullets else 0,
                    "is_attachment_like": 1 if "attachment" in source_kind else 0,
                    "is_non_html_file": 1 if suffix not in {".php", ".html", ".htm"} and "text/html" not in content_type else 0,
                    "date_distance_sec": date_distance_sec,
                    "updated_ts": updated_dt.timestamp() if updated_dt else 0.0,
                    "filename": filename,
                    "metadata": metadata,
                    "brief": dict(brief) if isinstance(brief, dict) else None,
                    "clean_bullets": clean_bullets,
                    "clean_question": clean_question,
                }
            )
    matches.sort(
        key=lambda row: (
            -int(row["week_rank"]),
            -int(row["is_attachment_like"]),
            -int(row["is_non_html_file"]),
            -int(row["score"]),
            float(row["date_distance_sec"]),
            -int(row["has_brief"]),
            -float(row["updated_ts"]),
            str(row["filename"]).lower(),
        )
    )
    return matches[: max(int(limit), 1)]


def _matched_notifications_for_class(
    db: Database,
    class_item: dict[str, Any],
    notifications: list[Any],
    limit: int = 2,
    *,
    user_id: int | None = None,
    match_context: _DayBriefMatchContext | None = None,
) -> list[str]:
    class_titles = _class_name_candidates(class_item)
    class_course_id = str(class_item.get("canonical_course_id") or "").strip()
    alias_map = (
        match_context.alias_map
        if match_context is not None
        else db.course_alias_resolution_map(user_id=user_id)
    )
    matches: list[tuple[int, str]] = []
    notification_candidates = (
        match_context.notification_candidates_for_course(class_course_id)
        if match_context is not None
        else ()
    )
    for notification_candidate in notification_candidates:
        notification_course_id = notification_candidate.canonical_course_id
        score = (
            200
            if class_course_id and notification_course_id and class_course_id == notification_course_id
            else 0
        )
        if score <= 0:
            score = max(
                _best_title_match_score(class_titles, notification_candidate.title),
                _best_title_match_score(
                    class_titles,
                    f"{notification_candidate.title} {notification_candidate.body}",
                ),
            )
        if score <= 0:
            continue
        matches.append((score, notification_candidate.title or "Notice"))
    if match_context is None:
        for item in notifications:
            source = str(getattr(item, "source", "")).strip().lower()
            if source not in {"uclass", "conflict"}:
                continue
            metadata = _json_load(getattr(item, "metadata_json", None))
            notification_course_id = _resolve_canonical_course_id(metadata, alias_map=alias_map)
            if class_course_id and notification_course_id and class_course_id != notification_course_id:
                continue
            title = str(getattr(item, "title", "")).strip()
            body = str(getattr(item, "body", "") or "").strip()
            score = 200 if class_course_id and notification_course_id and class_course_id == notification_course_id else 0
            if score <= 0:
                score = max(
                    _best_title_match_score(class_titles, title),
                    _best_title_match_score(class_titles, f"{title} {body}"),
                )
            if score <= 0:
                continue
            matches.append((score, title or "Notice"))
    matches.sort(key=lambda row: (-row[0], row[1].lower()))
    output: list[str] = []
    for _, title in matches:
        output.append(title)
        if len(output) >= max(int(limit), 1):
            break
    return output


def _task_course_name(task: Any) -> str:
    metadata = _json_load(getattr(task, "metadata_json", None))
    return _display_course_name(str(metadata.get("course_name") or ""))


def _clean_task_summary_text(value: str, *, task_title: str = "") -> str:
    text = html.unescape(str(value or "").strip())
    if not text:
        return ""
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"^\s*[-\d\.\s]+", " ", text)
    if task_title:
        text = re.sub(re.escape(task_title), " ", text, flags=re.IGNORECASE)
    text = re.sub(r"\([^)]*(?:월|화|수|목|금|토|일)[^)]*\)", " ", text)
    for pattern in MATERIAL_DATE_PATTERNS:
        text = pattern.sub(" ", text)
    for pattern in MATERIAL_TIME_PATTERNS:
        text = pattern.sub(" ", text)
    text = re.sub(r"(?i)\b(?:due|deadline)\b", " ", text)
    text = re.sub(r"(?i)\b(?:assignment|homework|project|report|quiz|hw)\b", " ", text)
    text = re.sub(r"(제출기한|마감일|마감|기한|과제|숙제|프로젝트|보고서|퀴즈)", " ", text)
    text = re.sub(r"\b\d+(?:\.\d+)?\b", " ", text)
    text = re.sub(r"[\[\]\(\):;,\-|]+", " ", text)
    text = re.sub(r"\s+", " ", text).strip(" .")
    return text[:120]


def _task_summary_text(task: Any) -> str:
    metadata = _json_load(getattr(task, "metadata_json", None))
    task_title = str(getattr(task, "title", "") or "").strip()
    for key in ("summary", "description", "details", "evidence"):
        text = _clean_task_summary_text(str(metadata.get(key) or ""), task_title=task_title)
        if text:
            return text
    artifact_title = str(metadata.get("artifact_title") or "").strip()
    artifact_stem = Path(artifact_title).stem.strip() if artifact_title else ""
    if artifact_stem and _normalize_task_title_key(artifact_stem) != _normalize_task_title_key(task_title):
        return artifact_stem[:120]
    return ""


def _task_merge_cache_fingerprint(settings: Settings) -> str:
    provider = str(getattr(settings, "llm_provider", "local") or "local").strip().lower()
    model = str(getattr(settings, "llm_model", "gemma4") or "gemma4").strip()
    enabled = "1" if bool(getattr(settings, "llm_enabled", False)) else "0"
    return f"v{TASK_MERGE_CACHE_VERSION}:{enabled}:{provider}:{model}"


def _task_merge_course_key(task: Any) -> str:
    metadata = _json_load(getattr(task, "metadata_json", None))
    canonical = str(metadata.get("canonical_course_id") or "").strip()
    if canonical:
        return canonical
    return _canonical_course_key(str(metadata.get("course_name") or ""))


def _task_merge_tokens(*values: str) -> set[str]:
    output: set[str] = set()
    for value in values:
        for raw_token in TASK_MERGE_TOKEN_RE.findall(str(value or "")):
            token = raw_token.strip().lower()
            if not token or token.isdigit() or token in TASK_MERGE_STOPWORDS:
                continue
            output.add(token)
    return output


def _task_merge_payload_item(task: Any) -> dict[str, Any]:
    metadata = _json_load(getattr(task, "metadata_json", None))
    return {
        "external_id": str(getattr(task, "external_id", "") or "").strip(),
        "source": str(getattr(task, "source", "") or "").strip(),
        "title": str(getattr(task, "title", "") or "").strip(),
        "due_at": str(getattr(task, "due_at", "") or "").strip() or None,
        "course_name": str(metadata.get("course_name") or "").strip(),
        "canonical_course_id": str(metadata.get("canonical_course_id") or "").strip() or None,
        "summary": _task_summary_text(task),
        "detected_via": str(metadata.get("detected_via") or "").strip() or None,
    }


def _task_merge_payload_hash(tasks: list[Any]) -> str:
    payload = [
        _task_merge_payload_item(task)
        for task in sorted(
            tasks,
            key=lambda item: (
                str(getattr(item, "due_at", "") or ""),
                str(getattr(item, "title", "") or "").lower(),
                str(getattr(item, "external_id", "") or ""),
            ),
        )
    ]
    return sha1(json.dumps(payload, ensure_ascii=True, sort_keys=True).encode("utf-8")).hexdigest()[:24]


def _task_merge_candidate_eligible(left: Any, right: Any) -> bool:
    if (
        str(getattr(left, "external_id", "") or "") == str(getattr(right, "external_id", "") or "")
        and str(getattr(left, "source", "") or "") == str(getattr(right, "source", "") or "")
    ):
        return False
    left_due = _parse_dt(getattr(left, "due_at", None))
    right_due = _parse_dt(getattr(right, "due_at", None))
    if left_due is None or right_due is None:
        return False
    if abs((left_due - right_due).total_seconds()) > TASK_MERGE_CANDIDATE_DUE_WINDOW_HOURS * 3600:
        return False
    left_course = _task_merge_course_key(left)
    right_course = _task_merge_course_key(right)
    if not left_course or not right_course or left_course != right_course:
        return False
    if _tasks_equivalent(left, right):
        return True
    left_title = str(getattr(left, "title", "") or "").strip()
    right_title = str(getattr(right, "title", "") or "").strip()
    if _material_task_titles_similar(left_title, right_title):
        return True
    left_summary = _task_summary_text(left)
    right_summary = _task_summary_text(right)
    overlap = _task_merge_tokens(left_title, left_summary) & _task_merge_tokens(right_title, right_summary)
    if len(overlap) >= 2:
        return True
    left_meta = _json_load(getattr(left, "metadata_json", None))
    right_meta = _json_load(getattr(right, "metadata_json", None))
    if len(overlap) >= 1:
        if abs((left_due - right_due).total_seconds()) <= 15 * 60:
            return True
        if (
            str(left_meta.get("detected_via") or "").strip()
            and str(left_meta.get("detected_via") or "").strip()
            == str(right_meta.get("detected_via") or "").strip()
        ):
            return True
    return False


def _task_merge_candidate_components(tasks: list[Any]) -> list[list[Any]]:
    by_course: dict[str, list[Any]] = {}
    for task in tasks:
        course_key = _task_merge_course_key(task)
        if not course_key:
            continue
        by_course.setdefault(course_key, []).append(task)

    components: list[list[Any]] = []
    for course_tasks in by_course.values():
        if len(course_tasks) < 2:
            continue
        ordered = sorted(
            course_tasks,
            key=lambda item: (
                str(getattr(item, "due_at", "") or ""),
                str(getattr(item, "title", "") or "").lower(),
                str(getattr(item, "external_id", "") or ""),
            ),
        )
        parents = list(range(len(ordered)))

        def _find(index: int) -> int:
            while parents[index] != index:
                parents[index] = parents[parents[index]]
                index = parents[index]
            return index

        def _union(left_index: int, right_index: int) -> None:
            left_root = _find(left_index)
            right_root = _find(right_index)
            if left_root != right_root:
                parents[right_root] = left_root

        for left_index, left_task in enumerate(ordered):
            for right_index in range(left_index + 1, len(ordered)):
                if _task_merge_candidate_eligible(left_task, ordered[right_index]):
                    _union(left_index, right_index)

        buckets: dict[int, list[Any]] = {}
        for index, task in enumerate(ordered):
            buckets.setdefault(_find(index), []).append(task)
        for bucket in buckets.values():
            if len(bucket) < 2:
                continue
            if len(bucket) > TASK_MERGE_LLM_MAX_COMPONENT_SIZE:
                continue
            components.append(bucket)
            if len(components) >= TASK_MERGE_LLM_MAX_COMPONENTS:
                return components
    return components


def _llm_task_merge_groups(
    settings: Settings,
    db: Database,
    *,
    tasks: list[Any],
) -> list[dict[str, Any]]:
    if len(tasks) < 2:
        return []
    if not bool(getattr(settings, "llm_enabled", False)):
        return []
    gate = _identity_warning_gate(
        settings=settings,
        db=db,
        step="task_merge_dedupe",
        destination="llm",
    )
    if gate is not None:
        return []
    payload = {
        "tasks": [
            {
                "id": str(getattr(task, "external_id", "") or "").strip(),
                "title": str(getattr(task, "title", "") or "").strip(),
                "due_at": str(getattr(task, "due_at", "") or "").strip() or None,
                "course_name": _task_course_name(task),
                "summary": _task_summary_text(task),
                "source": str(getattr(task, "source", "") or "").strip(),
                "detected_via": str(
                    _json_load(getattr(task, "metadata_json", None)).get("detected_via") or ""
                ).strip()
                or None,
            }
            for task in tasks
        ]
    }
    try:
        client = _llm_client(settings)
        raw = client.generate_text(
            system_prompt=TASK_MERGE_SYSTEM_PROMPT,
            prompt=json.dumps(payload, ensure_ascii=False, sort_keys=True),
        )
        parsed = _parse_llm_json_payload(raw) or {}
    except Exception as exc:
        logger.warning("task merge llm fallback", extra={"error": str(exc)})
        return []
    groups = parsed.get("groups")
    if not isinstance(groups, list):
        return []
    valid_ids = {str(getattr(task, "external_id", "") or "").strip() for task in tasks}
    confidence_rank = {"high": 2, "medium": 1, "low": 0}
    normalized: list[dict[str, Any]] = []
    for group in groups:
        if not isinstance(group, dict):
            continue
        ids = [
            str(item).strip()
            for item in list(group.get("ids") or [])
            if str(item).strip() in valid_ids
        ]
        deduped_ids: list[str] = []
        seen_ids: set[str] = set()
        for item in ids:
            if item in seen_ids:
                continue
            seen_ids.add(item)
            deduped_ids.append(item)
        if len(deduped_ids) < 2:
            continue
        confidence = str(group.get("confidence") or "low").strip().lower()
        if confidence not in confidence_rank:
            confidence = "low"
        normalized.append(
            {
                "ids": deduped_ids,
                "merged_title": str(group.get("merged_title") or "").strip()[:120],
                "confidence": confidence,
                "reason": str(group.get("reason") or "").strip()[:160],
            }
        )
    normalized.sort(
        key=lambda item: (
            -confidence_rank.get(str(item.get("confidence") or "low"), 0),
            -len(list(item.get("ids") or [])),
            str(item.get("merged_title") or "").lower(),
        )
    )
    output: list[dict[str, Any]] = []
    used_ids: set[str] = set()
    for item in normalized:
        ids = [task_id for task_id in list(item.get("ids") or []) if task_id not in used_ids]
        if len(ids) < 2:
            continue
        used_ids.update(ids)
        output.append({**item, "ids": ids})
    return output


def _precompute_task_merge_cache(
    settings: Settings,
    db: Database,
    *,
    user_id: int | None = None,
    tasks: list[Any] | None = None,
) -> dict[str, Any]:
    owner_id = _safe_int(user_id) or 0
    fingerprint = _task_merge_cache_fingerprint(settings)
    baseline = _dedupe_tasks_for_briefing(
        list(tasks)
        if tasks is not None
        else db.list_open_tasks(limit=600, user_id=owner_id or None)
    )
    payload_hash = _task_merge_payload_hash(baseline)
    state = db.get_sync_state(TASK_MERGE_CACHE_JOB_NAME, user_id=owner_id or None)
    existing = _json_load(state.last_cursor_json)
    if (
        int(existing.get("version") or 0) == TASK_MERGE_CACHE_VERSION
        and str(existing.get("payload_hash") or "") == payload_hash
        and str(existing.get("fingerprint") or "") == fingerprint
    ):
        return existing

    candidate_components = _task_merge_candidate_components(baseline)
    if not candidate_components:
        cursor = {
            "version": TASK_MERGE_CACHE_VERSION,
            "fingerprint": fingerprint,
            "payload_hash": payload_hash,
            "updated_at": now_utc_iso(),
            "task_count": len(baseline),
            "candidate_component_count": 0,
            "groups": [],
        }
        db.update_sync_state(
            TASK_MERGE_CACHE_JOB_NAME,
            last_run_at=now_utc_iso(),
            last_cursor_json=cursor,
            user_id=owner_id or None,
        )
        return cursor

    if bool(getattr(settings, "llm_enabled", False)):
        gate = _identity_warning_gate(
            settings=settings,
            db=db,
            step="task_merge_dedupe",
            destination="llm",
        )
        if gate is not None:
            return existing if isinstance(existing, dict) else {}

    groups: list[dict[str, Any]] = []
    for component in candidate_components[:TASK_MERGE_LLM_MAX_COMPONENTS]:
        groups.extend(_llm_task_merge_groups(settings, db, tasks=component))
    cursor = {
        "version": TASK_MERGE_CACHE_VERSION,
        "fingerprint": fingerprint,
        "payload_hash": payload_hash,
        "updated_at": now_utc_iso(),
        "task_count": len(baseline),
        "candidate_component_count": len(candidate_components),
        "groups": groups,
    }
    db.update_sync_state(
        TASK_MERGE_CACHE_JOB_NAME,
        last_run_at=now_utc_iso(),
        last_cursor_json=cursor,
        user_id=owner_id or None,
    )
    return cursor


def _tasks_equivalent(left: Any, right: Any) -> bool:
    if (
        str(getattr(left, "external_id", "") or "") == str(getattr(right, "external_id", "") or "")
        and str(getattr(left, "source", "") or "") == str(getattr(right, "source", "") or "")
    ):
        return True
    if not _material_task_titles_similar(
        str(getattr(left, "title", "") or ""),
        str(getattr(right, "title", "") or ""),
    ):
        return False
    left_meta = _json_load(getattr(left, "metadata_json", None))
    right_meta = _json_load(getattr(right, "metadata_json", None))
    left_canonical = str(left_meta.get("canonical_course_id") or "").strip()
    right_canonical = str(right_meta.get("canonical_course_id") or "").strip()
    if left_canonical and right_canonical and left_canonical != right_canonical:
        return False
    left_course = _canonical_course_key(str(left_meta.get("course_name") or ""))
    right_course = _canonical_course_key(str(right_meta.get("course_name") or ""))
    if left_course and right_course and left_course != right_course:
        return False
    left_due = _parse_dt(getattr(left, "due_at", None))
    right_due = _parse_dt(getattr(right, "due_at", None))
    if left_due is None or right_due is None:
        return False
    return abs((left_due - right_due).total_seconds()) <= 15 * 60


def _task_preference_score(task: Any) -> tuple[int, int, int]:
    metadata = _json_load(getattr(task, "metadata_json", None))
    course_score = len(_task_course_name(task))
    summary_score = len(_task_summary_text(task))
    metadata_score = sum(
        1
        for key in (
            "canonical_course_id",
            "course_name",
            "artifact_title",
            "summary",
            "description",
            "details",
            "evidence",
        )
        if str(metadata.get(key) or "").strip()
    )
    return (summary_score, course_score, metadata_score)


def _dedupe_tasks_for_briefing(tasks: list[Any]) -> list[Any]:
    deduped: list[Any] = []
    for task in tasks:
        duplicate_index: int | None = None
        for index, existing in enumerate(deduped):
            if _tasks_equivalent(existing, task):
                duplicate_index = index
                break
        if duplicate_index is None:
            deduped.append(task)
            continue
        if _task_preference_score(task) > _task_preference_score(deduped[duplicate_index]):
            deduped[duplicate_index] = task
    return deduped


def _merged_task_summary(tasks: list[Any]) -> str:
    seen: set[str] = set()
    pieces: list[str] = []
    for task in sorted(tasks, key=_task_preference_score, reverse=True):
        summary = _task_summary_text(task)
        summary_key = _normalize_task_title_key(summary)
        if not summary or not summary_key or summary_key in seen:
            continue
        seen.add(summary_key)
        pieces.append(summary)
        if len(pieces) >= 2:
            break
    return " / ".join(pieces)[:220]


def _build_task_merge_proxy(tasks: list[Any], group: dict[str, Any]) -> Task:
    representative = max(tasks, key=_task_preference_score)
    title = str(group.get("merged_title") or "").strip() or str(getattr(representative, "title", "") or "").strip()
    merged_due = None
    due_candidates = [
        _parse_dt(getattr(task, "due_at", None))
        for task in tasks
        if _parse_dt(getattr(task, "due_at", None)) is not None
    ]
    if due_candidates:
        merged_due = min(due_candidates).isoformat()
    metadata = _json_load(getattr(representative, "metadata_json", None))
    merged_summary = _merged_task_summary(tasks)
    if merged_summary:
        metadata["summary"] = merged_summary
    metadata["presentation_merge"] = {
        "group_ids": [
            str(getattr(task, "external_id", "") or "").strip()
            for task in tasks
            if str(getattr(task, "external_id", "") or "").strip()
        ],
        "confidence": str(group.get("confidence") or "").strip() or "medium",
        "reason": str(group.get("reason") or "").strip() or None,
        "merged_title": title,
    }
    return Task(
        external_id=str(getattr(representative, "external_id", "") or "").strip(),
        source=str(getattr(representative, "source", "") or "").strip(),
        due_at=merged_due or str(getattr(representative, "due_at", "") or "").strip() or None,
        title=title,
        status=str(getattr(representative, "status", "open") or "open").strip() or "open",
        metadata_json=json.dumps(metadata, ensure_ascii=True, sort_keys=True),
        user_id=_safe_int(getattr(representative, "user_id", None)),
    )


def _apply_task_merge_cache(
    settings: Settings,
    db: Database,
    *,
    tasks: list[Any],
    user_id: int | None = None,
) -> list[Any]:
    baseline = _dedupe_tasks_for_briefing(tasks)
    owner_id = _safe_int(user_id) or 0
    state = db.get_sync_state(TASK_MERGE_CACHE_JOB_NAME, user_id=owner_id or None)
    cursor = _json_load(state.last_cursor_json)
    if int(cursor.get("version") or 0) != TASK_MERGE_CACHE_VERSION:
        return baseline
    if str(cursor.get("fingerprint") or "") != _task_merge_cache_fingerprint(settings):
        return baseline
    task_by_id = {
        str(getattr(task, "external_id", "") or "").strip(): task
        for task in baseline
        if str(getattr(task, "external_id", "") or "").strip()
    }
    confidence_rank = {"high": 2, "medium": 1, "low": 0}
    groups = [
        item
        for item in list(cursor.get("groups") or [])
        if isinstance(item, dict) and confidence_rank.get(str(item.get("confidence") or "low"), 0) >= 1
    ]
    groups.sort(
        key=lambda item: (
            -confidence_rank.get(str(item.get("confidence") or "low"), 0),
            -len(list(item.get("ids") or [])),
        )
    )
    merged_by_id: dict[str, Any] = {}
    consumed: set[str] = set()
    for group in groups:
        ids = [
            str(item).strip()
            for item in list(group.get("ids") or [])
            if str(item).strip() in task_by_id and str(item).strip() not in consumed
        ]
        if len(ids) < 2:
            continue
        members = [task_by_id[item] for item in ids]
        proxy = _build_task_merge_proxy(members, group)
        representative_id = str(getattr(proxy, "external_id", "") or "").strip()
        if representative_id:
            merged_by_id[representative_id] = proxy
        consumed.update(ids)
    output: list[Any] = []
    for task in baseline:
        task_id = str(getattr(task, "external_id", "") or "").strip()
        if task_id in consumed:
            if task_id in merged_by_id:
                output.append(merged_by_id[task_id])
            continue
        output.append(task)
    return output


def _due_label(due_dt_local: datetime, reference_day_local: datetime) -> str:
    delta_days = (due_dt_local.date() - reference_day_local.date()).days
    if delta_days < 0:
        return f"D+{abs(delta_days)} overdue"
    if delta_days == 0:
        return "D-day"
    return f"D-{delta_days}"


def _format_briefing_task_line(
    task: Any,
    *,
    reference_day_local: datetime,
    include_course: bool,
) -> str:
    metadata = _json_load(getattr(task, "metadata_json", None))
    title = str(getattr(task, "title", "Task")).strip() or "Task"
    provenance_tag = _provenance_brief_tag(
        metadata,
        fallback_source=str(getattr(task, "source", "") or ""),
    )
    prefix = "[파일] " if str(metadata.get("detected_via") or "") == "material_deadline" else ""
    course_name = _task_course_name(task)
    summary = _task_summary_text(task)
    pieces = [f"{prefix}{title}"]
    if provenance_tag:
        pieces.insert(0, provenance_tag)
    if include_course and course_name:
        pieces.append(f"[{course_name}]")
    line = " ".join(pieces)
    due_dt = _parse_dt(getattr(task, "due_at", None))
    if due_dt is not None:
        due_local = due_dt.astimezone(reference_day_local.tzinfo or timezone.utc)
        label = _due_label(due_local, reference_day_local)
        line += f" ({label}, {due_local.strftime('%m-%d %H:%M')})"
    if summary:
        line += f" - {summary}"
    return line


def _matched_tasks_for_class(
    db: Database,
    class_item: dict[str, Any],
    tasks: list[Any],
    *,
    reference_day_local: datetime,
    limit: int = 2,
    user_id: int | None = None,
    match_context: _DayBriefMatchContext | None = None,
) -> tuple[list[str], list[str]]:
    class_titles = _class_name_candidates(class_item)
    class_course_id = str(class_item.get("canonical_course_id") or "").strip()
    alias_map = (
        match_context.alias_map
        if match_context is not None
        else db.course_alias_resolution_map(user_id=user_id)
    )
    regular_matches: list[tuple[int, Any]] = []
    file_matches: list[tuple[int, Any]] = []
    task_candidates = (
        match_context.task_candidates_for_course(class_course_id)
        if match_context is not None
        else ()
    )
    for task_candidate in task_candidates:
        due_local = task_candidate.due_dt.astimezone(reference_day_local.tzinfo or timezone.utc)
        if due_local.date() < reference_day_local.date():
            continue
        task_course_id = task_candidate.canonical_course_id
        score = 200 if class_course_id and task_course_id and class_course_id == task_course_id else 0
        if score <= 0:
            score = _best_title_match_score(class_titles, task_candidate.course_name) if task_candidate.course_name else 0
        if score <= 0:
            score = _best_title_match_score(class_titles, task_candidate.title)
        if score <= 0:
            continue
        target = file_matches if task_candidate.detected_via == "material_deadline" else regular_matches
        target.append((score, task_candidate.task))
    if match_context is None:
        for task in _dedupe_tasks_for_briefing(tasks):
            source = str(getattr(task, "source", "")).strip().lower()
            if source not in {"uclass", "inbox"}:
                continue
            due_dt = _parse_dt(getattr(task, "due_at", None))
            if due_dt is None:
                continue
            due_local = due_dt.astimezone(reference_day_local.tzinfo or timezone.utc)
            if due_local.date() < reference_day_local.date():
                continue
            metadata = _json_load(getattr(task, "metadata_json", None))
            task_course_id = _resolve_canonical_course_id(metadata, alias_map=alias_map)
            if class_course_id and task_course_id and class_course_id != task_course_id:
                continue
            title = str(getattr(task, "title", "Task")).strip() or "Task"
            course_name = str(metadata.get("course_name") or "").strip()
            score = 200 if class_course_id and task_course_id and class_course_id == task_course_id else 0
            if score <= 0:
                score = _best_title_match_score(class_titles, course_name) if course_name else 0
            if score <= 0:
                score = _best_title_match_score(class_titles, title)
            if score <= 0:
                continue
            target = file_matches if str(metadata.get("detected_via") or "") == "material_deadline" else regular_matches
            target.append((score, task))
    regular_matches.sort(
        key=lambda row: (
            -row[0],
            str(getattr(row[1], "title", "") or "").lower(),
        )
    )
    file_matches.sort(
        key=lambda row: (
            -row[0],
            str(getattr(row[1], "title", "") or "").lower(),
        )
    )

    def _take(items: list[tuple[int, Any]]) -> list[str]:
        output: list[str] = []
        for _, task in items:
            output.append(
                _format_briefing_task_line(
                    task,
                    reference_day_local=reference_day_local,
                    include_course=False,
                )
            )
            if len(output) >= max(int(limit), 1):
                break
        return output

    return _take(regular_matches), _take(file_matches)


def _collect_due_reviews_for_day(
    settings: Settings,
    db: Database,
    *,
    target_day_local: datetime,
    limit: int = 8,
) -> list[dict[str, Any]]:
    reviews: list[dict[str, Any]] = []
    for event in db.list_events(limit=3000):
        if str(event.source) != "review" and not str(event.external_id).startswith("review:"):
            continue
        occurrences = _event_occurrences_on_date(
            event=event,
            target_date_local=target_day_local,
            timezone_name=str(getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul"),
        )
        for start_local, end_local in occurrences:
            reviews.append(
                {
                    "title": event.title,
                    "start_local": start_local,
                    "end_local": end_local,
                }
            )
    reviews.sort(key=lambda row: (row["start_local"], str(row["title"]).lower()))
    return reviews[: max(int(limit), 1)]


def _build_briefing_llm_guidance(
    settings: Settings,
    db: Database,
    payload: dict[str, Any],
    *,
    cache: dict[str, list[str]] | None = None,
    enabled: bool = True,
) -> list[str]:
    if not enabled or not bool(getattr(settings, "llm_enabled", False)):
        return []
    cache_key = sha1(
        json.dumps(payload, ensure_ascii=False, sort_keys=True).encode("utf-8")
    ).hexdigest()
    if cache is not None and cache_key in cache:
        return list(cache[cache_key])
    gate = _identity_warning_gate(
        settings=settings,
        db=db,
        step="scheduled_briefings_llm",
        destination="llm",
    )
    if gate is not None:
        return []
    try:
        client = _llm_client_with_timeout(
            settings,
            timeout_sec=_llm_timeout_sec(
                settings,
                cap=BRIEFING_LLM_TIMEOUT_SEC_CAP,
            ),
        )
        system_prompt = (
            "You are a study assistant. Summarize the briefing payload into at most 3 short bullet lines. "
            "Do not include markdown headings."
        )
        raw = client.generate_text(
            system_prompt=system_prompt,
            prompt=json.dumps(payload, ensure_ascii=False),
        )
    except Exception as exc:
        logger.warning("briefing llm guidance failed", extra={"error": str(exc)})
        if cache is not None:
            cache[cache_key] = []
        return []
    lines: list[str] = []
    for line in str(raw or "").splitlines():
        text = line.strip(" -*\t")
        if not text:
            continue
        lines.append(text)
        if len(lines) >= 3:
            break
    if cache is not None:
        cache[cache_key] = list(lines)
    return lines


def _build_scheduled_briefing(
    settings: Settings,
    db: Database,
    *,
    slot: str,
    now_local: datetime,
    user_id: int | None = None,
    weather_snapshot_cache: dict[str, dict[str, Any] | None] | None = None,
    llm_guidance_cache: dict[str, list[str]] | None = None,
    enable_llm_guidance: bool = True,
) -> tuple[str, dict[str, Any]]:
    target_day = now_local if slot == "morning" else (now_local + timedelta(days=1))
    target_day = target_day.replace(hour=12, minute=0, second=0, microsecond=0)
    task_lookahead_days = max(int(getattr(settings, "briefing_task_lookahead_days", 7) or 7), 1)
    day_brief = DayBriefService(settings, db, user_id=user_id).build_day_brief(
        target_day_local=target_day,
        reference_day_local=now_local,
        max_classes=max(int(getattr(settings, "briefing_max_classes", 6) or 6), 1),
        artifact_limit=SCHEDULED_BRIEFING_ARTIFACT_LIMIT,
        notification_limit=SCHEDULED_BRIEFING_NOTIFICATION_LIMIT,
        open_task_limit=SCHEDULED_BRIEFING_OPEN_TASK_LIMIT,
        lookahead_days=task_lookahead_days,
        lookahead_limit=30,
        lookahead_now_iso=now_local.astimezone(timezone.utc).isoformat(),
        refresh_task_merge_cache=True,
    )
    meetings_result = day_brief.meetings_result
    meeting_items = list(day_brief.meeting_items)
    course_briefs = list(day_brief.course_briefs)
    due_tasks = list(day_brief.tasks_due_within_window)

    header_label = "아침" if slot == "morning" else "저녁"
    target_label = "오늘" if slot == "morning" else "내일"
    lines: list[str] = [f"[KU] {header_label} 브리핑 ({now_local.date().isoformat()})"]
    if slot != "morning":
        lines.append(f"{target_label} 기준일: {target_day.date().isoformat()}")
    lines.append("")
    weather_lines = _build_briefing_weather_lines(
        settings,
        db,
        now_local=now_local,
        user_id=user_id,
        snapshot_cache=weather_snapshot_cache,
    )
    if weather_lines:
        lines.extend(weather_lines)
        lines.append("")
    if bool(meetings_result.get("ok")):
        lines.append(f"약속/일정: {len(meeting_items)}")
        if not meeting_items:
            lines.append("- none")
        for item in meeting_items[:6]:
            title = str(item.get("title") or "Meeting")
            if bool(item.get("all_day")):
                when = "하루 종일"
            else:
                when = _format_time_range_local(item["start_local"], item["end_local"])
            location = str(item.get("location") or "").strip()
            if location:
                lines.append(f"- {when} {title} @ {location}")
            else:
                lines.append(f"- {when} {title}")
    else:
        lines.append("약속/일정: 확인 불가")
    skipped_reason = str(meetings_result.get("skipped_reason") or "").strip()
    if skipped_reason:
        lines.append(f"- 참고: {skipped_reason}")

    lines.append("")
    lines.append(f"수업: {len(course_briefs)}")
    if not course_briefs:
        lines.append("- none")
    for course_brief in course_briefs:
        item = course_brief.class_item
        start_time = item["start_local"].strftime("%H:%M")
        title = str(item["title"])
        location_text = str(item["location_text"])
        if location_text != "TBD":
            lines.append(f"- {start_time}에 {location_text}에서 {title} 수업")
        else:
            lines.append(f"- {start_time} {title} 수업")
        if course_brief.preparation:
            lines.append(f"  준비: {course_brief.preparation}")

        if course_brief.notice_titles:
            lines.append(f"  공지: {'; '.join(course_brief.notice_titles)}")
        if course_brief.task_lines:
            lines.append(f"  수업 과제: {'; '.join(course_brief.task_lines)}")
        if course_brief.file_task_lines:
            lines.append(f"  파일 감지 과제: {'; '.join(course_brief.file_task_lines)}")

    lines.append("")
    lines.append(f"다가오는 과제 ({task_lookahead_days}일): {len(due_tasks)}")
    if not due_tasks:
        lines.append("- none")
    for task in due_tasks[:8]:
        lines.append(
            "- "
            + _format_briefing_task_line(
                task,
                reference_day_local=now_local,
                include_course=True,
            )
        )

    llm_payload = {
        "slot": slot,
        "target_day": target_day.date().isoformat(),
        "meetings": [
            {
                "title": str(item.get("title") or ""),
                "all_day": bool(item.get("all_day")),
                "start": item.get("start_local").isoformat() if item.get("start_local") else None,
                "end": item.get("end_local").isoformat() if item.get("end_local") else None,
            }
            for item in meeting_items[:6]
        ],
        "classes": [
            {
                "title": str(course_brief.class_item.get("title") or ""),
                "start": (
                    course_brief.class_item.get("start_local").isoformat()
                    if course_brief.class_item.get("start_local")
                    else None
                ),
                "location": str(course_brief.class_item.get("location_text") or ""),
                "summary_short": str(course_brief.preparation or ""),
                "summary_long": list(course_brief.material_summary_bits[:3]),
            }
            for course_brief in course_briefs
        ],
        "tasks": [
            {
                "title": task.title,
                "due_at": task.due_at,
                "course_name": _task_course_name(task),
                "summary": _task_summary_text(task),
            }
            for task in due_tasks[:8]
        ],
    }
    guidance_lines = _build_briefing_llm_guidance(
        settings=settings,
        db=db,
        payload=llm_payload,
        cache=llm_guidance_cache,
        enabled=enable_llm_guidance,
    )
    if guidance_lines:
        lines.append("")
        lines.append("가이드 [AI]:")
        for item in guidance_lines:
            lines.append(f"- {item}")

    message = "\n".join(lines).strip()
    if len(message) > 3400:
        message = message[:3390].rstrip() + "\n..."
    return message, llm_payload


def _briefing_delivery_mode(settings: Settings) -> str:
    raw = str(getattr(settings, "briefing_delivery_mode", "direct") or "direct").strip().lower()
    if raw in {"precompute_only", "direct"}:
        return raw
    return "direct"


def _scheduled_briefing_occurrences(
    *,
    base_local: datetime,
    morning_hour: int,
    morning_minute: int,
    evening_hour: int,
    evening_minute: int,
    days: int = 2,
) -> list[tuple[str, datetime]]:
    tz = base_local.tzinfo
    occurrences: list[tuple[str, datetime]] = []
    for day_offset in range(max(int(days), 1)):
        target_day = (base_local + timedelta(days=day_offset)).date()
        morning_dt = datetime(
            target_day.year,
            target_day.month,
            target_day.day,
            morning_hour,
            morning_minute,
            tzinfo=tz,
        )
        evening_dt = datetime(
            target_day.year,
            target_day.month,
            target_day.day,
            evening_hour,
            evening_minute,
            tzinfo=tz,
        )
        occurrences.append((f"{target_day.isoformat()}-morning", morning_dt))
        occurrences.append((f"{target_day.isoformat()}-evening", evening_dt))
    return occurrences


def build_precomputed_telegram_briefings(
    settings: Settings,
    db: Database,
    *,
    now_local: datetime | None = None,
    enable_llm_guidance: bool = True,
) -> dict[str, Any]:
    if not bool(getattr(settings, "briefing_enabled", False)):
        return {"ok": False, "skipped": True, "reason": "BRIEFING_ENABLED is false"}
    channel = str(getattr(settings, "briefing_channel", "telegram") or "telegram").strip().lower()
    if channel != "telegram":
        return {"ok": False, "skipped": True, "reason": "BRIEFING_CHANNEL is not telegram"}
    if not bool(getattr(settings, "telegram_enabled", False)):
        return {"ok": False, "skipped": True, "reason": "TELEGRAM_ENABLED is false"}

    try:
        morning_hour, morning_minute = _parse_clock_time(
            str(getattr(settings, "briefing_morning_time_local", "09:00") or "09:00"),
            "BRIEFING_MORNING_TIME_LOCAL",
        )
        evening_hour, evening_minute = _parse_clock_time(
            str(getattr(settings, "briefing_evening_time_local", "21:00") or "21:00"),
            "BRIEFING_EVENING_TIME_LOCAL",
        )
    except ValueError as exc:
        return {"ok": False, "error": str(exc)}

    timezone_name = str(getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul")
    tz = ZoneInfo(timezone_name)
    anchor_local = (
        now_local.astimezone(tz)
        if now_local is not None
        else datetime.now(tz)
    )
    slot_schedule = _scheduled_briefing_occurrences(
        base_local=anchor_local,
        morning_hour=morning_hour,
        morning_minute=morning_minute,
        evening_hour=evening_hour,
        evening_minute=evening_minute,
        days=2,
    )

    items: dict[str, Any] = {}
    relay_endpoint = str(
        getattr(settings, "briefing_relay_endpoint", "") or ""
    ).strip()
    relay_shared_secret = str(
        getattr(settings, "briefing_relay_shared_secret", "") or ""
    ).strip()
    relay_ready = bool(relay_endpoint and relay_shared_secret)
    weather_snapshot_cache: dict[str, dict[str, Any] | None] = {}
    llm_guidance_cache: dict[str, list[str]] | None = {} if enable_llm_guidance else None
    all_chat_ids: set[str] = set()
    for item_key, send_at_local in slot_schedule:
        slot = "morning" if item_key.endswith("-morning") else "evening"
        slot_policy_kind = (
            NOTIFICATION_POLICY_KIND_BRIEFING_MORNING
            if slot == "morning"
            else NOTIFICATION_POLICY_KIND_BRIEFING_EVENING
        )
        slot_chat_ids = _chat_ids_for_notification_dispatch(
            settings,
            db,
            preference=USER_PREFERENCE_SCHEDULED_BRIEFINGS_ENABLED,
            policy_kinds=(slot_policy_kind,),
            reference_local=send_at_local,
        )
        all_chat_ids.update(slot_chat_ids)
        deliveries: list[dict[str, Any]] = []
        messages_by_chat: dict[str, str] = {}
        payloads_by_chat: dict[str, Any] = {}
        for chat_id in slot_chat_ids:
            user_scope = _resolve_user_scope(
                settings,
                db,
                chat_id=chat_id,
                create_if_missing=True,
                metadata_source="precomputed_briefings",
            )
            owner_id = int(user_scope["user_id"] or 0)
            message, payload = _build_scheduled_briefing(
                settings=settings,
                db=db,
                slot=slot,
                now_local=send_at_local,
                user_id=owner_id,
                weather_snapshot_cache=weather_snapshot_cache,
                llm_guidance_cache=llm_guidance_cache,
                enable_llm_guidance=enable_llm_guidance,
            )
            deliveries.append(
                {
                    "chat_id": chat_id,
                    "user_id": owner_id,
                    "message": message,
                    "message_length": len(message),
                    "payload": payload,
                }
            )
            messages_by_chat[chat_id] = message
            payloads_by_chat[chat_id] = payload
        default_delivery = deliveries[0] if deliveries else {}
        item_payload = {
            "item_key": item_key,
            "slot": slot,
            "send_at_local": send_at_local.isoformat(),
            "message": str(default_delivery.get("message") or ""),
            "message_length": int(default_delivery.get("message_length") or 0),
            "payload": default_delivery.get("payload") or {},
            "chat_ids": list(slot_chat_ids),
            "deliveries": deliveries,
            "messages_by_chat": messages_by_chat,
            "payloads_by_chat": payloads_by_chat,
        }
        if relay_ready:
            item_payload["relay_request"] = build_signed_briefing_delivery_request(
                endpoint=relay_endpoint,
                shared_secret=relay_shared_secret,
                payload=item_payload,
            )
        items[item_key] = item_payload

    if not all_chat_ids:
        return {
            "ok": False,
            "skipped": True,
            "reason": "No eligible chat_ids",
        }

    return {
        "ok": True,
        "generated_at": now_utc_iso(),
        "generated_from_local": anchor_local.isoformat(),
        "channel": "telegram",
        "chat_ids": sorted(all_chat_ids),
        "relay": {
            "configured": relay_ready,
            "endpoint": relay_endpoint if relay_ready else None,
            "reason": (
                None
                if relay_ready
                else "BRIEFING_RELAY_ENDPOINT or BRIEFING_RELAY_SHARED_SECRET is missing"
            ),
        },
        "items": items,
    }


def send_scheduled_briefings(settings: Settings, db: Database) -> dict[str, Any]:
    if not bool(getattr(settings, "briefing_enabled", False)):
        return {"skipped": True, "reason": "BRIEFING_ENABLED is false"}
    channel = str(getattr(settings, "briefing_channel", "telegram") or "telegram").strip().lower()
    if channel != "telegram":
        return {"skipped": True, "reason": "BRIEFING_CHANNEL is not telegram"}
    delivery_mode = _briefing_delivery_mode(settings)
    if delivery_mode != "direct":
        return {
            "skipped": True,
            "reason": f"BRIEFING_DELIVERY_MODE={delivery_mode}",
        }
    if not bool(getattr(settings, "telegram_enabled", False)) or not str(
        getattr(settings, "telegram_bot_token", "") or ""
    ).strip():
        return {"skipped": True, "reason": "Telegram is not configured"}
    try:
        morning_hour, morning_minute = _parse_clock_time(
            str(getattr(settings, "briefing_morning_time_local", "09:00") or "09:00"),
            "BRIEFING_MORNING_TIME_LOCAL",
        )
        evening_hour, evening_minute = _parse_clock_time(
            str(getattr(settings, "briefing_evening_time_local", "21:00") or "21:00"),
            "BRIEFING_EVENING_TIME_LOCAL",
        )
    except ValueError as exc:
        return {"error": str(exc)}

    timezone_name = str(getattr(settings, "timezone", "Asia/Seoul") or "Asia/Seoul")
    now_local = datetime.now(ZoneInfo(timezone_name))
    today = now_local.date().isoformat()
    morning_trigger = now_local.replace(hour=morning_hour, minute=morning_minute, second=0, microsecond=0)
    evening_trigger = now_local.replace(hour=evening_hour, minute=evening_minute, second=0, microsecond=0)
    target_chat_ids_by_slot = {
        "morning": set(
            _chat_ids_for_notification_dispatch(
                settings,
                db,
                preference=USER_PREFERENCE_SCHEDULED_BRIEFINGS_ENABLED,
                policy_kinds=(NOTIFICATION_POLICY_KIND_BRIEFING_MORNING,),
                reference_local=now_local,
            )
        ),
        "evening": set(
            _chat_ids_for_notification_dispatch(
                settings,
                db,
                preference=USER_PREFERENCE_SCHEDULED_BRIEFINGS_ENABLED,
                policy_kinds=(NOTIFICATION_POLICY_KIND_BRIEFING_EVENING,),
                reference_local=now_local,
            )
        ),
    }
    chat_ids = sorted(target_chat_ids_by_slot["morning"] | target_chat_ids_by_slot["evening"])
    if not chat_ids:
        return {"skipped": True, "reason": "No eligible chat_ids"}

    gate = _identity_warning_gate(
        settings=settings,
        db=db,
        step="scheduled_briefings",
        destination="telegram",
    )
    if gate is not None:
        return gate

    client = TelegramBotClient(str(getattr(settings, "telegram_bot_token")))
    results: dict[str, Any] = {}
    sent_any = False
    weather_snapshot_cache: dict[str, dict[str, Any] | None] = {}
    llm_guidance_cache: dict[str, list[str]] = {}
    for chat_id in chat_ids:
        user_scope = _resolve_user_scope(
            settings,
            db,
            chat_id=chat_id,
            create_if_missing=True,
            metadata_source="scheduled_briefings",
        )
        owner_id = int(user_scope["user_id"] or 0)
        state = db.get_sync_state("scheduled_briefings", user_id=owner_id)
        cursor = _json_load(state.last_cursor_json)
        pending_slots: list[str] = []
        if (
            chat_id in target_chat_ids_by_slot["morning"]
            and now_local >= morning_trigger
            and cursor.get("morning_sent_date") != today
        ):
            pending_slots.append("morning")
        if (
            chat_id in target_chat_ids_by_slot["evening"]
            and now_local >= evening_trigger
            and cursor.get("evening_sent_date") != today
        ):
            pending_slots.append("evening")
        if not pending_slots:
            continue
        slot_results: dict[str, Any] = {}
        sent_slots: list[str] = []
        for slot in pending_slots:
            message, payload = _build_scheduled_briefing(
                settings=settings,
                db=db,
                slot=slot,
                now_local=now_local,
                user_id=owner_id,
                weather_snapshot_cache=weather_snapshot_cache,
                llm_guidance_cache=llm_guidance_cache,
            )
            sent = False
            try:
                sent = bool(client.send_message(chat_id=chat_id, text=message))
            except Exception as exc:
                logger.warning(
                    "failed to send scheduled briefing message",
                    extra={"slot": slot, "chat_id": chat_id, "error": str(exc)},
                )
            slot_results[slot] = {
                "sent_to": [chat_id] if sent else [],
                "message_length": len(message),
                "payload": payload,
            }
            if sent:
                sent_slots.append(slot)
                cursor[f"{slot}_sent_date"] = today
                sent_any = True
        db.update_sync_state(
            "scheduled_briefings",
            last_run_at=now_utc_iso(),
            last_cursor_json={
                **cursor,
                "pending_slots": pending_slots,
                "sent_slots": sent_slots,
                "last_results": {
                    key: {"sent_to": value["sent_to"], "message_length": value["message_length"]}
                    for key, value in slot_results.items()
                },
            },
            user_id=owner_id,
        )
        results[str(chat_id)] = {
            "user_id": owner_id,
            "pending_slots": pending_slots,
            "sent_slots": sent_slots,
            "results": slot_results,
        }
    if not results:
        return {"skipped": True, "reason": "scheduled time not reached or already sent", "pending_slots": []}
    if not sent_any:
        return {"error": "briefing send failed", "results": results}
    return {"sent_slots": sorted({slot for row in results.values() for slot in row.get("sent_slots", [])}), "results": results}


def _parse_digest_time(value: str) -> tuple[int, int]:
    return _parse_clock_time(value, "DIGEST_TIME_LOCAL")


def send_daily_digest(settings: Settings, db: Database) -> dict[str, Any]:
    if not settings.digest_enabled:
        return {"skipped": True, "reason": "DIGEST_ENABLED is false"}
    if settings.digest_channel != "telegram":
        return {"skipped": True, "reason": "DIGEST_CHANNEL is not telegram"}
    if not settings.telegram_enabled or not settings.telegram_bot_token:
        return {"skipped": True, "reason": "Telegram is not configured"}
    try:
        hour, minute = _parse_digest_time(settings.digest_time_local)
    except ValueError as exc:
        return {"error": str(exc)}

    tz = ZoneInfo(settings.timezone)
    now_local = datetime.now(tz)
    scheduled_local = now_local.replace(hour=hour, minute=minute, second=0, microsecond=0)
    if now_local < scheduled_local:
        return {"skipped": True, "reason": "scheduled time not reached"}
    chat_ids = _chat_ids_for_notification_dispatch(
        settings,
        db,
        preference=USER_PREFERENCE_DAILY_DIGEST_ENABLED,
        policy_kinds=(NOTIFICATION_POLICY_KIND_DAILY_DIGEST,),
        reference_local=now_local,
    )
    if not chat_ids:
        return {"skipped": True, "reason": "No eligible chat_ids"}

    gate = _identity_warning_gate(
        settings=settings,
        db=db,
        step="daily_digest",
        destination="telegram",
    )
    if gate is not None:
        return gate

    client = TelegramBotClient(settings.telegram_bot_token)
    results: dict[str, Any] = {}
    sent_to: list[str] = []
    pending_chat_ids: list[str] = []
    skipped_chat_ids: list[str] = []
    today = now_local.date().isoformat()
    for chat_id in chat_ids:
        user_scope = _resolve_user_scope(
            settings,
            db,
            chat_id=str(chat_id),
            create_if_missing=True,
            metadata_source="daily_digest",
        )
        owner_id = int(user_scope["user_id"] or 0)
        state = db.get_sync_state("daily_digest", user_id=owner_id)
        cursor = _json_load(state.last_cursor_json)
        if cursor.get("sent_date") == today:
            skipped_chat_ids.append(str(chat_id))
            continue
        pending_chat_ids.append(str(chat_id))
        since_iso = cursor.get("last_sent_at")
        since_dt = _parse_dt(str(since_iso)) if since_iso else None
        if since_dt is None:
            since_dt = datetime.now(timezone.utc) - timedelta(days=1)
        since_value = since_dt.replace(microsecond=0).isoformat()

        notifications = db.list_notifications_since(since_value, limit=100, user_id=owner_id)
        tasks = _dedupe_tasks_for_briefing(
            db.list_tasks_due_within(
                days=settings.digest_task_lookahead_days,
                now_iso=now_utc_iso(),
                limit=20,
                user_id=owner_id,
            )
        )
        materials: list[Any] = []
        for artifact in db.list_artifacts(limit=300, user_id=owner_id):
            meta = _json_load(artifact.metadata_json)
            downloaded_at = _parse_dt(str(meta.get("downloaded_at") or ""))
            if downloaded_at is None or downloaded_at <= since_dt:
                continue
            materials.append(artifact)
        materials = materials[:20]
        lines: list[str] = [f"[KU] Daily digest ({today})", ""]
        lines.append(f"New notifications: {len(notifications)}")
        for item in notifications[:3]:
            lines.append(f"- {item.title}")
        lines.append("")
        lines.append(f"Upcoming tasks ({settings.digest_task_lookahead_days}d): {len(tasks)}")
        for task in tasks[:5]:
            due = task.due_at or "no due date"
            lines.append(f"- {task.title} ({due})")
        lines.append("")
        lines.append(f"New materials: {len(materials)}")
        for artifact in materials[:5]:
            lines.append(f"- {artifact.filename}")
        message = "\n".join(lines).strip()

        sent = False
        try:
            sent = bool(client.send_message(chat_id=chat_id, text=message))
        except Exception as exc:
            logger.warning(
                "failed to send digest message",
                extra={"chat_id": str(chat_id), "error": str(exc)},
            )
        if not sent:
            continue
        sent_to.append(str(chat_id))
        db.update_sync_state(
            "daily_digest",
            last_run_at=now_utc_iso(),
            last_cursor_json={
                "sent_date": today,
                "last_sent_at": now_utc_iso(),
                "sent_to": [str(chat_id)],
                "notifications": len(notifications),
                "tasks": len(tasks),
                "materials": len(materials),
            },
            user_id=owner_id,
        )
        results[str(chat_id)] = {
            "user_id": owner_id,
            "notifications": len(notifications),
            "tasks": len(tasks),
            "materials": len(materials),
        }
    if not sent_to:
        if not pending_chat_ids and skipped_chat_ids:
            return {
                "skipped": True,
                "reason": "already sent today",
                "skipped_chat_ids": skipped_chat_ids,
            }
        return {
            "error": "digest send failed",
            "pending_chat_ids": pending_chat_ids,
            "skipped_chat_ids": skipped_chat_ids,
            "results": results,
        }
    return {
        "sent_to": sent_to,
        "skipped_chat_ids": skipped_chat_ids,
        "results": results,
    }


def sync_llm_summaries(settings: Settings, db: Database) -> dict[str, Any]:
    if not settings.llm_enabled:
        return {"skipped": True, "reason": "LLM_ENABLED is false"}

    notifications = db.list_notifications(limit=10)
    artifacts = db.list_artifacts(limit=10)
    payload = {
        "notifications": [
            {
                "title": item.title,
                "created_at": item.created_at,
                "body": item.body,
                "url": item.url,
            }
            for item in notifications
        ],
        "materials": [
            {
                "filename": item.filename,
                "icloud_path": item.icloud_path,
                "content_hash": item.content_hash,
            }
            for item in artifacts
        ],
    }
    if not payload["notifications"] and not payload["materials"]:
        return {"skipped": True, "reason": "No data to summarize"}

    payload_hash = sha1(
        json.dumps(payload, ensure_ascii=True, sort_keys=True).encode("utf-8")
    ).hexdigest()[:24]
    summary_external_id = f"llm:summary:{payload_hash}"
    if db.has_summary(summary_external_id, source="llm"):
        return {"skipped": True, "reason": "No new data", "external_id": summary_external_id}
    gate = _identity_warning_gate(
        settings=settings,
        db=db,
        step="sync_llm_summaries",
        destination="llm",
    )
    if gate is not None:
        return gate

    try:
        client = _llm_client(settings)
        summary = client.summarize(payload)
    except Exception as exc:
        logger.warning("llm summary failed", extra={"error": str(exc)})
        return {"error": str(exc)}

    body = "\n".join([f"- {line}" for line in summary.bullets])
    db.record_summary(
        external_id=summary_external_id,
        source="llm",
        created_at=now_utc_iso(),
        title="Updates Summary",
        body=body,
        action_item=summary.action_item,
        metadata_json=attach_provenance(
            {"payload_hash": payload_hash},
            source="llm_inferred",
            confidence="medium",
            last_verified_at=now_utc_iso(),
            raw_source_ids=[payload_hash],
            derivation="llm_summary",
        ),
    )
    _record_sync_dashboard_state(
        db,
        "sync_llm_summaries",
        status="success",
        new_items=1,
        cursor_payload={"external_id": summary_external_id},
    )
    return {"created": 1, "external_id": summary_external_id}


def publish_dashboard(settings: Settings, db: Database) -> dict[str, Any]:
    storage_root = resolve_storage_root(settings)
    if storage_root is None:
        raise ValueError("STORAGE_ROOT_DIR is required for publish")
    precomputed_briefings = build_precomputed_telegram_briefings(
        settings=settings,
        db=db,
        enable_llm_guidance=False,
    )
    result = render_dashboard_snapshot(
        db=db,
        storage_root_dir=storage_root,
        extra_data={
            "precomputed_telegram_briefings": precomputed_briefings,
        },
    )
    result["precomputed_telegram_briefings"] = precomputed_briefings
    _record_sync_dashboard_state(
        db,
        "publish_dashboard",
        status="success",
        new_items=1,
        cursor_payload={
            "output": result["dashboard_dir"],
            "precomputed_briefings_ok": bool(precomputed_briefings.get("ok")),
            "precomputed_briefing_slots": {
                slot: str(item.get("send_at_local") or "")
                for slot, item in (
                    precomputed_briefings.get("items", {}).items()
                    if isinstance(precomputed_briefings.get("items"), dict)
                    else []
                )
                if isinstance(item, dict)
            },
        },
    )
    return result


def import_portal_events(
    settings: Settings,
    db: Database,
    ics_url: str | None = None,
    ics_file: Path | None = None,
    csv_file: Path | None = None,
) -> dict[str, Any]:
    inputs = [bool(ics_url), bool(ics_file), bool(csv_file)]
    if sum(inputs) != 1:
        raise ValueError("provide exactly one of ics_url, ics_file, csv_file")

    portal_provenance_source = "portal_csv"
    import_origin = "csv"
    if ics_url:
        events = parse_ics_url(ics_url, timezone_name=settings.timezone)
        portal_provenance_source = "portal_ics_url"
        import_origin = "ics_url"
    elif ics_file:
        events = parse_ics_file(ics_file, timezone_name=settings.timezone)
        portal_provenance_source = "portal_ics"
        import_origin = "ics_file"
    else:
        if not csv_file:
            raise ValueError("csv_file is required")
        events = parse_csv_file(csv_file, timezone_name=settings.timezone)

    upserted = 0
    for event in events:
        metadata = dict(event.metadata or {})
        metadata["import_origin"] = import_origin
        metadata = attach_provenance(
            metadata,
            source=portal_provenance_source,
            confidence="high",
            last_verified_at=event.start_at,
            raw_source_ids=[event.external_id],
            derivation="portal_import",
        )
        db.upsert_event(
            external_id=event.external_id,
            source="portal",
            start=event.start_at,
            end=event.end_at,
            title=event.title,
            location=event.location,
            rrule=event.rrule,
            metadata_json=metadata,
        )
        upserted += 1

    _record_sync_dashboard_state(
        db,
        "import_portal",
        status="success",
        new_items=upserted,
        cursor_payload={
            "upserted_events": upserted,
            "import_origin": import_origin,
            "provenance_source": portal_provenance_source,
        },
    )
    return {"upserted_events": upserted}


def _ops_scope_summary(
    settings: Settings,
    db: Database,
    *,
    chat_id: str | None = None,
    user_id: int | None = None,
) -> dict[str, Any]:
    requested_user_id = _safe_int(user_id)
    requested_chat_id = str(chat_id or "").strip() or None
    resolved = _resolve_user_scope(
        settings,
        db,
        chat_id=requested_chat_id,
        user_id=requested_user_id,
        create_if_missing=False,
        metadata_source="ops",
    )
    resolved_user_id = _safe_int(resolved.get("user_id")) or 0
    return {
        "requested_user_id": requested_user_id,
        "requested_chat_id": requested_chat_id,
        "user_id": resolved_user_id,
        "chat_id": str(resolved.get("chat_id") or requested_chat_id or "").strip() or None,
        "status": str(resolved.get("status") or "unknown").strip() or "unknown",
        "found": resolved_user_id > 0,
        "timezone": str(
            resolved.get("timezone")
            or getattr(settings, "timezone", "Asia/Seoul")
            or "Asia/Seoul"
        ).strip()
        or "Asia/Seoul",
    }


def _ops_surface_state(
    db: Database,
    job_name: str,
    *,
    user_id: int | None = None,
    allow_global_fallback: bool = True,
) -> tuple[dict[str, Any], Any, dict[str, Any]]:
    return _ops_surface_state_impl(
        db,
        job_name,
        user_id=user_id,
        allow_global_fallback=allow_global_fallback,
    )


def _build_ku_official_api_health(
    settings: Settings,
    db: Database,
    *,
    user_id: int | None = None,
) -> dict[str, Any]:
    return _build_ku_official_api_health_impl(
        settings,
        db,
        user_id=user_id,
        resolve_ku_portal_timetable_targets=_resolve_ku_portal_timetable_targets,
    )


def _build_uclass_sync_health(
    settings: Settings,
    db: Database,
    *,
    user_id: int | None = None,
) -> dict[str, Any]:
    return _build_uclass_sync_health_impl(
        settings,
        db,
        user_id=user_id,
    )


def _build_telegram_listener_health(settings: Settings, db: Database) -> dict[str, Any]:
    return _build_telegram_listener_health_impl(
        settings,
        db,
        effective_telegram_allowed_chat_ids=_effective_telegram_allowed_chat_ids,
    )


def _build_telegram_send_health(settings: Settings, db: Database) -> dict[str, Any]:
    return _build_telegram_send_health_impl(
        settings,
        db,
        effective_telegram_allowed_chat_ids=_effective_telegram_allowed_chat_ids,
    )


def _build_weather_sync_health(
    settings: Settings,
    db: Database,
    *,
    user_id: int | None = None,
) -> dict[str, Any]:
    return _build_weather_sync_health_impl(
        settings,
        db,
        user_id=user_id,
    )


def _build_notice_feed_health(db: Database, kind: str) -> dict[str, Any]:
    return _build_notice_feed_health_impl(db, kind)


def build_beta_ops_health_report(
    settings: Settings,
    db: Database,
    *,
    chat_id: str | None = None,
    user_id: int | None = None,
) -> dict[str, Any]:
    return _build_beta_ops_health_report_state(
        settings,
        db,
        chat_id=chat_id,
        user_id=user_id,
        effective_telegram_allowed_chat_ids=_effective_telegram_allowed_chat_ids,
        resolve_ku_portal_timetable_targets=_resolve_ku_portal_timetable_targets,
    )


def refresh_uclass_for_user(
    settings: Settings,
    db: Database,
    *,
    chat_id: str | None = None,
    user_id: int | None = None,
    send_material_brief_push_enabled: bool = False,
) -> dict[str, Any]:
    scope = _ops_scope_summary(settings, db, chat_id=chat_id, user_id=user_id)
    owner_id = int(scope.get("user_id") or 0)
    chat = str(scope.get("chat_id") or "").strip() or None
    if owner_id <= 0:
        return {"skipped": True, "reason": "user not found", "scope": scope}
    targets = [
        item
        for item in _resolve_uclass_sync_targets(settings, db)
        if int(item.get("user_id") or 0) == owner_id
    ]
    if chat:
        targets = [item for item in targets if str(item.get("chat_id") or "").strip() == chat]
    if not targets:
        reason = "no active moodle connection"
        _record_sync_dashboard_state(
            db,
            "sync_uclass",
            status="skipped",
            last_error=reason,
            cursor_payload={"skipped": True, "reason": reason},
            user_id=owner_id,
        )
        return {"skipped": True, "reason": reason, "scope": scope}
    try:
        result = _sync_uclass_target(
            settings=settings,
            db=db,
            target=targets[0],
            send_material_brief_push_enabled=send_material_brief_push_enabled,
        )
    except Exception as exc:
        _record_sync_dashboard_state(
            db,
            "sync_uclass",
            status="error",
            action_required=1,
            last_error=str(exc),
            cursor_payload={"error": str(exc)},
            user_id=owner_id,
        )
        return {"ok": False, "error": str(exc), "scope": scope}
    return {
        "ok": True,
        "status": "success",
        "scope": scope,
        "upserted_notifications": int(result.get("upserted_notifications") or 0),
        "upserted_tasks": int(result.get("upserted_tasks") or 0),
        "upserted_events": int(result.get("upserted_events") or 0),
        "recorded_artifacts": int(result.get("recorded_artifacts") or 0),
        "failed_artifact_downloads": int(result.get("failed_artifact_downloads") or 0),
        "failed_artifact_extractions": int(result.get("failed_artifact_extractions") or 0),
    }


def refresh_weather_for_user(
    settings: Settings,
    db: Database,
    *,
    chat_id: str | None = None,
    user_id: int | None = None,
) -> dict[str, Any]:
    scope = _ops_scope_summary(settings, db, chat_id=chat_id, user_id=user_id)
    owner_id = int(scope.get("user_id") or 0)
    chat = str(scope.get("chat_id") or "").strip() or None
    if owner_id <= 0:
        return {"skipped": True, "reason": "user not found", "scope": scope}
    snapshot = _get_or_refresh_weather_snapshot(
        settings,
        db,
        user_id=owner_id,
        chat_id=chat,
        force_refresh=True,
    )
    if snapshot is None:
        return {"skipped": True, "reason": "weather snapshot unavailable", "scope": scope}
    current = snapshot.get("current") if isinstance(snapshot.get("current"), dict) else {}
    if str(snapshot.get("error") or "").strip():
        return {
            "ok": False,
            "error": str(snapshot.get("error")).strip(),
            "scope": scope,
            "location_label": str(snapshot.get("location_label") or "").strip() or None,
        }
    return {
        "ok": True,
        "status": "success",
        "scope": scope,
        "location_label": str(snapshot.get("location_label") or "").strip() or None,
        "observed_at": str(snapshot.get("observed_at") or "").strip() or None,
        "temperature_c": current.get("temperature_c"),
    }


def _ops_result_is_error(result: dict[str, Any]) -> bool:
    if not isinstance(result, dict):
        return True
    if bool(result.get("skipped")):
        return False
    if str(result.get("status") or "").strip().lower() == "error":
        return True
    if result.get("error"):
        return True
    if result.get("ok") is False:
        return True
    return False


def refresh_beta_user(
    settings: Settings,
    db: Database,
    *,
    chat_id: str | None = None,
    user_id: int | None = None,
) -> dict[str, Any]:
    scope = _ops_scope_summary(settings, db, chat_id=chat_id, user_id=user_id)
    owner_id = int(scope.get("user_id") or 0)
    chat = str(scope.get("chat_id") or "").strip() or None
    if owner_id <= 0:
        return {"ok": False, "error": "user not found", "scope": scope}
    components = {
        "ku_portal_timetable": prime_ku_portal_timetable_for_user(
            settings,
            db,
            chat_id=chat,
            user_id=owner_id,
            force=True,
        ),
        "uclass_sync": refresh_uclass_for_user(
            settings,
            db,
            chat_id=chat,
            user_id=owner_id,
            send_material_brief_push_enabled=False,
        ),
        "weather_sync": refresh_weather_for_user(
            settings,
            db,
            chat_id=chat,
            user_id=owner_id,
        ),
    }
    error_components = [
        name for name, result in components.items() if _ops_result_is_error(result)
    ]
    skipped_components = [
        name for name, result in components.items() if bool(result.get("skipped"))
    ]
    return {
        "ok": not error_components,
        "scope": scope,
        "components": components,
        "error_components": error_components,
        "skipped_components": skipped_components,
        "health": build_beta_ops_health_report(
            settings,
            db,
            chat_id=chat,
            user_id=owner_id,
        ),
    }


def repair_missing_material_briefs_for_user(
    settings: Settings,
    db: Database,
    *,
    chat_id: str | None = None,
    user_id: int | None = None,
    limit: int = 200,
) -> dict[str, Any]:
    scope = _ops_scope_summary(settings, db, chat_id=chat_id, user_id=user_id)
    owner_id = int(scope.get("user_id") or 0)
    if owner_id <= 0:
        return {"ok": False, "error": "user not found", "scope": scope}
    artifacts = list(db.list_artifacts(limit=max(int(limit), 1), user_id=owner_id))
    result: dict[str, Any] = {
        "ok": True,
        "scope": scope,
        "scanned": len(artifacts),
        "generated": 0,
        "generated_llm": 0,
        "generated_heuristic": 0,
        "skipped_existing": 0,
        "skipped_missing_file": 0,
        "skipped_no_text": 0,
        "failed": 0,
        "items": [],
    }
    for artifact in artifacts:
        metadata = json.loads(str(artifact.metadata_json or "{}"))
        brief = metadata.get("brief")
        if isinstance(brief, dict):
            result["skipped_existing"] += 1
            continue

        local_path = str(artifact.icloud_path or "").strip() or None
        local_file = Path(local_path).expanduser() if local_path else None
        if local_file is not None and not local_file.exists():
            local_file = None
            local_path = None

        text_extract = metadata.get("text_extract") if isinstance(metadata.get("text_extract"), dict) else {}
        extracted_text = str(text_extract.get("excerpt") or "").strip() or None
        if not extracted_text and local_file is not None:
            extracted_text, text_error, text_extract_type = extract_material_text(
                local_file,
                max_chars=max(int(getattr(settings, "material_extract_max_chars", 12000) or 12000), 1),
            )
            if extracted_text:
                metadata["text_extract"] = {
                    "ok": True,
                    "type": text_extract_type,
                    "hash": sha1(extracted_text.encode("utf-8")).hexdigest(),
                    "chars": len(extracted_text),
                    "excerpt": extracted_text[: min(2000, len(extracted_text))],
                }
            elif text_error and not str(text_error).startswith("unsupported"):
                metadata["text_extract"] = {
                    "ok": False,
                    "type": text_extract_type,
                    "error": text_error,
                }
        if not extracted_text:
            if local_file is None:
                result["skipped_missing_file"] += 1
            else:
                result["skipped_no_text"] += 1
            continue

        try:
            provenance = normalize_provenance(metadata, fallback_source=str(artifact.source or "uclass"))
            metadata["brief"] = _build_material_brief(
                settings=settings,
                db=db,
                title=str(artifact.filename or "material"),
                extracted_text=extracted_text,
                local_path=str(local_file) if local_file is not None else local_path,
                artifact_provenance_source=str(provenance.get("source") or artifact.source or "uclass"),
            )
            stored = db.record_artifact(
                external_id=str(artifact.external_id or ""),
                source=str(artifact.source or "uclass"),
                filename=str(artifact.filename or "material"),
                icloud_path=str(local_file) if local_file is not None else local_path,
                content_hash=str(artifact.content_hash or "").strip() or None,
                metadata_json=metadata,
                user_id=owner_id,
            )
            brief_payload = metadata.get("brief") if isinstance(metadata.get("brief"), dict) else {}
            mode = str(brief_payload.get("mode") or "unknown").strip().lower() or "unknown"
            result["generated"] += 1
            if mode == "llm":
                result["generated_llm"] += 1
            else:
                result["generated_heuristic"] += 1
            if len(result["items"]) < 8:
                result["items"].append(
                    {
                        "external_id": str(stored.external_id or ""),
                        "filename": str(stored.filename or ""),
                        "updated_at": str(stored.updated_at or "").strip() or now_utc_iso(),
                        "mode": mode,
                    }
                )
        except Exception as exc:
            logger.warning(
                "failed to repair material brief",
                extra={
                    "external_id": artifact.external_id,
                    "user_id": owner_id,
                    "error": str(exc),
                },
            )
            result["failed"] += 1
    return result


def _failed_stage_candidate(
    *,
    component: str,
    job_name: str,
    stage: str,
    message: str | None,
    status: str,
    last_run_at: str | None,
    scope: dict[str, Any],
    details: dict[str, Any] | None = None,
) -> dict[str, Any]:
    return {
        "component": component,
        "job_name": job_name,
        "stage": stage,
        "status": status,
        "message": str(message or "").strip() or None,
        "last_run_at": str(last_run_at or "").strip() or None,
        "scope": {
            "user_id": _safe_int(scope.get("user_id")),
            "chat_id": str(scope.get("chat_id") or "").strip() or None,
        },
        "details": dict(details or {}),
    }


def _portal_failed_stage_candidate(
    settings: Settings,
    db: Database,
    *,
    scope: dict[str, Any],
) -> dict[str, Any] | None:
    owner_id = _safe_int(scope.get("user_id")) or None
    _, state, cursor = _ops_surface_state(
        db,
        "sync_ku_portal_timetable",
        user_id=owner_id,
        allow_global_fallback=False,
    )
    attempts = _normalize_timetable_source_attempts(cursor.get("source_attempts"))
    for attempt in attempts:
        attempt_status = str(attempt.get("status") or "").strip().lower()
        if attempt_status not in {"error", "fallback", "unsupported"}:
            continue
        stage = (
            "ku_official_api"
            if str(attempt.get("source") or "").strip() == KU_OPENAPI_TIMETABLE_SOURCE
            else "ku_portal_browser"
        )
        return _failed_stage_candidate(
            component="ku_official_api",
            job_name="sync_ku_portal_timetable",
            stage=stage,
            message=attempt.get("reason"),
            status=attempt_status,
            last_run_at=state.last_run_at,
            scope=scope,
            details=attempt,
        )
    if bool(cursor.get("auth_required")):
        return _failed_stage_candidate(
            component="ku_official_api",
            job_name="sync_ku_portal_timetable",
            stage="portal_session",
            message=cursor.get("reason"),
            status="error",
            last_run_at=state.last_run_at,
            scope=scope,
            details={"auth_required": True},
        )
    failed_targets = list(cursor.get("failed_targets") or [])
    if failed_targets:
        item = failed_targets[0] if isinstance(failed_targets[0], dict) else {}
        return _failed_stage_candidate(
            component="ku_official_api",
            job_name="sync_ku_portal_timetable",
            stage="portal_target_sync",
            message=item.get("error"),
            status="error",
            last_run_at=state.last_run_at,
            scope=scope,
            details=item,
        )
    return None


def _uclass_failed_stage_candidate(
    settings: Settings,
    db: Database,
    *,
    scope: dict[str, Any],
) -> dict[str, Any] | None:
    owner_id = _safe_int(scope.get("user_id")) or None
    _, state, cursor = _ops_surface_state(db, "sync_uclass", user_id=owner_id)
    wsfunctions = cursor.get("wsfunctions") if isinstance(cursor.get("wsfunctions"), dict) else {}
    for name, item in wsfunctions.items():
        if not isinstance(item, dict):
            continue
        if int(item.get("failed") or 0) <= 0 and not str(item.get("last_error") or "").strip():
            continue
        return _failed_stage_candidate(
            component="uclass_sync",
            job_name="sync_uclass",
            stage=f"uclass_ws:{name}",
            message=item.get("last_error") or f"{int(item.get('failed') or 0)} failures",
            status="error",
            last_run_at=state.last_run_at,
            scope=scope,
            details={
                "wsfunction": name,
                "failed": int(item.get("failed") or 0),
            },
        )
    html_material_error = str(cursor.get("html_material_error") or "").strip()
    if html_material_error:
        return _failed_stage_candidate(
            component="uclass_sync",
            job_name="sync_uclass",
            stage="uclass_html_materials",
            message=html_material_error,
            status="error",
            last_run_at=state.last_run_at,
            scope=scope,
        )
    if int(cursor.get("material_download_failures") or 0) > 0:
        return _failed_stage_candidate(
            component="uclass_sync",
            job_name="sync_uclass",
            stage="uclass_material_download",
            message=f"{int(cursor.get('material_download_failures') or 0)} downloads failed",
            status="degraded",
            last_run_at=state.last_run_at,
            scope=scope,
        )
    if int(cursor.get("material_extract_failures") or 0) > 0:
        return _failed_stage_candidate(
            component="uclass_sync",
            job_name="sync_uclass",
            stage="uclass_material_extract",
            message=f"{int(cursor.get('material_extract_failures') or 0)} extractions failed",
            status="degraded",
            last_run_at=state.last_run_at,
            scope=scope,
        )
    if str(cursor.get("error") or "").strip():
        return _failed_stage_candidate(
            component="uclass_sync",
            job_name="sync_uclass",
            stage="uclass_sync",
            message=cursor.get("error"),
            status="error",
            last_run_at=state.last_run_at,
            scope=scope,
        )
    return None


def _telegram_failed_stage_candidate(
    settings: Settings,
    db: Database,
    *,
    scope: dict[str, Any],
) -> dict[str, Any] | None:
    _, state, cursor = _ops_surface_state(db, "sync_telegram")
    menu = cursor.get("menu") if isinstance(cursor.get("menu"), dict) else {}
    commands = cursor.get("commands") if isinstance(cursor.get("commands"), dict) else {}
    reminders = cursor.get("reminders") if isinstance(cursor.get("reminders"), dict) else {}
    if str(cursor.get("error") or "").strip():
        return _failed_stage_candidate(
            component="telegram_listener",
            job_name="sync_telegram",
            stage="telegram_polling",
            message=cursor.get("error"),
            status="error",
            last_run_at=state.last_run_at,
            scope=scope,
        )
    if menu and menu.get("ok") is False:
        return _failed_stage_candidate(
            component="telegram_listener",
            job_name="sync_telegram",
            stage="telegram_menu_registration",
            message=menu.get("error") or menu.get("reason"),
            status="error",
            last_run_at=state.last_run_at,
            scope=scope,
            details=menu,
        )
    if int(commands.get("blocked_sends") or 0) > 0:
        return _failed_stage_candidate(
            component="telegram_send",
            job_name="sync_telegram",
            stage="telegram_command_send",
            message=f"{int(commands.get('blocked_sends') or 0)} sends were blocked",
            status="degraded",
            last_run_at=state.last_run_at,
            scope=scope,
            details=commands,
        )
    if int(reminders.get("failed") or 0) > 0:
        return _failed_stage_candidate(
            component="telegram_send",
            job_name="sync_telegram",
            stage="telegram_reminder_send",
            message=f"{int(reminders.get('failed') or 0)} reminder sends failed",
            status="degraded",
            last_run_at=state.last_run_at,
            scope=scope,
            details=reminders,
        )
    if int(commands.get("failed") or 0) > 0:
        return _failed_stage_candidate(
            component="telegram_listener",
            job_name="sync_telegram",
            stage="telegram_command_processing",
            message=f"{int(commands.get('failed') or 0)} command handlers failed",
            status="degraded",
            last_run_at=state.last_run_at,
            scope=scope,
            details=commands,
        )
    return None


def _weather_failed_stage_candidate(
    settings: Settings,
    db: Database,
    *,
    scope: dict[str, Any],
) -> dict[str, Any] | None:
    owner_id = _safe_int(scope.get("user_id")) or None
    _, state, cursor = _ops_surface_state(db, "sync_weather", user_id=owner_id)
    if str(cursor.get("error") or "").strip():
        return _failed_stage_candidate(
            component="weather_sync",
            job_name="sync_weather",
            stage="weather_fetch",
            message=cursor.get("error"),
            status="error",
            last_run_at=state.last_run_at,
            scope=scope,
        )
    air_quality = cursor.get("air_quality") if isinstance(cursor.get("air_quality"), dict) else {}
    if air_quality and air_quality.get("ok") is False and str(air_quality.get("error") or "").strip():
        return _failed_stage_candidate(
            component="weather_sync",
            job_name="sync_weather",
            stage="air_quality_fetch",
            message=air_quality.get("error"),
            status="degraded",
            last_run_at=state.last_run_at,
            scope=scope,
            details=air_quality,
        )
    return None


def _notice_failed_stage_candidates(db: Database, *, scope: dict[str, Any]) -> list[dict[str, Any]]:
    output: list[dict[str, Any]] = []
    for kind in ("general", "academic"):
        state = db.get_sync_state(_portal_notice_snapshot_job(kind))
        cursor = _json_load(state.last_cursor_json)
        attempt = cursor.get("last_attempt") if isinstance(cursor.get("last_attempt"), dict) else {}
        if attempt.get("ok") is not False:
            continue
        output.append(
            _failed_stage_candidate(
                component="notice_fetch",
                job_name=_portal_notice_snapshot_job(kind),
                stage=f"notice_fetch:{kind}",
                message=attempt.get("error"),
                status="error",
                last_run_at=attempt.get("attempted_at") or state.last_run_at,
                scope=scope,
                details={
                    "kind": kind,
                    "http_status": _safe_int(attempt.get("http_status")),
                },
            )
        )
    return output


def inspect_last_failed_stage(
    settings: Settings,
    db: Database,
    *,
    chat_id: str | None = None,
    user_id: int | None = None,
    component: str | None = None,
) -> dict[str, Any]:
    scope = _ops_scope_summary(settings, db, chat_id=chat_id, user_id=user_id)
    if (chat_id or user_id) and not bool(scope.get("found")):
        return {"ok": False, "error": "user not found", "scope": scope, "match": None, "candidates": []}
    requested = str(component or "").strip().lower() or None
    candidates: list[dict[str, Any]] = []
    if requested in {None, "ku_official_api", "portal"}:
        candidate = _portal_failed_stage_candidate(settings, db, scope=scope)
        if candidate:
            candidates.append(candidate)
    if requested in {None, "uclass_sync", "uclass"}:
        candidate = _uclass_failed_stage_candidate(settings, db, scope=scope)
        if candidate:
            candidates.append(candidate)
    if requested in {None, "telegram", "telegram_listener", "telegram_send"}:
        candidate = _telegram_failed_stage_candidate(settings, db, scope=scope)
        if candidate:
            candidates.append(candidate)
    if requested in {None, "weather", "weather_sync"}:
        candidate = _weather_failed_stage_candidate(settings, db, scope=scope)
        if candidate:
            candidates.append(candidate)
    if requested in {None, "notice", "notice_fetch"}:
        candidates.extend(_notice_failed_stage_candidates(db, scope=scope))
    candidates.sort(
        key=lambda item: _parse_dt(str(item.get("last_run_at") or "")) or datetime.fromtimestamp(0, tz=timezone.utc),
        reverse=True,
    )
    return {
        "ok": not candidates,
        "scope": scope,
        "match": candidates[0] if candidates else None,
        "candidates": candidates,
    }


def run_all_jobs(
    settings: Settings,
    db: Database,
    *,
    job_runner: Callable[[str, Callable[[Settings, Database], Any]], Any] | None = None,
) -> PipelineSummary:
    stats: dict[str, Any] = {}
    errors: list[str] = []
    jobs = [
        ("sync_uclass", sync_uclass),
        ("sync_ku_portal_timetable", sync_ku_portal_timetable),
        ("sync_weather", sync_weather),
        ("sync_telegram", sync_telegram),
        ("scheduled_briefings", send_scheduled_briefings),
        ("sync_llm_summaries", sync_llm_summaries),
        ("daily_digest", send_daily_digest),
        ("publish_dashboard", publish_dashboard),
    ]
    for name, fn in jobs:
        try:
            if job_runner is None:
                stats[name] = fn(settings, db)
            else:
                stats[name] = job_runner(name, fn)
            logger.info("job completed", extra={"job": name, "result": stats[name]})
        except Exception as exc:
            error = f"{name}: {exc}"
            errors.append(error)
            logger.exception("job failed", extra={"job": name})
    return PipelineSummary(ok=not errors, stats=stats, errors=errors)
